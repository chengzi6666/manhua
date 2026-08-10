"""
大阅读精灵漫画家 - Flask主程序

请配置以下环境变量：
- API_KEY: 文生图API密钥
- SECRET_KEY: 图片API密钥（已不再使用）
- MODEL_API_URL: 大模型API地址（通用）
- IMAGE_API_URL: 文生图API地址
- ARK_API_KEY: 豆包API密钥（推荐，用于生成高质量教育剧本）
- DOUBAO_MODEL_ID: 豆包模型ID（推荐: doubao-seed-evolving，默认: doubao-1-5-vision-pro-32k-250115）

推荐配置：
1. 豆包API（推荐用于剧本生成）: https://console.volcengine.com/ark/
2. Pollinations.AI生图API: https://pollinations.ai/
3. 其他兼容OpenAI格式的大模型API

注意：要支持图片版PDF（扫描件）的OCR识别，需要安装Tesseract OCR引擎：
Windows: 下载安装 https://github.com/UB-Mannheim/tesseract/wiki
Linux: sudo apt-get install tesseract-ocr tesseract-ocr-chi-sim
"""

import os
import re
import uuid
import logging
import random
import math
import json
import socket
import subprocess
import hashlib
import base64
from contextvars import ContextVar, copy_context
from datetime import datetime
import time
import concurrent.futures
from datetime import datetime
from flask import Flask, render_template, request, jsonify, send_file, send_from_directory, redirect, url_for, session, abort, has_request_context
from flask_cors import CORS
from werkzeug.security import generate_password_hash, check_password_hash
from cryptography.fernet import Fernet, InvalidToken

import sys

log_file_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'app_debug.log')

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(sys.stdout),
        logging.FileHandler(log_file_path, encoding='utf-8', mode='a')
    ]
)
logging.getLogger().handlers[0].stream = sys.stdout
logger = logging.getLogger(__name__)

logger.info(f"日志文件路径: {log_file_path}")

def load_env_file():
    env_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '.env')
    if os.path.exists(env_path):
        try:
            with open(env_path, 'r', encoding='utf-8') as f:
                for line in f:
                    line = line.strip()
                    if line and not line.startswith('#') and '=' in line:
                        key, value = line.split('=', 1)
                        os.environ.setdefault(key.strip(), value.strip())
            logger.info(f"已从.env文件加载配置")
        except Exception as e:
            logger.warning(f"加载.env文件失败: {e}")
    else:
        logger.info("未找到.env文件，将使用系统环境变量")

load_env_file()

from PIL import Image, ImageDraw, ImageFont, ImageFilter
import requests

app = Flask(__name__)
CORS(app)
# 用户与会话：签名 Cookie 只保存用户 ID 与会话版本，不保存密码或 API Key。
app.secret_key = os.environ.get('COMIC_SECRET_KEY', 'd3a7f1c2b9e8450a7c6f2d1e8b4a3c5f7e9d0a1b2c3f4e5d6c7b8a9f0e1d2c3')
app.config.update(
    SESSION_COOKIE_HTTPONLY=True,
    SESSION_COOKIE_SAMESITE='Lax',
    SESSION_COOKIE_SECURE=bool(os.environ.get('RAILWAY_ENVIRONMENT')),
)


def _redact_log_secrets(value):
    """递归隐藏请求日志中的密码、恢复信息和各种 API Key。"""
    sensitive_tokens = ('password', 'api_key', 'apikey', 'secret', 'token', 'authorization')
    if isinstance(value, dict):
        return {
            key: ('[REDACTED]' if any(token in str(key).lower() for token in sensitive_tokens)
                  else _redact_log_secrets(item))
            for key, item in value.items()
        }
    if isinstance(value, list):
        return [_redact_log_secrets(item) for item in value]
    return value

@app.before_request
def log_request():
    import json
    request_id = f"req_{id(request) % 100000}"
    request.start_time = datetime.now()
    request.request_id = request_id
    
    body_summary = ""
    try:
        if request.method in ['POST', 'PUT', 'PATCH']:
            content_type = request.content_type or ''
            if 'application/json' in content_type:
                data = _redact_log_secrets(request.get_json(silent=True) or {})
                body_summary = json.dumps(data, ensure_ascii=False)[:500]
            else:
                body_summary = f"Content-Type: {content_type}, length: {request.content_length or 0}"
    except Exception as e:
        body_summary = f"读取请求体失败: {str(e)}"
    
    logger.info(f"[REQ-{request_id}] {request.method} {request.path} | IP: {request.remote_addr} | Body: {body_summary}")


@app.after_request
def log_response(response):
    try:
        elapsed = (datetime.now() - request.start_time).total_seconds() * 1000
        request_id = getattr(request, 'request_id', 'unknown')
        
        response_summary = ""
        try:
            content_type = response.content_type or ''
            if 'application/json' in content_type:
                try:
                    response_data = response.get_json()
                    response_summary = json.dumps(response_data, ensure_ascii=False)[:300]
                except:
                    response_summary = f"JSON解析失败, length: {response.content_length or 0}"
            else:
                response_summary = f"Content-Type: {content_type}, length: {response.content_length or 0}"
        except Exception as e:
            response_summary = f"读取响应失败: {str(e)}"
        
        logger.info(f"[RES-{request_id}] {request.method} {request.path} | Status: {response.status_code} | Time: {elapsed:.2f}ms | Response: {response_summary}")
    except Exception as e:
        logger.error(f"日志记录异常: {str(e)}")
    
    return response

# 配置上传文件夹
UPLOAD_FOLDER = os.path.join(app.root_path, 'static', 'uploads')
ALLOWED_EXTENSIONS = {'pdf', 'png', 'jpg', 'jpeg', 'gif', 'bmp', 'docx', 'doc', 'txt', 'md', 'rtf'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 1024 * 1024 * 1024  # 最大1GB
# 单表单字段内存上限（multipart 解析时，单个文本字段如 base64 图片超过该值会抛 413）
# 编辑器导出的 dataURL 可达数 MB，需调大，否则保存接口报 RequestEntityTooLarge
app.config['MAX_FORM_MEMORY_SIZE'] = 64 * 1024 * 1024  # 64MB
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'dev-secret-key')

# 上传体积超限时返回明确的 413（而非被当 500 重试），便于前端提示
from werkzeug.exceptions import RequestEntityTooLarge
@app.errorhandler(RequestEntityTooLarge)
def handle_too_large(e):
    return jsonify({'error': '上传内容过大，请减小图片尺寸或联系管理员'}), 413

# 确保上传目录存在
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# API配置
API_KEY = os.environ.get('API_KEY')
SECRET_ACCESS_KEY = os.environ.get('SECRET_ACCESS_KEY')
MODEL_API_URL = os.environ.get('MODEL_API_URL', '')
MODEL_API_KEY = os.environ.get('MODEL_API_KEY', '')
IMAGE_API_URL = os.environ.get('IMAGE_API_URL', 'https://api.example.com/image')

# IP图片目录
IP_DIR = os.path.join(app.root_path, 'static', 'ip')

# 任务会话管理（服务器端状态）
class TaskSession:
    def __init__(self, task_id):
        self.task_id = task_id
        self.pdf_text = ""
        self.script = None
        self.backgrounds = []
        self.comics = []
        self.ip_paths = []
        self.comic_type = 'preview'
        self.created_at = datetime.now()
    
    def to_dict(self):
        return {
            'task_id': self.task_id,
            'pdf_text': self.pdf_text,
            'script': self.script,
            'backgrounds': self.backgrounds,
            'comics': self.comics,
            'ip_paths': self.ip_paths,
            'comic_type': self.comic_type,
            'has_script': self.script is not None,
            'has_backgrounds': len(self.backgrounds) > 0,
            'has_comics': len(self.comics) > 0
        }

tasks = {}

def create_task():
    task_id = str(uuid.uuid4())[:8]
    tasks[task_id] = TaskSession(task_id)
    return task_id
OUTPUT_DIR = os.path.join(app.root_path, 'static', 'output')

# 确保输出目录存在
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Logo 水印图片目录
LOGO_DIR = os.path.join(app.root_path, 'static', 'logos')
os.makedirs(LOGO_DIR, exist_ok=True)

# 全局变量存储上传的文件内容
uploaded_files = {}

# 尝试初始化OCR相关库
try:
    import pytesseract
    OCR_TESSERACT_AVAILABLE = True
    if os.name == 'nt':
        tesseract_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            r'D:\Program Files\Tesseract-OCR\tesseract.exe'
        ]
        for path in tesseract_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                logger.info(f"Tesseract OCR路径设置成功: {path}")
                break
    logger.info("Tesseract OCR已启用")
except ImportError:
    OCR_TESSERACT_AVAILABLE = False
    logger.warning("pytesseract未安装")

try:
    import easyocr
    OCR_EASYOCR_AVAILABLE = True
    easyocr_reader = None
    logger.info("EasyOCR已安装（按需加载模型）")
except ImportError:
    OCR_EASYOCR_AVAILABLE = False
    logger.warning("easyocr未安装")

try:
    import fitz
    PDF2IMAGE_AVAILABLE = True
    logger.info("PyMuPDF(fitz)已安装")
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("PyMuPDF未安装")

BAIDU_OCR_API_KEY = os.environ.get('BAIDU_OCR_API_KEY', '')
BAIDU_OCR_SECRET_KEY = os.environ.get('BAIDU_OCR_SECRET_KEY', '')
OCR_BAIDU_AVAILABLE = bool(BAIDU_OCR_API_KEY and BAIDU_OCR_SECRET_KEY)

if OCR_BAIDU_AVAILABLE:
    logger.info("百度OCR API已配置")
else:
    logger.warning("百度OCR API未配置（可选）")

SERVER_ARK_API_KEY = os.environ.get('ARK_API_KEY', '')
DOUBAO_MODEL_ID = os.environ.get('DOUBAO_MODEL_ID', 'doubao-seed-evolving')
# 图像生成与文本模型分开配置。ARK_IMAGE_MODEL 可在 .env 中覆盖，默认使用已开通的 Seedream 5.0 lite。
ARK_IMAGE_MODEL = os.environ.get('ARK_IMAGE_MODEL', 'doubao-seedream-5-0-lite-260128').strip()
_request_ark_key = ContextVar('request_ark_key', default=None)
_request_image_provider = ContextVar('request_image_provider', default=None)


class UserScopedArkKey:
    """兼容旧调用点的用户级密钥代理，避免并发请求之间串用 API Key。"""
    def value(self):
        contextual = _request_ark_key.get()
        if contextual is not None:
            return contextual
        if has_request_context():
            uid = session.get('user_id')
            if uid:
                try:
                    return get_user_ark_api_key(uid) or ''
                except Exception as exc:
                    logger.warning(f'[user-api-key] 读取当前用户密钥失败: {exc}')
                    return ''
            # Railway 公网环境不允许匿名用户借用服务器公共密钥。
            if os.environ.get('RAILWAY_ENVIRONMENT'):
                return ''
        return SERVER_ARK_API_KEY

    def __str__(self):
        return self.value()

    def __bool__(self):
        return bool(self.value())


ARK_API_KEY = UserScopedArkKey()


def get_request_image_provider():
    """Return the current user's image provider without leaking credentials."""
    configured = _request_image_provider.get()
    if configured:
        return configured
    key = str(ARK_API_KEY or '')
    return {
        'provider': 'doubao', 'api_key': key,
        'model': ARK_IMAGE_MODEL,
        'base_url': 'https://ark.cn-beijing.volces.com/api/v3/images/generations',
    } if key else None
OCR_DOUBAO_AVAILABLE = ARK_API_KEY

if OCR_DOUBAO_AVAILABLE:
    logger.info(f"豆包API已配置，模型: {DOUBAO_MODEL_ID}")
else:
    logger.warning("豆包API未配置（可选，设置ARK_API_KEY启用）")


def call_openai_api(prompt, model=None, temperature=0.7, max_tokens=2000):
    """调用OpenAI兼容的LLM API（支持Ollama、free.ai等）"""
    if not MODEL_API_URL or MODEL_API_URL == 'https://api.example.com/model':
        logger.warning("未配置MODEL_API_URL")
        return None
    
    try:
        headers = {
            'Content-Type': 'application/json',
        }
        
        if MODEL_API_KEY:
            headers['Authorization'] = f'Bearer {MODEL_API_KEY}'
        
        data = {
            'model': model or 'glm-4-flash',
            'messages': [
                {'role': 'system', 'content': '你是一位专业的教育漫画编剧，擅长为小学生设计有趣的漫画故事。'},
                {'role': 'user', 'content': prompt}
            ],
            'temperature': temperature,
            'max_tokens': max_tokens
        }
        
        logger.info(f"正在调用LLM API: {MODEL_API_URL[:50]}...")
        response = requests.post(MODEL_API_URL, headers=headers, json=data, timeout=120)
        response.raise_for_status()
        
        result = response.json()
        
        if 'choices' in result and len(result['choices']) > 0:
            content = result['choices'][0]['message']['content']
            return content
        else:
            logger.warning("LLM API返回格式异常")
            return None
            
    except Exception as e:
        logger.error(f"调用LLM API失败: {str(e)}")
        return None

def get_random_ip_image():
    """从static/ip/目录随机选取一张IP图片"""
    try:
        ip_files = [f for f in os.listdir(IP_DIR) if f.lower().endswith('.png')]
        if not ip_files:
            logger.error("static/ip/目录下没有找到PNG图片")
            return None
        random_ip = random.choice(ip_files)
        return os.path.join(IP_DIR, random_ip)
    except Exception as e:
        logger.error(f"获取IP图片失败: {str(e)}")
        return None


def get_selected_guide_image(characters, selected_sprite=None):
    """返回用户为知识小精灵明确选择的本地图片；仅在没有选择时才随机回退。"""
    candidates = []
    if isinstance(characters, dict):
        for name in ('知识小精灵', '小精灵', '精灵'):
            if characters.get(name):
                candidates.append(characters[name])
    if selected_sprite:
        candidates.append(selected_sprite)

    root = os.path.abspath(app.root_path)
    for url in candidates:
        if not isinstance(url, str) or not url.startswith('/static/'):
            continue
        path = os.path.abspath(os.path.join(app.root_path, url.lstrip('/')))
        if path.startswith(root + os.sep) and os.path.isfile(path):
            return path
    return get_random_ip_image()


SCENARIO_STYLES = {
    'humorous': '幽默搞笑风格，夸张的表情和动作，让小朋友开怀大笑',
    'adventure': '冒险探索风格，神秘的场景，充满好奇心和探索精神',
    'friendship': '友情温暖风格，温馨的场景，体现朋友之间的互助和关爱',
    'magic': '奇幻魔法风格，魔法元素，星星和闪光，充满想象力',
    'school': '校园生活风格，教室和校园场景，贴近学生生活',
    'animal': '动物可爱风格，动物角色，可爱的表情和动作',
    'learning': '学习教育风格，清晰的知识点讲解，寓教于乐',
    'family': '温馨家庭风格，家庭场景，亲子互动'
}

SOUND_EFFECTS = []

PANEL_ANGLES = [
    'eye level', 'low angle', 'high angle', 'close-up', 
    'medium shot', 'wide shot', 'extreme close-up', 'over the shoulder'
]

PANEL_COMPOSITION = [
    'centered', 'rule of thirds', 'leading lines', 'frame within frame',
    'symmetric', 'asymmetric', 'negative space', 'dynamic diagonal'
]

EDUCATION_PROMPTS = {
    'preview': """你是一位专业的教育漫画编剧，擅长将枯燥的知识点转化为有趣的漫画故事。

你的任务是根据课程重点，为小学生设计一个引人入胜的漫画剧本，激发他们的好奇心和学习兴趣。

请遵循以下原则：
1. 故事性强：要有有趣的情节和生动的角色互动
2. 知识准确：确保知识点的准确性和科学性
3. 年龄适配：语言要适合小学生理解，生动有趣
4. 启发思考：通过提问和引导，激发孩子主动思考
5. 情感共鸣：故事要有温暖、积极向上的基调

请输出JSON格式，包含以下字段：
- scenes: 场景中文描述（用于生成背景图）
- scenes_zh: 场景中文描述
- dialogues: 对话内容
- speakers: 说话者
- hints: 思考提示
""",
    'summary': """你是一位专业的教育漫画编剧，擅长将复杂的知识点转化为易于理解的漫画故事。

你的任务是根据课程重点，为小学生设计一个帮助复习和巩固知识的漫画剧本。

请遵循以下原则：
1. 结构清晰：逻辑清晰，层次分明
2. 重点突出：突出核心知识点和易错点
3. 联系实际：将知识与生活实际联系起来
4. 总结归纳：帮助孩子形成知识体系
5. 鼓励激励：给予积极的鼓励和肯定

请输出JSON格式，包含以下字段：
- scenes: 场景中文描述（用于生成背景图）
- scenes_zh: 场景中文描述
- dialogues: 对话内容
- speakers: 说话者
- summary: 知识点总结
"""
}

CHARACTER_SYSTEMS = {
    'student': {
        'name': '校园学生',
        'roles': ['调皮学生', '学霸同桌', '普通学生'],
        'visual_style': 'Q版卡通，明亮色彩，校园风格',
        'interaction_patterns': ['斗嘴', '互助', '学习', '玩耍']
    },
    'explorer': {
        'name': '探险家',
        'roles': ['勇敢队长', '好奇队员', '知识队员'],
        'visual_style': '冒险风格，探险装备，自然场景',
        'interaction_patterns': ['探索', '发现', '合作', '解决问题']
    },
    'time_traveler': {
        'name': '时空旅行者',
        'roles': ['小穿越者', '向导', '历史人物'],
        'visual_style': '古风/复古风格，历史场景',
        'interaction_patterns': ['对话', '体验', '学习', '感悟']
    },
    'classic': {
        'name': '经典名著',
        'roles': ['名著角色A', '名著角色B', '旁白'],
        'visual_style': '经典文学风格，符合原著描述',
        'interaction_patterns': ['对话', '冲突', '成长', '领悟']
    },
    'elves': {
        'name': '小精灵',
        'roles': ['好奇的探索者', '知识渊博者', '幽默搞笑者'],
        'visual_style': 'Q版卡通，可爱精灵形象',
        'interaction_patterns': ['提问', '解答', '探索', '互助']
    }
}

ROLE_GUIDE_CONFIG = {
    'name': '知识小精灵',
    'role': '知识引导者',
    'personality': ['智慧', '耐心', '幽默', '鼓励'],
    'appearance': '可爱精灵形象，发光效果，小巧玲珑，透明翅膀',
    'position': 'corner',
    'appearance_panel': 'last',
    'dialogue_style': '启发式提问，简明讲解'
}

PANEL_COUNT_OPTIONS = [1, 2, 3, 4, 5, 6, 7, 8]

def calculate_panel_structure(panel_count):
    if panel_count == 1:
        return {
            'phase1': [1],
            'phase2': [],
            'phase3': [],
            'phase4': [],
            'guide_panel': 1
        }
    elif panel_count == 2:
        return {
            'phase1': [1],
            'phase2': [],
            'phase3': [],
            'phase4': [2],
            'guide_panel': 2
        }
    elif panel_count == 3:
        return {
            'phase1': [1],
            'phase2': [2],
            'phase3': [],
            'phase4': [3],
            'guide_panel': 3
        }
    elif panel_count == 4:
        return {
            'phase1': [1],
            'phase2': [2],
            'phase3': [3],
            'phase4': [4],
            'guide_panel': 4
        }
    elif panel_count == 5:
        return {
            'phase1': [1],
            'phase2': [2, 3],
            'phase3': [4],
            'phase4': [5],
            'guide_panel': 5
        }
    elif panel_count == 6:
        return {
            'phase1': [1],
            'phase2': [2, 3],
            'phase3': [4, 5],
            'phase4': [6],
            'guide_panel': 6
        }
    elif panel_count == 7:
        return {
            'phase1': [1, 2],
            'phase2': [3, 4],
            'phase3': [5, 6],
            'phase4': [7],
            'guide_panel': 7
        }
    elif panel_count == 8:
        return {
            'phase1': [1, 2],
            'phase2': [3, 4],
            'phase3': [5, 6],
            'phase4': [7, 8],
            'guide_panel': 8
        }
    else:
        return {
            'phase1': [1],
            'phase2': [2],
            'phase3': [3],
            'phase4': [4],
            'guide_panel': 4
        }

def generate_panel_structure_desc(panel_count):
    if panel_count == 1:
        return """   - 第1格：完整故事展示，画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""
    elif panel_count == 2:
        return """   - 第1格：故事开端，人物登场，交代故事基础背景
   - 第2格：情节收尾，画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""
    elif panel_count == 3:
        return """   - 第1格：故事开端，人物登场，交代故事基础背景
   - 第2格：情节推进，埋下矛盾伏笔
   - 第3格：情节收尾，画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""
    elif panel_count == 4:
        return """   - 第1格：故事开端，人物登场，交代故事基础背景
   - 第2格：情节推进，埋下矛盾伏笔
   - 第3格：剧情高潮，出现冲突、危机或反转
   - 第4格：情节暂时收尾，画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""
    elif panel_count == 5:
        return """   - 第1格：故事开端，人物登场，交代故事基础背景
   - 第2-3格：情节推进，详细展示矛盾伏笔和角色互动
   - 第4格：剧情高潮，出现冲突、危机或反转
   - 第5格：情节收尾，画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""
    elif panel_count == 6:
        return """   - 第1格：故事开端，人物登场，交代故事基础背景
   - 第2-3格：情节推进，详细展示矛盾伏笔和角色互动
   - 第4-5格：剧情高潮，详细展示冲突、危机或反转过程
   - 第6格：情节暂时收尾，画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""
    elif panel_count == 7:
        return """   - 第1-2格：故事开端，详细展示人物登场和故事背景
   - 第3-4格：情节推进，详细展示矛盾伏笔和角色互动过程
   - 第5-6格：剧情高潮，详细展示冲突、危机或反转的完整过程
   - 第7格：情节收尾，画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""
    elif panel_count == 8:
        return """   - 第1-2格：故事开端，详细展示人物登场和故事背景
   - 第3-4格：情节推进，详细展示矛盾伏笔和角色互动过程
   - 第5-6格：剧情高潮，详细展示冲突、危机或反转的完整过程
   - 第7-8格：情节收尾，第7格展示冲突解决过程，第8格画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""
    else:
        return """   - 第1格：故事开端，人物登场，交代故事基础背景
   - 第2格：情节推进，埋下矛盾伏笔
   - 第3格：剧情高潮，出现冲突、危机或反转
   - 第4格：情节暂时收尾，画面角落悬浮知识小精灵，小精灵对话气泡提出开放式思考问题"""

def generate_characters(character_system='student', count=3):
    system = CHARACTER_SYSTEMS.get(character_system, CHARACTER_SYSTEMS['student'])
    roles = system['roles'][:count]
    characters = []
    for i, role in enumerate(roles):
        characters.append({
            'name': f'{role}{chr(ord("A") + i)}',
            'role': role,
            'personality': ['活泼', '好奇', '勇敢'][i % 3],
            'visual_style': system['visual_style']
        })
    return characters

def clean_script_narration(script, characters):
    """清洗剧本中的旁白，确保说话者都是角色名"""
    if not script:
        return script
    
    speakers = script.get('speakers', [])
    dialogues = script.get('dialogues', [])
    
    for i, speaker in enumerate(speakers):
        if speaker and ('旁白' in speaker or '叙述' in speaker or '解说' in speaker or '旁白者' in speaker):
            speakers[i] = characters[i % len(characters)]
        
        if dialogues[i] and ('旁白：' in dialogues[i] or '旁白:' in dialogues[i]):
            dialogues[i] = dialogues[i].replace('旁白：', '').replace('旁白:', '').strip()
    
    script['speakers'] = speakers
    script['dialogues'] = dialogues
    
    if 'sound_effects' in script:
        del script['sound_effects']
    
    return script

def force_speaker_rotation(script, characters):
    """强制小精灵轮流发言，确保每个角色都有发言机会"""
    if not script:
        return script
    
    speakers = script.get('speakers', [])
    if not speakers or len(characters) <= 1:
        return script
    
    panel_count = len(speakers)
    
    for i in range(panel_count):
        speakers[i] = characters[i % len(characters)]
    
    script['speakers'] = speakers
    return script

def analyze_pdf_content(content, scenario_style='humorous'):
    """Step1：分析PDF内容，提取知识点和教学目标"""
    default_result = {
        "topic": "课程重点",
        "key_points": ["知识点1", "知识点2", "知识点3"],
        "teaching_goal": "帮助孩子理解课程内容",
        "difficulty": "小学",
        "story_premise": f"小精灵们在{SCENARIO_STYLES.get(scenario_style, '有趣')}的场景中学习知识"
    }
    
    prompt = f"""你是一位专业的教育内容分析师，擅长从学术文档中提取适合小学生学习的知识点。

请分析以下内容，提取关键信息：

【内容】
{content[:2000]}

【输出要求】
请输出JSON格式，包含以下字段：
- topic: 课程主题（简洁明了）
- key_points: 数组，包含3-5个核心知识点（每个知识点不超过20字）
- teaching_goal: 教学目标（让孩子掌握什么）
- difficulty: 适合的年级（如：小学低年级、小学高年级）
- story_premise: 故事创意（基于知识点构思一个适合小精灵的故事背景，20-30字）

注意：输出必须是合法的JSON格式，不要包含任何其他文字。"""
    
    if OCR_DOUBAO_AVAILABLE:
        try:
            url = "https://ark.cn-beijing.volces.com/api/v3/responses"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {ARK_API_KEY}"
            }
            
            data = {
                "model": DOUBAO_MODEL_ID,
                "input": prompt,
                "temperature": 0.7
            }
            
            logger.info("正在分析PDF内容...")
            response = requests.post(url, headers=headers, json=data, timeout=120)
            response.raise_for_status()
            
            result = response.json()
            
            content_text = ""
            for output_item in result.get('output', []):
                if output_item.get('type') == 'message':
                    for content_item in output_item.get('content', []):
                        if content_item.get('type') == 'output_text':
                            content_text = content_item.get('text', '')
                            break
                    if content_text:
                        break
            
            if content_text:
                content_text = content_text.replace('```json', '').replace('```', '').strip()
                analysis = json.loads(content_text)
                return analysis
            else:
                logger.warning("PDF分析返回格式异常")
        except Exception as e:
            logger.error(f"豆包API分析失败: {str(e)}，尝试使用通用LLM API")
    
    if MODEL_API_URL and MODEL_API_URL != 'https://api.example.com/model':
        logger.info("使用通用LLM API分析PDF内容")
        content_text = call_openai_api(prompt)
        if content_text:
            content_text = content_text.replace('```json', '').replace('```', '').strip()
            try:
                analysis = json.loads(content_text)
                return analysis
            except:
                logger.warning("通用LLM API返回格式异常")
    
    logger.warning("未配置可用的LLM API，使用默认提取结果")
    return default_result

def generate_story_plan(analysis, character_count=3, panel_count=6, scenario_style='humorous', character_system='student'):
    """Step2：生成故事大纲（世界设定、角色分配、情节发展）"""
    system = CHARACTER_SYSTEMS.get(character_system, CHARACTER_SYSTEMS['student'])
    system_roles = system['roles'][:character_count]
    characters = [f'{system_roles[i]}{chr(ord("A") + i)}' for i in range(character_count)]
    
    default_plan = {
        "world_setting": {
            "main_location": "神奇的知识森林",
            "time_period": "现代",
            "visual_style": system['visual_style'],
            "atmosphere": "有趣、探索"
        },
        "characters": [{"name": char, "role": system_roles[i], "personality": "可爱"} for i, char in enumerate(characters)],
        "story_arc": [{"panel": i+1, "location": "神奇的知识森林", "action": f"第{i+1}格场景", "mood": "开心"} for i in range(panel_count)]
    }
    
    style_instruction = {
        'humorous': '搞笑风格：角色有滑稽行为，对话幽默，通过搞笑场景学习',
        'adventure': '冒险风格：在探险场景中（森林/洞穴/神秘遗迹）发现知识',
        'friendship': '友情风格：通过友情故事学习，角色互相帮助',
        'magic': '魔法风格：用魔法比喻解释知识，使用魔法道具',
        'school': '校园风格：课堂互动，角色扮演学生和老师',
        'learning': '学习风格：通过探索发现知识，强调知识的力量',
        'animal': '动物风格：角色拟人化为可爱动物',
        'family': '家庭风格：家庭场景互动学习'
    }.get(scenario_style, '轻松有趣的风格')
    
    panel_structure_desc = generate_panel_structure_desc(panel_count)
    
    prompt = f"""你是一位专业的儿童漫画编剧，擅长为小学生设计有趣的漫画故事。

请根据以下知识点，为{character_count}个{system['name']}设计一个完整的漫画故事大纲：

【知识点分析】
主题：{analysis.get('topic', '课程重点')}
核心知识点：{', '.join(analysis.get('key_points', []))}
教学目标：{analysis.get('teaching_goal', '')}
故事创意：{analysis.get('story_premise', '')}

【剧本要求】
- 角色：{', '.join(characters)}（角色类型：{system['name']}）
- 格子数：{panel_count}格
- 风格：{style_instruction}
- 视觉风格：{system['visual_style']}
- 互动模式：{', '.join(system['interaction_patterns'])}

【输出要求】
请输出JSON格式，包含以下字段：

1. world_setting（世界设定）：
   - main_location: 主要场景地点（如：古代中国山洞、神奇的知识森林、校园教室）
   - time_period: 时间设定（如：远古时期、现代）
   - visual_style: 视觉风格描述（{system['visual_style']}）
   - atmosphere: 整体氛围（如：神秘、探索、欢乐）

2. characters（角色设定）：
   数组，每个角色包含：
   - name: 角色名（{', '.join(characters)}）
   - role: 角色定位（如：{', '.join(system_roles)}）
   - personality: 性格特点（如：爱提问、爱解答、爱吐槽）

3. story_arc（情节发展）：
   数组，包含{panel_count}个元素，每个元素包含：
   - panel: 格子序号（1-{panel_count}）
   - location: 该格发生的具体地点
   - action: 角色在做什么
   - mood: 该格的情绪氛围
   - speaker: 主要说话者（必须是{', '.join(characters)}之一）
   - dialogue_summary: 对话内容概要

【四段式故事结构要求】
{panel_structure_desc}

注意：输出必须是合法的JSON格式，不要包含任何其他文字。"""
    
    if OCR_DOUBAO_AVAILABLE:
        try:
            url = "https://ark.cn-beijing.volces.com/api/v3/responses"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {ARK_API_KEY}"
            }
            
            data = {
                "model": DOUBAO_MODEL_ID,
                "input": prompt,
                "temperature": 0.7
            }
            
            logger.info("正在生成故事大纲...")
            response = requests.post(url, headers=headers, json=data, timeout=120)
            response.raise_for_status()
            
            result = response.json()
            
            content = ""
            for output_item in result.get('output', []):
                if output_item.get('type') == 'message':
                    for content_item in output_item.get('content', []):
                        if content_item.get('type') == 'output_text':
                            content = content_item.get('text', '')
                            break
                    if content:
                        break
            
            if content:
                content = content.replace('```json', '').replace('```', '').strip()
                story_plan = json.loads(content)
                return story_plan
            else:
                logger.warning("故事大纲返回格式异常，尝试使用通用LLM API")
        except Exception as e:
            logger.error(f"豆包API生成故事大纲失败: {str(e)}，尝试使用通用LLM API")
    
    if MODEL_API_URL and MODEL_API_URL != 'https://api.example.com/model':
        logger.info("使用通用LLM API生成故事大纲")
        content = call_openai_api(prompt)
        if content:
            content = content.replace('```json', '').replace('```', '').strip()
            try:
                story_plan = json.loads(content)
                return story_plan
            except:
                logger.warning("通用LLM API返回格式异常")
    
    logger.warning("未配置可用的LLM API，使用默认故事大纲")
    return default_plan

def generate_panel_content(story_plan, panel_index, scenario_style='humorous'):
    """Step3：基于Story Plan生成单格内容（对话+背景描述）"""
    world_setting = story_plan.get('world_setting', {})
    characters = story_plan.get('characters', [])
    story_arc = story_plan.get('story_arc', [])
    
    panel_info = story_arc[panel_index] if panel_index < len(story_arc) else {}
    character_names = [c['name'] for c in characters] if characters else ['小精灵A']
    speaker = panel_info.get('speaker', character_names[0]) if character_names else '小精灵A'
    
    default_content = {
        "scene": f"cartoon background, Q version style, bright colors, {world_setting.get('main_location', '')}",
        "scene_zh": panel_info.get('location', '卡通背景'),
        "dialogue": panel_info.get('dialogue_summary', '小精灵在说话'),
        "speaker": speaker,
        "hint": "思考一下这个问题"
    }
    
    character_desc = ', '.join([f"{c['name']}（{c['role']}，{c['personality']}）" for c in characters]) if characters else ''
    
    prompt = f"""你是一位专业的儿童漫画编剧。请根据以下故事大纲，生成第{panel_index+1}格的具体内容。

【世界设定】
地点：{world_setting.get('main_location', '')}
时间：{world_setting.get('time_period', '')}
视觉风格：{world_setting.get('visual_style', '')}
氛围：{world_setting.get('atmosphere', '')}

【角色设定】
{character_desc}

【当前格子规划】
格子：第{panel_index+1}格
地点：{panel_info.get('location', '')}
动作：{panel_info.get('action', '')}
情绪：{panel_info.get('mood', '')}
说话者：{panel_info.get('speaker', '')}
对话概要：{panel_info.get('dialogue_summary', '')}

【输出要求】
请输出JSON格式，包含以下字段：
- scene: 英文背景描述（用于文生图API，详细描述场景环境，不要包含人物，Q版卡通风格，与世界设定一致）
- scene_zh: 中文背景描述
- dialogue: 中文对白（符合角色性格，有趣生动，包含具体知识内容）
- speaker: 说话者（必须是{', '.join(character_names)}之一）
- hint: 小提示（帮助孩子思考，与知识点相关）

注意：背景描述要详细，适合Q版卡通风格，不要包含人物。对话要生动有趣，体现角色性格。输出必须是合法的JSON格式，不要包含任何其他文字。"""
    
    if OCR_DOUBAO_AVAILABLE:
        try:
            url = "https://ark.cn-beijing.volces.com/api/v3/responses"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {ARK_API_KEY}"
            }
            
            data = {
                "model": DOUBAO_MODEL_ID,
                "input": prompt,
                "temperature": 0.7
            }
            
            logger.info(f"正在生成第{panel_index+1}格内容...")
            response = requests.post(url, headers=headers, json=data, timeout=120)
            response.raise_for_status()
            
            result = response.json()
            
            content = ""
            for output_item in result.get('output', []):
                if output_item.get('type') == 'message':
                    for content_item in output_item.get('content', []):
                        if content_item.get('type') == 'output_text':
                            content = content_item.get('text', '')
                            break
                    if content:
                        break
            
            if content:
                content = content.replace('```json', '').replace('```', '').strip()
                panel_content = json.loads(content)
                return panel_content
            else:
                logger.warning(f"第{panel_index+1}格内容返回格式异常，尝试使用通用LLM API")
        except Exception as e:
            logger.error(f"豆包API生成第{panel_index+1}格内容失败: {str(e)}，尝试使用通用LLM API")
    
    if MODEL_API_URL and MODEL_API_URL != 'https://api.example.com/model':
        logger.info(f"使用通用LLM API生成第{panel_index+1}格内容")
        content = call_openai_api(prompt)
        if content:
            content = content.replace('```json', '').replace('```', '').strip()
            try:
                panel_content = json.loads(content)
                return panel_content
            except:
                logger.warning("通用LLM API返回格式异常")
    
    logger.warning(f"未配置可用的LLM API，使用第{panel_index+1}格默认内容")
    return default_content

def generate_comic_script(course_topic, comic_type='preview', panel_count=4, character_count=1, scenario_style='humorous', ip_images=None, character_system='student'):
    """生成漫画剧本（支持多格数、多角色、多种情景风格、图片输入）
    
    支持多种LLM API：豆包API、通用OpenAI兼容API（Ollama、free.ai等）
    
    comic_type: 'preview' - 埋下伏笔类（课前预习）, 'summary' - 点题类（课后复习）
    panel_count: 漫画格子数（4/6/8）
    character_count: 主角数量（1-4）
    scenario_style: 情景风格
    ip_images: IP图片路径列表（可选，用于让模型识别角色形象）
    character_system: 角色体系（student/explorer/time_traveler/classic/elves）
    """
    
    system = CHARACTER_SYSTEMS.get(character_system, CHARACTER_SYSTEMS['student'])
    system_roles = system['roles'][:character_count]
    characters = [f'{system_roles[i]}{chr(ord("A") + i)}' for i in range(character_count)]
    
    story_plan = {
        "world_setting": {
            "main_location": "漫画故事场景",
            "time_period": "现代",
            "visual_style": system['visual_style'],
            "atmosphere": "有趣、探索"
        },
        "characters": [{"name": characters[i], "role": system_roles[i], "personality": "活泼"} for i in range(character_count)],
        "story_arc": [{"panel": i+1, "location": "故事场景", "action": f"第{i+1}格场景", "mood": "开心"} for i in range(panel_count)]
    }
    
    style_desc = SCENARIO_STYLES.get(scenario_style, '轻松有趣的风格')
    
    if comic_type == 'preview':
        prompt = f"""你是一位专业的教育漫画编剧，擅长将枯燥的知识点转化为有趣的漫画故事。

你的任务是根据课程重点，为小学生设计一个引人入胜的漫画剧本，用于**课前预习**。

【课程重点】
{course_topic}

【剧本要求】
- 类型：埋下伏笔型（课前预习）
- 目的：激发好奇心，让孩子带着问题听课，绝对不要直接给出答案！
- 风格：{style_desc}
- 角色：{character_count}个{system['name']}（{', '.join(characters)}）
- 格子数：{panel_count}格

【创作原则】
1. ❌ 绝对不要直接讲知识点！不要解释原理！不要给出答案！
2. ✅ 只提出有趣的问题，制造悬念和好奇心
3. ✅ 故事性强：要有有趣的情节和生动的角色互动，让孩子产生代入感
4. ✅ 年龄适配：语言要适合小学生理解，生动有趣，避免使用过于复杂的词汇
5. ✅ 启发思考：通过提问和引导，激发孩子主动思考
6. ✅ 情感共鸣：故事要有温暖、积极向上的基调

【{panel_count}格漫画结构 - 悬念模式】
第1格【引入】：从日常场景引入，出现一个奇怪的现象或物体（制造好奇）
第2格【疑问】：小精灵们提出有趣的问题，表达困惑（引发思考）
第3格【探索】：小精灵们尝试猜测或探索，但得不到答案（加深悬念）
第{panel_count}格【悬念】：留下一个大大的问号，引导孩子期待课堂揭晓答案

请输出JSON格式，包含以下字段：
- scenes: 数组，包含{panel_count}个场景的英文背景描述（用于生成背景图，不要包含人物，要详细描述场景环境，包含视角角度如low angle/close-up等）
- scenes_zh: 数组，包含{panel_count}个场景的中文描述
- dialogues: 数组，包含{panel_count}句中文对白（以疑问为主，引发好奇心，不要给出答案）
- speakers: 数组，包含{panel_count}个说话者（必须是{', '.join(characters)}之一）
- hints: 数组，包含{panel_count}个小提示（引导孩子思考，不要给出答案）
- angles: 数组，包含{panel_count}个视角描述（如"俯视"、"仰视"、"特写"等）

示例格式：
{{
    "scenes": ["公园里，一个气球自己往前飞", "小精灵们惊讶地看着气球", ...],
    "scenes_zh": ["公园里，一个气球自己往前飞", "小精灵们惊讶地看着气球", ...],
    "dialogues": ["气球怎么自己飞？", "谁在推它呀？", "没看到人推啊！", "到底是什么力量？"],
    "speakers": ["小精灵A", "小精灵B", ...],
    "hints": ["想想看，气球为什么会动？", "什么东西能让气球飞起来？", ...]
}}

注意：背景描述要详细，适合Q版卡通风格，不要包含人物。对话要以疑问为主，引发好奇心。
"""
    else:
        prompt = f"""你是一位专业的教育漫画编剧，擅长将复杂的知识点转化为易于理解的漫画故事。

你的任务是根据课程重点，为小学生设计一个帮助**复习和巩固知识**的漫画剧本。

【课程重点】
{course_topic}

【剧本要求】
- 类型：点题型（课后复习）
- 目的：总结知识点，强化记忆，帮助迁移应用，必须给出明确答案！
- 风格：{style_desc}
- 角色：{character_count}个{system['name']}（{', '.join(characters)}）
- 格子数：{panel_count}格

【创作原则】
1. ✅ 必须明确讲解知识点！给出清晰答案！帮助孩子巩固记忆！
2. ✅ 结构清晰：逻辑清晰，层次分明，帮助孩子建立知识框架
3. ✅ 重点突出：突出核心知识点和易错点，强化记忆
4. ✅ 联系实际：将知识与生活实际联系起来，让孩子知道知识的应用场景
5. ✅ 总结归纳：帮助孩子形成知识体系，学会举一反三
6. ✅ 鼓励激励：给予积极的鼓励和肯定，增强学习信心

【{panel_count}格漫画结构 - 复习总结模式】
第1格【回顾】：小精灵开心地回忆今天学了什么
第2格【讲解】：一句话明确解释核心知识点，给出答案
第3格【应用】：用生活中的场景举例说明知识的用途
第4格【测试】：一个有趣的小问题，检验学习成果（给出答案）
第5-{panel_count-1}格【拓展】：有趣的知识延伸或生活应用场景
第{panel_count}格【鼓励】：一句简短的鼓励，增强学习信心

请输出JSON格式，包含以下字段：
- scenes: 数组，包含{panel_count}个场景的英文背景描述（用于生成背景图，不要包含人物，要详细描述场景环境，包含视角角度如low angle/close-up等）
- scenes_zh: 数组，包含{panel_count}个场景的中文描述
- dialogues: 数组，包含{panel_count}句中文对白（明确讲解知识，给出答案，帮助记忆）
- speakers: 数组，包含{panel_count}个说话者（必须是{', '.join(characters)}之一）
- summary: 字符串，包含知识点总结和迁移建议（详细总结核心知识点，提供生活中的应用建议）
- angles: 数组，包含{panel_count}个视角描述（如"俯视"、"仰视"、"特写"等）

示例格式：
{{
    "scenes": ["明亮的教室，彩色黑板上写着作用力和反作用力", "公园里，小朋友在划船", ...],
    "scenes_zh": ["明亮的教室，彩色黑板上写着作用力和反作用力", "公园里，小朋友在划船", ...],
    "dialogues": ["今天学了什么？", "作用力和反作用力！", "划船就是这个原理！", "推墙时墙也推你！"],
    "speakers": ["小精灵A", "小精灵B", ...],
    "summary": "【知识点总结】作用力和反作用力\n\n重点内容：\n1. 作用力和反作用力大小相等\n2. 方向相反\n3. 作用在不同物体上\n\n生活中的应用：划船、走路、火箭发射..."
}}

注意：背景描述要详细，适合Q版卡通风格，不要包含人物。对话要明确给出答案，帮助孩子理解和记忆知识点。
"""
    
    if OCR_DOUBAO_AVAILABLE:
        logger.info("=== 调用豆包API生成漫画剧本 ===")
        script = generate_comic_script_with_doubao_v2(course_topic, comic_type, panel_count, character_count, scenario_style, ip_images, None, character_system)
        
        if script:
            script = clean_script_narration(script, characters)
            script = force_speaker_rotation(script, characters)
            script['story_plan'] = story_plan
            logger.info("=== 剧本生成完成 ===")
            return script
        else:
            logger.warning("豆包API生成失败，尝试使用通用LLM API")
    
    if MODEL_API_URL and MODEL_API_URL != 'https://api.example.com/model':
        logger.info("=== 使用通用LLM API生成漫画剧本 ===")
        content = call_openai_api(prompt)
        
        if content:
            json_match = re.search(r'\{[\s\S]*\}', content)
            if json_match:
                try:
                    script = json.loads(json_match.group())
                    script = clean_script_narration(script, characters)
                    script = force_speaker_rotation(script, characters)
                    script['story_plan'] = story_plan
                    logger.info("=== 剧本生成完成 ===")
                    return script
                except:
                    logger.warning("通用LLM API返回JSON解析失败")
            else:
                logger.warning("通用LLM API返回非JSON格式")
    
    logger.info("未配置可用的LLM API，使用默认剧本")
    script = get_default_script(course_topic, comic_type, panel_count, character_count, scenario_style, character_system)
    script['story_plan'] = story_plan
    return script


def extract_key_concepts(text):
    """从文本中提取关键概念"""
    if not text:
        return []
    
    text = text.replace('\\n', '\n').replace('\\r', '')
    
    keywords = []
    
    lines = text.split('\n')
    for line in lines:
        line = line.strip()
        if len(line) > 5:
            keywords.append(line[:20])
    
    for pattern in [r'第[\u4e00-\u9fa5\d]+章', r'第[\u4e00-\u9fa5\d]+节', r'[\u4e00-\u9fa5]{2,8}法', r'[\u4e00-\u9fa5]{2,8}技巧', r'[\u4e00-\u9fa5]{2,8}方法']:
        matches = re.findall(pattern, text)
        keywords.extend(matches)
    
    return list(set(keywords))[:5]

def get_default_script(course_topic, comic_type='preview', panel_count=4, character_count=1, scenario_style='humorous', character_system='student'):
    """生成默认剧本（支持多格数、多角色、多种情景风格），包含丰富的教育内容"""
    key_concepts = extract_key_concepts(course_topic)
    main_topic = course_topic[:30] if len(course_topic) > 30 else course_topic
    
    system = CHARACTER_SYSTEMS.get(character_system, CHARACTER_SYSTEMS['student'])
    system_roles = system['roles'][:character_count]
    characters = [f'{system_roles[i]}{chr(ord("A") + i)}' for i in range(character_count)]
    
    style_prompts = {
        'humorous': 'funny, humorous, exaggerated expressions, laughter',
        'adventure': 'adventure, exploration, mysterious, curious',
        'friendship': 'warm, friendly, caring, helping each other',
        'magic': 'magic, fantasy, stars, sparkles, magical elements',
        'school': 'classroom, school, educational, studying',
        'animal': 'cute animals, adorable, animal characters',
        'learning': 'educational, teaching, knowledge, explaining',
        'family': 'family, home, parents, warm atmosphere'
    }
    
    style_en = style_prompts.get(scenario_style, 'cute, cartoon, bright colors')
    
    scenes = []
    scenes_zh = []
    dialogues = []
    speakers = []
    hints = []
    
    educational_scenarios = {
        'humorous': {
            'scenes': [
                ("cute cartoon classroom with silly funny books doing cartwheels on desk, Q version style, bright colorful, chaotic fun", "教室里，书本们在课桌上翻跟头，好有趣！"),
                ("cute cartoon elves laughing and pointing at floating numbers, Q version style, bright colors, surprised expressions", "小精灵们指着飘浮的数字哈哈大笑！"),
                ("cute cartoon chalkboard with numbers dancing, Q version style, bright colors, magical", "黑板上数字在跳舞！"),
                ("cute cartoon classroom with giant question mark balloon, Q version style, bright colors, funny", "一个巨大的问号气球飘浮在教室！"),
                ("cute cartoon elves trying to catch flying numbers, Q version style, bright colors, playful", "小精灵们在抓飞来飞去的数字！"),
                ("cute cartoon books stacked in funny shapes, Q version style, bright colors, silly", "书本堆成了奇怪的形状！"),
                ("cute cartoon numbers having a party, Q version style, bright colors, festive", "数字们在开派对！"),
                ("cute cartoon elves with confused faces looking at math problems, Q version style, bright colors, puzzled", "小精灵们看着数学题一脸困惑！")
            ],
            'dialogues': [
                f"{main_topic}是什么？书在翻跟头！" if main_topic else "这是什么？书在翻跟头！",
                f"{key_concepts[0] if key_concepts else '数字'}怎么会飞呢？",
                "它们在玩什么游戏呀？",
                "为什么数字会跳舞？",
                "谁在指挥它们呢？",
                "这是什么魔法？",
                "答案藏在哪里？",
                "课堂上揭晓谜底！"
            ],
            'hints': [
                "想想看，什么能让数字动起来？",
                "观察一下数字之间有什么关系？",
                "和小伙伴讨论你的发现！",
                "带着问题去听课吧！",
                "猜猜数字们在做什么？",
                "数字的秘密是什么？",
                "动动脑筋想一想！",
                "答案就在课堂里！"
            ]
        },
        'adventure': {
            'scenes': [
                ("cute cartoon magical forest with mysterious glowing path leading to unknown, Q version style, adventurous, twilight", "神秘的魔法森林，发光小路通向未知！"),
                ("cute cartoon elves discovering ancient treasure chest covered in vines, Q version style, adventurous, curious", "小精灵们发现了被藤蔓覆盖的古老宝箱！"),
                ("cute cartoon magical creatures guiding explorers through forest, Q version style, adventurous, friendly", "神奇生物在森林中指引探险者！"),
                ("cute cartoon cave entrance with glowing math symbols on walls, Q version style, adventurous, mysterious", "洞穴入口的墙上有发光的数学符号！"),
                ("cute cartoon elves crossing rickety rope bridge over river, Q version style, adventurous, brave", "小精灵们勇敢地走过摇晃的绳索桥！"),
                ("cute cartoon treasure map with X marking spot, Q version style, adventurous, exciting", "藏宝图上有一个大大的X！"),
                ("cute cartoon ancient temple with floating numbers, Q version style, adventurous, magical", "古老神庙里飘浮着数字！"),
                ("cute cartoon elves celebrating discovery with fireworks, Q version style, adventurous, joyful", "小精灵们庆祝发现，烟花绽放！")
            ],
            'dialogues': [
                f"{main_topic}是什么？森林里藏着秘密！" if main_topic else "森林里藏着什么秘密？",
                f"{key_concepts[0] if key_concepts else '知识'}藏在哪里？",
                "宝箱里会有答案吗？",
                "这些符号是什么意思？",
                "我们能找到宝藏吗？",
                "地图上的X在哪里？",
                "神庙里有什么？",
                "一起来揭开谜题！"
            ],
            'hints': [
                "探险需要什么装备？",
                "遇到困难时要怎么办？",
                "团队合作能做什么？",
                "勇敢迈出第一步！",
                "地图会告诉我们什么？",
                "符号背后有什么秘密？",
                "仔细观察周围环境！",
                "知识就是宝藏！"
            ]
        },
        'friendship': {
            'scenes': [
                ("cute cartoon park with friends playing together, Q version style, warm colors", "公园里，朋友们在一起玩耍"),
                ("cute cartoon elves helping each other climb tree, Q version style, warm colors", "小精灵们互相帮助爬树"),
                ("cute cartoon picnic scene with happy friends, Q version style, warm colors", "快乐的野餐场景"),
                ("cute cartoon sunset with friends holding hands, Q version style, warm colors", "夕阳下，朋友们手牵手")
            ],
            'dialogues': [
                f"{main_topic}是什么？我们一起找找看！" if main_topic else "这是什么？我们一起找找看！",
                f"{key_concepts[0] if key_concepts else '这个知识'}在哪里呢？",
                "好朋友一起找答案！",
                "课堂上揭晓谜底！"
            ],
            'hints': [
                "好朋友之间应该怎么做？",
                "分享想法能让友谊更牢固吗？",
                "互相帮助是什么感觉？",
                "和好朋友一起去课堂找答案！"
            ]
        },
        'magic': {
            'scenes': [
                ("cute cartoon magical wizard lab with floating potions, Q version style, magical", "魔法巫师的实验室，药水在空中飘浮"),
                ("cute cartoon elves casting colorful magic spells, Q version style, magical", "小精灵们在施展彩色魔法"),
                ("cute cartoon magical library with talking books, Q version style, magical", "会说话的书本的魔法图书馆"),
                ("cute cartoon rainbow bridge leading to knowledge castle, Q version style, magical", "通向知识城堡的彩虹桥")
            ],
            'dialogues': [
                f"{main_topic}是什么魔法？药水怎么会飘？" if main_topic else "这是什么魔法？药水怎么会飘？",
                f"{key_concepts[0] if key_concepts else '这个现象'}是什么咒语？",
                "谁在施魔法呀？",
                "魔法的秘密在课堂！"
            ],
            'hints': [
                "魔法和知识有什么相似之处？",
                "如果你有魔法，会怎么探索？",
                "相信知识的力量！",
                "一起来揭开魔法的秘密！"
            ]
        },
        'school': {
            'scenes': [
                ("cute cartoon classroom with colorful chalkboard, Q version style, educational", "明亮的教室，彩色黑板"),
                ("cute cartoon teacher explaining to curious students, Q version style, educational", "老师在给好奇的同学们讲解"),
                ("cute cartoon group study session, Q version style, educational", "小组学习场景"),
                ("cute cartoon graduation ceremony with proud students, Q version style, educational", "毕业典礼，自豪的同学们")
            ],
            'dialogues': [
                f"{main_topic}是什么？黑板上写着什么？" if main_topic else "这是什么？黑板上写着什么？",
                f"{key_concepts[0] if key_concepts else '这个知识'}是什么意思呀？",
                "老师会告诉我们吗？",
                "认真听课找答案！"
            ],
            'hints': [
                "课堂上应该怎么做？",
                "带着问题听课效果更好哦！",
                "不懂的问题要及时问！",
                "一起来课堂找答案！"
            ]
        },
        'learning': {
            'scenes': [
                ("cute cartoon knowledge tree with books as leaves, Q version style, educational", "知识树，书本是它的叶子"),
                ("cute cartoon elves climbing knowledge ladder, Q version style, educational", "小精灵们在攀登知识的阶梯"),
                ("cute cartoon treasure map of knowledge, Q version style, educational", "知识宝藏地图"),
                ("cute cartoon graduation cap with stars, Q version style, educational", "带着星星的毕业帽")
            ],
            'dialogues': [
                f"{main_topic}是什么？树上的书在说话！" if main_topic else "这是什么？树上的书在说话！",
                f"{key_concepts[0] if key_concepts else '知识'}藏在哪里呢？",
                "地图能找到答案吗？",
                "一起来探索知识吧！"
            ],
            'hints': [
                "你最喜欢探索什么？",
                "发现新知识的感觉怎么样？",
                "知识能帮我们解决什么问题？",
                "带着问题去探索吧！"
            ]
        },
        'animal': {
            'scenes': [
                ("cute cartoon forest with adorable animals playing, Q version style, cute", "森林里，可爱的小动物们在玩耍"),
                ("cute cartoon little animals learning together, Q version style, cute", "小动物们在一起学习"),
                ("cute cartoon animal teacher explaining to students, Q version style, cute", "动物老师在给学生们讲解"),
                ("cute cartoon animals celebrating with flowers, Q version style, cute", "小动物们拿着花庆祝")
            ],
            'dialogues': [
                f"{main_topic}是什么？小动物们在讨论什么？" if main_topic else "这是什么？小动物们在讨论什么？",
                f"{key_concepts[0] if key_concepts else '这个知识'}在哪里呀？",
                "动物老师会知道吗？",
                "和小动物一起找答案！"
            ],
            'hints': [
                "小动物们是怎么探索的？",
                "你最喜欢哪种小动物？",
                "动物世界里有哪些有趣的现象？",
                "和小动物一起去课堂找答案！"
            ]
        },
        'family': {
            'scenes': [
                ("cute cartoon warm family room with parents and kids, Q version style, warm", "温馨的客厅，爸爸妈妈和孩子们在一起"),
                ("cute cartoon family reading together, Q version style, warm", "一家人在一起阅读"),
                ("cute cartoon parents helping children with homework, Q version style, warm", "爸爸妈妈在帮孩子做作业"),
                ("cute cartoon happy family dinner, Q version style, warm", "快乐的家庭晚餐")
            ],
            'dialogues': [
                f"{main_topic}是什么？我们一起找找看！" if main_topic else "这是什么？我们一起找找看！",
                f"{key_concepts[0] if key_concepts else '这个知识'}在哪里呀？",
                "爸爸妈妈知道答案吗？",
                "一起去课堂找答案！"
            ],
            'hints': [
                "和家人一起探索是什么感觉？",
                "爸爸妈妈能帮你做什么？",
                "家庭里有哪些有趣的现象？",
                "和家人一起去课堂找答案！"
            ]
        }
    }
    
    if comic_type == 'preview':
        scenario = educational_scenarios.get(scenario_style, educational_scenarios['humorous'])
        
        scene_templates = scenario['scenes']
        dialogue_templates = scenario['dialogues']
        hint_templates = scenario['hints']
        
        for i in range(panel_count):
            scene_en, scene_zh = scene_templates[i % len(scene_templates)]
            scenes.append(scene_en)
            scenes_zh.append(scene_zh)
            dialogues.append(dialogue_templates[i % len(dialogue_templates)])
            if i == panel_count - 1:
                speakers.append('知识小精灵')
            else:
                speakers.append(characters[i % character_count])
            hints.append(hint_templates[i % len(hint_templates)])
        
        sound_effects = []
        angles = []
        for i in range(panel_count):
            if i == 0:
                sound_effects.append('')
            elif i == panel_count - 1:
                sound_effects.append('')
            else:
                sound_effects.append('')
            
            angles.append(random.choice(['俯视', '仰视', '平视', '特写', '全景']))
        
        return {
            "scenes": scenes,
            "scenes_zh": scenes_zh,
            "dialogues": dialogues,
            "speakers": speakers,
            "hints": hints,
            "sound_effects": sound_effects,
            "angles": angles
        }
    else:
        scene_templates = [
            (f"cute cartoon classroom with colorful chalkboard showing important notes, Q version style, {style_en}", "明亮的教室，彩色黑板上写着重要的笔记"),
            (f"cute cartoon chalkboard with diagrams and explanations, Q version style, {style_en}", "黑板上画着图表和解释"),
            (f"cute cartoon children applying knowledge in daily life, Q version style, {style_en}", "小朋友们在生活中应用知识"),
            (f"cute cartoon celebration with stars and confetti, Q version style, {style_en}", "庆祝成功，星星和彩带飞舞"),
            (f"cute cartoon characters reviewing together, Q version style, {style_en}", "小精灵们一起复习"),
            (f"cute cartoon solving problems happily, Q version style, {style_en}", "开心地解决问题"),
            (f"cute cartoon sharing knowledge with friends, Q version style, {style_en}", "和朋友分享知识"),
            (f"cute cartoon proud achievement with trophy, Q version style, {style_en}", "获得奖杯，自豪的小精灵")
        ]
        
        dialogue_templates = [
            f"今天学了{main_topic}！" if main_topic else "今天学了新知识！",
            f"{key_concepts[0] if key_concepts else '核心知识'}就是答案！",
            f"生活中{key_concepts[1] if len(key_concepts) > 1 else '这样用'}！",
            "记住了！太棒了！",
            "复习巩固，加深记忆！",
            "这个你学会了吗？",
            "和小伙伴分享心得！",
            "继续加油，你最棒！"
        ]
        
        hint_templates = [
            "回顾今天学的内容", "找出重点和难点", "练习应用知识", "总结学习方法",
            "和同学互相提问", "做一些练习题", "记录学习笔记", "保持学习热情"
        ]
        
        for i in range(panel_count):
            scene_en, scene_zh = scene_templates[i % len(scene_templates)]
            scenes.append(scene_en)
            scenes_zh.append(scene_zh)
            dialogues.append(dialogue_templates[i % len(dialogue_templates)])
            if i == panel_count - 1:
                speakers.append('知识小精灵')
            else:
                speakers.append(characters[i % character_count])
            hints.append(hint_templates[i % len(hint_templates)])
        
        summary_content = f"""【知识点总结】{main_topic}

重点内容：
1. {key_concepts[0] if key_concepts else '理解核心概念'}
2. {key_concepts[1] if len(key_concepts) > 1 else '掌握关键方法'}
3. 学会举一反三

生活中的应用：
- {key_concepts[0] if key_concepts else '知识'}可以帮助我们解决很多问题
- 在日常生活中多加练习，加深理解
- 和家人朋友分享你学到的知识

鼓励寄语：
学习是一个不断积累的过程，每天进步一点点，你就是最棒的！继续加油哦！"""
        
        sound_effects = []
        angles = []
        for i in range(panel_count):
            if i == panel_count - 1:
                sound_effects.append('')
            else:
                sound_effects.append('')
            
            angles.append(random.choice(['俯视', '仰视', '平视', '特写', '全景']))
        
        return {
            "scenes": scenes,
            "scenes_zh": scenes_zh,
            "dialogues": dialogues,
            "speakers": speakers,
            "summary": summary_content,
            "sound_effects": sound_effects,
            "angles": angles
        }

def translate_to_english(text):
    """将中文场景描述翻译成英文（用于文生图API）"""
    if not text:
        return text
    
    if not any('\u4e00' <= char <= '\u9fff' for char in text):
        return text
    
    if not OCR_DOUBAO_AVAILABLE:
        logger.warning("豆包API未配置，无法翻译中文场景描述")
        return text
    
    try:
        url = "https://ark.cn-beijing.volces.com/api/v3/responses"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {ARK_API_KEY}"
        }
        
        prompt = f"""请将以下中文场景描述翻译成英文，用于文生图API。
保持场景细节，不要添加额外内容，确保翻译准确且适合图像生成。

中文描述：
{text}

请只输出英文翻译结果，不要包含任何其他文字。"""
        
        data = {
            "model": DOUBAO_MODEL_ID,
            "input": prompt
        }
        
        response = requests.post(url, headers=headers, json=data, timeout=30)
        response.raise_for_status()
        
        result = response.json()
        translated = ""
        for output_item in result.get('output', []):
            if output_item.get('type') == 'message':
                for content_item in output_item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        translated = content_item.get('text', '').strip()
                        break
        
        logger.info(f"中文场景描述翻译完成: '{text[:30]}...' -> '{translated[:30]}...'")
        return translated if translated else text
    except Exception as e:
        logger.error(f"翻译失败: {str(e)}")
        return text


def generate_jimeng_signature(access_key, secret_key, method, url, body_json):
    """生成Pollinations.AI请求头（无需签名，返回空字典）"""
    return {}


def submit_jimeng_task(prompt, width=1024, height=1024, seed=-1):
    """使用Pollinations.AI提交图片生成任务（直接返回图片内容）"""
    import urllib.parse
    
    max_retries = 5
    retry_delay = 3
    
    for attempt in range(max_retries):
        try:
            encoded_prompt = urllib.parse.quote(prompt)
            url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width={width}&height={height}&nologo=true&model=flux"
            
            logger.info(f"[Pollinations] 请求URL: {url[:200]}...")
            
            response = requests.get(url, timeout=120)
            response.raise_for_status()
            
            content_type = response.headers.get('Content-Type', '')
            content_length = len(response.content)
            
            if 'image' in content_type and content_length > 1000:
                logger.info(f"Pollinations.AI图片生成成功！Content-Type: {content_type}, 大小: {content_length}字节")
                return response.content
            else:
                logger.warning(f"Pollinations返回的不是有效图片: Content-Type={content_type}, 长度={content_length}字节")
                
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)
                    retry_delay = min(retry_delay * 2, 15)
                    continue
                return None
                
        except requests.exceptions.RequestException as e:
            status_code = getattr(e.response, 'status_code', 'unknown')
            logger.error(f"Pollinations.AI请求异常（第{attempt+1}次）: {str(e)}, 状态码: {status_code}")
            
            if status_code == 429:
                time.sleep(retry_delay * 3)
                retry_delay = min(retry_delay * 2, 30)
            elif attempt < max_retries - 1:
                time.sleep(retry_delay)
                retry_delay = min(retry_delay * 2, 15)
                continue
        except Exception as e:
            logger.error(f"Pollinations.AI生成失败（第{attempt+1}次）: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
                retry_delay = min(retry_delay * 2, 15)
                continue
    
    logger.warning("Pollinations.AI生成失败")
    return None


def get_jimeng_result(task_id):
    """查询Pollinations.AI任务结果（同步模式，直接返回结果）"""
    return {"code": 10000, "data": {"status": "SUCCESS", "image_urls": [task_id]}}


COMMON_CHARACTER_NAMES = [
    '小明', '小红', '小刚', '小丽', '小华', '小强', '小芳', '小林',
    '小宇', '朵朵', '小炎犬', '火焰犬', '烈焰犬', '叶影猫', '叶绿猫',
    '草叶猫', '小冰熊', '冰冻熊', '急冻熊', '暴风鸟', '疾风鸟', '迷你雀',
    '暮光龙', '精灵龙', '蓝焰龙', '电光鼠', '电球鼠', '雷电鼠',
    '知识小精灵', '小精灵', '老师', '同学', '朋友', '爸爸', '妈妈',
    '爷爷', '奶奶', '哥哥', '姐姐', '弟弟', '妹妹'
]


def extract_environment_only(scene_prompt, character_names=None):
    """从场景描述中提取纯环境信息，剥离角色行为描述和角色名称
    
    将包含角色动作的场景描述转换为纯环境描述，避免AI在背景中生成人物。
    
    示例:
    "小明在教室里回答老师的问题" -> "明亮的教室，黑板，课桌椅"
    "小红在操场上跑步" -> "操场，跑道，蓝天白云"
    
    character_names: 可选的角色名称列表，用于更精确地移除角色名称
    """
    if not scene_prompt:
        return ""
    
    import re
    
    cleaned = scene_prompt.strip()
    
    character_indicators = [
        r'[他她它]在', r'[他她它]正在', r'[他她它]站在', r'[他她它]坐在',
        r'[他她它]跑', r'[他她它]走', r'[他她它]说', r'[他她它]看',
        r'[他她它]做', r'[他她它]想', r'[他她它]问', r'[他她它]回答',
        r'[他她它]们在', r'[他她它]们正在', r'[他她它]们站在', r'[他她它]们坐在',
        r'[他她它]们跑', r'[他她它]们走', r'[他她它]们说', r'[他她它]们看',
    ]
    
    for indicator in character_indicators:
        match = re.search(indicator, cleaned)
        if match:
            cleaned = cleaned[:match.start()].strip()
            break
    
    action_words = ['在', '正在', '站在', '坐在', '跑', '走', '说', '看', '做', '想', '问', '回答']
    
    parts = re.split(r'[,，。；;]', cleaned)
    filtered_parts = []
    for part in parts:
        part = part.strip()
        if not part:
            continue
        
        has_action = False
        for action_word in action_words:
            if action_word in part:
                has_action = True
                break
        
        if not has_action:
            filtered_parts.append(part)
        else:
            match = re.search(r'(在|正在|站在|坐在|位于|处于)\s*(.+)', part)
            if match:
                env_part = match.group(2).strip()
                if env_part and len(env_part) > 2:
                    filtered_parts.append(env_part)
    
    result = '，'.join(filtered_parts)
    
    if not result:
        result = scene_prompt
    
    names_to_remove = set(COMMON_CHARACTER_NAMES)
    if character_names and isinstance(character_names, list):
        names_to_remove.update([str(name).strip() for name in character_names if str(name).strip()])
    
    for name in sorted(names_to_remove, key=len, reverse=True):
        if name and name in result:
            result = result.replace(name, '').replace('，，', '，').replace('。。', '。').strip()
    
    # 显式人物名词（中文本层再剥一层，避免翻译后英文名漏检）：男孩/女孩/小孩/老师/学生/人物/人群等
    _zh_person = [
        '人物', '人群', '人们', '人类', '人影', '人形', '人偶',
        '小孩', '孩子', '儿童', '小朋友', '少儿', '幼儿',
        '男孩', '女孩', '少男', '少女', '男童', '女童',
        '男人', '女人', '男士', '女士', '青年', '少年', '中年', '老年', '婴儿', '宝宝',
        '学生', '老师', '教师', '同学', '路人', '群众', '行人', '民众', '大众', '市民', '观众',
        '背影', '剪影', '侧影',
        '角色', '主角', '配角', '主人公',
        '伙伴', '朋友',
    ]
    for w in _zh_person:
        if w in result:
            result = result.replace(w, '').replace('，，', '，').replace('，，', '，')
    # 常见诗人/先贤名（大阅读古诗文高频人物词）
    for w in _FAMOUS_NAMES_ZH:
        if w in result:
            result = result.replace(w, '').replace('，，', '，')
    # 中文「X人」结构（一个人/有人/一群人/两个人）
    result = re.sub(r'(一|两|二|三|四|五|六|七|八|九|十|几|多|全|每|各|某|这|那|众)\s*个?\s*(名|位|个)?\s*人', '', result)
    result = re.sub(r'(有|无|没有)\s*人', '', result)
    result = re.sub(r'一\s*群\s*人', '', result)
    result = re.sub(r'几\s*个?\s*人', '', result)

    result = re.sub(r'[，。;；]+$', '', result).strip()

    if not result:
        result = scene_prompt

    return result.strip()


# 背景（纯场景）模式下，必须屏蔽的中/英文人物词。Pollinations 接收翻译后的英文，
# 因此中英文都要覆盖。名单偏保守，只删明确指「人/角色」的词，避免误伤环境词
# （如「人行道」「人工」「人口」不在名单内，不会被删）。
_PERSON_NOUNS_ZH = [
    '人物', '人群', '人们', '人类', '人影', '人形', '人偶', '人像',
    '小孩', '孩子', '儿童', '小朋友', '少儿', '幼儿', '婴幼儿',
    '男孩', '女孩', '少男', '少女', '男童', '女童',
    '男人', '女人', '男士', '女士', '青年', '少年', '中年', '老年', '婴儿', '宝宝',
    '学生', '老师', '教师', '同学', '路人', '群众', '行人', '民众', '大众', '市民', '观众', '顾客',
    '背影', '剪影', '侧影', '黑影',
    '角色', '主角', '配角', '主人公', '人物角色',
    '伙伴', '朋友', '小伙伴们',
]
_PERSON_NOUNS_EN = [
    'person', 'people', 'humans', 'human', 'child', 'children', 'kid', 'kids',
    'boy', 'girl', 'man', 'men', 'woman', 'women', 'baby', 'babies', 'infant',
    'student', 'students', 'teacher', 'teachers', 'classmate', 'classmates',
    'crowd', 'crowds', 'figure', 'figures', 'silhouette', 'silhouettes',
    'passerby', 'passersby', 'audience', 'character', 'characters',
    'protagonist', 'hero', 'heroine', 'spectator', 'spectators', 'pedestrian',
    'friend', 'friends',
]

# 大阅读产品古诗文场景极多，这些常见诗人/先贤/文中人物名也是「人物提示词」，
# 必须屏蔽（背景里不该出现具体人物）。中文层 extract_environment_only 与最终
# sanitize_no_people 都会用。英文列为其拼音/惯用英文，防翻译后漏检。
_FAMOUS_NAMES_ZH = [
    '李白', '杜甫', '王维', '苏轼', '苏东坡', '白居易', '孟浩然', '陶渊明', '贺知章',
    '王之涣', '杜牧', '柳宗元', '杨万里', '王安石', '陆游', '辛弃疾', '李清照',
    '屈原', '王昌龄', '岑参', '高适', '王勃', '李商隐', '杜牧', '刘禹锡', '韩愈',
    '孟郊', '贾岛', '元稹', '张继', '王翰', '卢纶', '韦应物', '刘长卿', '常建',
    '曹操', '曹植', '陶渊明', '谢灵运', '屈原', '庄子', '孔子', '孟子', '老子',
    '司马迁', '司马光', '诸葛亮', '周瑜', '刘备', '关羽', '张飞', '曹操', '孙权',
    '唐僧', '孙悟空', '猪八戒', '沙僧', '梁山伯', '祝英台', '牛郎', '织女',
    '夸父', '后羿', '嫦娥', '精卫', '愚公', '木兰',
]
_FAMOUS_NAMES_EN = [
    'Li Bai', 'Du Fu', 'Wang Wei', 'Su Shi', 'Su Dongpo', 'Bai Juyi', 'Meng Haoran',
    'Tao Yuanming', 'Du Mu', 'Li Shangyin', 'Li Qingzhao', 'Lu You', 'Xin Qiji',
    'Wang Anshi', 'Yang Wanli', 'Qu Yuan', 'Confucius', 'Mencius', 'Laozi',
    'Zhuangzi', 'Sima Qian', 'Zhuge Liang', 'Sun Wukong', 'Tang Monk', 'Monkey King',
]


def sanitize_no_people(prompt, extra_names=None):
    """纯背景（纯场景）模式下，剥离提示词里任何人物/角色相关词，杜绝模型生成人。

    同时处理中文与翻译后的英文（Pollinations 实际接收英文）。extra_names 可传入
    角色名列表（中/英文皆可），一并删除（如脚本里的「李白」→ 翻译后「Li Bai」）。
    返回清洗后的 prompt。注意：本函数只删「人」相关词，不改动后面的正面空环境约束。
    """
    if not prompt:
        return prompt
    import re
    p = prompt

    # ---- 中文显式人物名词（多字，直接删除）----
    for w in _PERSON_NOUNS_ZH:
        if w in p:
            p = p.replace(w, '')

    # ---- 中文「X人」结构：一个人/有人/一群人/两个人/许多人/几个人/众人 ----
    p = re.sub(r'(一|两|二|三|四|五|六|七|八|九|十|几|多|全|每|各|某|这|那|众)\s*个?\s*(名|位|个)?\s*人', '', p)
    p = re.sub(r'(有|无|没有|无认)\s*人', '', p)
    p = re.sub(r'一\s*群\s*人', '', p)
    p = re.sub(r'几\s*个?\s*人', '', p)

    # ---- 英文人物词（翻译后，词边界匹配，忽略大小写）----
    for w in _PERSON_NOUNS_EN:
        p = re.sub(r'\b' + re.escape(w) + r'\b', '', p, flags=re.IGNORECASE)

    # ---- 常见诗人/先贤名（大阅读古诗文高频人物词）----
    for w in _FAMOUS_NAMES_EN:
        p = re.sub(r'\b' + re.escape(w) + r'\b', '', p, flags=re.IGNORECASE)
    for w in _FAMOUS_NAMES_ZH:
        if w in p:
            p = p.replace(w, '')

    # ---- 额外角色名（如「李白」「Li Bai」）----
    if extra_names:
        names = extra_names if isinstance(extra_names, (list, tuple, set)) else [extra_names]
        for nm in names:
            nm = str(nm).strip()
            if not nm:
                continue
            if nm in p:
                p = p.replace(nm, '')
            # 英文首字母大写形式（Li Bai）
            cap = ' '.join(seg.capitalize() for seg in nm.split())
            if cap and cap != nm and cap in p:
                p = p.replace(cap, '')

    # ---- 清理残留标点/空格 ----
    p = re.sub(r'[,，、;；:：]+\s*', ', ', p)
    p = re.sub(r'\s{2,}', ' ', p)
    p = re.sub(r'^[\s,，、;；:：]+|[\s,，、;；:：]+$', '', p).strip()
    return p


def generate_image(prompt, index, world_setting=None, style_seed=None, is_character=False, characters=None, angle=None, composition=None, is_translated=False, layout_prompt=None, width=None, height=None, panel_spec=None, is_freeform=False, force_pure_background=False, seed=None, block_names=None, allow_default_fallback=True):
    """生成图片（优先使用豆包方舟API，其次使用Pollinations.AI生图API，未配置则使用默认图）
    
    world_setting: 世界设定字典，包含视觉风格、地点等信息，用于生成连贯背景
    style_seed: 风格种子字符串，确保同一故事的所有面板使用相同的视觉标识
    is_character: 是否为角色图片，为True时API未配置会生成角色头像而不是背景
    characters: 角色描述列表，用于将角色融入背景图中（当is_character=False且characters不为空时）
    angle: 视角描述（如"low angle", "close-up"等）
    composition: 构图描述（如"rule of thirds", "dynamic diagonal"等）
    is_translated: 是否已翻译过，True则跳过翻译步骤（用于批量生成优化）
    layout_prompt: 布局提示词，描述分镜的布局、视角、构图要求
    width: 图片宽度（像素），None则使用默认768
    height: 图片高度（像素），None则使用默认768
    panel_spec: 分镜规格字典，包含aspect_ratio、shot_type、angle等
    is_freeform: 是否为自由布局模式，True时会将border_prompt嵌入提示词
    """
    # 纯背景必须在翻译前完成中文场景清洗；翻译后中文动作模式无法再识别，
    # 会把“某人正在做什么”残留为英文动作提示，从而使背景偏离场景要求。
    pure_background = force_pure_background or (
        not is_character and not (characters and isinstance(characters, list) and len(characters) > 0)
    )
    if pure_background:
        raw_prompt = prompt
        raw_names = list(block_names) if isinstance(block_names, (list, tuple, set)) else ([block_names] if block_names else [])
        cleaned_prompt = extract_environment_only(raw_prompt, character_names=raw_names)
        if cleaned_prompt:
            prompt = cleaned_prompt
        logger.info(f"[纯背景预清洗] {raw_prompt[:100]} -> {prompt[:100]}")

    if not is_translated:
        try:
            translated = translate_to_english(prompt)
            if translated and len(translated) > 0:
                prompt = translated
                logger.info(f"提示词已翻译为英文")
            else:
                logger.warning("翻译返回空结果，使用原文")
        except Exception as e:
            logger.warning(f"翻译失败，使用原文: {str(e)}")
    
    style_elements = []
    if world_setting:
        visual_style = world_setting.get('visual_style', '')
        main_location = world_setting.get('main_location', '')
        atmosphere = world_setting.get('atmosphere', '')
        time_period = world_setting.get('time_period', '')
        
        if visual_style:
            style_elements.append(visual_style)
        if main_location:
            style_elements.append(f"same location: {main_location}")
        if atmosphere:
            style_elements.append(f"atmosphere: {atmosphere}")
        if time_period:
            style_elements.append(f"time period: {time_period}")
    
    if style_seed:
        style_elements.append(f"unique story id: {style_seed}")
    
    if angle:
        style_elements.append(f"camera angle: {angle}")
    
    # “人物特写”等分镜构图对纯背景没有意义，还会让模型偏离场景；
    # 纯背景只保留明确的环境广角构图。
    if composition and (not pure_background or str(composition).lower() in ('wide shot', 'establishing shot', 'eye-level')):
        style_elements.append(f"composition: {composition}")
    
    # layout_prompt 来自分镜说明，常包含“人物特写/动作/表情”等内容。
    # 纯背景不能把它重新拼进提示词，否则会抵消上面的角色清洗并诱导错误画面。
    if layout_prompt and not pure_background:
        try:
            layout_prompt_en = translate_to_english(layout_prompt)
            if layout_prompt_en and len(layout_prompt_en) > 0:
                style_elements.append(f"layout: {layout_prompt_en}")
            else:
                style_elements.append(f"layout: {layout_prompt}")
        except Exception as e:
            logger.warning(f"布局提示词翻译失败，使用原文: {str(e)}")
            style_elements.append(f"layout: {layout_prompt}")
    
    style_constraint = ", ".join(style_elements) if style_elements else ""
    
    has_characters = characters and isinstance(characters, list) and len(characters) > 0
    is_crowd = False  # 是否为人群场景（背影人群，避免免费引擎生成的脸怪异）

    # ⚠️ 纯背景（纯场景）模式：翻译后的 prompt 仍可能残留人物词（如「男孩」→「boy」、
    # 「李白」→「Li Bai」），且用户自定义角色名（许多/莉莉/脚本里的人名）未必在静态名单。
    # Flux 不理解否定词，只要出现人物 token 就可能生成人。
    # 这里做最终兜底——①先用 extract_environment_only 按中文动作结构+角色名剥离场景里的人；
    # ②再用 sanitize_no_people 剥离中英文人物词（含指定角色名），再拼正面空环境描述。
    # 即使背景提示词里明确写了人物，也要屏蔽其出现。
    if (not is_character and not has_characters) or force_pure_background:
        _bn = []
        if block_names:
            _bn += list(block_names) if isinstance(block_names, (list, tuple, set)) else [block_names]
        if characters and isinstance(characters, list):
            _bn += [str(c.get('name') if isinstance(c, dict) else c) for c in characters]
        _bn = [str(x).strip() for x in _bn if str(x).strip()]
        # ① 先按中文动作结构/角色名剥离（extract_environment_only 处理「X在…」「李白」等）
        _cleaned = extract_environment_only(prompt, character_names=_bn)
        if _cleaned:
            prompt = _cleaned
        # ② 再剥中英文人物词 + 角色名（覆盖翻译后的英文层）
        prompt = sanitize_no_people(prompt, extra_names=_bn)
        logger.info(f"[纯背景清洗] 最终背景 prompt: {prompt[:120]}")
    
    common_style = "comic panel style, bright vibrant colors, clean line art, manga panel layout, no text, no watermark, high quality, professional illustration"
    
    if is_character:
        # 精简但强有效的约束：防止多腿/多臂/多人问题
        # 将关键约束放在提示词末尾，避免被截断
        char_constraint = (
            "solo, exactly one person only, no extra limbs, no deformed anatomy, "
            "perfect proportions, clean outlines, professional character design"
        )
        base_prompt = (
            f"{prompt}, Q version cartoon character, full body, {common_style}, "
            f"solid white background, isolated character, cute adorable, {char_constraint}"
        )
    elif has_characters:
        character_desc = ", ".join([str(c) for c in characters])
        base_prompt = (
            f"{prompt}, {common_style}, consistent visual style, same color palette, "
            f"full scene illustration, characters naturally integrated into environment, "
            f"same lighting and shadows on characters and background, natural poses and expressions, "
            f"characters: {character_desc}, consistent character design, same appearance across scenes, "
            f"no transparent cutout feeling, cohesive single-layer illustration"
        )
        logger.info(f"[背景+角色模式] 已融入角色描述: {character_desc[:100]}")
    else:
        base_prompt = (
            f"{prompt}, {common_style}, consistent visual style, same color palette"
        )
        # 关键词分组，按场景意图决定背景是否含人
        # 同时支持英文和中文关键词，避免翻译失败时漏检
        platform_keywords = ['boat', 'ship', 'raft', 'pier', 'dock', 'deck', 'harbor', 'shore',
                             '船', '筏', '码头', '栈桥', '甲板', '港', '岸']
        water_keywords = ['water', 'river', 'lake', 'sea', 'ocean', 'pond', 'stream', 'wave',
                          '水', '河', '湖', '海', '池塘', '溪', '浪']
        crowd_keywords = ['crowd', 'people', 'street', 'market', 'playground', 'audience',
                          'classroom', 'square', 'festival', 'parade', 'group of people',
                          '人群', '街道', '集市', '市场', '操场', '教室', '广场', '会场',
                          '观众', '学生', '同学们', '放学', '热闹', '大街']
        # 只检查用户场景描述，避免 common_style 里的 "watermark" 误触发
        lower_prompt = prompt.lower()
        is_crowd = any(k in lower_prompt for k in crowd_keywords)
        is_water = (any(k in lower_prompt for k in platform_keywords)
                    or any(k in lower_prompt for k in water_keywords))
        if is_water:
            # 水面/船只类场景，强制前景生成 solid deck/pier，避免后续角色合成时站在水上
            base_prompt += (", empty solid wooden deck/pier/platform in the foreground, "
                            "vacant unoccupied platform, water behind and below the platform, "
                            "空旷的木质平台，纯环境背景")
    
    if style_constraint:
        base_prompt += f", {style_constraint}"
    
    if is_freeform and panel_spec and panel_spec.get('border_prompt'):
        border_prompt = panel_spec.get('border_prompt')
        base_prompt += f", frame border style: {border_prompt}"
        logger.info(f"[自由布局模式] 已追加边框样式: {border_prompt[:100]}...")
    
    if (not is_character and not has_characters) or force_pure_background:
        # ⚠️ BugFix：此前把「miniature diorama model / scale model / architectural maquette /
        # empty film set / still life scene / architectural and landscape only」等强风格词
        # 无条件追加到 prompt 末尾。Flux 类模型对末尾 token 权重更高，这批词会直接覆盖
        # 用户自己写的风格描述（如「3D卡通黏土渲染、圆润软质、马卡龙色系」），
        # 导致生成建筑沙盘 / 玻璃球 / 写实场景等与预期完全不符的结果。
        #
        # 新策略：
        #   1. 删除上述强风格污染词，只保留轻量的「空环境 / 无人 / 给角色留出站立地面」约束；
        #   2. 把约束词前置到 prompt 最前面，用户原始场景描述留在后面，
        #      使用户风格词获得更高权重，约束词仅起辅助作用。
        # 注：仍然刻意不使用 "no people" / "no human figures" 这类否定表达——
        # Flux 不理解否定词，出现 people/human token 反而会诱导生成人物；
        # 人物词的剥离由上游的 extract_environment_only + sanitize_no_people 负责。
        pure_background_constraint = (
            "empty environment, vacant unoccupied space, "
            "wide unobstructed walkable floor across the lower third, "
            "clear continuous ground plane reserved for standing characters, "
            "pure background scenery only"
        )
        base_prompt = f"{pure_background_constraint}, {base_prompt}"
        logger.info("[纯背景模式] 空环境约束已前置，用户场景描述保留在后（不再追加微缩模型/沙盘等强风格词）")
    
    max_prompt_length = 1200
    if len(base_prompt) > max_prompt_length:
        logger.warning(f"提示词过长({len(base_prompt)}字符)，截断至{max_prompt_length}字符")
        base_prompt = base_prompt[:max_prompt_length]
    
    logger.info(f"[最终Prompt] len={len(base_prompt)} | {base_prompt}")
    
    # 从panel_spec推导尺寸
    if panel_spec and (width is None or height is None):
        w, h = calculate_image_size_for_panel(panel_spec)
        if width is None:
            width = w
        if height is None:
            height = h
    
    if width is None:
        width = 768
    if height is None:
        height = 768
    
    logger.info(f"生图尺寸: {width}x{height}")
    
    # 已配置方舟时必须优先使用用户选择的付费模型。此前免费模型在这里先返回，
    # 导致即使配置了豆包，背景仍会由不稳定的免费服务生成并与提示词不匹配。
    image_config = get_request_image_provider()
    if image_config:
        provider = image_config.get('provider')
        if provider == 'doubao':
            result = generate_image_ark(base_prompt, index, is_character, width, height, image_config)
        elif provider == 'aliyun':
            result = generate_image_tongyi_text(
                base_prompt, seed=seed, width=width, height=height,
                api_key=image_config.get('api_key'), model=image_config.get('model'))
        else:
            result = generate_image_openai_compatible(
                base_prompt, index, width, height, image_config)
        if result:
            return result
        # 纯背景接口不能把生成失败伪装成“已按提示词生成”的免费结果。
        if not allow_default_fallback:
            logger.error("方舟图像生成失败；已禁止纯背景回退至免费模型")
            return None

    # 未配置方舟，或允许普通功能采用免费回退时，才使用 Pollinations。
    result = generate_image_pollinations(base_prompt, index, is_character, width, height, seed=seed)
    if result:
        return result
    
    # 优先级3: Pollinations.AI生图API（完全免费，与优先级1相同）
    result = generate_image_jimeng(base_prompt, index, style_seed, is_character, width, height)
    if result:
        return result
    
    # 兜底: 使用默认图片。纯背景接口可以禁用此行为，避免限流时把随机示意图
    # 伪装成“根据提示词生成”的背景返回给用户。
    if not allow_default_fallback:
        logger.error("图像服务均失败，已禁用默认占位图回退")
        return None
    if is_character:
        logger.info(f"所有图片API均失败，生成默认角色头像（第{index+1}个）")
        return get_default_character_image(index, prompt)
    else:
        logger.info(f"所有图片API均失败，使用默认背景图（第{index+1}张）")
        return get_default_background(index, prompt, world_setting)


def background_visual_style_matches_prompt(image_data, scene_prompt):
    """拦截与明确暖色卡通提示词明显冲突的冷灰写实背景。

    免费生图模型会偶发返回 HTTP 200 但内容偏题；这不是网络失败，不能靠状态码发现。
    这里只处理特征非常明确的情况，避免把用户主动要求的冷色场景误判掉。
    """
    text = str(scene_prompt or '').lower()
    warm_style_tokens = (
        '暖', '暖色', '马卡龙', '黏土', '粘土', '可爱', '治愈',
        'warm', 'pastel', 'macaron', 'clay', 'cute', 'cozy', 'cartoon', '卡通'
    )
    if not any(token in text for token in warm_style_tokens):
        return True

    try:
        from io import BytesIO
        image = Image.open(BytesIO(image_data)).convert('RGB').resize((64, 64))
        pixels = list(image.getdata())
        total = max(1, len(pixels))
        avg_r = sum(p[0] for p in pixels) / total
        avg_g = sum(p[1] for p in pixels) / total
        avg_b = sum(p[2] for p in pixels) / total
        # 简化饱和度：通道差相对亮度。冷灰工业照片通常同时满足“低饱和 + 偏蓝”。
        avg_spread = sum(max(p) - min(p) for p in pixels) / total
        avg_light = max(1.0, (avg_r + avg_g + avg_b) / 3)
        saturation_proxy = avg_spread / avg_light
        is_cold_gray = avg_b > avg_r + 7 and saturation_proxy < 0.28
        is_nearly_gray = saturation_proxy < 0.07
        if is_cold_gray or is_nearly_gray:
            logger.warning(
                '[背景质检] 拒绝与暖色卡通提示词冲突的结果: '
                f'RGB=({avg_r:.0f},{avg_g:.0f},{avg_b:.0f}), sat={saturation_proxy:.2f}'
            )
            return False
        return True
    except Exception as error:
        # 质检异常不能阻断正常生图，但保留日志便于追踪。
        logger.warning(f'[背景质检] 无法分析图片色彩，跳过质检: {error}')
        return True


def remove_background(image_data):
    """移除图片背景，返回带透明背景的PNG图片数据
    
    优先使用rembg库进行AI背景移除，失败时回退到颜色阈值方法。
    
    Args:
        image_data: 原始图片数据（bytes）
        
    Returns:
        移除背景后的图片数据（bytes），处理失败返回None
    """
    import time
    from io import BytesIO
    
    start_time = time.time()
    
    try:
        logger.info("开始背景移除处理...")
        
        import os
        os.environ['U2NET_HOME'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), '.u2net_cache')
        
        from rembg import remove
        
        result_data = remove(image_data, force_return_bytes=True)
        
        elapsed_time = time.time() - start_time
        logger.info(f"背景移除处理完成（rembg），耗时: {elapsed_time:.2f}秒")
        
        return result_data
        
    except ImportError:
        logger.warning("rembg库未安装，使用颜色阈值方法")
        
    except Exception as e:
        elapsed_time = time.time() - start_time
        logger.warning(f"rembg处理失败（耗时: {elapsed_time:.2f}秒）: {str(e)}，使用颜色阈值方法")
    
    try:
        from PIL import Image
        from collections import Counter
        
        img = Image.open(BytesIO(image_data)).convert('RGBA')
        w, h = img.size
        if w < 10 or h < 10:
            return image_data
        
        pixels = img.load()
        total_pixels = w * h
        transparent_pixels = 0
        for y in range(h):
            for x in range(w):
                if pixels[x, y][3] < 128:
                    transparent_pixels += 1
        
        # 如果已有大量透明像素，说明 rembg 可能已经生效，直接返回
        if transparent_pixels > total_pixels * 0.05:
            output_bytes = BytesIO()
            img.save(output_bytes, format='PNG')
            elapsed_time = time.time() - start_time
            logger.info(f"背景移除处理完成（保持已有透明），耗时: {elapsed_time:.2f}秒")
            return output_bytes.getvalue()
        
        # 从四角采样背景色（角色通常不会贴边到角落）
        corner_size = max(10, min(w, h) // 20)
        corners = [
            (0, 0, corner_size, corner_size),
            (w - corner_size, 0, w, corner_size),
            (0, h - corner_size, corner_size, h),
            (w - corner_size, h - corner_size, w, h)
        ]
        bg_colors = []
        for x0, y0, x1, y1 in corners:
            crop = img.crop((x0, y0, x1, y1))
            crop_pixels = crop.load()
            crop_w, crop_h = crop.size
            most_common = Counter([(crop_pixels[x, y][0], crop_pixels[x, y][1], crop_pixels[x, y][2])
                                    for y in range(crop_h) for x in range(crop_w)]).most_common(1)[0][0]
            bg_colors.append(most_common)
        bg_color = Counter(bg_colors).most_common(1)[0][0]
        
        tolerance = 35  # 容忍 JPEG 压缩/反锯齿导致的色差
        white_tolerance = 25  # 专门处理纯白背景
        
        for y in range(h):
            for x in range(w):
                r, g, b, a = pixels[x, y]
                # 情况1：接近四角背景色（处理米色/灰底等 AI 常见背景）
                near_bg = (abs(r - bg_color[0]) <= tolerance and
                           abs(g - bg_color[1]) <= tolerance and
                           abs(b - bg_color[2]) <= tolerance)
                # 情况2：接近纯白（处理白底）
                near_white = (r > 240 - white_tolerance and
                            g > 240 - white_tolerance and
                            b > 240 - white_tolerance)
                if near_bg or near_white:
                    pixels[x, y] = (r, g, b, 0)
        
        output_bytes = BytesIO()
        img.save(output_bytes, format='PNG')
        
        elapsed_time = time.time() - start_time
        logger.info(f"背景移除处理完成（四角背景色检测），耗时: {elapsed_time:.2f}秒")
        
        return output_bytes.getvalue()
        
    except Exception as e:
        elapsed_time = time.time() - start_time
        logger.error(f"背景移除处理失败（四角背景色检测，耗时: {elapsed_time:.2f}秒）: {str(e)}")
        return None


def generate_image_ark(prompt, index, is_character=False, width=512, height=512, config=None):
    """使用豆包方舟平台图片生成API"""
    max_retries = 3
    retry_delay = 2

    # Seedream 5.0 lite 要求单张图像至少 3,686,400 像素（例如 1920×1920）。
    # 前端所需的 768 尺寸仅用于漫画格展示，不能直接作为方舟生图尺寸；这里按原比例
    # 放大后由既有合成流程缩放回格子，既满足模型限制，也不会改变布局比例。
    ark_min_pixels = 3_686_400
    requested_width = max(1, int(width or 768))
    requested_height = max(1, int(height or 768))
    if requested_width * requested_height < ark_min_pixels:
        scale = math.sqrt(ark_min_pixels / (requested_width * requested_height))
        width = int(math.ceil(requested_width * scale / 8) * 8)
        height = int(math.ceil(requested_height * scale / 8) * 8)
        logger.info(
            f"Seedream 尺寸已按原比例放大: {requested_width}x{requested_height} -> {width}x{height}"
        )
    else:
        width, height = requested_width, requested_height
    
    # 方舟API支持的尺寸
    ark_sizes = {
        (512, 512): '512x512',
        (768, 768): '768x768',
        (1024, 1024): '1024x1024',
        (768, 512): '768x512',
        (512, 768): '512x768',
        (1024, 576): '1024x576',
        (576, 1024): '576x1024',
        (1024, 768): '1024x768',
        (768, 1024): '768x1024',
    }
    size_str = ark_sizes.get((width, height), f'{width}x{height}')
    
    for attempt in range(max_retries):
        try:
            logger.info(f"正在调用豆包方舟图片API生成第{index+1}张图片（第{attempt+1}次尝试）...")
            logger.info(f"提示词: {prompt[:200]}...")
            logger.info(f"图片尺寸: {size_str}")
            
            config = config or get_request_image_provider() or {}
            url = config.get('base_url') or "https://ark.cn-beijing.volces.com/api/v3/images/generations"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {config.get('api_key') or ARK_API_KEY}"
            }
            
            data = {
                "model": config.get('model') or ARK_IMAGE_MODEL,
                "prompt": prompt,
                "size": size_str,
                "n": 1,
                "response_format": "url",
                "watermark": False,
                "sequential_image_generation": "disabled"
            }
            
            response = requests.post(url, headers=headers, json=data, timeout=120)
            if response.status_code >= 400:
                # 仅记录服务端原因，绝不记录 Authorization/API Key，便于区分密钥、模型开通和参数问题。
                logger.error(
                    f"方舟图片请求被拒绝 HTTP {response.status_code}: {response.text[:500]}"
                )
            response.raise_for_status()
            
            result = response.json()
            
            if 'data' in result and len(result['data']) > 0:
                image_url = result['data'][0].get('url', '')
                if image_url:
                    image_response = requests.get(image_url, timeout=30)
                    image_response.raise_for_status()
                    logger.info(f"豆包方舟图片生成成功！")
                    return image_response.content
            
            logger.error("豆包方舟API返回格式不正确")
            if attempt < max_retries - 1:
                logger.info(f"等待{retry_delay}秒后重试...")
                time.sleep(retry_delay)
                retry_delay *= 2
                continue
                
        except requests.exceptions.RequestException as e:
            logger.error(f"调用豆包方舟API失败（第{attempt+1}次）: {str(e)}")
            if attempt < max_retries - 1:
                logger.info(f"等待{retry_delay}秒后重试...")
                time.sleep(retry_delay)
                retry_delay *= 2
                continue
        except Exception as e:
            logger.error(f"生成图片时发生错误（第{attempt+1}次）: {str(e)}")
            if attempt < max_retries - 1:
                logger.info(f"等待{retry_delay}秒后重试...")
                time.sleep(retry_delay)
                retry_delay *= 2
                continue
    
    logger.warning("豆包方舟API调用失败，尝试Pollinations.AI生图API")
    return None


def _download_or_decode_generated_image(payload, require_public_url=False):
    """Read the common OpenAI-compatible image response (URL or base64)."""
    import base64
    items = payload.get('data') if isinstance(payload, dict) else None
    if not items or not isinstance(items, list):
        return None
    item = items[0] or {}
    encoded = item.get('b64_json') or item.get('base64')
    if encoded:
        return base64.b64decode(encoded)
    image_url = item.get('url')
    if image_url:
        if require_public_url:
            image_url = _validate_public_https_url(image_url)
        response = requests.get(image_url, timeout=60)
        response.raise_for_status()
        return response.content
    return None


def generate_image_openai_compatible(prompt, index, width, height, config):
    """OpenAI Images API and compatible /images/generations endpoints."""
    provider = config.get('provider') or 'custom'
    url = config.get('base_url') or 'https://api.openai.com/v1/images/generations'
    model = config.get('model') or ('gpt-image-1' if provider == 'openai' else '')
    # OpenAI image models accept a documented set of square/landscape/portrait sizes.
    ratio = float(width or 1) / max(1.0, float(height or 1))
    if provider == 'openai':
        size = '1536x1024' if ratio > 1.12 else ('1024x1536' if ratio < 0.89 else '1024x1024')
    else:
        size = f'{int(width)}x{int(height)}'
    # Do not force response_format: modern compatible services may return either URL or b64_json.
    # The parser below accepts both, which is more portable across gateways.
    body = {'model': model, 'prompt': prompt, 'n': 1, 'size': size}
    try:
        if provider == 'custom':
            # Revalidate on every call to reduce DNS-rebinding/changed-domain risk.
            url = _validate_custom_image_url(url)
        response = requests.post(
            url,
            headers={'Authorization': f"Bearer {config.get('api_key')}",
                     'Content-Type': 'application/json'},
            json=body, timeout=180,
        )
        if response.status_code >= 400:
            logger.error('[image-provider:%s] HTTP %s: %s', provider,
                         response.status_code, response.text[:500])
        response.raise_for_status()
        result = _download_or_decode_generated_image(
            response.json(), require_public_url=(provider == 'custom'))
        if result:
            return result
        logger.error('[image-provider:%s] 返回中没有可用图片', provider)
    except Exception as exc:
        logger.error('[image-provider:%s] 生图失败: %s', provider, exc)
    return None


def generate_image_jimeng(prompt, index, style_seed=None, is_character=False, width=768, height=768):
    """使用Pollinations.AI生图API生成图片"""
    return generate_image_pollinations(prompt, index, is_character, width, height)


def _poll_tongyi_task(task_id, api_key, max_wait=240):
    """轮询通义万相异步任务，成功后立即下载图片（URL 有效期极短），返回图片 bytes 或 None"""
    import requests, time
    url = f'https://dashscope.aliyuncs.com/api/v1/tasks/{task_id}'
    headers = {'Authorization': f'Bearer {api_key}'}
    start = time.time()
    while time.time() - start < max_wait:
        try:
            r = requests.get(url, headers=headers, timeout=30)
            if r.status_code != 200:
                time.sleep(3)
                continue
            data = r.json()
            status = data.get('output', {}).get('task_status')
            if status == 'SUCCEEDED':
                results = data.get('output', {}).get('results', [])
                if results and results[0].get('url'):
                    img_url = results[0]['url']
                    img_r = requests.get(img_url, timeout=30)
                    if img_r.status_code == 200:
                        return img_r.content
                logger.error('万相任务成功但无图片URL')
                return None
            elif status == 'FAILED':
                logger.error(f'万相任务失败: {data.get("message")} | {data.get("output", {}).get("task_status")}')
                return None
        except Exception as e:
            logger.warning(f'万相轮询异常: {e}')
        time.sleep(4)
    logger.error('万相任务轮询超时')
    return None


def generate_image_tongyi_text(prompt, seed=None, width=1024, height=1024,
                               api_key=None, model='wanx2.1-t2i-turbo'):
    """通义万相文生图（wanx2.1-t2i-turbo），用于生成标准角色。返回图片 bytes 或 None。
    注意：通义万相对于中文 prompt 效果最好，调用方应直接传中文、不要翻译成英文。
    """
    import requests, base64, time
    api_key = api_key or os.environ.get('DASHSCOPE_API_KEY')
    if not api_key:
        return None
    try:
        ratio = float(width or 1) / max(1.0, float(height or 1))
        if ratio > 1.18:
            width, height = 1280, 720
        elif ratio < 0.85:
            width, height = 720, 1280
        else:
            width, height = 1024, 1024
        url = 'https://dashscope.aliyuncs.com/api/v1/services/aigc/text2image/image-synthesis'
        headers = {
            'Authorization': f'Bearer {api_key}',
            'X-DashScope-Async': 'enable',
            'Content-Type': 'application/json'
        }
        params = {'n': 1, 'size': f'{width}*{height}'}
        if seed is not None:
            params['seed'] = int(seed)
        body = {'model': model or 'wanx2.1-t2i-turbo', 'input': {'prompt': prompt}, 'parameters': params}
        r = requests.post(url, headers=headers, json=body, timeout=40)
        if r.status_code != 200:
            logger.error(f'万相文生图提交失败: HTTP {r.status_code} {r.text[:200]}')
            return None
        task_id = r.json().get('output', {}).get('task_id')
        if not task_id:
            logger.error(f'万相文生图无 task_id: {r.text[:200]}')
            return None
        return _poll_tongyi_task(task_id, api_key)
    except Exception as e:
        logger.error(f'万相文生图异常: {e}')
        return None


def generate_image_tongyi_img2img(init_image_bytes, prompt, seed=None, function='control_cartoon_feature'):
    """通义万相图生图（wanx2.1-imageedit），以 init_image 角色图为参考生成同形象不同姿势。
    function=control_cartoon_feature（参考卡通形象生图）最适合漫画角色换姿势。返回图片 bytes 或 None。
    """
    import requests, base64
    api_key = os.environ.get('DASHSCOPE_API_KEY')
    if not api_key:
        return None
    try:
        b64 = base64.b64encode(init_image_bytes).decode('ascii')
        base_image_url = f'data:image/png;base64,{b64}'
        url = 'https://dashscope.aliyuncs.com/api/v1/services/aigc/image2image/image-synthesis'
        headers = {
            'Authorization': f'Bearer {api_key}',
            'X-DashScope-Async': 'enable',
            'Content-Type': 'application/json'
        }
        input_obj = {'function': function, 'prompt': prompt, 'base_image_url': base_image_url}
        params = {'n': 1}
        if seed is not None:
            params['seed'] = int(seed)
        body = {'model': 'wanx2.1-imageedit', 'input': input_obj, 'parameters': params}
        r = requests.post(url, headers=headers, json=body, timeout=40)
        if r.status_code != 200:
            logger.error(f'万相图生图提交失败: HTTP {r.status_code} {r.text[:200]}')
            return None
        task_id = r.json().get('output', {}).get('task_id')
        if not task_id:
            logger.error(f'万相图生图无 task_id: {r.text[:200]}')
            return None
        return _poll_tongyi_task(task_id, api_key)
    except Exception as e:
        logger.error(f'万相图生图异常: {e}')
        return None


def _poll_horde_task(req_id, max_wait=300):
    """轮询 AI Horde 异步任务，返回生成图片的 bytes 或 None。
    匿名 key 优先级最低，可能排队较久，故 max_wait 给到 5 分钟。
    """
    import requests, time, base64
    check_url = f'https://aihorde.net/api/v2/generate/check/{req_id}'
    status_url = f'https://aihorde.net/api/v2/generate/status/{req_id}'
    headers = {'apikey': os.environ.get('HORDE_API_KEY', '0000000000'), 'Client-Agent': 'dydj_manhua:1.0:local'}
    start = time.time()
    while time.time() - start < max_wait:
        try:
            chk = requests.get(check_url, headers=headers, timeout=30)
            if chk.status_code != 200:
                time.sleep(4)
                continue
            cd = chk.json()
            # 没有 worker 能接这个任务，立即放弃，避免空等 5 分钟导致前端超时
            if cd.get('is_possible') is False:
                logger.error('horde 当前无可用 worker，无法完成该请求')
                return None
            if cd.get('done') or cd.get('finished'):
                break
        except Exception as e:
            logger.warning(f'horde 轮询异常: {e}')
        time.sleep(4)
    else:
        logger.error('horde 任务轮询超时')
        return None
    try:
        st = requests.get(status_url, headers=headers, timeout=30)
        if st.status_code != 200:
            logger.error(f'horde 状态获取失败: HTTP {st.status_code}')
            return None
        sd = st.json()
        if sd.get('faulted'):
            logger.error('horde 任务 faulted（生成失败）')
            return None
        gens = sd.get('generations', [])
        if not gens:
            logger.error('horde 无生成结果')
            return None
        g = gens[0]
        if g.get('url'):
            img_r = requests.get(g['url'], timeout=30)
            if img_r.status_code == 200:
                return img_r.content
        if g.get('img'):
            return base64.b64decode(g['img'])
        logger.error('horde 结果既无 url 也无 img')
        return None
    except Exception as e:
        logger.error(f'horde 结果解析异常: {e}')
        return None


def generate_image_horde(text_prompt, init_image_bytes=None, denoising_strength=0.5,
                         seed=None, width=576, height=576, models=None):
    """AI Horde 免费生图（匿名 key 0000000000，无需注册、无需 API Key）。

    - init_image_bytes 为 None → 文生图（text2img）
    - init_image_bytes 不为 None → 图生图（img2img），以该图为底图换姿势，保脸

    SD 系列模型对英文 prompt 远强于中文，故内部自动把 prompt 翻成英文。
    返回图片 bytes 或 None（失败时上层会降级 Pollinations）。
    """
    import requests, time, base64, io
    from PIL import Image
    try:
        # SD 模型偏好英文，自动翻译（translate_to_english 在外部已定义）
        en_prompt = translate_to_english(text_prompt) if text_prompt else text_prompt
        if not en_prompt:
            en_prompt = text_prompt
        # AI Horde 匿名/低优先级用户限制：尺寸不能超过 576x576，steps 不能过高
        # 否则直接返回 403 "Due to heavy demand"
        width = min(int(width), 576)
        height = min(int(height), 576)
        params = {
            'width': width,
            'height': height,
            'seed': str(int(seed)) if seed is not None else str(random.randint(1, 2147483647)),
            'steps': 25,
            'cfg_scale': 7.5,
            'sampler_name': 'k_euler',
            'karras': True,
            'n': 1,
            'denoising_strength': float(denoising_strength),
        }
        payload = {
            'prompt': en_prompt,
            'params': params,
            'nsfw': False,
            'censor_nsfw': True,
            'trusted_workers': False,
            'slow_workers': True,
            'shared': False,
            'replacement_filter': True,
            'r2': True,
            'models': models or ['stable_diffusion_xl'],
            'Client-Agent': 'dydj_manhua:1.0:local',
        }
        if init_image_bytes:
            try:
                src = Image.open(io.BytesIO(init_image_bytes)).convert('RGBA')
                buf = io.BytesIO()
                src.save(buf, format='WebP', quality=95)
                payload['source_image'] = base64.b64encode(buf.getvalue()).decode('ascii')
            except Exception as e:
                logger.warning(f'horde 底图转 WebP 失败: {e}')
        headers = {'apikey': os.environ.get('HORDE_API_KEY', '0000000000'), 'Client-Agent': 'dydj_manhua:1.0:local',
                   'Content-Type': 'application/json'}
        r = requests.post('https://aihorde.net/api/v2/generate/async',
                          headers=headers, json=payload, timeout=40)
        if r.status_code != 200:
            logger.error(f'horde 提交失败: HTTP {r.status_code} {r.text[:200]}')
            return None
        res = r.json()
        req_id = res.get('id')
        # 如果服务端立即提示无 worker，直接放弃，让上层降级 Pollinations
        warnings = res.get('warnings', [])
        if any(w.get('code') == 'NoAvailableWorker' for w in warnings):
            logger.error('horde 提交后提示无可用 worker，将降级 Pollinations')
            return None
        if not req_id:
            logger.error(f'horde 无 req_id: {r.text[:200]}')
            return None
        return _poll_horde_task(req_id)
    except Exception as e:
        logger.error(f'horde 生图异常: {e}')
        return None


def generate_image_pollinations(prompt, index=0, is_character=False, width=768, height=768, seed=None):
    """使用Pollinations.AI免费API生成图片（无需API Key，完全免费）
    
    Pollinations.AI: https://image.pollinations.ai/prompt/{prompt}
    - 完全免费，不需要注册
    - 不需要API Key
    - 支持中文prompt（会自动翻译）
    - 支持模型选择: flux（高质量）, turbo（快速）
    - 限制: 每秒1次调用，图片质量略低于付费API
    """
    import urllib.parse
    
    max_retries = 3
    retry_delay = 3
    rate_limited = False  # 标记是否遭遇限流(429)，用于提前跳出重试循环改走兜底

    for attempt in range(max_retries):
        try:
            logger.info(f"正在调用Pollinations.AI生成第{index+1}张图片（第{attempt+1}次尝试）...")
            logger.info(f"提示词: {prompt[:200]}...")
            logger.info(f"图片尺寸: {width}x{height}")
            
            encoded_prompt = urllib.parse.quote(prompt)
            
            if seed is not None:
                random_seed = int(seed)
            else:
                random_seed = random.randint(100000, 999999)
            
            # 零密钥可用模型里 zimage 实测会把“3D 黏土卡通”生成成写实照片；
            # 默认保持 flux，且允许通过环境变量配置已授权的高质量模型。
            image_model = os.environ.get('POLLINATIONS_IMAGE_MODEL', 'flux').strip() or 'flux'
            url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width={width}&height={height}&nologo=true&seed={random_seed}&model={urllib.parse.quote(image_model)}&nofeed=true"
            logger.info(f"[Pollinations] 使用图像模型: {image_model}")
            
            logger.info(f"[Pollinations] 请求URL长度: {len(url)}")
            
            response = requests.get(url, timeout=120)
            response.raise_for_status()
            
            content_type = response.headers.get('Content-Type', '')
            content_length = len(response.content)
            
            if 'image' in content_type and content_length > 1000:
                logger.info(f"Pollinations.AI图片生成成功！Content-Type: {content_type}, 大小: {content_length}字节")
                return response.content
            else:
                logger.warning(f"Pollinations返回的不是有效图片: Content-Type={content_type}, 长度={content_length}字节")
                
                if content_length < 1000:
                    logger.warning(f"[Pollinations] 响应内容过小，可能是限流或错误响应，内容: {response.content[:200]}")
                
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)
                    retry_delay = min(retry_delay * 2, 15)
                    continue
                return None
                
        except requests.exceptions.Timeout as e:
            logger.error(f"Pollinations.AI请求超时（第{attempt+1}次）: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
                retry_delay = min(retry_delay * 2, 15)
                continue
        except requests.exceptions.RequestException as e:
            status_code = getattr(e.response, 'status_code', 'unknown')
            logger.error(f"Pollinations.AI请求异常（第{attempt+1}次）: {str(e)}, 状态码: {status_code}")
            
            if status_code == 429:
                rate_limited = True
                if attempt >= 1:
                    logger.warning(f"[Pollinations] 连续限流(429)，提前跳出重试循环")
                    break
                wait = 3
                logger.warning(f"[Pollinations] 限流(429)，等待{wait}秒后重试（第{attempt+1}次）")
                time.sleep(wait)
                retry_delay = min(retry_delay * 2, 15)
                continue
            elif attempt < max_retries - 1:
                time.sleep(retry_delay)
                retry_delay = min(retry_delay * 2, 15)
                continue
        except Exception as e:
            logger.error(f"Pollinations.AI生成失败（第{attempt+1}次）: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
                retry_delay = min(retry_delay * 2, 15)
                continue
    
    if rate_limited:
        logger.warning("[fallback] Pollinations 限流，AI Horde 也不可用，跳过")
    else:
        logger.warning("[fallback] Pollinations 失败，尝试降级 AI Horde 匿名通道（免Key）")
        try:
            horde_bytes = generate_image_horde(prompt, width=width, height=height, seed=seed)
        except Exception as _e:
            logger.error(f"[fallback] AI Horde 兜底调用异常: {_e}")
            horde_bytes = None
        if horde_bytes:
            logger.info("[fallback] AI Horde 兜底成功")
            return horde_bytes
    logger.error("Pollinations.AI 与 AI Horde 兜底均失败，返回 None 使用默认背景")
    return None


def generate_image_craiyon(prompt, index=0, width=512, height=512):
    """使用Craiyon免费API生成图片（完全免费，无需注册，无API Key）
    
    Craiyon (原名DALL·E Mini): https://api.craiyon.com/v3
    - 完全免费，不需要注册
    - 不需要API Key
    - 支持中文prompt
    - 限制: 生成速度较慢，图片质量略低，但稳定性好
    """
    import json
    
    max_retries = 1
    retry_delay = 5
    
    for attempt in range(max_retries):
        try:
            logger.info(f"正在调用Craiyon生成第{index+1}张图片（第{attempt+1}次尝试）...")
            logger.info(f"提示词: {prompt[:100]}...")
            
            payload = {
                "prompt": prompt,
                "model": "none",
                "negative_prompt": "",
                "width": width,
                "height": height
            }
            
            response = requests.post(
                "https://api.craiyon.com/v3",
                json=payload,
                timeout=30,
                headers={"Content-Type": "application/json"}
            )
            response.raise_for_status()
            
            data = response.json()
            
            if "images" in data and len(data["images"]) > 0:
                image_data = data["images"][0]
                if image_data:
                    import base64
                    image_bytes = base64.b64decode(image_data)
                    logger.info(f"Craiyon图片生成成功！大小: {len(image_bytes)}字节")
                    return image_bytes
            
            logger.warning(f"Craiyon返回的不是有效图片: {data}")
            break
            
        except requests.exceptions.Timeout as e:
            logger.error(f"Craiyon请求超时（第{attempt+1}次）: {str(e)}")
            break
        except requests.exceptions.RequestException as e:
            status_code = getattr(e.response, 'status_code', 'unknown')
            logger.error(f"Craiyon请求异常（第{attempt+1}次）: {str(e)}，状态码: {status_code}")
            break
        except Exception as e:
            logger.error(f"Craiyon处理异常（第{attempt+1}次）: {str(e)}")
            break
    
    logger.warning("Craiyon生图失败")
    return None


def _reencode_reference_image(init_image_bytes, max_side=256, quality=70):
    """把参考图重新编码为小尺寸、去元数据的 JPEG，避免 data URI 过长触发 414/431。

    Pollinations 图生图端点把参考图放在 URL 的 image= 参数里，整条 URL（请求行）
    必须 < ~8KB（nginx large_client_header_buffers 限制），否则返回
    414/431。因此自适应缩小：base64 后若仍 > 阈值则把边长减半重试，
    直到足够小或降到 floor。

    返回重新编码后的 bytes；若失败则原样返回（让上层继续尝试）。
    """
    try:
        import base64 as _b64
        from io import BytesIO as _BytesIO
        from PIL import Image as _Image
        _img = _Image.open(_BytesIO(init_image_bytes))
        # 把透明区域填充为白色，再统一为 RGB，避免参考图残留米色/杂色背景
        if _img.mode in ('RGBA', 'P', 'LA'):
            if _img.mode == 'RGBA':
                _bg = _Image.new('RGB', _img.size, (255, 255, 255))
                _bg.paste(_img, mask=_img.split()[3])
                _img = _bg
            else:
                _img = _img.convert('RGB')
        _w, _h = _img.size
        _side = max_side
        _floor = 96
        _threshold = 6000  # base64 字节上限（URL 约 6.4KB，稳在 8KB 限制内）
        while True:
            if max(_w, _h) > _side:
                _scale = _side / max(_w, _h)
                _tmp = _img.resize((int(_w * _scale), int(_h * _scale)), _Image.BICUBIC)
            else:
                _tmp = _img
            _buf = _BytesIO()
            _tmp.save(_buf, format='JPEG', quality=quality, optimize=True)
            _out = _buf.getvalue()
            if len(_b64.b64encode(_out)) <= _threshold or _side <= _floor:
                break
            _side = max(_floor, _side // 2)
        logger.info(f"[Pollinations-img2img] 参考图重编码: {len(init_image_bytes)}→{len(_out)} 字节 (边长≤{_side})")
        return _out
    except Exception as _e:
        logger.warning(f"[Pollinations-img2img] 参考图重编码失败，使用原图: {_e}")
        return init_image_bytes


def generate_image_pollinations_img2img(init_image_bytes, prompt, seed=None, width=768, height=768,
                                        model='flux'):
    """使用 Pollinations 图生图端点做 img2img（免 Key、可保脸换姿势）。

    端点：https://image.pollinations.ai/image/{prompt}?image={参考图URL或data URI}
    - 图像端点本身免 Key（仅文本框强制 key），比 AI Horde 匿名模式稳（自家 GPU 队列）。
    - 参考图（标准角色）以 base64 data URI 内联传入，无需公网可访问的 URL。
    - 配合固定 seed + 英文"保持脸/服装/发型一致、只换姿势"提示词，实现换姿势保脸。
    返回图片 bytes 或 None（失败时上层降级普通文生图）。
    """
    import urllib.parse
    import base64
    from io import BytesIO as _BytesIO

    if not init_image_bytes:
        logger.warning("[Pollinations-img2img] 无参考图，返回 None")
        return None

    # 重新编码参考图：去掉元数据 + 缩小边长，避免 base64 data URI 过长触发 414 URI Too Long。
    # 标准角色图（来自 Pollinations t2i）常内嵌大段 JSON 元数据，直接编码会让 URL 超过 100KB。
    init_image_bytes = _reencode_reference_image(init_image_bytes, max_side=448)

    max_retries = 3
    retry_delay = 4

    for attempt in range(max_retries):
        try:
            if seed is not None:
                random_seed = int(seed)
            else:
                random_seed = random.randint(100000, 999999)

            data_uri = 'data:image/jpeg;base64,' + base64.b64encode(init_image_bytes).decode('ascii')
            encoded_prompt = urllib.parse.quote(prompt)
            # 参考图作为 image 参数；safe='' 允许 data URI 中的 : / = 等字符
            encoded_image = urllib.parse.quote(data_uri, safe='')
            url = (f"https://image.pollinations.ai/image/{encoded_prompt}"
                   f"?image={encoded_image}&width={width}&height={height}"
                   f"&seed={random_seed}&nologo=true&model={model}")

            logger.info(f"[Pollinations-img2img] 第{attempt+1}次尝试，提示词: {prompt[:120]}")
            response = requests.get(url, timeout=120)
            response.raise_for_status()

            content_type = response.headers.get('Content-Type', '')
            if 'image' in content_type and len(response.content) > 1000:
                logger.info(f"[Pollinations-img2img] 成功，{len(response.content)} 字节")
                return response.content
            else:
                logger.warning(f"[Pollinations-img2img] 返回非图片: {content_type}, {len(response.content)}字节")
                if attempt < max_retries - 1:
                    time.sleep(retry_delay)
                    retry_delay = min(retry_delay * 2, 15)
                    continue
                return None
        except requests.exceptions.RequestException as e:
            status_code = getattr(e.response, 'status_code', 'unknown')
            logger.error(f"[Pollinations-img2img] 异常(第{attempt+1}次): {e}, 状态码 {status_code}")
            if status_code == 429:
                wait = min(retry_delay * (attempt + 1) * 2, 20)
                logger.warning(f"[Pollinations-img2img] 限流，等待{wait}s")
                time.sleep(wait)
                retry_delay = min(retry_delay * 2, 15)
                continue
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
                retry_delay = min(retry_delay * 2, 15)
                continue
        except Exception as e:
            logger.error(f"[Pollinations-img2img] 失败(第{attempt+1}次): {e}")
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
                retry_delay = min(retry_delay * 2, 15)
                continue
    logger.warning("[Pollinations-img2img] 全部重试失败，返回 None")
    return None


def get_default_background(index=0, scene_description=None, world_setting=None):
    """生成默认背景图（当API调用失败时使用），根据场景描述和世界设定智能匹配背景主题"""
    from io import BytesIO
    
    theme_keywords = {
        'classroom': ['classroom', 'school', 'desk', 'chalkboard', 'blackboard', 'student', 'teacher', 'lesson'],
        'nature': ['nature', 'outdoor', 'field', 'meadow', 'grass', 'hill'],
        'space': ['space', 'star', 'moon', 'planet', 'universe', 'galaxy', 'cosmos'],
        'ocean': ['ocean', 'sea', 'water', 'beach', 'wave', 'fish', 'boat'],
        'forest': ['forest', 'tree', 'wood', 'jungle', 'bush', 'leaf'],
        'magic': ['magic', 'wizard', 'spell', 'castle', 'enchanted', 'fantasy', 'fairy'],
        'city': ['city', 'building', 'street', 'house', 'town', 'urban'],
        'garden': ['garden', 'flower', 'plant', 'floral', 'park'],
        'library': ['library', 'book', 'reading', 'study', 'knowledge'],
        'bedroom': ['bedroom', 'room', 'home', 'house', 'living'],
        'adventure': ['adventure', 'cave', 'treasure', 'explore', 'mountain'],
        'laboratory': ['lab', 'laboratory', 'science', 'experiment', 'test'],
        'ancient': ['ancient', 'temple', 'pyramid', 'cave', 'history', 'primitive', 'tribal'],
        'future': ['future', 'robot', 'cyber', 'neon', 'hologram'],
        'snow': ['snow', 'winter', 'ice', 'cold', 'frozen'],
    }
    
    themes = {
        'classroom': {'bg_start': (255, 245, 230), 'bg_end': (255, 220, 200), 'elements': 'books'},
        'nature': {'bg_start': (200, 240, 200), 'bg_end': (150, 200, 150), 'elements': 'trees'},
        'space': {'bg_start': (100, 100, 150), 'bg_end': (50, 50, 100), 'elements': 'stars'},
        'ocean': {'bg_start': (150, 200, 255), 'bg_end': (50, 100, 200), 'elements': 'waves'},
        'forest': {'bg_start': (180, 220, 180), 'bg_end': (100, 160, 100), 'elements': 'trees'},
        'magic': {'bg_start': (220, 180, 255), 'bg_end': (150, 100, 200), 'elements': 'sparkles'},
        'city': {'bg_start': (220, 220, 220), 'bg_end': (150, 150, 150), 'elements': 'buildings'},
        'garden': {'bg_start': (200, 240, 200), 'bg_end': (150, 200, 150), 'elements': 'flowers'},
        'library': {'bg_start': (245, 240, 235), 'bg_end': (200, 190, 180), 'elements': 'books'},
        'bedroom': {'bg_start': (255, 248, 252), 'bg_end': (220, 200, 210), 'elements': 'stars'},
        'adventure': {'bg_start': (160, 180, 140), 'bg_end': (80, 100, 60), 'elements': 'trees'},
        'laboratory': {'bg_start': (200, 230, 240), 'bg_end': (100, 150, 180), 'elements': 'sparkles'},
        'ancient': {'bg_start': (180, 150, 100), 'bg_end': (100, 70, 40), 'elements': 'trees'},
        'future': {'bg_start': (50, 50, 80), 'bg_end': (20, 20, 40), 'elements': 'sparkles'},
        'snow': {'bg_start': (200, 220, 255), 'bg_end': (150, 180, 220), 'elements': 'stars'},
    }
    
    matched_theme = None
    
    if world_setting:
        visual_style = world_setting.get('visual_style', '')
        main_location = world_setting.get('main_location', '')
        
        combined_desc = f"{visual_style} {main_location}"
        desc_lower = combined_desc.lower()
        for theme_name, keywords in theme_keywords.items():
            for keyword in keywords:
                if keyword.lower() in desc_lower:
                    matched_theme = theme_name
                    logger.info(f"世界设定匹配到主题: {theme_name}")
                    break
            if matched_theme:
                break
    
    if matched_theme is None and scene_description and isinstance(scene_description, str):
        desc_lower = scene_description.lower()
        for theme_name, keywords in theme_keywords.items():
            for keyword in keywords:
                if keyword.lower() in desc_lower:
                    matched_theme = theme_name
                    logger.info(f"场景描述 '{scene_description[:50]}...' 匹配到主题: {theme_name}")
                    break
            if matched_theme:
                break
    
    if matched_theme is None or matched_theme not in themes:
        theme_list = list(themes.values())
        theme = theme_list[index % len(theme_list)]
    else:
        theme = themes[matched_theme]
    img = Image.new('RGB', (512, 512), color=theme['bg_start'])
    draw = ImageDraw.Draw(img)
    
    start_r, start_g, start_b = theme['bg_start']
    end_r, end_g, end_b = theme['bg_end']
    
    panel_offset = index * 10
    start_r = min(255, max(0, start_r + panel_offset))
    start_g = min(255, max(0, start_g + panel_offset))
    start_b = min(255, max(0, start_b + panel_offset))
    
    for y in range(512):
        r = start_r + int((end_r - start_r) * y / 512)
        g = start_g + int((end_g - start_g) * y / 512)
        b = start_b + int((end_b - start_b) * y / 512)
        draw.line([(0, y), (512, y)], fill=(r, g, b))
    
    elements = theme['elements']
    seed_value = hash(f"{scene_description}_{index}_{world_setting}") % 1000000
    random.seed(seed_value)
    
    if elements == 'books':
        book_color = (255, 179, 198)
        for i in range(6):
            bx = random.randint(20, 480)
            by = random.randint(250, 480)
            bw = random.randint(30, 50)
            bh = random.randint(20, 35)
            draw.rectangle([bx, by, bx + bw, by + bh], fill=book_color)
            draw.line([bx + bw * 0.6, by, bx + bw * 0.6, by + bh], fill=(200, 100, 120), width=2)
        for i in range(3):
            cx = random.randint(50, 450)
            cy = random.randint(20, 150)
            for j in range(5):
                cloud_x = cx + (j - 2) * 30
                cloud_y = cy + random.randint(-10, 10)
                draw.ellipse([cloud_x, cloud_y, cloud_x + 40, cloud_y + 30], fill=(255, 255, 255))
    
    elif elements == 'trees':
        for i in range(4):
            tx = random.randint(50, 450)
            ty = random.randint(200, 480)
            draw.rectangle([tx - 5, ty, tx + 5, ty + 60], fill=(139, 69, 19))
            draw.ellipse([tx - 30, ty - 40, tx + 30, ty], fill=(34, 139, 34))
            draw.ellipse([tx - 25, ty - 50, tx + 25, ty - 10], fill=(50, 205, 50))
    
    elif elements == 'stars':
        star_color = (255, 255, 255)
        for i in range(30):
            sx = random.randint(20, 490)
            sy = random.randint(30, 480)
            size = random.randint(3, 8)
            draw.regular_polygon((sx, sy, size), n_sides=5, rotation=0, fill=star_color)
        moon_x = random.randint(350, 450)
        moon_y = random.randint(50, 120)
        draw.ellipse([moon_x, moon_y, moon_x + 60, moon_y + 60], fill=(255, 255, 200))
    
    elif elements == 'waves':
        for i in range(5):
            y = 300 + i * 40
            for x in range(0, 512, 20):
                draw.arc([x, y - 10, x + 20, y + 10], 0, 3.14, fill=(100, 180, 255), width=2)
        for i in range(3):
            bx = random.randint(50, 450)
            by = random.randint(20, 100)
            for j in range(5):
                cloud_x = bx + (j - 2) * 30
                cloud_y = by + random.randint(-10, 10)
                draw.ellipse([cloud_x, cloud_y, cloud_x + 40, cloud_y + 30], fill=(255, 255, 255))
    
    elif elements == 'sparkles':
        sparkle_color = (255, 255, 200)
        for i in range(25):
            sx = random.randint(20, 490)
            sy = random.randint(30, 480)
            size = random.randint(4, 10)
            draw.regular_polygon((sx, sy, size), n_sides=5, rotation=0, fill=sparkle_color)
        for i in range(8):
            mx = random.randint(20, 490)
            my = random.randint(30, 480)
            draw.polygon([(mx, my-8), (mx+4, my+4), (mx-8, my)], fill=(255, 200, 255))
    
    elif elements == 'buildings':
        for i in range(5):
            bx = random.randint(30, 480)
            by = random.randint(250, 480)
            bw = random.randint(40, 80)
            bh = random.randint(80, 200)
            draw.rectangle([bx, by - bh, bx + bw, by], fill=(200, 200, 200))
            for j in range(4):
                wx = bx + 10 + j * 15
                wy = by - 20 - j * 30
                if wy > by - bh:
                    draw.rectangle([wx, wy - 10, wx + 8, wy], fill=(255, 200, 100))
    
    elif elements == 'flowers':
        for i in range(8):
            fx = random.randint(30, 480)
            fy = random.randint(250, 480)
            draw.line([fx, fy, fx, fy - 20], fill=(34, 139, 34), width=2)
            flower_colors = [(255, 179, 198), (255, 200, 150), (200, 150, 255), (255, 150, 200)]
            for j in range(5):
                px = fx + int(15 * math.cos(j * 0.8))
                py = fy - 20 + int(15 * math.sin(j * 0.8))
                draw.ellipse([px - 5, py - 5, px + 5, py + 5], fill=random.choice(flower_colors))
            draw.ellipse([fx - 5, fy - 25, fx + 5, fy - 15], fill=(255, 255, 200))
    
    # 绘制彩色圆圈装饰
    colors = [(255, 179, 198), (168, 230, 207), (255, 200, 150), (180, 180, 255), (255, 180, 200)]
    for i in range(15):
        x = random.randint(0, 512)
        y = random.randint(0, 512)
        size = random.randint(8, 25)
        color = random.choice(colors)
        draw.ellipse([x, y, x+size, y+size], fill=color)
    
    # 绘制底部草地
    grass_color = (168, 230, 207)
    draw.rectangle([0, 460, 512, 512], fill=grass_color)
    
    # 绘制小草
    for i in range(20):
        gx = random.randint(0, 512)
        gy = random.randint(440, 470)
        draw.polygon([(gx, gy), (gx + 5, gy - 15), (gx + 10, gy)], fill=(140, 200, 170))
    
    buffer = BytesIO()
    img.save(buffer, format='PNG')
    return buffer.getvalue()


def get_default_character_image(index=0, character_description=None):
    """生成默认角色头像（优先使用IP目录中的图片，当API调用失败时使用）"""
    from io import BytesIO
    
    try:
        ip_files = [f for f in os.listdir(IP_DIR) if f.lower().endswith('.png')]
        if ip_files:
            ip_files.sort()
            selected_ip = ip_files[index % len(ip_files)]
            ip_path = os.path.join(IP_DIR, selected_ip)
            
            with open(ip_path, 'rb') as f:
                ip_image_data = f.read()
            
            img = Image.open(BytesIO(ip_image_data))
            
            if img.mode != 'RGBA':
                img = img.convert('RGBA')
            
            target_size = (512, 512)
            bg = Image.new('RGBA', target_size, (255, 255, 255, 255))
            
            scale = min(target_size[0] / img.width, target_size[1] / img.height)
            new_width = int(img.width * scale)
            new_height = int(img.height * scale)
            img = img.resize((new_width, new_height), Image.Resampling.BICUBIC)
            
            offset_x = (target_size[0] - new_width) // 2
            offset_y = (target_size[1] - new_height) // 2
            bg.paste(img, (offset_x, offset_y), img)
            
            buffer = BytesIO()
            bg.save(buffer, format='PNG')
            logger.info(f"使用IP图片作为角色头像: {selected_ip}")
            return buffer.getvalue()
    except Exception as e:
        logger.warning(f"使用IP图片失败，回退到默认头像: {str(e)}")
    
    avatar_colors = [
        (255, 100, 150), (50, 200, 150), (255, 150, 50), 
        (100, 150, 255), (255, 50, 100), (50, 150, 255),
        (200, 100, 255), (255, 200, 50), (150, 255, 100)
    ]
    
    bg_colors = [
        (255, 200, 220), (200, 255, 220), (255, 220, 200),
        (200, 220, 255), (255, 200, 255), (200, 255, 255)
    ]
    
    color_index = index % len(avatar_colors)
    bg_index = index % len(bg_colors)
    
    img = Image.new('RGBA', (512, 512), (255, 255, 255, 255))
    draw = ImageDraw.Draw(img)
    
    center_x, center_y = 256, 256
    avatar_radius = 200
    
    gradient_colors = avatar_colors[color_index]
    r_start, g_start, b_start = gradient_colors
    
    for r in range(avatar_radius, 0, -1):
        factor = 1 - (avatar_radius - r) / avatar_radius
        r_color = int(r_start * factor + 255 * (1 - factor) * 0.5)
        g_color = int(g_start * factor + 255 * (1 - factor) * 0.5)
        b_color = int(b_start * factor + 255 * (1 - factor) * 0.5)
        draw.ellipse([
            center_x - r, center_y - r,
            center_x + r, center_y + r
        ], fill=(r_color, g_color, b_color))
    
    try:
        font = ImageFont.truetype('msyh.ttc', 150)
    except:
        font = ImageFont.load_default()
    
    if character_description and isinstance(character_description, str):
        first_char = character_description[0]
        if first_char.isalpha() or '\u4e00' <= first_char <= '\u9fff':
            text_char = first_char
        else:
            text_char = chr(ord('A') + (index % 26))
    else:
        text_char = chr(ord('A') + (index % 26))
    
    text_bbox = draw.textbbox((0, 0), text_char, font=font)
    text_width = text_bbox[2] - text_bbox[0]
    text_height = text_bbox[3] - text_bbox[1]
    
    text_x = center_x - text_width // 2
    text_y = center_y - text_height // 2
    
    draw.text((text_x, text_y), text_char, font=font, fill=(255, 255, 255), stroke_width=5, stroke_fill=(0, 0, 0))
    
    stroke_color = tuple(max(0, c - 80) for c in avatar_colors[color_index])
    draw.ellipse([
        center_x - avatar_radius - 10, center_y - avatar_radius - 10,
        center_x + avatar_radius + 10, center_y + avatar_radius + 10
    ], outline=(0, 0, 0), width=12)
    
    for i in range(8):
        sparkle_x = center_x + int(avatar_radius * 0.7 * math.cos(i * 2 * math.pi / 8))
        sparkle_y = center_y + int(avatar_radius * 0.7 * math.sin(i * 2 * math.pi / 8))
        draw.polygon([
            (sparkle_x, sparkle_y - 8),
            (sparkle_x + 4, sparkle_y + 4),
            (sparkle_x - 8, sparkle_y)
        ], fill=(255, 255, 255))
    
    buffer = BytesIO()
    img.save(buffer, format='PNG')
    return buffer.getvalue()


POSE_PRESETS = {
    'standing': {'flip': False, 'rotate': 0, 'scale': 1.0, 'offset_y': 0, 'description': '站立'},
    'sitting': {'flip': False, 'rotate': 0, 'scale': 0.9, 'offset_y': 15, 'description': '坐着'},
    'waving': {'flip': False, 'rotate': 0, 'scale': 1.0, 'offset_y': 0, 'description': '挥手'},
    'thinking': {'flip': False, 'rotate': 0, 'scale': 1.0, 'offset_y': 5, 'description': '思考'},
    'happy': {'flip': False, 'rotate': 0, 'scale': 1.05, 'offset_y': -5, 'description': '开心跳跃'},
    'sad': {'flip': False, 'rotate': 0, 'scale': 0.95, 'offset_y': 10, 'description': '难过'},
    'surprised': {'flip': False, 'rotate': 0, 'scale': 1.1, 'offset_y': -10, 'description': '惊讶'},
    'walking': {'flip': False, 'rotate': 0, 'scale': 1.0, 'offset_y': 0, 'description': '走路'},
    'running': {'flip': False, 'rotate': 0, 'scale': 1.05, 'offset_y': -5, 'description': '奔跑'},
    'pointing': {'flip': False, 'rotate': 0, 'scale': 1.0, 'offset_y': 0, 'description': '指向'},
    'reading': {'flip': False, 'rotate': 0, 'scale': 0.95, 'offset_y': 10, 'description': '阅读'},
    'sleeping': {'flip': False, 'rotate': 0, 'scale': 0.85, 'offset_y': 25, 'description': '睡觉'},
}


LAYOUT_CONFIG = {
    'single': {
        'name': '单格',
        'rows': 1,
        'cols': 1,
        'panel_count': 1,
        'grid_specs': [{'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1}],
        'css_grid': 'grid-template-columns: 1fr; grid-template-rows: minmax(250px, 1fr);'
    },
    'two-horizontal': {
        'name': '两格横排',
        'rows': 1,
        'cols': 2,
        'panel_count': 2,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr; grid-template-rows: minmax(250px, 1fr);'
    },
    'two-vertical': {
        'name': '两格竖排',
        'rows': 2,
        'cols': 1,
        'panel_count': 2,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'four-grid': {
        'name': '四格经典',
        'rows': 2,
        'cols': 2,
        'panel_count': 4,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'four-grid-l-shape': {
        'name': '四格L型',
        'rows': 3,
        'cols': 2,
        'panel_count': 4,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 2, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 2}
        ],
        'css_grid': 'grid-template-columns: 2fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'four-grid-top-big': {
        'name': '四格上大下小',
        'rows': 2,
        'cols': 3,
        'panel_count': 4,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 3},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr 1fr; grid-template-rows: minmax(250px, 2fr) minmax(250px, 1fr);'
    },
    'four-grid-left-big': {
        'name': '四格左大右小',
        'rows': 3,
        'cols': 2,
        'panel_count': 4,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 3, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 2, 'col': 1, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 2fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'three-horizontal': {
        'name': '三格横排',
        'rows': 1,
        'cols': 3,
        'panel_count': 3,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr 1fr; grid-template-rows: minmax(250px, 1fr);'
    },
    'one-big-two-small': {
        'name': '一大两小',
        'rows': 2,
        'cols': 2,
        'panel_count': 3,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 2, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 2fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'three-vertical': {
        'name': '三格竖排',
        'rows': 3,
        'cols': 1,
        'panel_count': 3,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'six-grid': {
        'name': '六格布局',
        'rows': 2,
        'cols': 3,
        'panel_count': 6,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'eight-grid': {
        'name': '八格布局',
        'rows': 2,
        'cols': 4,
        'panel_count': 8,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 3, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 3, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: repeat(4, 1fr); grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'five-grid': {
        'name': '五格布局',
        'rows': 2,
        'cols': 3,
        'panel_count': 5,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'seven-grid': {
        'name': '七格布局',
        'rows': 2,
        'cols': 4,
        'panel_count': 7,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1},
            {'row': 0, 'col': 3, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1}
        ],
        'css_grid': 'grid-template-columns: repeat(4, 1fr); grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'four-cinematic': {
        'name': '四格电影感',
        'description': '顶部满幅横版 + 左右双小格 + 底部跨幅收尾，经典电影分镜节奏',
        'rows': 3,
        'cols': 2,
        'panel_count': 4,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'action-close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'expression-close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '4:3', 'shot_type': 'medium', 'angle': 'low-angle', 'sound_effect': False, 'guide_sprite': True, 'sprite_position': 'bottom-right'}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr; grid-template-rows: minmax(250px, 1.2fr) minmax(250px, 1fr) minmax(250px, 1.3fr);'
    },
    'vertical-strip-6': {
        'name': '竖版长条六格',
        'description': '第1行通栏大横格(全景) + 第2行两个等宽竖格(特写) + 第3行通栏大横格(高潮) + 第4行左宽右窄格(收尾+悬念)',
        'rows': 4,
        'cols': 2,
        'panel_count': 6,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '9:16', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '9:16', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '4:3', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 3, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:3', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'bottom-right'}
        ],
        'css_grid': 'grid-template-columns: 2fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1.2fr) minmax(250px, 1.1fr) minmax(250px, 0.8fr);'
    },
    'freeform-strip-6': {
        'name': '异形错落拼贴六格',
        'description': '竖版长条多格分镜漫画，无对称网格，无整齐行列，分镜大小不一、高低错落、带自然倾斜角度，相邻格子半错位叠压，混搭6种完全不同的异形边框，手绘漫画拼贴感',
        'panel_count': 6,
        'freeform': True,
        'grid_specs': [
            {
                'id': 0, 'name': '开篇全景格',
                'width': 700, 'height': 280,
                'x': 20, 'y': 30,
                'tilt_angle': 1.5,
                'z_order': 1,
                'border_type': 'paper-tear',
                'border_prompt': 'warm brown torn rice paper edge, no solid outline, natural fiber texture and roughness on edges, irregular arc corners, ancient book torn page texture, warm beige paper background texture',
                'shot_type': 'wide', 'angle': 'eye-level'
            },
            {
                'id': 1, 'name': '特写疑问格',
                'width': 220, 'height': 340,
                'x': 30, 'y': 300,
                'tilt_angle': -0.8,
                'z_order': 3,
                'border_type': 'water-ripple',
                'border_prompt': 'light blue dashed water ripple lines border, line thickness varies with wave undulation, concave curved corners, no right angles or regular rounded corners, soft flowing organic border',
                'shot_type': 'close-up', 'angle': 'eye-level'
            },
            {
                'id': 2, 'name': '思考近景格',
                'width': 260, 'height': 380,
                'x': 320, 'y': 350,
                'tilt_angle': 1.2,
                'z_order': 2,
                'border_type': 'pencil-sketch',
                'border_prompt': 'dark brown pencil sketch rough border, hand-drawn jitter and grain texture, irregular small bubble outline bulging at top-right corner, hand-drawn thought bubble style edge',
                'shot_type': 'medium', 'angle': 'eye-level'
            },
            {
                'id': 3, 'name': '高潮反转格',
                'width': 720, 'height': 320,
                'x': 10, 'y': 720,
                'tilt_angle': -2.0,
                'z_order': 4,
                'border_type': 'wave-crash',
                'border_prompt': 'navy blue thick solid base border, left and right sides middle section broken by water current and not closed, water elements overflowing outside the border frame, top and bottom borders preserved with rough edge texture',
                'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False
            },
            {
                'id': 4, 'name': '动感收尾格',
                'width': 480, 'height': 360,
                'x': 40, 'y': 1020,
                'tilt_angle': -2.5,
                'z_order': 5,
                'border_type': 'speed-lines',
                'border_prompt': 'no closed solid border, black radiating speed lines extending outward from center, jagged irregular outer edge naturally formed at line ends, soft blur gradient at edges',
                'shot_type': 'dynamic', 'angle': 'low-angle'
            },
            {
                'id': 5, 'name': '悬念小精灵格',
                'width': 200, 'height': 200,
                'x': 500, 'y': 1080,
                'tilt_angle': 3.0,
                'z_order': 6,
                'border_type': 'glow-bubble',
                'border_prompt': 'no solid outline stroke, pale gold soft light gradient forming fuzzy boundary, light gradually transparent from inside to outside, scattered tiny starlight dots around edge, floating glowing bubble effect',
                'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'
            }
        ]
    },

    'five-cinematic': {
        'name': '五格电影感',
        'description': '顶部满幅横版(全景) + 第二行左右双小格(特写) + 第三行左侧大格(高潮) + 右下角小精灵格',
        'rows': 3,
        'cols': 2,
        'panel_count': 5,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '4:3', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 2, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:3', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: 2fr 1fr; grid-template-rows: minmax(250px, 1.2fr) minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'five-l-shape': {
        'name': '五格L型',
        'description': '左侧竖版大格(全景) + 右上三小格(特写) + 底部横版收尾',
        'rows': 3,
        'cols': 2,
        'panel_count': 5,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 2, 'col_span': 1, 'aspect_ratio': '9:16', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 2, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:3', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'bottom-right'}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 0.8fr);'
    },
    'five-top-big': {
        'name': '五格上大下小',
        'description': '顶部满幅横版(全景) + 底部四小格并列',
        'rows': 2,
        'cols': 4,
        'panel_count': 5,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 4, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 1, 'col': 3, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: repeat(4, 1fr); grid-template-rows: minmax(250px, 1.5fr) minmax(250px, 1fr);'
    },

    'six-cinematic': {
        'name': '六格电影感',
        'description': '顶部满幅(全景) + 第二行三小格(特写) + 第三行跨幅(高潮) + 底部双小格(收尾)',
        'rows': 4,
        'cols': 3,
        'panel_count': 6,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '4:3', 'shot_type': 'medium', 'angle': 'eye-level', 'guide_sprite': True, 'sprite_position': 'bottom-right'}
        ],
        'css_grid': 'grid-template-columns: repeat(3, 1fr); grid-template-rows: minmax(250px, 1.2fr) minmax(250px, 1fr) minmax(250px, 1.1fr) minmax(250px, 0.8fr);'
    },
    'six-diamond': {
        'name': '六格菱形',
        'description': '中心大格(高潮) + 周围五小格围绕',
        'rows': 3,
        'cols': 3,
        'panel_count': 6,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 2, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:3', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: 1fr 2fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 2fr) minmax(250px, 1fr);'
    },

    'seven-cinematic': {
        'name': '七格电影感',
        'description': '顶部满幅(全景) + 第二行三小格 + 第三行双小格(高潮) + 底部小精灵格',
        'rows': 4,
        'cols': 3,
        'panel_count': 7,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 2, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:3', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '4:3', 'shot_type': 'medium', 'angle': 'eye-level', 'guide_sprite': True, 'sprite_position': 'bottom-right'}
        ],
        'css_grid': 'grid-template-columns: repeat(3, 1fr); grid-template-rows: minmax(250px, 1.2fr) minmax(250px, 1fr) minmax(250px, 1.1fr) minmax(250px, 0.7fr);'
    },
    'seven-staggered': {
        'name': '七格错落',
        'description': '错落排列的七格布局，大小不一，富有动感',
        'rows': 3,
        'cols': 3,
        'panel_count': 7,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 2, 'col_span': 1, 'aspect_ratio': '9:16', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 2, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:3', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 0.8fr);'
    },

    'eight-cinematic': {
        'name': '八格电影感',
        'description': '两行四格 + 底部两行跨幅，电影分镜节奏',
        'rows': 4,
        'cols': 4,
        'panel_count': 8,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 3, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 4, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 4, 'aspect_ratio': '4:3', 'shot_type': 'medium', 'angle': 'eye-level', 'guide_sprite': True, 'sprite_position': 'bottom-right'}
        ],
        'css_grid': 'grid-template-columns: repeat(4, 1fr); grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 1.2fr) minmax(250px, 0.8fr);'
    },
    'eight-vertical': {
        'name': '八格竖版',
        'description': '两列四行竖版布局，适合手机阅读',
        'rows': 4,
        'cols': 2,
        'panel_count': 8,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '9:16', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '9:16', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '9:16', 'shot_type': 'medium', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 2, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '9:16', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'bottom-right'}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr; grid-template-rows: minmax(250px, 0.8fr) minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 0.8fr);'
    },

    'nine-grid': {
        'name': '九格经典',
        'description': '三行三列标准网格布局',
        'rows': 3,
        'cols': 3,
        'panel_count': 9,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 2, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 2, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: repeat(3, 1fr); grid-template-rows: repeat(3, minmax(250px, 1fr));'
    },
    'nine-cinematic': {
        'name': '九格电影感',
        'description': '顶部满幅 + 中间六小格 + 底部跨幅收尾',
        'rows': 3,
        'cols': 3,
        'panel_count': 9,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 2, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 2, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '16:9', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 4, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '4:3', 'shot_type': 'medium', 'angle': 'eye-level', 'guide_sprite': True, 'sprite_position': 'bottom-right'}
        ],
        'css_grid': 'grid-template-columns: repeat(3, 1fr); grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 0.8fr);'
    },
    'nine-cross': {
        'name': '九格十字',
        'description': '十字形布局，中心大格 + 四周小格',
        'rows': 3,
        'cols': 3,
        'panel_count': 9,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:1', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:2', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:2', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 2, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 2, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: 1fr 2fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 2fr) minmax(250px, 1fr);'
    },

    'ten-grid': {
        'name': '十格经典',
        'description': '两行五列标准网格布局',
        'rows': 2,
        'cols': 5,
        'panel_count': 10,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 0, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 0, 'col': 3, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 0, 'col': 4, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 3, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 4, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: repeat(5, 1fr); grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr);'
    },
    'ten-cinematic': {
        'name': '十格电影感',
        'description': '顶部满幅 + 中间两行四小格 + 底部双小格收尾',
        'rows': 4,
        'cols': 4,
        'panel_count': 10,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 4, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 3, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 2, 'col': 2, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '4:3', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 3, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'high-angle'},
            {'row': 3, 'col': 3, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: repeat(4, 1fr); grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 1.2fr) minmax(250px, 0.8fr);'
    },
    'ten-staggered': {
        'name': '十格错落',
        'description': '错落排列的十格布局，大小不一，富有动感',
        'rows': 4,
        'cols': 3,
        'panel_count': 10,
        'grid_specs': [
            {'row': 0, 'col': 0, 'row_span': 1, 'col_span': 2, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'eye-level'},
            {'row': 0, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:3', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 1, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 1, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'eye-level'},
            {'row': 2, 'col': 0, 'row_span': 1, 'col_span': 3, 'aspect_ratio': '16:9', 'shot_type': 'wide', 'angle': 'low-angle', 'sound_effect': False},
            {'row': 3, 'col': 0, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 3, 'col': 1, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'medium', 'angle': 'eye-level'},
            {'row': 3, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '1:1', 'shot_type': 'close-up', 'angle': 'high-angle'},
            {'row': 3, 'col': 2, 'row_span': 1, 'col_span': 1, 'aspect_ratio': '2:3', 'shot_type': 'close-up', 'angle': 'high-angle', 'guide_sprite': True, 'sprite_position': 'center'}
        ],
        'css_grid': 'grid-template-columns: 1fr 1fr 1fr; grid-template-rows: minmax(250px, 1fr) minmax(250px, 1fr) minmax(250px, 1.2fr) minmax(250px, 1fr);'
    }
}


def get_layout_config(layout_name):
    return LAYOUT_CONFIG.get(layout_name, LAYOUT_CONFIG['four-grid'])


def get_layout_variants(panel_count):
    """获取指定格子数的所有布局变体"""
    variants = []
    for key, config in LAYOUT_CONFIG.items():
        if config.get('panel_count') == panel_count:
            variants.append({
                'id': key,
                'name': config.get('name', key),
                'panel_count': config.get('panel_count'),
                'css_grid': config.get('css_grid', '')
            })
    return variants


def parse_layout_prompts(layout_prompts):
    """从自然语言布局提示词解析出布局配置
    
    输入：布局提示词列表，每个元素对应一格的自然语言描述
    输出：结构化布局配置（包含grid_specs、css_grid等）
    
    解析策略：
    1. 先根据分镜数量匹配默认布局模板
    2. 根据每格的提示词关键词，细化该格的参数（宽高比、视角、景别、是否有小精灵等）
    """
    if not layout_prompts or not isinstance(layout_prompts, list):
        panel_count = 4
        return get_layout_config('four-grid')
    
    panel_count = len(layout_prompts)
    
    # 第一步：根据整体特征匹配基础布局模板
    base_layout = detect_base_layout(layout_prompts, panel_count)
    config = get_layout_config(base_layout)
    
    # 第二步：逐格细化参数
    grid_specs = config.get('grid_specs', [])
    enriched_specs = []
    
    for i, prompt in enumerate(layout_prompts):
        if i < len(grid_specs):
            spec = dict(grid_specs[i])
        else:
            spec = {'row': i, 'col': 0, 'row_span': 1, 'col_span': 1}
        
        # 解析宽高比
        aspect = detect_aspect_ratio(prompt)
        if aspect:
            spec['aspect_ratio'] = aspect
        
        # 解析视角
        angle = detect_angle(prompt)
        if angle:
            spec['angle'] = angle
        
        # 解析景别
        shot = detect_shot_type(prompt)
        if shot:
            spec['shot_type'] = shot
        
        # 检测音效
        if detect_sound_effect(prompt):
            spec['sound_effect'] = True
        
        # 检测小精灵
        sprite_info = detect_guide_sprite(prompt)
        if sprite_info:
            spec['guide_sprite'] = True
            spec['sprite_position'] = sprite_info
        
        spec['layout_prompt'] = prompt
        enriched_specs.append(spec)
    
    result = dict(config)
    result['grid_specs'] = enriched_specs
    result['panel_count'] = panel_count
    result['detected_layout'] = base_layout
    
    return result


def detect_base_layout(layout_prompts, panel_count):
    """根据布局提示词的整体特征判断基础布局模板"""
    if not layout_prompts:
        return 'four-grid' if panel_count == 4 else f'{panel_count}-grid'
    
    full_text = ' '.join(layout_prompts)
    
    # 4格电影感布局关键词
    cinematic_keywords = ['满幅', '横版', '左右等分', '双小格', '并列', '跨幅', '错落', '电影', '全景']
    cinematic_score = sum(1 for kw in cinematic_keywords if kw in full_text)
    
    if panel_count == 4 and cinematic_score >= 3:
        return 'four-cinematic'
    
    # 按格子数匹配默认布局
    layout_map = {
        1: 'single',
        2: 'two-horizontal',
        3: 'three-horizontal',
        4: 'four-grid',
        5: 'five-grid',
        6: 'six-grid',
        7: 'seven-grid',
        8: 'eight-grid'
    }
    
    return layout_map.get(panel_count, 'four-grid')


def detect_aspect_ratio(prompt):
    """从提示词中检测宽高比"""
    if not prompt or not isinstance(prompt, str):
        return None
    
    prompt_lower = prompt.lower()
    
    # 横版关键词
    if any(kw in prompt for kw in ['横版', '宽幅', '满幅', '全景', '横向', '宽屏', 'landscape', 'wide']):
        if any(kw in prompt for kw in ['16:9', '十六比九']):
            return '16:9'
        return '16:9'
    
    # 竖版关键词
    if any(kw in prompt for kw in ['竖版', '竖幅', '纵向', 'portrait', 'tall']):
        return '9:16'
    
    # 方形关键词
    if any(kw in prompt for kw in ['方形', '正方', '等大', '等分', 'square']):
        return '1:1'
    
    # 3:2关键词
    if any(kw in prompt for kw in ['3:2', '三比二']):
        return '3:2'
    
    # 4:3关键词
    if any(kw in prompt for kw in ['4:3', '四比三']):
        return '4:3'
    
    return None


def detect_angle(prompt):
    """从提示词中检测视角/角度"""
    if not prompt or not isinstance(prompt, str):
        return None
    
    angle_keywords = {
        'low-angle': ['仰视', '低角度', '仰拍', 'low angle', '仰视视角'],
        'high-angle': ['俯视', '高角度', '俯拍', 'high angle', '俯视视角'],
        'eye-level': ['平视', 'eye level', '水平视角', '正面'],
        'dutch-angle': ['倾斜', '斜角', '荷兰角', 'dutch angle']
    }
    
    for angle, keywords in angle_keywords.items():
        for kw in keywords:
            if kw in prompt:
                return angle
    
    return None


def detect_shot_type(prompt):
    """从提示词中检测景别/镜头类型"""
    if not prompt or not isinstance(prompt, str):
        return None
    
    shot_keywords = {
        'wide': ['全景', '远景', '广角', 'wide shot', 'full shot', '开阔'],
        'medium': ['中景', 'medium shot', '半身'],
        'close-up': ['特写', '近景', 'close-up', '近拍'],
        'action-close-up': ['动作特写', '动态特写', 'action'],
        'expression-close-up': ['表情特写', '表情近景', '面部特写', 'expression'],
        'extreme-close-up': ['大特写', '细节特写', 'extreme close-up']
    }
    
    for shot, keywords in shot_keywords.items():
        for kw in keywords:
            if kw in prompt:
                return shot
    
    return None


def detect_sound_effect(prompt):
    """检测是否包含音效/拟声词要求"""
    if not prompt or not isinstance(prompt, str):
        return False
    
    sound_keywords = ['拟声', '音效', '声效', '艺术字', '砰', '轰', '哗', '咚', '锵', '哇哦', 'boom', 'sound effect']
    return any(kw in prompt for kw in sound_keywords)


def detect_guide_sprite(prompt):
    """检测小精灵位置"""
    if not prompt or not isinstance(prompt, str):
        return None
    
    if any(kw in prompt for kw in ['小精灵', '知识精灵', '精灵', 'guide', 'sprite']):
        if '左下' in prompt:
            return 'bottom-left'
        elif '右下' in prompt or '右下角' in prompt:
            return 'bottom-right'
        elif '左上' in prompt:
            return 'top-left'
        elif '右上' in prompt or '右上角' in prompt:
            return 'top-right'
        else:
            return 'bottom-right'
    
    return None


def calculate_image_size_for_panel(spec, base_size=1024):
    """根据分镜规格计算生图尺寸（宽高）。base_size 提高到 1024 以改善清晰度"""
    aspect_ratio = spec.get('aspect_ratio', '1:1')
    
    try:
        w_ratio, h_ratio = [int(x) for x in aspect_ratio.split(':')]
    except:
        w_ratio, h_ratio = 1, 1
    
    # 以base_size为基准，按比例缩放
    if w_ratio >= h_ratio:
        width = base_size
        height = int(base_size * h_ratio / w_ratio)
    else:
        height = base_size
        width = int(base_size * w_ratio / h_ratio)
    
    # 确保是偶数（API要求）
    width = width - (width % 2)
    height = height - (height % 2)
    
    return width, height


def normalize_background_for_panel(image, target_width, target_height):
    """去除生成器附带的纯黑电影边，并按漫画格尺寸 cover 裁切。

    生图服务偶尔会忽略请求比例，在图片内部加入上下或左右黑边。这里先只裁掉
    从画面边缘连续延伸的近黑色带，再用 ImageOps.fit 填满目标格，绝不补黑边。
    """
    from PIL import ImageOps

    img = image.convert('RGB')
    width, height = img.size
    if width > 8 and height > 8:
        gray = img.convert('L')
        sample_w = min(width, 256)
        sample_h = min(height, 256)
        sample = gray.resize((sample_w, sample_h), Image.Resampling.BILINEAR)
        pixels = sample.load()

        # 生图模型的“黑边”不总是纯黑；实测也会输出亮度约 34 的深灰电影边。
        # 阈值放到 50，但仍要求一整行/列至少 96% 都很暗，避免误裁普通暗色场景。
        def row_is_bar(y):
            dark = sum(1 for x in range(sample_w) if pixels[x, y] <= 50)
            return dark >= sample_w * 0.96

        def col_is_bar(x):
            dark = sum(1 for y in range(sample_h) if pixels[x, y] <= 50)
            return dark >= sample_h * 0.96

        top = 0
        while top < sample_h and row_is_bar(top):
            top += 1
        bottom = sample_h
        while bottom > top and row_is_bar(bottom - 1):
            bottom -= 1
        left = 0
        while left < sample_w and col_is_bar(left):
            left += 1
        right = sample_w
        while right > left and col_is_bar(right - 1):
            right -= 1

        # 至少保留一半画面，防止真正的夜景被误判为边框。
        if right - left >= sample_w * 0.5 and bottom - top >= sample_h * 0.5:
            scale_x = width / sample_w
            scale_y = height / sample_h
            # 缩略检测在黑边交界处会产生双线性灰色过渡；向内容区多收两个
            # 采样像素，避免 cover 后仍残留 1～2px 黑线。
            inset_x = int(round(scale_x * 2))
            inset_y = int(round(scale_y * 2))
            crop_box = (
                int(round(left * scale_x)) + (inset_x if left > 0 else 0),
                int(round(top * scale_y)) + (inset_y if top > 0 else 0),
                int(round(right * scale_x)) - (inset_x if right < sample_w else 0),
                int(round(bottom * scale_y)) - (inset_y if bottom < sample_h else 0),
            )
            if crop_box != (0, 0, width, height):
                logger.info(
                    f"[背景规范化] 裁除黑边: {width}x{height} crop={crop_box}"
                )
                img = img.crop(crop_box)

    target_width = max(2, int(target_width))
    target_height = max(2, int(target_height))
    return ImageOps.fit(
        img, (target_width, target_height),
        method=Image.Resampling.LANCZOS,
        centering=(0.5, 0.5)
    )


def calculate_cell_aspect_ratio(layout_config, grid_spec):
    """根据CSS Grid的fr值和grid_spec的跨度计算单元格实际宽高比。

    返回最简整数比字符串，如 '16:9'、'1:1'、'9:16'。
    """
    css_grid = layout_config.get('css_grid', '')
    rows, cols = parse_fr_values(css_grid)
    
    if not rows:
        rows = [1.0] * layout_config.get('rows', 1)
    if not cols:
        cols = [1.0] * layout_config.get('cols', 1)
    
    row_span = grid_spec.get('row_span', 1)
    col_span = grid_spec.get('col_span', 1)
    row = grid_spec.get('row', 0)
    col = grid_spec.get('col', 0)
    
    col_fr = sum(cols[col:col + col_span])
    row_fr = sum(rows[row:row + row_span])
    
    if row_fr == 0:
        row_fr = 1.0
    if col_fr == 0:
        col_fr = 1.0
    
    ratio = col_fr / row_fr
    # 转换为最简整数比
    def float_to_ratio(value, max_denominator=100):
        from fractions import Fraction
        try:
            frac = Fraction(value).limit_denominator(max_denominator)
            return frac.numerator, frac.denominator
        except Exception:
            # 兜底：按比例放大到整数
            w, h = int(value * 100), 100
            from math import gcd
            g = gcd(w, h)
            return w // g, h // g
    
    w, h = float_to_ratio(ratio)
    return f"{w}:{h}"
def parse_fr_values(css_grid):
    """解析CSS grid模板中的fr值，返回行列的比例列表。

    支持：
    - 普通 Xfr
    - minmax(250px, Xfr)
    - repeat(N, Xfr)
    """
    rows = []
    cols = []

    def _parse_values(template_str):
        values = []
        pos = 0
        while pos < len(template_str):
            # 跳过空白
            while pos < len(template_str) and template_str[pos].isspace():
                pos += 1
            if pos >= len(template_str):
                break
            # repeat(N, Xfr)
            m = re.match(r'repeat\(\s*(\d+)\s*,\s*(\d*\.?\d+)fr\s*\)', template_str[pos:])
            if m:
                n = int(m.group(1))
                fr = float(m.group(2))
                values.extend([fr] * n)
                pos += m.end()
                continue
            # minmax(..., Xfr)
            m = re.match(r'minmax\([^,]+,\s*(\d*\.?\d+)fr\s*\)', template_str[pos:])
            if m:
                values.append(float(m.group(1)))
                pos += m.end()
                continue
            # 普通 Xfr
            m = re.match(r'(\d*\.?\d+)fr', template_str[pos:])
            if m:
                values.append(float(m.group(1)))
                pos += m.end()
                continue
            # 未知字符，跳过
            pos += 1
        return values

    if 'grid-template-columns' in css_grid:
        col_match = re.search(r'grid-template-columns:\s*([^;]+)', css_grid)
        if col_match:
            cols = _parse_values(col_match.group(1))

    if 'grid-template-rows' in css_grid:
        row_match = re.search(r'grid-template-rows:\s*([^;]+)', css_grid)
        if row_match:
            rows = _parse_values(row_match.group(1))

    return rows, cols


def calculate_panel_dimensions(layout_config, total_width, total_height, gap=4):
    """根据布局配置计算每个面板的实际像素尺寸"""
    grid_specs = layout_config.get('grid_specs', [])
    css_grid = layout_config.get('css_grid', '')
    
    rows, cols = parse_fr_values(css_grid)
    
    if not rows:
        rows = [1.0] * layout_config.get('rows', 2)
    if not cols:
        cols = [1.0] * layout_config.get('cols', 2)
    
    total_row_fr = sum(rows)
    total_col_fr = sum(cols)
    
    row_heights = []
    remaining_height = total_height - (len(rows) - 1) * gap
    for r in rows:
        row_heights.append(int(remaining_height * r / total_row_fr))
    
    col_widths = []
    remaining_width = total_width - (len(cols) - 1) * gap
    for c in cols:
        col_widths.append(int(remaining_width * c / total_col_fr))
    
    panel_dimensions = []
    for spec in grid_specs:
        row = spec.get('row', 0)
        col = spec.get('col', 0)
        row_span = spec.get('row_span', 1)
        col_span = spec.get('col_span', 1)
        
        start_row = row
        end_row = row + row_span
        start_col = col
        end_col = col + col_span
        
        panel_width = sum(col_widths[start_col:end_col]) + (col_span - 1) * gap
        panel_height = sum(row_heights[start_row:end_row]) + (row_span - 1) * gap
        
        panel_dimensions.append({
            'width': panel_width,
            'height': panel_height,
            'row': row,
            'col': col,
            'row_span': row_span,
            'col_span': col_span
        })
    
    return panel_dimensions


BUBBLE_STYLES = {
    'rounded': {'name': '圆角矩形', 'shape': 'rounded_rectangle', 'radius': 12},
    'cloud': {'name': '云朵', 'shape': 'cloud', 'radius': 0.5},
    'oval': {'name': '椭圆形', 'shape': 'ellipse', 'radius': 0.5},
    'sharp': {'name': '尖角矩形', 'shape': 'sharp_rectangle', 'radius': 4},
    'heart': {'name': '心形', 'shape': 'heart', 'radius': 0.5},
    'star': {'name': '星形', 'shape': 'star', 'radius': 5},
    'burst': {'name': '爆炸形', 'shape': 'burst', 'radius': 8},
    'thought': {'name': '思考气泡', 'shape': 'thought', 'radius': 0.4},
    'whisper': {'name': '悄悄话', 'shape': 'whisper', 'radius': 0.3},
    'shout': {'name': '大喊', 'shape': 'shout', 'radius': 6},
    'wave': {'name': '波浪边', 'shape': 'wave', 'radius': 8},
    'diamond': {'name': '菱形', 'shape': 'diamond', 'radius': 0},
    'bubble': {'name': '圆形气泡', 'shape': 'circle', 'radius': 0.5},
    'scroll': {'name': '卷轴', 'shape': 'scroll', 'radius': 10},
    'banner': {'name': '横幅', 'shape': 'banner', 'radius': 6}
}


FONT_FAMILIES = {
    'msyh': {'name': '微软雅黑', 'path': 'C:/Windows/Fonts/msyh.ttc'},
    'msyhbd': {'name': '微软雅黑加粗', 'path': 'C:/Windows/Fonts/msyhbd.ttc'},
    'simkai': {'name': '楷体', 'path': 'C:/Windows/Fonts/simkai.ttf'},
    'simhei': {'name': '黑体', 'path': 'C:/Windows/Fonts/simhei.ttf'},
    'simsun': {'name': '宋体', 'path': 'C:/Windows/Fonts/simsun.ttc'},
    'simfang': {'name': '仿宋', 'path': 'C:/Windows/Fonts/simfang.ttf'},
    'stkaiti': {'name': '华文楷体', 'path': 'C:/Windows/Fonts/stkaiti.ttf'},
    'stxingkai': {'name': '华文行楷', 'path': 'C:/Windows/Fonts/stxingkai.ttf'},
    'stfangsong': {'name': '华文仿宋', 'path': 'C:/Windows/Fonts/stfangsong.ttf'},
    'stxinwei': {'name': '华文新魏', 'path': 'C:/Windows/Fonts/stxinwei.ttf'},
    'stzhongsong': {'name': '华文中宋', 'path': 'C:/Windows/Fonts/stzhongsong.ttf'}
}


def parse_pose_description(description):
    """解析姿势描述文本，返回姿势参数"""
    if not description or not isinstance(description, str):
        return None
    
    desc_lower = description.lower().strip()
    
    pose_keywords = {
        'standing': ['站立', '站着', '直立', 'standing', 'stand'],
        'sitting': ['坐着', '坐', '坐姿', 'sitting', 'sit'],
        'waving': ['挥手', '招手', '打招呼', 'waving', 'wave'],
        'thinking': ['思考', '想', '疑惑', 'thinking', 'think'],
        'happy': ['开心', '高兴', '快乐', '跳跃', 'happy', 'jump'],
        'sad': ['难过', '悲伤', '伤心', 'sad', 'cry'],
        'surprised': ['惊讶', '吃惊', 'surprised', 'shocked'],
        'walking': ['走路', '行走', '步行', 'walking', 'walk'],
        'running': ['奔跑', '跑', 'running', 'run'],
        'pointing': ['指向', '指', 'pointing', 'point'],
        'reading': ['阅读', '看书', '读书', 'reading', 'read'],
        'sleeping': ['睡觉', '睡', 'sleeping', 'sleep'],
    }
    
    for pose_name, keywords in pose_keywords.items():
        for keyword in keywords:
            if keyword.lower() in desc_lower:
                logger.info(f"姿势描述 '{description}' 匹配到预设姿势: {pose_name}")
                return POSE_PRESETS.get(pose_name)


EMOTION_KEYWORDS = {
    'happy': [
        '开心', '高兴', '快乐', '欢乐', '喜悦', '大笑', '笑', '哈哈哈', '嘻嘻', '嘿嘿',
        '兴奋', '激动', '兴奋', '愉快', '乐', '欢', '喜', '笑眯眯', '美滋滋',
        '开心极了', '高兴极了', '太棒了', '太好了', '精彩', '完美', '赞',
        'happy', 'glad', 'joy', 'joyful', 'delighted', 'excited', 'laugh', 'smile'
    ],
    'sad': [
        '难过', '伤心', '悲伤', '沮丧', '失落', '失望', '伤心欲绝', '难过极了',
        '哭泣', '哭', '流泪', '泪眼', '泪汪汪', '呜呜', '泣',
        '愁', '忧郁', '惆怅', '愁眉', '叹气', '叹息',
        'sad', 'sorrow', 'grief', 'upset', 'depressed', 'cry', 'tear'
    ],
    'angry': [
        '生气', '愤怒', '发火', '发脾气', '怒', '恼火', '怒气冲冲', '火冒三丈',
        '气', '恼', '暴跳如雷', '气坏了', '气死我了',
        '愤怒的', '生气的', '凶狠', '凶',
        'angry', 'furious', 'annoyed', 'mad', 'rage', 'irritated'
    ],
    'surprised': [
        '惊讶', '吃惊', '震惊', '意外', '忽然', '突然', '吓了一跳', '目瞪口呆',
        '哇', '哇哦', '啊', '天哪', '天呐', '竟然', '居然',
        '惊讶的', '吃惊的', '震惊的',
        'surprised', 'shocked', 'amazed', 'astonished', 'wow'
    ],
    'neutral': [
        '说', '讲', '问', '答', '回答', '道', '说', '告诉', '说道',
        'think', 'thought', 'say', 'tell', 'ask', 'answer', 'speak'
    ]
}


def analyze_emotion(text):
    """分析文本情绪，返回情绪类型（happy/sad/angry/surprised/neutral）"""
    if not text or not isinstance(text, str):
        return 'neutral'
    
    text_lower = text.lower().strip()
    
    scores = {'happy': 0, 'sad': 0, 'angry': 0, 'surprised': 0}
    
    for emotion, keywords in EMOTION_KEYWORDS.items():
        if emotion == 'neutral':
            continue
        for keyword in keywords:
            if keyword.lower() in text_lower:
                scores[emotion] += 1
    
    max_emotion = max(scores, key=scores.get)
    if scores[max_emotion] > 0:
        logger.info(f"情绪分析结果: '{text}' -> {max_emotion} (分数: {scores[max_emotion]})")
        return max_emotion
    
    return 'neutral'


EMOTION_BUBBLE_MAP = {
    'happy': ['cloud', 'burst', 'wave', 'star', 'heart'],
    'sad': ['ellipse', 'whisper', 'thought', 'rounded', 'bubble'],
    'angry': ['burst', 'sharp', 'shout', 'diamond', 'wave'],
    'surprised': ['star', 'burst', 'cloud', 'diamond', 'oval'],
    'neutral': ['rounded', 'oval', 'bubble', 'cloud']
}


def get_bubble_style_for_emotion(text):
    """根据文本情绪获取气泡样式"""
    emotion = analyze_emotion(text)
    styles = EMOTION_BUBBLE_MAP.get(emotion, EMOTION_BUBBLE_MAP['neutral'])
    return random.choice(styles)


# 自动选择时排除的装饰性/过小气泡——这些形状包不住全部对白文字
_DECORATIVE_BUBBLE_SHAPES = {'star', 'heart', 'burst'}


def get_random_bubble_image(emotion, allow_decorative=False):
    """根据情绪获取随机气泡图片路径。
    allow_decorative=False（默认）时，自动选择会排除星星/爱心/爆炸等包不住文字的小气泡；
    用户在前端手动指定星星/爱心时不受影响（手动走 bubble_url 直传）。"""
    base_dir = os.path.join(app.root_path, 'static', 'bubbles', emotion)

    if not os.path.exists(base_dir):
        base_dir = os.path.join(app.root_path, 'static', 'bubbles', 'neutral')
        if not os.path.exists(base_dir):
            return None

    try:
        all_files = [f for f in os.listdir(base_dir) if f.lower().endswith('.png')]
        if not all_files:
            return None

        if not allow_decorative:
            # 文件名形如 happy_03_star.png，形状在倒数第一段
            candidates = []
            for f in all_files:
                shape = f[:-4].split('_')[-1].lower() if f.lower().endswith('.png') else ''
                if shape not in _DECORATIVE_BUBBLE_SHAPES:
                    candidates.append(f)
            # 兜底：若某情绪目录下过滤后无图，则放开装饰性，避免无图
            if candidates:
                all_files = candidates

        random_file = random.choice(all_files)
        return os.path.join(base_dir, random_file)
    except Exception as e:
        logger.error(f"获取气泡图片失败: {str(e)}")
        return None


def get_font_path(font_family):
    font_info = FONT_FAMILIES.get(font_family)
    if font_info and os.path.exists(font_info['path']):
        return font_info['path']

    # Railway 使用 Linux 容器，不能使用 Windows 随系统授权的微软雅黑文件。
    # Docker 已安装 fonts-noto-cjk：对于“微软雅黑加粗”必须优先取 Noto 的
    # Bold 字重，避免原先回退到 Regular 后公网成品看起来没有加粗。
    bold_noto_paths = [
        '/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc',
        '/usr/share/fonts/truetype/noto/NotoSansCJK-Bold.ttc',
        '/usr/share/fonts/noto-cjk/NotoSansCJK-Bold.ttc',
    ]
    regular_noto_paths = [
        '/usr/share/fonts/truetype/noto/NotoSansCJK-SC.ttc',
        '/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc',
        '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc',
        '/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc',
    ]
    fallback_paths = (
        bold_noto_paths + regular_noto_paths
        if font_family == 'msyhbd'
        else [
        'C:/Windows/Fonts/msyh.ttc',
        'C:/Windows/Fonts/msyhbd.ttc',
        '/System/Library/Fonts/STHeiti Light.ttc',
        '/usr/share/fonts/truetype/wqy/wqy-microhei.ttc',
        '/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc',
        *regular_noto_paths,
        *bold_noto_paths,
    ]
    )
    
    for path in fallback_paths:
        if os.path.exists(path):
            return path
    return None


def calculate_character_positions(ip_images, bg_width, bg_height, speaker_index=-1, is_guide=False,
                                  character_names=None):
    if not ip_images:
        return [], [], []
    
    num_ips = len(ip_images)
    
    scaled_ips = []
    ip_positions = []
    ip_transform_params = []
    
    for idx, ip_img in enumerate(ip_images):
        is_speaker = idx == speaker_index

        # 自动识别格子尺寸：同时按宽度和高度约束人物最大尺寸，
        # 避免在矮/宽格子中人物顶天立地或头被裁掉
        # Bug2 修复：≥3 人时显著缩小人物高度占比，给头顶气泡留足空间
        height_caps = {1: 0.75, 2: 0.70, 3: 0.58, 4: 0.40}
        height_cap = height_caps.get(num_ips, 0.48)

        if num_ips == 1:
            width_size = int(bg_width * 0.55)
        elif num_ips == 2:
            width_size = int(bg_width * 0.48) if is_speaker else int(bg_width * 0.38)
        elif num_ips == 3:
            # Bug2 修复：显著缩小人物并加大间距，使三人明显分散、给头顶气泡留足空间
            gap = 35
            width_size = max(1, min(int(bg_width * 0.28), int(bg_height * 0.58)))
        elif num_ips == 4:
            # Bug2 修复：缩小人物，配合下方两排布局形成清晰分散
            gap = 30
            width_size = max(1, min(int(bg_width * 0.20), int(bg_height * 0.48)))
        else:
            width_size = int(bg_width * 0.34) if is_speaker else int(bg_width * 0.27)

        height_size = int(bg_height * height_cap)
        ip_size = min(width_size, height_size)
        # 老师是成年人，不能和小学生/小精灵使用同一个视觉身高。
        # 名称来自本格真正参与合成的角色，避免仅依赖图片文件名而误判。
        character_name = ''
        if isinstance(character_names, (list, tuple)) and idx < len(character_names):
            character_name = str(character_names[idx] or '')
        if character_name.endswith('老师'):
            ip_size = int(ip_size * 1.24)

        ip_width, ip_height = ip_img.size
        ratio = min(ip_size / ip_width, ip_size / ip_height)
        new_size = (int(ip_width * ratio), int(ip_height * ratio))
        
        # flip 先占位，位置排好后再按“朝向对方/画面中心”统一决定（见函数末尾）
        rotate_angle = 0
        scale_factor = 1.0
        offset_y = 0
        
        scaled_ips.append(ip_img)
        ip_transform_params.append({
            'size': new_size, 
            'flip': False, 
            'rotate': rotate_angle, 
            'scale': scale_factor, 
            'offset_y': offset_y,
            'is_speaker': is_speaker
        })
    
    if num_ips == 1:
        ip_x = (bg_width - ip_transform_params[0]['size'][0]) // 2
        ip_y = bg_height - ip_transform_params[0]['size'][1] - 15
        ip_positions.append((ip_x, ip_y))
    elif num_ips == 2:
        total_width = sum(params['size'][0] for params in ip_transform_params) + 25
        start_x = (bg_width - total_width) // 2
        current_x = start_x

        for params in ip_transform_params:
            ip_w, ip_h = params['size']
            # 成人老师放大后，仍要完整落在画面内，不能因为两人居中计算而被裁切。
            ip_x = max(10, min(current_x, bg_width - ip_w - 10))
            ip_y = bg_height - ip_h - 15
            ip_positions.append((ip_x, ip_y))
            current_x += ip_w + 25
    elif num_ips == 3:
        # 三人横向排列：中间角色后排抬高，两侧前排下沉，
        # 拉开间距并避免气泡互相遮挡人脸
        gap = 35
        total_width = sum(params['size'][0] for params in ip_transform_params) + gap * 2
        start_x = max(10, (bg_width - total_width) // 2)
        center_idx = 1
        back_offset = int(bg_height * 0.12)

        # 三套构图随机轮换，避免每格人物都像同一模板复制。
        variant = random.randrange(3)
        x_ratios = ([0.08, 0.40, 0.72], [0.16, 0.48, 0.75], [0.05, 0.34, 0.67])[variant]
        y_lifts = ([0.00, 0.12, 0.02], [0.10, 0.00, 0.13], [0.03, 0.15, 0.00])[variant]
        for idx, params in enumerate(ip_transform_params):
            ip_w, ip_h = params['size']
            ip_x = int(bg_width * x_ratios[idx])
            ip_x = max(10, min(ip_x, bg_width - ip_w - 10))
            ip_y = bg_height - ip_h - 15 - int(bg_height * y_lifts[idx])
            ip_positions.append((ip_x, ip_y))
    elif num_ips == 4:
        gap = 30
        row1_ips = ip_transform_params[:2]
        row2_ips = ip_transform_params[2:]

        row1_width = sum(p['size'][0] for p in row1_ips) + 30
        row2_width = sum(p['size'][0] for p in row2_ips) + 30

        start_x1 = (bg_width - row1_width) // 2
        start_x2 = (bg_width - row2_width) // 2

        # Bug2 修复：真正分两排——上排整体抬高约一个人物高度，下排贴底，
        # 使两层在垂直方向明显分开（不再只差 35px 导致重叠）。
        for r, (row_ips, start_x) in enumerate([(row1_ips, start_x1), (row2_ips, start_x2)]):
            current_x = start_x
            for params in row_ips:
                ip_w, ip_h = params['size']
                ip_x = current_x
                if r == 0:
                    # Bug3 修复：上排抬高约一个人物高度 + gap(30) 间距，
                    # 使上排底 < 下排顶（固定间隔 30），两排不重叠、不顶到天花板
                    ip_y = bg_height - 2 * ip_h - gap
                else:
                    ip_y = bg_height - ip_h - 15       # 下排：贴底
                ip_positions.append((ip_x, ip_y))
                current_x += ip_w + 30
    else:
        rows = 2
        cols = (num_ips + 1) // 2

        for row_idx in range(rows):
            row_params = ip_transform_params[row_idx * cols:(row_idx + 1) * cols]
            row_width = sum(p['size'][0] for p in row_params) + (len(row_params) - 1) * 15
            start_x = (bg_width - row_width) // 2
            current_x = start_x

            for params in row_params:
                ip_w, ip_h = params['size']
                ip_x = current_x
                if row_idx == 0:
                    ip_y = bg_height - ip_h - 35
                else:
                    ip_y = bg_height - ip_h - 10
                ip_positions.append((ip_x, ip_y))
                current_x += ip_w + 15

    # 所有人物脚底必须落在画面下方的同一“地面带”。
    # 旧版用最高 15% 画高的纵向错位制造变化，会让人物看起来站上桌子、柜台或橱窗。
    # 现在只保留左右位置、朝向和轻微大小变化，脚底高度最多相差画高的 2%。
    if ip_positions:
        ground_y = bg_height - max(8, int(bg_height * 0.025))
        grounded_positions = []
        for idx, ((ip_x, _), params) in enumerate(zip(ip_positions, ip_transform_params)):
            _, ip_h = params['size']
            # 极小错位用于避免机械齐线，但仍处于同一地面带。
            foot_jitter = int(bg_height * random.uniform(0.0, 0.02))
            ip_y = max(0, ground_y - ip_h - foot_jitter)
            grounded_positions.append((ip_x, ip_y))
        ip_positions = grounded_positions

    # ---- 人物重叠避让：3+ 人场景下若边界框重叠（含 margin），做最小平移微调 ----
    if num_ips >= 3:
        ip_positions = _resolve_character_overlaps(
            ip_positions,
            [p['size'] for p in ip_transform_params],
            bg_width, bg_height, margin=35
        )

    # ---- 智能朝向：让人物面向对方/说话者，制造对话互动感 ----
    # 约定：姿势图默认朝右（flip=False 朝右，flip=True 朝左）。
    # 单人不翻转；多人时每个人朝向“目标”：优先朝向说话者，否则朝向画面中心。
    if num_ips >= 2 and len(ip_positions) == num_ips:
        centers = []
        for idx in range(num_ips):
            px, _py = ip_positions[idx]
            pw = ip_transform_params[idx]['size'][0]
            centers.append(px + pw / 2.0)

        # 目标 x：有说话者则所有人朝向说话者，说话者本人朝向离它最远的听众（面向听众群）
        speaker_center = centers[speaker_index] if 0 <= speaker_index < num_ips else None
        for idx in range(num_ips):
            my_c = centers[idx]
            if speaker_center is not None and idx != speaker_index:
                target_c = speaker_center
            elif speaker_center is not None and idx == speaker_index:
                # 说话者朝向听众重心（其余人的平均中心）
                others = [c for j, c in enumerate(centers) if j != idx]
                target_c = sum(others) / len(others) if others else bg_width / 2.0
            else:
                target_c = bg_width / 2.0  # 无说话者：一律朝画面中心
            # 目标在右侧→朝右(flip=False)；目标在左侧→朝左(flip=True)
            ip_transform_params[idx]['flip'] = (target_c < my_c)

    return scaled_ips, ip_positions, ip_transform_params


def _resolve_character_overlaps(positions, sizes, bg_width, bg_height, margin=25, edge=10):
    """确保人物边界框之间不重叠（含 margin）。
    1~2 人基本保持原样；3+ 人在检测到重叠/过近时，沿水平方向做最小平移微调，
    并约束在画布内。返回调整后的 positions（list of (x, y)）。
    margin 取 20~30px（人物之间间距）。
    """
    n = len(positions)
    if n <= 2:
        return positions
    boxes = [[positions[i][0], positions[i][1], sizes[i][0], sizes[i][1]] for i in range(n)]
    for _ in range(30):
        moved = False
        for i in range(n):
            for j in range(i + 1, n):
                xi, yi, wi, hi = boxes[i]
                xj, yj, wj, hj = boxes[j]
                sep_x = min(xi + wi, xj + wj) - max(xi, xj)
                sep_y = min(yi + hi, yj + hj) - max(yi, yj)
                # 仅当两个方向都过近（即处于同一排）时才做水平推开
                if sep_x > 0 and sep_y > 0:
                    need = sep_x
                    ci = xi + wi / 2.0
                    cj = xj + wj / 2.0
                    if ci <= cj:
                        dx, dxj = -need / 2.0, need / 2.0
                    else:
                        dx, dxj = need / 2.0, -need / 2.0
                    # 若纵向也真正重叠（不同排错位），同时做垂直微调
                    if sep_y > 0:
                        need_y = sep_y
                        if yi <= yj:
                            dy, dyj = -need_y / 2.0, need_y / 2.0
                        else:
                            dy, dyj = need_y / 2.0, -need_y / 2.0
                        boxes[i][1] += dy
                        boxes[j][1] += dyj
                    boxes[i][0] += dx
                    boxes[j][0] += dxj
                    moved = True
        if not moved:
            break
    # 约束在画布内
    for i in range(n):
        w, h = boxes[i][2], boxes[i][3]
        boxes[i][0] = max(edge, min(boxes[i][0], bg_width - w - edge))
        boxes[i][1] = max(edge, min(boxes[i][1], bg_height - h - edge))
    return [(int(round(b[0])), int(round(b[1]))) for b in boxes]


def apply_character_lighting(img, strength=1.0):
    """给 RGBA 人物图打光：叠加“顶亮底暗”的环境光渐变 + 轻微对比增强，
    让 3D 人物有统一的明暗关系，减少浮于表面的贴纸感。
    光效只作用于不透明像素（用 alpha 通道做遮罩）。返回新的 RGBA 图。
    strength: 光照强度系数（1.0 标准，可调）。
    """
    try:
        from PIL import ImageChops, ImageEnhance
        img = img.convert('RGBA')
        w, h = img.size
        if w < 4 or h < 4:
            return img
        alpha = img.split()[3]

        # 1) 垂直渐变：顶部偏亮(>128)，底部偏暗(<128)，128 为中性
        grad = Image.new('L', (1, h))
        for y in range(h):
            t = y / max(1, h - 1)
            grad.putpixel((0, y), int(158 - t * 96))  # 顶≈158 → 底≈62
        grad = grad.resize((w, h))

        # 2) 拆成高光层(白)与阴影层(黑)，强度随渐变偏离中性程度变化
        hi = grad.point(lambda p: int(max(0, p - 128) * 0.42 * strength))
        lo = grad.point(lambda p: int(max(0, 128 - p) * 0.52 * strength))
        # 只在人物范围内生效
        hi = ImageChops.multiply(hi, alpha)
        lo = ImageChops.multiply(lo, alpha)

        light = Image.new('RGBA', (w, h), (255, 255, 255, 0)); light.putalpha(hi)
        dark  = Image.new('RGBA', (w, h), (0, 0, 0, 0));       dark.putalpha(lo)

        out = img.copy()
        out = Image.alpha_composite(out, dark)
        out = Image.alpha_composite(out, light)

        # 3) 轻微提升对比与饱和，让人物更“实”、更贴合场景
        rgb = out.convert('RGB')
        rgb = ImageEnhance.Contrast(rgb).enhance(1.06)
        rgb = ImageEnhance.Color(rgb).enhance(1.05)
        out = Image.merge('RGBA', (*rgb.split(), out.split()[3]))
        return out
    except Exception as e:
        logger.warning(f"人物打光失败，返回原图: {e}")
        return img


def bubble_native_tail_kind(bubble_image_path):
    """识别素材库中已经画好尾巴的气泡，避免再叠加程序三角尾巴。"""
    if not bubble_image_path:
        return None
    name = os.path.basename(str(bubble_image_path)).lower()
    if '_02_cloud' in name:
        return 'bottom_center'
    if '_07_diamond' in name:
        return 'bottom_left'
    return None


def bubble_content_scale(bubble_image_path):
    """返回素材气泡为容纳同量文字所需的额外外框缩放比例。

    07 思考气泡的圆点尾巴占据整张 PNG 下方较大区域，真正可写字的主体
    明显小于普通圆角气泡。若仍按 PNG 外接矩形计算，文字会贴边或被迫缩小。
    """
    if not bubble_image_path:
        return 1.0, 1.0
    name = os.path.basename(str(bubble_image_path)).lower()
    if '_07_diamond' in name:
        return 1.24, 1.28
    return 1.0, 1.0


def render_bubble(draw, bubble_x, bubble_y, bubble_width, bubble_height, style='rounded', 
                  fill_color=(255, 255, 255, 245), outline_color=(255, 179, 198), 
                  outline_width=2, opacity=0.96, tail_pos=None, bubble_image=None):
    alpha = int(opacity * 255)
    fill_color = (fill_color[0], fill_color[1], fill_color[2], min(255, int(fill_color[3] * opacity)))
    
    if bubble_image is not None:
        bubble_img = Image.open(bubble_image).convert('RGBA')
        bubble_img = bubble_img.resize((int(bubble_width), int(bubble_height)), Image.LANCZOS)
        
        alpha_channel = bubble_img.split()[3]
        alpha_channel = alpha_channel.point(lambda p: int(p * opacity))
        bubble_img.putalpha(alpha_channel)
        
        # 素材自带尾巴时直接使用原图，绝不能再叠加程序生成的三角尾巴。
        # 没有原生尾巴的素材才补画一体式尾巴。
        if tail_pos and not bubble_native_tail_kind(bubble_image):
            if len(tail_pos) >= 4:
                base_x, base_y, target_x, target_y = tail_pos[:4]
                side = tail_pos[4] if len(tail_pos) > 4 else ('bottom' if target_y >= bubble_y + bubble_height / 2 else 'top')
                tail_fill = bubble_img.getpixel((bubble_img.width // 2, bubble_img.height // 2))
                if tail_fill[3] < 40:
                    tail_fill = fill_color
                opaque_pixels = [px for px in bubble_img.getdata() if px[3] > 100]
                tail_outline = min(opaque_pixels, key=lambda px: px[0] + px[1] + px[2]) if opaque_pixels else outline_color
                if side in ('left', 'right'):
                    base_x = int(bubble_x + 18) if side == 'left' else int(bubble_x + bubble_width - 18)
                    base_y = int(max(bubble_y + 22, min(base_y, bubble_y + bubble_height - 22)))
                    tip_x = (
                        int(max(bubble_x - 30, target_x))
                        if side == 'left'
                        else int(min(bubble_x + bubble_width + 30, target_x))
                    )
                    tip_y = base_y + max(-26, min(26, int(target_y - base_y)))
                    points = [(base_x, base_y - 14), (base_x, base_y + 14), (tip_x, tip_y)]
                else:
                    base_x = int(max(bubble_x + 22, min(base_x, bubble_x + bubble_width - 22)))
                    base_y = int(bubble_y + 18) if side == 'top' else int(bubble_y + bubble_height - 18)
                    tip_x = base_x + max(-30, min(30, int(target_x - base_x)))
                    tip_y = (
                        int(max(bubble_y - 30, target_y))
                        if side == 'top'
                        else int(min(bubble_y + bubble_height + 30, target_y))
                    )
                    points = [(base_x - 14, base_y), (base_x + 14, base_y), (tip_x, tip_y)]
                draw.polygon(points, fill=tail_fill, outline=tail_outline, width=max(2, outline_width))
            else:
                tail_x, tail_y, tail_direction = tail_pos
                tip_y = tail_y - 18 if tail_direction == 'up' else tail_y + 18
                draw.polygon([(tail_x - 10, tail_y), (tail_x + 10, tail_y), (tail_x, tip_y)],
                             fill=fill_color, outline=outline_color, width=max(2, outline_width))

        draw._image.paste(bubble_img, (int(bubble_x), int(bubble_y)), bubble_img)
        return
    
    bubble_config = BUBBLE_STYLES.get(style, BUBBLE_STYLES['rounded'])
    
    if bubble_config['shape'] == 'rounded_rectangle':
        radius = bubble_config['radius']
        draw.rounded_rectangle(
            [bubble_x, bubble_y, bubble_x + bubble_width, bubble_y + bubble_height],
            radius=radius,
            fill=fill_color,
            outline=outline_color,
            width=outline_width
        )
    elif bubble_config['shape'] == 'ellipse':
        draw.ellipse(
            [bubble_x, bubble_y, bubble_x + bubble_width, bubble_y + bubble_height],
            fill=fill_color,
            outline=outline_color,
            width=outline_width
        )
    elif bubble_config['shape'] == 'sharp_rectangle':
        radius = bubble_config['radius']
        draw.rounded_rectangle(
            [bubble_x, bubble_y, bubble_x + bubble_width, bubble_y + bubble_height],
            radius=radius,
            fill=fill_color,
            outline=outline_color,
            width=outline_width
        )
    elif bubble_config['shape'] == 'cloud':
        cloud_radius = bubble_height // 2
        center_x = bubble_x + bubble_width // 2
        center_y = bubble_y + bubble_height // 2
        
        draw.ellipse([bubble_x + cloud_radius, bubble_y, bubble_x + bubble_width - cloud_radius, bubble_y + bubble_height], fill=fill_color)
        draw.ellipse([bubble_x, bubble_y + cloud_radius, bubble_x + bubble_width, bubble_y + bubble_height - cloud_radius], fill=fill_color)
        
        draw.ellipse([center_x - cloud_radius, bubble_y - cloud_radius // 2, center_x + cloud_radius, bubble_y + cloud_radius // 2], fill=fill_color)
        draw.ellipse([bubble_x - cloud_radius // 2, center_y - cloud_radius, bubble_x + cloud_radius // 2, center_y + cloud_radius], fill=fill_color)
        draw.ellipse([bubble_x + bubble_width - cloud_radius // 2, center_y - cloud_radius, bubble_x + bubble_width + cloud_radius // 2, center_y + cloud_radius], fill=fill_color)
        
        draw.ellipse([bubble_x + cloud_radius, bubble_y, bubble_x + bubble_width - cloud_radius, bubble_y + bubble_height], outline=outline_color, width=outline_width)
        draw.ellipse([bubble_x, bubble_y + cloud_radius, bubble_x + bubble_width, bubble_y + bubble_height - cloud_radius], outline=outline_color, width=outline_width)
        draw.ellipse([center_x - cloud_radius, bubble_y - cloud_radius // 2, center_x + cloud_radius, bubble_y + cloud_radius // 2], outline=outline_color, width=outline_width)
    elif bubble_config['shape'] == 'circle':
        min_dim = min(bubble_width, bubble_height)
        center_x = bubble_x + bubble_width // 2
        center_y = bubble_y + bubble_height // 2
        radius = min_dim // 2
        draw.ellipse([center_x - radius, center_y - radius, center_x + radius, center_y + radius], fill=fill_color, outline=outline_color, width=outline_width)
    elif bubble_config['shape'] == 'diamond':
        cx, cy = bubble_x + bubble_width // 2, bubble_y + bubble_height // 2
        points = [
            (cx, bubble_y),
            (bubble_x + bubble_width, cy),
            (cx, bubble_y + bubble_height),
            (bubble_x, cy)
        ]
        draw.polygon(points, fill=fill_color, outline=outline_color)
        draw.polygon(points, outline=outline_color, width=outline_width)
    elif bubble_config['shape'] == 'burst':
        spikes = bubble_config['radius']
        cx, cy = bubble_x + bubble_width // 2, bubble_y + bubble_height // 2
        rx, ry = bubble_width // 2, bubble_height // 2
        
        points = []
        for i in range(spikes * 2):
            angle = math.pi * i / spikes
            if i % 2 == 0:
                r = rx
            else:
                r = rx * 0.6
            px = cx + r * math.cos(angle)
            py = cy + ry * math.sin(angle)
            points.append((px, py))
        
        draw.polygon(points, fill=fill_color, outline=outline_color)
        draw.polygon(points, outline=outline_color, width=outline_width)
    elif bubble_config['shape'] == 'wave':
        segments = 8
        points = []
        for i in range(segments + 1):
            x = bubble_x + (bubble_width / segments) * i
            if i % 2 == 0:
                y = bubble_y
            else:
                y = bubble_y + 6
            points.append((x, y))
        
        for i in range(segments + 1):
            x = bubble_x + (bubble_width / segments) * i
            if i % 2 == 0:
                y = bubble_y + bubble_height
            else:
                y = bubble_y + bubble_height - 6
            points.append((x, y))
        
        draw.polygon(points, fill=fill_color, outline=outline_color)
        draw.polygon(points, outline=outline_color, width=outline_width)
    else:
        draw.rounded_rectangle(
            [bubble_x, bubble_y, bubble_x + bubble_width, bubble_y + bubble_height],
            radius=8,
            fill=fill_color,
            outline=outline_color,
            width=outline_width
        )
    
    if tail_pos:
        if len(tail_pos) >= 4:
            base_x, base_y, target_x, target_y = tail_pos[:4]
            side = tail_pos[4] if len(tail_pos) > 4 else ('bottom' if target_y >= base_y else 'top')
            if side in ('left', 'right'):
                tip_x = (
                    max(base_x - 30, target_x)
                    if side == 'left'
                    else min(base_x + 30, target_x)
                )
                tip_y = base_y + max(-26, min(26, target_y - base_y))
                points = [(base_x, base_y - 13), (base_x, base_y + 13), (tip_x, tip_y)]
            else:
                tip_x = base_x + max(-28, min(28, target_x - base_x))
                tip_y = (
                    max(base_y - 30, target_y)
                    if side == 'top'
                    else min(base_y + 30, target_y)
                )
                points = [(base_x - 13, base_y), (base_x + 13, base_y), (tip_x, tip_y)]
            draw.polygon(points, fill=fill_color, outline=outline_color, width=max(2, outline_width))
        else:
            tail_x, tail_y, tail_direction = tail_pos
            tip_y = tail_y - 12 if tail_direction == 'up' else tail_y + 12
            draw.polygon([(tail_x - 8, tail_y), (tail_x + 8, tail_y), (tail_x, tip_y)], fill=fill_color, outline=outline_color)


def bubble_text_colors(bubble_image_path=None, fill_color=None):
    """根据气泡底色自动选择文字和描边颜色，确保可读性。
    深色气泡 -> 白字（无描边）；浅色气泡 -> 黑字（无描边）。
    注：原为浅色气泡返回白色描边，会在文字外围形成一圈白边，编辑器无法单独编辑，
    现已移除描边（返回 None），由字符本体颜色保证可读性。
    """
    try:
        if bubble_image_path:
            img = Image.open(bubble_image_path).convert('RGBA')
            r, g, b, a = img.split()
            mask = a.point(lambda p: 255 if p > 10 else 0)
            bbox = mask.getbbox()
            if bbox:
                rgb = Image.merge('RGB', (r, g, b)).crop(bbox)
                m = mask.crop(bbox)
                total = count = 0
                for (pr, pg, pb), mv in zip(rgb.getdata(), m.getdata()):
                    if mv > 10:
                        total += 0.299 * pr + 0.587 * pg + 0.114 * pb
                        count += 1
                if count:
                    avg_lum = total / count
                    if avg_lum < 140:
                        return (255, 255, 255), None
        elif fill_color is not None:
            rgb = fill_color[:3]
            lum = 0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]
            if lum < 140:
                return (255, 255, 255), None
    except Exception as e:
        logger.debug(f"气泡底色亮度分析失败: {e}")
    return (50, 50, 50), None


def render_text(draw, text, bubble_x, bubble_y, bubble_width, bubble_height, font,
                line_height, padding=12, align='center', text_color=(50, 50, 50), font_path=None,
                stroke_fill=None, content_bottom_inset=0):
    text = text.replace('\r', '').strip()
    if not text:
        return None

    max_width = bubble_width - padding * 2
    # 原生尾巴气泡的尾巴不属于正文区域。为它预留空间，避免文字被下方尾巴挤出主体。
    content_bottom_inset = max(0, int(content_bottom_inset or 0))
    max_height = bubble_height - padding * 2 - content_bottom_inset

    def try_layout(current_font, current_size):
        """尝试用指定字号排版，返回 (是否放得下, lines, actual_line_height)"""
        lines = _wrap_dialogue_lines(draw, text, current_font, max_width, 11)
        lh = current_size + 4
        total_h = len(lines) * lh
        return total_h <= max_height, lines, lh, current_font

    # 优先使用传入字号，若放不下则逐步缩小
    current_font = font
    current_size = getattr(font, 'size', 18)
    fits, lines, line_height, current_font = try_layout(current_font, current_size)

    if font_path and not fits:
        for size in [current_size - 2, current_size - 4, current_size - 6, 16, 14]:
            if size < 14:
                break
            try:
                current_font = ImageFont.truetype(font_path, size)
            except Exception:
                current_font = font
            fits, lines, line_height, current_font = try_layout(current_font, size)
            if fits:
                current_size = size
                break

    if not lines:
        # 兜底：强行按原字号一行显示
        lines = [text]

    # 垂直居中
    total_text_height = len(lines) * line_height
    text_y = bubble_y + padding + (max_height - total_text_height) // 2
    if text_y < bubble_y + padding:
        text_y = bubble_y + padding

    stroke_width = max(1, int(current_size / 12))

    for line in lines:
        line_width = draw.textlength(line, font=current_font)

        if align == 'left':
            text_x = bubble_x + padding
        elif align == 'right':
            text_x = bubble_x + bubble_width - padding - line_width
        else:
            text_x = bubble_x + (bubble_width - line_width) // 2

        if stroke_fill is not None:
            draw.text((text_x, text_y), line, font=current_font, fill=text_color,
                      stroke_width=stroke_width, stroke_fill=stroke_fill)
        else:
            draw.text((text_x, text_y), line, font=current_font, fill=text_color)
        text_y += line_height

    # 编辑器需要复用最终实际生效的排版值；长文本可能在这里被缩小，不能只传
    # 气泡初始计算时的字号，否则编辑画面会比预览更大、更细。
    return {
        'font_size': current_size,
        'line_height': line_height,
        'stroke_width': stroke_width if stroke_fill is not None else 0,
        # 把 Pillow 已验证的最终断行一并保存。编辑器直接复用它，就不会把中文
        # 标点、书名号和粗体字重新按另一套浏览器规则拆开。
        'rendered_text': '\n'.join(lines),
    }


def render_sound_effect(draw, text, x, y, bg_width, bg_height, font_path):
    """渲染音效文字（拟声词），如"呼呼"、"滴滴"、"啊啊"等"""
    if not text or not isinstance(text, str):
        return
    
    text = text.strip()
    if not text:
        return
    
    try:
        font_size = min(40, max(24, int(len(text) * -4 + 48)))
        font = ImageFont.truetype(font_path, font_size) if font_path else ImageFont.load_default()
        
        text_bbox = draw.textbbox((0, 0), text, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        text_height = text_bbox[3] - text_bbox[1]
        
        bubble_width = text_width + 20
        bubble_height = text_height + 16
        
        bubble_x = x - bubble_width // 2
        bubble_y = y - bubble_height - 10
        
        safe_m = max(15, int(bg_width * 0.06))
        bubble_x = max(safe_m, min(bubble_x, bg_width - bubble_width - safe_m))
        bubble_y = max(safe_m, min(bubble_y, bg_height - bubble_height - safe_m))
        
        spike_count = 8
        cx, cy = bubble_x + bubble_width // 2, bubble_y + bubble_height // 2
        rx, ry = bubble_width // 2, bubble_height // 2
        
        points = []
        for i in range(spike_count * 2):
            angle = math.pi * i / spike_count
            if i % 2 == 0:
                r = rx * 1.2
            else:
                r = rx * 0.7
            px = cx + r * math.cos(angle)
            py = cy + ry * math.sin(angle)
            points.append((px, py))
        
        draw.polygon(points, fill=(255, 255, 255, 240), outline=(255, 100, 150), width=3)
        
        glow_layer = Image.new('RGBA', (bg_width, bg_height), (0, 0, 0, 0))
        glow_draw = ImageDraw.Draw(glow_layer)
        for r in range(30, 0, -5):
            alpha_val = max(0, int(20 * (r / 30)))
            glow_draw.polygon(points, fill=(255, 200, 220, alpha_val))
        
        text_x = cx - text_width // 2
        text_y = cy - text_height // 2
        
        draw.text((text_x, text_y), text, font=font, fill=(255, 50, 100))
        
    except Exception as e:
        logger.warning(f"渲染音效文字失败: {str(e)}")


def _bubble_path_to_url(path):
    """把气泡图片的绝对路径转成 /static/bubbles/... URL（用于编辑面板加载）"""
    if not path:
        return None
    try:
        static_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static')
        rel = os.path.relpath(path, static_dir).replace('\\', '/')
        return f'/static/{rel}'
    except Exception:
        return path


def _download_bubble_to_local(url):
    """把远程气泡图片 url 下载到本地临时文件，返回本地绝对路径；失败返回 None。

    用于「按人物指定气泡」中用户直接传入 http(s) 气泡图片 url 的场景。
    本地 /static/... 路径无需下载，由调用方直接转绝对路径。
    """
    try:
        import hashlib
        import urllib.request
        ext = '.png'
        low = url.lower()
        if low.endswith('.jpg') or low.endswith('.jpeg'):
            ext = '.jpg'
        digest = hashlib.md5(url.encode('utf-8')).hexdigest()[:12]
        out_dir = os.path.join(app.root_path, 'static', 'bubbles', '_tmp')
        os.makedirs(out_dir, exist_ok=True)
        out_path = os.path.join(out_dir, f'remote_{digest}{ext}')
        if os.path.exists(out_path):
            return out_path
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=15) as resp:
            data = resp.read()
        with open(out_path, 'wb') as f:
            f.write(data)
        return out_path
    except Exception as e:
        logger.warning(f"下载远程气泡图片失败: {e}")
        return None


# ===== 角色姿势库 =====
# 每个角色在第二步用「相同 seed + 相同外貌描述、只改姿势词」生成以下姿势变体，
# 同 seed 可让脸型/配色保持基本一致（免费 Pollinations 下最稳的做法）。
POSE_LIBRARY = {
    'stand': '自然站立',
    'sit': '坐在地上',
    'laugh': '开心大笑、嘴巴张开',
    'smile': '神秘温柔地微笑',
    'raise_hand': '兴奋地举起一只手',
    'nod': '点头表示赞同',
    'scratch_head': '困惑地挠头',
    'chin_rest': '手托下巴思考',
    'point': '用手指着某处',
}

# 给 AI Horde（SD 系列模型）用的英文姿势描述，因为 SD 对英文远强于中文。
# horde 模式 = Pollinations 出标准角色(中文友好) + horde img2img 换姿势(英文姿势词)。
POSE_LIBRARY_EN = {
    'stand': 'standing naturally, full body',
    'sit': 'sitting on the ground, full body',
    'laugh': 'laughing happily with mouth open, full body',
    'smile': 'gentle mysterious smile, full body',
    'raise_hand': 'raising one hand excitedly, full body',
    'nod': 'nodding in agreement, full body',
    'scratch_head': 'scratching head in confusion, full body',
    'chin_rest': 'resting chin on hand, thinking, full body',
    'point': 'pointing at something with a finger, full body',
}

# ===== 新版角色专属姿势库（高远/许多/莉莉/知识小精灵）=====
# 每个角色有自己的一组姿势 key，与旧版 POSE_LIBRARY（stand/sit/laugh/...）共存。
# 旧版老师（哈哈/张驰/大唐/徐剑/可可/布布）继续使用 POSE_LIBRARY 的 key。

GAOYUAN_POSES = {
    'gather': '集合引导式',
    'reveal': '神秘揭秘式',
    'guide_knowledge': '知识点指引式',
    'gentle_remind': '温和叮嘱式',
    'ask_question': '课堂提问式',
    'preview_method': '方法预告式',
}

XUDUO_POSES = {
    'curious_lean': '好奇探身式',
    'raise_answer': '举手抢答式',
    'shock_back': '震惊后仰式',
    'scratch_think': '挠头思索式',
    'excited_fist': '兴奋握拳式',
    'point_laugh': '指物笑弯式',
}

LILI_POSES = {
    'book_lookup': '持本抬望式',
    'chin_think': '托腮沉思式',
    'point_analyze': '点指分析式',
    'cover_surprise': '捂嘴惊讶式',
    'bow_write': '低头书写式',
    'smile_confirm': '浅笑印证式',
}

SPRITE_POSES = {
    'float_ask': '悬浮提问式',
    'wave_interact': '挥手互动式',
    'static_float': '静立漂浮式',
}

# 角色目录 → 该角色的姿势 key 列表
CHARACTER_POSE_KEYS = {
    'gaoyuan': list(GAOYUAN_POSES.keys()),
    'xuduo': list(XUDUO_POSES.keys()),
    'lili': list(LILI_POSES.keys()),
    'sprite': list(SPRITE_POSES.keys()),
}

# 角色目录 → 默认姿势 key（无动作标记时使用）
CHARACTER_DEFAULT_POSE = {
    'gaoyuan': 'ask_question',
    'xuduo': 'curious_lean',
    'lili': 'book_lookup',
    'sprite': 'static_float',
}

# 角色目录 → 动作关键词 → 姿势 key 映射
CHARACTER_ACTION_MAP = {
    'gaoyuan': [
        ('gather', ['举手', '集合', '引导', '招呼', '过来', '招手', '集合啦', '过来这边', '聚拢']),
        ('reveal', ['揭秘', '神秘', '展示', '惊喜', '看这个', '揭秘时刻', '看这里', '揭晓', '放出大招']),
        ('guide_knowledge', ['指向', '指引', '知识点', '指', '这边', '指向知识点', '往这看', '这里很重要']),
        ('gentle_remind', ['叮嘱', '温和', '认真', '放松', '记住', '别忘了', '叮嘱一句', '温柔提醒', '慢点']),
        ('ask_question', ['提问', '课堂', '问', '考考', '想想', '谁来答', '考考你', '想一想', '提问时间']),
        ('preview_method', ['预告', '方法', '三步', '介绍', '接下来', '接下来这样', '方法三步', '预告一下', '看步骤']),
    ],
    'xuduo': [
        ('curious_lean', ['好奇', '探身', '前倾', '探', '凑近', '凑近看', '探着身子', '好奇脸', '往前倾']),
        ('raise_answer', ['举手', '抢答', '回答', '抢', '我知道', '我会我会', '抢着说', '举手回答', '蹦着举手']),
        ('shock_back', ['震惊', '惊讶', '后仰', '难以置信', '不会吧', '吓一跳', '目瞪口呆', '大跌眼镜', '不敢相信']),
        ('scratch_think', ['挠头', '思考', '思索', '困惑', '挠', '想不通', '挠头想', '犯愁', '皱眉思索']),
        ('excited_fist', ['兴奋', '握拳', '开心', '成就', '太好了', '耶', '太棒了', '开心到跳', '成就感']),
        ('point_laugh', ['指物', '笑弯', '弯腰', '乐', '哈哈', '指', '指着他笑', '笑得弯下腰', '乐呵呵指']),
    ],
    'lili': [
        ('book_lookup', ['抬头', '望', '期待', '注视', '看着', '抬头看', '满怀期待', '望向远方', '注视远']),
        ('chin_think', ['托腮', '沉思', '思考', '想', '托', '琢磨', '托腮想', '陷入沉思', '认真琢磨', '歪头思考']),
        ('point_analyze', ['分析', '点指', '认真', '讲解', '指', '其实', '分析一下', '点着说', '认真讲解', '其实呢']),
        ('cover_surprise', ['惊讶', '捂嘴', '意外', '不敢相信', '天哪', '惊讶状', '天哪', '不敢相信']),
        ('bow_write', ['书写', '低头', '写字', '记录', '写', '记', '低头写', '奋笔疾书', '边听边记']),
        ('smile_confirm', ['浅笑', '印证', '收获', '点头', '笑', '对了', '印证了', '收获满满', '点头笑']),
    ],
    'sprite': [
        ('float_ask', ['提问', '问', '疑惑', '好奇', '？', '为什么']),
        ('wave_interact', ['挥手', '互动', '招呼', '欢迎', '嗨']),
    ],
}

# 预设 IP 角色姿势库：用 WorkBuddy ImageGen 生成的固定角色，不走 Pollinations 动态生成。
# 角色名 -> 静态姿势图目录（static/poses/<folder>/ 下存放 stand/sit/.../point.png，透明背景）。
# 许多=男孩，莉莉=女孩；哈哈/张驰/大唐/徐剑为四位老师（均已生成并抠透明背景）。
PRESET_CHARACTERS = {
    '许多': 'xuduo',
    '莉莉': 'lili',
    '高远': 'gaoyuan',
    '哈哈': 'haha',
    '高远老师': 'gaoyuan',    # 高远老师专属姿势库在 gaoyuan/ 目录
    '张驰': 'zhangchi',
    '驰哥': 'zhangchi',    # 张驰老师别名
    '大唐': 'datang',
    '唐润然': 'datang',    # 大唐老师别名
    '唐润然老师': 'datang',    # 大唐老师别名
    '徐剑': 'xujian',
    '雪糕': 'xujian',    # 徐剑老师别名
    '雪糕老师': 'xujian',    # 徐剑老师别名
    '可可': 'keke',
    '布布': 'bubu',
}
# 知识小精灵：非预设角色的默认姿势库
SPRITE_FOLDER = 'sprite'


def get_preset_poses(key):
    """若 key 是预设 IP 角色名或目录名，返回 {poseKey: url}（来自 static/poses/<folder>）。否则返回 None。
    同时包含正面姿势（stand/sit/...）与侧身姿势（stand_left/stand_right/...），供漫画按朝向选择。
    新版角色（高远/许多新姿势/莉莉新姿势/精灵）扫描 CHARACTER_POSE_KEYS 中的专属 key，
    同时也扫描旧版 POSE_LIBRARY key 作为 fallback（许多/莉莉保留有旧文件）。
    """
    if not key:
        return None
    folder = PRESET_CHARACTERS.get(key) or key  # 允许直接传目录名（xuduo/lili）
    base = os.path.join(app.root_path, 'static', 'poses', folder)
    if not os.path.isdir(base):
        return None
    poses = {}

    # 1. 扫描角色专属姿势 key（新版）
    if folder in CHARACTER_POSE_KEYS:
        for pk in CHARACTER_POSE_KEYS[folder]:
            p = os.path.join(base, f'{pk}.png')
            if os.path.exists(p):
                poses[pk] = f'/static/poses/{folder}/{pk}.png'

    # 2. 扫描旧版 POSE_LIBRARY key（作为 fallback，对所有角色都尝试）
    for pk in POSE_LIBRARY:
        p = os.path.join(base, f'{pk}.png')
        if os.path.exists(p) and pk not in poses:
            poses[pk] = f'/static/poses/{folder}/{pk}.png'
        # 侧身图：{pose}_left.png / {pose}_right.png
        for side in ('left', 'right'):
            sp = os.path.join(base, f'{pk}_{side}.png')
            if os.path.exists(sp):
                poses[f'{pk}_{side}'] = f'/static/poses/{folder}/{pk}_{side}.png'

    # 3. 标准图兜底：如果有 {folder}_standard.png 或 standard.png，加入 'stand' key
    for std_name in [f'{folder}_standard.png', 'standard.png', 'stand.png']:
        sp = os.path.join(base, std_name)
        if os.path.exists(sp) and 'stand' not in poses:
            poses['stand'] = f'/static/poses/{folder}/{std_name}'
            break

    return poses if poses else None



def extract_action(text):
    """从对话文本中提取括号里的动作说明，并返回(清理后的文本, 动作或None)。

    支持：
      - 中英文括号：「（举手抢着说）我知道！」
      - 全角与半角混合：「（指着黑板)」或「(指着黑板）」
      - 一句话多个动作：「（举手）我知道！（开心地笑）」
      - 嵌套括号：「（（挥手）打招呼）」
      - 空括号：「（）」

    清理后的文本不再包含任何括号及其中内容，用于气泡渲染；
    提取出的动作以「，」连接，用于选择角色姿势。
    """
    if not text:
        return ('', None)

    open_chars = {'（', '('}
    close_chars = {'）', ')'}
    cleaned_chars = []
    action_buffer = []
    actions = []
    depth = 0

    for ch in text:
        if ch in open_chars:
            depth += 1
        elif ch in close_chars:
            if depth > 0:
                depth -= 1
                if depth == 0:
                    action_text = ''.join(action_buffer).strip()
                    if action_text:
                        actions.append(action_text)
                    action_buffer = []
            # 没有对应开括号的 stray 闭括号直接丢弃
        else:
            if depth > 0:
                action_buffer.append(ch)
            else:
                cleaned_chars.append(ch)

    # 容错：如果文本末尾仍有未闭合的括号，把已收集的内容也作为动作
    if depth > 0 and action_buffer:
        action_text = ''.join(action_buffer).strip()
        if action_text:
            actions.append(action_text)

    cleaned = ''.join(cleaned_chars).strip()
    action = '，'.join(actions) if actions else None
    return (cleaned, action)


def _wrap_dialogue_lines(draw, text, font, max_width, max_chars_per_line=11):
    """按语义（标点优先）+ 单字宽度 + 单行星限制对台词断行。

    composite_image 与 render_text 共用此函数，确保气泡尺寸计算与最终文字排版
    使用完全一致的行数，避免文字溢出气泡边界。
    """
    if not text:
        return []
    text = text.replace('\r', '')
    # 11 个字只是舒适阅读的软上限；标点、闭合书名号必须能跟随前字越过这一格。
    forbidden_head = set('，。！？、；：”’）】〉》〕｝》…—,.!?;:)')
    forbidden_tail = set('（【〈《〔｛“‘(')

    # 短书名、引号、括号内容作为一个排版单元，优先整块留在同一行。
    pairs = {'《': '》', '“': '”', '‘': '’', '（': '）', '【': '】', '(': ')'}
    tokens, i = [], 0
    while i < len(text):
        ch = text[i]
        if ch == '\n':
            tokens.append('\n'); i += 1; continue
        close = pairs.get(ch)
        if close:
            end = text.find(close, i + 1)
            if end != -1:
                tokens.append(text[i:end + 1]); i = end + 1; continue
        tokens.append(ch); i += 1

    def fits(line, token, allow_soft_overflow=False):
        candidate = line + token
        if draw.textlength(candidate, font=font) > max_width:
            return False
        return allow_soft_overflow or len(candidate) <= max_chars_per_line

    lines, current_line = [], ''
    for token in tokens:
        if token == '\n':
            if current_line: lines.append(current_line); current_line = ''
            continue
        # 整块（例如《猜猜他是谁》）放不下时，先换行；空行仍放不下再安全地拆字。
        if len(token) > 1 and not fits(current_line, token):
            if current_line:
                lines.append(current_line); current_line = ''
            if fits('', token, allow_soft_overflow=True):
                current_line = token
                continue
            token_parts = list(token)
        else:
            token_parts = [token]
        for char in token_parts:
            # 标点不得成为新行首；可在不超出像素宽度时突破 11 字软上限。
            is_closing = char in forbidden_head
            if current_line and fits(current_line, char, allow_soft_overflow=is_closing):
                current_line += char
                continue
            if not current_line:
                current_line = char
                continue
            # 若闭合标点因像素宽度放不下，把上一行最后一个字连同标点一起
            # 挪到下一行。必须保持原文字符顺序；例如“看。”只能变成下一行
            # 的“看。”，绝不能被重排成上一行末尾“。”、下一行“看”。
            if is_closing and len(current_line) > 1:
                moved = current_line[-1]
                lines.append(current_line[:-1])
                current_line = moved + char
            else:
                # 开书名号/开括号不应孤立在行尾，尽量连同下一个字留给下一行。
                if current_line and current_line[-1] in forbidden_tail and len(current_line) > 1:
                    moved = current_line[-1]
                    current_line = current_line[:-1]
                    lines.append(current_line)
                    current_line = moved + char
                else:
                    lines.append(current_line)
                    current_line = char
    if current_line:
        lines.append(current_line)
    return lines


def map_action_to_pose(action):
    """把动作说明映射到姿势 key（见 POSE_LIBRARY）。匹配不到返回 'stand'。"""
    if not action:
        return 'stand'
    a = action
    rules = [
        ('raise_hand', ['举手', '抢', '挥手', '举起', '高高举', '招手', '抬起手', '手举高', '比耶', '举手抢答']),
        ('scratch_head', ['挠头', '摸头', '抓头', '挠', '挠脑袋', '拍脑袋', '抓耳挠腮']),
        ('chin_rest', ['托腮', '托着下巴', '琢磨', '思考', '沉思', '歪头想', '手撑下巴', '托着腮帮', '若有所思', '一手托腮']),
        ('nod', ['点头', '附和', '认同', '赞同', '点头称是', '首肯', '颔首']),
        ('smile', ['神秘一笑', '微笑', '笑', '开心', '乐', '窃笑', '抿嘴笑', '浅笑', '莞尔', '会心一笑']),
        ('point', ['指着', '指', '比划', '指向', '指着绘本', '伸手一指', '点指', '手指某处']),
        ('sit', ['坐', '盘腿', '蹲', '坐下', '跪坐', '席地而坐']),
        ('lie', ['躺', '趴', '卧', '倚', '靠']),
        ('jump', ['跳', '蹦', '跃', '弹跳', '跳起来']),
        ('sad', ['哭', '难过', '伤心', '委屈', '泪', '抽泣', '哽咽']),
        ('run', ['跑', '冲', '追', '奔', '疾走']),
        ('angry', ['生气', '怒', '急', '瞪', '气鼓鼓', '皱眉', '不满']),
    ]
    for pose, kws in rules:
        for kw in kws:
            if kw in a:
                return pose

    # 兜底1：英文姿势 key 匹配（如 raise_hand / nod / sit / laugh / smile / point / chin / scratch / stand）
    en_key_map = {
        'raise_hand': 'raise_hand',
        'raise': 'raise_hand',
        'nod': 'nod',
        'sit': 'sit',
        'laugh': 'laugh',
        'smile': 'smile',
        'point': 'point',
        'chin': 'chin_rest',
        'scratch': 'scratch_head',
        'stand': 'stand',
    }
    for ek in ['raise_hand', 'raise', 'nod', 'sit', 'laugh', 'smile', 'point', 'chin', 'scratch', 'stand']:
        if ek in a:
            return en_key_map[ek]

    # 兜底2：中文描述匹配（POSE_LIBRARY 值按顿号拆分，匹配到子串即返回对应 key）
    for pk, desc in POSE_LIBRARY.items():
        if any(seg in a for seg in desc.split('、')):
            return pk

    return 'stand'


def auto_pose_from_text(text):
    """台词没有显式(动作)标记时，根据语义/情绪自动选一个自然姿势。

    避免所有人一律 stand 导致的「直愣愣杵在原地」观感。结果确定（同句同人每次一致）。
    返回 POSE_LIBRARY 中的姿势 key。
    """
    if not text:
        return 'smile'
    t = text
    has_excl = ('！' in t or '!' in t)
    has_q = ('？' in t or '?' in t)

    # 惊喜 / 赞叹 / 兴奋
    if any(k in t for k in ['哇塞', '哇', '居然', '竟然', '太棒', '好棒', '好厉害', '厉害',
                             '天哪', '天啊', '真的吗', '难以置信', 'amazing', 'wow', 'WOW', '酷']):
        return 'raise_hand' if has_excl else 'smile'
    # 疑问 / 困惑
    if any(k in t for k in ['为什么', '怎么会', '怎么', '啥', '哪', '如何', '什么意思',
                             '不懂', '不明白', '困惑', '难道']):
        return 'scratch_head' if has_q else 'chin_rest'
    # 赞同 / 肯定 / 领悟
    if any(k in t for k in ['对呀', '是的', '没错', '原来如此', '明白了', '我懂', '说得对',
                             '确实', '认同', '赞同', '当然', '就是']):
        return 'nod'
    # 解释 / 说明 / 指向
    if any(k in t for k in ['这就是', '比如', '你看', '看这里', '那个', '指向', '比如说',
                             '因为', '所以', '其实', '换句话说']):
        return 'point'
    # 温柔 / 共情 / 感受
    if any(k in t for k in ['舍不得', '朋友', '温柔', '像', '一样', '爱', '喜欢', '美丽',
                             '美', '感动', '真好', '暖']):
        return 'smile'
    # 句子末尾感叹号（但没命中上面关键词）：用笑/举手体现情绪
    if has_excl:
        return 'laugh'
    # 疑问句兜底
    if has_q:
        return 'chin_rest'
    # 平静叙述：给个有交流的姿势而不是僵硬站立
    return 'smile'


def _infer_character_folder(available_poses):
    """从可用姿势 key 集合推断角色目录名。

    参数:
        available_poses: dict {pose_key: url}，该角色可用的所有姿势

    返回:
        角色目录名字符串（如 'gaoyuan'/'xuduo'/'lili'/'sprite'），无法判定时返回 None
    """
    keys_set = set(available_poses.keys())
    for folder, pose_keys in CHARACTER_POSE_KEYS.items():
        # 如果可用姿势包含该角色超过一半的专属 key，判定为该角色
        overlap = keys_set & set(pose_keys)
        if len(overlap) >= max(2, len(pose_keys) // 2):
            return folder
    return None


def _select_pose_for_character(action, text, available_poses):
    """根据角色的可用姿势，选择最合适的姿势 key。

    选择优先级：
    1. 显式动作 → 旧版 map_action_to_pose → 命中可用姿势则用
    2. 显式动作 → 角色专属动作映射 → 命中可用姿势则用
    3. 显式动作 → 跨角色映射 → 命中可用姿势则用
    4. 无动作 → auto_pose_from_text → 命中可用姿势则用
    5. 无动作 → 语义自动选择（疑问/感叹）→ 命中可用姿势则用
    6. 角色默认姿势
    7. 'stand' 兜底
    8. 第一个可用姿势

    参数:
        action: 从台词中提取的动作（可能为 None）
        text: 台词原文（用于 auto_pose_from_text 兜底）
        available_poses: dict {pose_key: url}，该角色可用的所有姿势

    返回:
        pose_key 字符串（保证在 available_poses 中，除非 available_poses 为空）
    """
    if not available_poses:
        return 'stand'

    char_folder = _infer_character_folder(available_poses)

    # 1. 如果有显式动作，先尝试旧版映射（仅当匹配到非 stand 的有效姿势才直接使用）
    if action:
        pose_key = map_action_to_pose(action)
        if pose_key != 'stand' and pose_key in available_poses:
            return pose_key
        # 动作存在但旧版未强匹配到：继续尝试角色专属/跨角色映射与语义选择，避免直接退化成 stand

        # 动作文本本身也做一次语义自动选（兼容"手插兜若有所思"这类无关键词但有情绪的动作）
        sem = auto_pose_from_text(action)
        if sem != 'stand' and sem in available_poses:
            return sem

        # 2. 尝试角色专属映射
        if char_folder and char_folder in CHARACTER_ACTION_MAP:
            a = action
            for pose, kws in CHARACTER_ACTION_MAP[char_folder]:
                for kw in kws:
                    if kw in a and pose in available_poses:
                        return pose

        # 3. 尝试其他角色的映射（跨角色匹配，找最接近的语义）
        for folder, rules in CHARACTER_ACTION_MAP.items():
            if folder == char_folder:
                continue
            for pose, kws in rules:
                for kw in kws:
                    if kw in action and pose in available_poses:
                        return pose

    # 4. 无动作标记时，尝试语义自动选择
    if text:
        pose_key = auto_pose_from_text(text)
        if pose_key in available_poses:
            return pose_key

        # 语义选择失败，尝试角色专属语义匹配
        t = text
        has_q = ('？' in t or '?' in t)
        has_excl = ('！' in t or '!' in t)
        if has_q:
            for folder, rules in CHARACTER_ACTION_MAP.items():
                for pose, kws in rules:
                    if any(kw in ['问', '？', '疑惑', '好奇', '提问'] for kw in kws) and pose in available_poses:
                        return pose
        if has_excl:
            for folder, rules in CHARACTER_ACTION_MAP.items():
                for pose, kws in rules:
                    if any(kw in ['兴奋', '开心', '震惊', '惊讶'] for kw in kws) and pose in available_poses:
                        return pose

    # 5. 返回角色默认姿势
    if char_folder and CHARACTER_DEFAULT_POSE.get(char_folder) in available_poses:
        return CHARACTER_DEFAULT_POSE[char_folder]

    # 6. 如果有 'stand'，用 stand
    if 'stand' in available_poses:
        return 'stand'

    # 7. 最终兜底：返回第一个可用姿势
    return next(iter(available_poses))


def _apply_logo_to_image(background, logo_url=None, logo_position='top-right', logo_size=0.08):
    """在合成图片上叠加 Logo 水印（RGBA 透明合成）。

    参数:
        background: PIL.Image (RGBA) — 已合成好的漫画格子图
        logo_url: str — Logo 图片的服务器路径（如 /static/logos/xxx.png）或本地绝对路径
        logo_position: str — 位置 top-right / top-left / bottom-right / bottom-left
        logo_size: float — Logo 宽度占格子宽度的比例（如 0.08 表示 8%）

    返回:
        PIL.Image (RGBA) — 叠加 Logo 后的图片（如果 logo_url 为空或加载失败则原样返回）
    """
    if not logo_url:
        return background

    try:
        # 把服务器路径转成本地文件系统路径
        if logo_url.startswith('/'):
            logo_path = os.path.join(app.root_path, logo_url.lstrip('/'))
        elif os.path.isabs(logo_url):
            logo_path = logo_url
        else:
            logo_path = os.path.join(app.root_path, logo_url)

        if not os.path.exists(logo_path):
            logger.warning(f"Logo 文件不存在: {logo_path}")
            return background

        logo_img = Image.open(logo_path).convert('RGBA')

        bg_width, bg_height = background.size
        logo_w, logo_h = logo_img.size

        # 按比例缩放 Logo
        target_logo_w = max(1, int(bg_width * logo_size))
        scale_ratio = target_logo_w / logo_w if logo_w > 0 else 1
        target_logo_h = max(1, int(logo_h * scale_ratio))
        logo_img = logo_img.resize((target_logo_w, target_logo_h), Image.LANCZOS)

        # 边距（格子宽度的 2%）
        margin = int(bg_width * 0.02)

        # 根据位置计算坐标
        if logo_position == 'top-left':
            pos_x, pos_y = margin, margin
        elif logo_position == 'bottom-right':
            pos_x = bg_width - target_logo_w - margin
            pos_y = bg_height - target_logo_h - margin
        elif logo_position == 'bottom-left':
            pos_x = margin
            pos_y = bg_height - target_logo_h - margin
        else:  # 默认 top-right
            pos_x = bg_width - target_logo_w - margin
            pos_y = margin

        # 创建与背景同尺寸的透明图层，粘贴 Logo 后再 alpha_composite
        logo_layer = Image.new('RGBA', (bg_width, bg_height), (0, 0, 0, 0))
        logo_layer.paste(logo_img, (pos_x, pos_y), logo_img)
        result = Image.alpha_composite(background, logo_layer)

        logger.info(f"Logo 水印叠加成功: position={logo_position}, size={logo_size}, "
                     f"logo尺寸={target_logo_w}x{target_logo_h}")
        return result

    except Exception as e:
        logger.error(f"Logo 水印叠加失败: {str(e)}")
        return background


def composite_image(background_data, ip_paths, dialogue, output_path, speaker='', speakers=None,
                    font_family='msyhbd', font_size=None, poses=None, grid_spec=None, layout=None,
                    bubble_style='rounded', text_align='center', opacity=0.96, dialogues=None,
                    bubble_images=None, character_bubble_map=None, skip_ip_overlay=False,
                    character_poses=None, target_width=None, target_height=None,
                    logo_url=None, logo_position='top-right', logo_size=0.08):
    """使用Pillow合成图片：背景 + 多个IP贴图 + 气泡对白 + 说话者

    参数:
        dialogues: 对话列表，支持多个说话者的对话
        speakers: 说话者列表，与dialogues对应
        skip_ip_overlay: 是否跳过IP贴图（背景已是完整场景时启用）
    """
    # 旧脚本同时保留 `speaker`（单值）和 `speakers`（逐句列表）。
    # 多人物格里只要主说话者碰巧是知识小精灵，不能把整格误判成
    # “单独精灵格”，否则普通人物会被丢掉，精灵还会走两次绘制分支。
    panel_speaker_names = []
    if isinstance(speakers, (list, tuple)):
        panel_speaker_names = [str(name).strip() for name in speakers if str(name).strip()]
    elif speakers:
        panel_speaker_names = [str(speakers).strip()]
    if not panel_speaker_names and speaker:
        panel_speaker_names = [str(speaker).strip()]
    unique_panel_speakers = list(dict.fromkeys(panel_speaker_names))
    is_guide = unique_panel_speakers == ['知识小精灵']
    # 收集合成元数据，用于前端编辑面板重建布局
    meta = {
        'success': True,
        'characters': [],
        'bubbles': [],
        'sound_effect': None,
        'background_url': None,
        'is_guide': is_guide
    }
    try:
        if not isinstance(ip_paths, list):
            ip_paths = [ip_paths]
        
        from io import BytesIO
        background = Image.open(BytesIO(background_data)).convert('RGBA')

        if target_width and target_height:
            # 不直接拉伸：先裁掉生图服务偶发的黑边，再按目标格比例 cover 填满。
            background = normalize_background_for_panel(
                background, target_width, target_height
            ).convert('RGBA')
            bg_width, bg_height = target_width, target_height
        else:
            bg_width, bg_height = background.size
        # 编辑器必须以合成画布尺寸为坐标基准，不能误用原始背景文件的像素尺寸。
        meta['canvas_width'] = bg_width
        meta['canvas_height'] = bg_height
        
        # ---- 提前解析对话：提取(动作)并清理文本，用于按姿势选角色图 ----
        if dialogues and isinstance(dialogues, list) and len(dialogues) > 0:
            _dlist = dialogues
            _slist = speakers if isinstance(speakers, list) else []
        else:
            _dlist = [dialogue]
            _slist = [speaker]

        valid_pairs = []      # (clean_text, speaker)
        valid_actions = []    # 与 valid_pairs 对齐的动作说明
        for i, d in enumerate(_dlist):
            raw = d.replace('\n', '').replace('\r', '').strip()
            if not raw:
                continue
            clean_text, action = extract_action(raw)
            clean_text = clean_text.strip()
            if not clean_text:
                # 只剩动作没台词则跳过（不画空气泡）
                continue
            s = _slist[i] if i < len(_slist) else (speaker or '')
            valid_pairs.append((clean_text, s))
            valid_actions.append(action)

        # ---- 决定用于绘制的角色图：有对话时按每句(动作)选姿势图，否则用标准图 ----
        if valid_pairs:
            # 人物图层按“角色”而不是按“台词”创建。同一人物连续说两句时，
            # 仍只出现一次；气泡继续使用 valid_pairs，数量不受影响。
            speaker_names = []
            _char_source_paths = []
            _char_source_urls = []
            _pose_keys = []  # 与去重后的 speaker_names 一一对应
            for i, (ct, s) in enumerate(valid_pairs):
                s = str(s or '').strip()
                if s in speaker_names:
                    continue
                speaker_names.append(s)
                action = valid_actions[i]
                # 获取该角色的可用姿势
                avail = {}
                if character_poses and s in character_poses and isinstance(character_poses[s], dict):
                    avail = character_poses[s]
                
                # 使用新的智能姿势选择（兼容旧版角色和新版角色）
                pose_key = _select_pose_for_character(action, ct, avail)
                _pose_keys.append(pose_key)
                
                pose_url = avail.get(pose_key) if avail else None
                if not pose_url and avail:
                    pose_url = next(iter(avail.values()))  # 最终兜底
                if not pose_url and s in (speakers or []):
                    try:
                        si = list(speakers or []).index(s)
                        if si < len(ip_paths):
                            pose_url = _bubble_path_to_url(ip_paths[si])
                    except Exception:
                        pass
                if not pose_url and ip_paths:
                    pose_url = _bubble_path_to_url(ip_paths[0])
                abs_p = os.path.join(app.root_path, pose_url.lstrip('/')) if pose_url else (ip_paths[0] if ip_paths else None)
                _char_source_paths.append(abs_p)
                _char_source_urls.append(pose_url)
        else:
            speaker_names = speakers if isinstance(speakers, list) else []
            _char_source_paths = list(ip_paths)
            _char_source_urls = [_bubble_path_to_url(p) if p else None for p in ip_paths]
            _pose_keys = ['stand'] * len(ip_paths)

        ip_images = []
        for p in _char_source_paths:
            if p and os.path.exists(p):
                ip_images.append(Image.open(p).convert('RGBA'))

        if not ip_images:
            logger.warning("没有可用的IP图片")
            background = _apply_logo_to_image(background, logo_url, logo_position, logo_size)
            background.save(output_path, format='PNG')
            meta['success'] = True
            return meta
        
        speaker_index = -1
        if speaker and not is_guide:
            if speaker in speaker_names:
                speaker_index = speaker_names.index(speaker)
            else:
                speaker_char = speaker[-1]
                if speaker_char >= 'A' and speaker_char <= 'Z':
                    speaker_index = ord(speaker_char) - ord('A')
        
        if speaker_index >= len(ip_images):
            speaker_index = -1
        
        display_indices = []
        if is_guide:
            display_indices = [0]
        elif speaker_index >= 0:
            display_indices.append(speaker_index)
            other_indices = [i for i in range(len(ip_images)) if i != speaker_index]
            for i in range(min(4, len(other_indices))):
                display_indices.append(other_indices[i])
        else:
            display_indices = list(range(min(5, len(ip_images))))
        
        display_images = [ip_images[i] for i in display_indices]
        display_pose_keys = [_pose_keys[i] for i in display_indices]
        speaker_idx_in_display = display_indices.index(speaker_index) if speaker_index >= 0 else -1
        
        display_character_names = [speaker_names[i] for i in display_indices]
        scaled_ips, ip_positions, ip_transform_params = calculate_character_positions(
            display_images, bg_width, bg_height, speaker_idx_in_display, is_guide,
            character_names=display_character_names,
        )
        
        for i, (ip_img, (ip_x, ip_y), params) in enumerate(zip(scaled_ips, ip_positions, ip_transform_params)):
            # 完整场景模式下背景已含人物，跳过额外IP贴图（保留知识小精灵的单独处理）
            if skip_ip_overlay:
                continue

            transformed_img = ip_img.copy()

            # 若存在真正的侧身图，优先按朝向替换正面图，避免正面图镜像后仍朝镜头的僵硬感
            try:
                speaker_idx_for_meta = display_indices[i]
                speaker_name_for_pose = speaker_names[speaker_idx_for_meta] if speaker_idx_for_meta < len(speaker_names) else ''
                pose_key = display_pose_keys[i]
                if character_poses and speaker_name_for_pose in character_poses and isinstance(character_poses[speaker_name_for_pose], dict):
                    side_key = f"{pose_key}_{'left' if params['flip'] else 'right'}"
                    side_url = character_poses[speaker_name_for_pose].get(side_key)
                    if side_url:
                        side_path = os.path.join(app.root_path, side_url.lstrip('/'))
                        if os.path.exists(side_path):
                            side_img = Image.open(side_path).convert('RGBA')
                            # 侧身图已经朝正确方向，取消后续镜像
                            transformed_img = side_img
                            params['flip'] = False
                            logger.info(f"角色[{speaker_name_for_pose}] 使用侧身图 {side_key}")
            except Exception as e:
                logger.warning(f"侧身图替换失败: {e}")

            final_size = params['size']
            scale_factor = params.get('scale', 1.0)

            target_w = int(final_size[0] * scale_factor)
            target_h = int(final_size[1] * scale_factor)

            if params['flip']:
                transformed_img = transformed_img.transpose(Image.FLIP_LEFT_RIGHT)

            # 旋转逻辑已移除：人物不再旋转
            # if params['rotate'] != 0:
            #     transformed_img = transformed_img.rotate(params['rotate'], expand=True, resample=Image.BICUBIC)

            if transformed_img.size != (target_w, target_h):
                transformed_img = transformed_img.resize((target_w, target_h), Image.BICUBIC)

            # 给人物打光：叠加环境光渐变（顶亮底暗）+对比增强，减少3D贴纸感
            if not is_guide:
                transformed_img = apply_character_lighting(transformed_img, strength=1.35)

            adjusted_x = ip_x
            adjusted_y = ip_y + params.get('offset_y', 0)

            # 兜底：确保人物底部不超出背景边界（防止旋转/scale 后越界截断）
            if adjusted_y + target_h > bg_height:
                adjusted_y = max(0, bg_height - target_h)
            # 兜底：确保人物右侧不超出背景边界
            if adjusted_x + target_w > bg_width:
                adjusted_x = max(0, bg_width - target_w)

            if params['rotate'] != 0:
                adjusted_x += (final_size[0] - target_w) // 2
                adjusted_y += (final_size[1] - target_h) // 2

            # 给角色脚下加柔和阴影，减少“漂浮”贴纸感
            if not is_guide:
                try:
                    shadow_layer = Image.new('RGBA', background.size, (0, 0, 0, 0))
                    shadow_draw = ImageDraw.Draw(shadow_layer)
                    foot_y = adjusted_y + target_h - 2
                    # 光源假设在左上，接地阴影往右下偏移一点，更有落地感
                    dir_offset = int(target_w * 0.06)
                    center_x = adjusted_x + target_w // 2 + dir_offset
                    shadow_w = max(target_w * 1.02, 24)
                    shadow_h = max(20, int(target_h * 0.10))
                    shadow_draw.ellipse(
                        [center_x - shadow_w // 2, foot_y,
                         center_x + shadow_w // 2, foot_y + shadow_h],
                        fill=(0, 0, 0, 150)
                    )
                    shadow_layer = shadow_layer.filter(ImageFilter.GaussianBlur(radius=10))
                    background = Image.alpha_composite(background, shadow_layer)
                except Exception as e:
                    logger.warning(f"绘制角色阴影失败: {e}")

            # 使用alpha_composite进行透明混合，让角色更自然地融入背景
            temp_layer = Image.new('RGBA', background.size, (0, 0, 0, 0))
            temp_layer.paste(transformed_img, (adjusted_x, adjusted_y), transformed_img)
            background = Image.alpha_composite(background, temp_layer)

            # 保存“实际参与合成”的角色图层。侧身替换、镜像、尺寸和环境光都已经
            # 烘焙在这张透明 PNG 中，编辑器直接加载它，避免再次根据原始姿势猜测，
            # 从而保证编辑面板与最终预览使用的是同一张人物图。
            editor_layer_url = None
            try:
                output_stem = os.path.splitext(os.path.basename(output_path))[0]
                editor_layer_name = f"{output_stem}_character_{i}.png"
                editor_layer_path = os.path.join(os.path.dirname(output_path), editor_layer_name)
                transformed_img.save(editor_layer_path, format='PNG')
                editor_layer_url = '/' + os.path.relpath(
                    editor_layer_path, app.root_path
                ).replace('\\', '/')
            except Exception as e:
                logger.warning(f"保存角色编辑图层失败: {e}")

            # 记录角色元数据（用于编辑面板重建）
            speaker_name = speaker_names[display_indices[i]] if display_indices[i] < len(speaker_names) else ''
            meta['characters'].append({
                'speaker': speaker_name,
                'x': adjusted_x,
                'y': adjusted_y,
                'width': target_w,
                'height': target_h,
                'rotate': 0,
                # editor_layer_url 已经包含镜像结果，编辑器不可再次翻转。
                'flip': False if editor_layer_url else params.get('flip', False),
                'scale': scale_factor,
                'path': _char_source_paths[display_indices[i]] if display_indices[i] < len(_char_source_paths) else None,
                'url': editor_layer_url or (
                    _char_source_urls[display_indices[i]]
                    if display_indices[i] < len(_char_source_urls) else None
                )
            })
        
        draw = ImageDraw.Draw(background)
        
        font_path = get_font_path(font_family)
        
        if not font_path:
            logger.warning("未找到可用的字体文件，使用默认字体")
        
        # 正文与边缘必须有明显呼吸感；不再使用旧版“紧凑”10px 留白。
        padding = 18
        # 多对白时收窄单个气泡，为横向并排预留空间。
        dialogue_count = max(1, len(valid_pairs))
        width_ratio = 0.60 if dialogue_count == 1 else (0.40 if dialogue_count == 2 else 0.22)
        max_width = min(int(bg_width * width_ratio), int(40 * (11 if dialogue_count == 1 else 8)))
        # 安全边距：按图片尺寸比例计算，避免气泡/文字贴近边缘被 CSS object-fit:cover 裁掉
        safe_margin_x = max(15, int(bg_width * 0.06))   # 左右安全边距 ≈ 6% 宽度
        safe_margin_top = max(15, int(bg_height * 0.06)) # 顶部安全边距 ≈ 6% 高度
        safe_margin_bottom = max(20, int(bg_height * 0.08)) # 底部安全边距更大（留人物空间）
        
        # 绘制多个气泡（valid_pairs 已在函数开头解析：clean_text + speaker）
        bubble_colors = [
            (255, 255, 255, 245),   # 白色
            (255, 240, 245, 245),   # 粉色
            (240, 255, 245, 245),   # 绿色
            (240, 245, 255, 245),   # 蓝色
            (255, 255, 230, 245),   # 黄色
        ]
        
        outline_colors = [
            (255, 179, 198),   # 粉色
            (168, 230, 207),   # 绿色
            (100, 170, 255),   # 蓝色
            (255, 200, 100),   # 橙色
            (180, 150, 255),   # 紫色
        ]
        
        drawn_bubbles = []
        
        # 气泡只需要避让人物头脸，不应把整具身体都当禁区。
        # 使用最终合成元数据，避免初始尺寸与实际缩放位置不一致。
        character_boxes = []
        if not skip_ip_overlay:
            for char_meta in meta['characters']:
                cx, cy = char_meta['x'], char_meta['y']
                cw, ch = char_meta['width'], char_meta['height']
                face_x = cx + int(cw * 0.14)
                face_y = cy + int(ch * 0.02)
                face_w = int(cw * 0.72)
                face_h = int(ch * 0.40)
                character_boxes.append((face_x, face_y, face_w, face_h))

        def check_overlap(new_x, new_y, new_w, new_h, margin=10, char_margin=15):
            """检测与已绘气泡、已放置人物的重叠（含 margin）。
            margin：气泡之间间距；char_margin：气泡与人物之间间距。
            """
            for bx, by, bw, bh in drawn_bubbles:
                if (new_x - margin < bx + bw and
                    new_x + new_w + margin > bx and
                    new_y - margin < by + bh and
                    new_y + new_h + margin > by):
                    return True
            for cx, cy, cw, ch in character_boxes:
                if (new_x - char_margin < cx + cw and
                    new_x + new_w + char_margin > cx and
                    new_y - char_margin < cy + ch and
                    new_y + new_h + char_margin > cy):
                    return True
            return False

        def _overlap_area(x, y, w, h, char_margin=15, margin=10):
            """计算与已绘气泡/人物的重叠总面积，用于回退策略挑选最优位置。"""
            area = 0
            for bx, by, bw, bh in drawn_bubbles:
                ox = max(0, min(x + w, bx + bw) - max(x, bx))
                oy = max(0, min(y + h, by + bh) - max(y, by))
                area += ox * oy
            for cx, cy, cw, ch in character_boxes:
                ox = max(0, min(x + w, cx + cw) - max(x, cx))
                oy = max(0, min(y + h, cy + ch) - max(y, cy))
                area += ox * oy
            return area

        def _tail_for(cpos, bx, by, bw, bh,
                      speaker_x, speaker_y, speaker_w, speaker_h):
            """返回一体式短尾巴锚点。

            尾巴只指到脸部外边缘，不进入脸部矩形。上方停在头顶外、下方停在
            下巴外、侧面停在脸颊外，避免三角尾巴遮住五官。
            """
            face_left = speaker_x + speaker_w * 0.14
            face_top = speaker_y + speaker_h * 0.02
            face_right = face_left + speaker_w * 0.72
            face_bottom = face_top + speaker_h * 0.40
            face_center_x = (face_left + face_right) / 2
            face_center_y = (face_top + face_bottom) / 2

            if by + bh <= face_top:
                side = 'bottom'
                target_x, target_y = face_center_x, face_top - 4
            elif by >= face_bottom:
                side = 'top'
                target_x, target_y = face_center_x, face_bottom + 4
            elif bx + bw <= face_left:
                side = 'right'
                target_x, target_y = face_left - 4, face_center_y
            else:
                side = 'left'
                target_x, target_y = face_right + 4, face_center_y
            target_x, target_y = int(target_x), int(target_y)
            if side == 'left':
                base_x = bx
                base_y = max(by + 18, min(target_y, by + bh - 18))
            elif side == 'right':
                base_x = bx + bw
                base_y = max(by + 18, min(target_y, by + bh - 18))
            elif side == 'top':
                base_x = max(bx + 18, min(target_x, bx + bw - 18))
                base_y = by
            else:
                base_x = max(bx + 18, min(target_x, bx + bw - 18))
                base_y = by + bh
            return (int(base_x), int(base_y), target_x, target_y, side)

        def _place_speaker_bubble(speaker_x, speaker_y, speaker_w, speaker_h,
                                  bubble_width, bubble_height, bubble_center_x):
            """把气泡放在说话者头顶附近，并让气泡主体做小幅自然错位。

            正常候选只改变气泡本身的位置，不把尾巴横向拉长；全部候选仍须通过
            人脸与其他气泡避让。大范围位置只作为空间不足时的安全兜底。
            """
            candidates = []
            base_x = bubble_center_x - bubble_width / 2
            base_y = speaker_y - bubble_height - 26
            # 侧移的是气泡主体。偏移量需要在成品缩放后仍然看得出来，因此采用
            # 气泡宽度约 22%～32% 的左上/右上构图，而不是几乎不可见的几像素。
            offset_variants = [
                (-min(62, bubble_width * 0.32), -8),
                (min(62, bubble_width * 0.32), 6),
                (-min(46, bubble_width * 0.22), 8),
                (min(46, bubble_width * 0.22), -5),
            ]
            first = random.randrange(len(offset_variants))
            offset_variants = offset_variants[first:] + offset_variants[:first]
            for dx, dy in offset_variants:
                candidates.append(('above', base_x + dx, base_y + dy))

            def _clamp(cx, cy):
                cx = max(safe_margin_x, min(cx, bg_width - bubble_width - safe_margin_x))
                cy = max(safe_margin_top, min(cy, bg_height - bubble_height - safe_margin_bottom))
                return cx, cy

            # 第一轮：严格无重叠（char_margin=15 保证远离人脸）
            for cpos, cx, cy in candidates:
                cx, cy = _clamp(cx, cy)
                if not check_overlap(cx, cy, bubble_width, bubble_height, char_margin=15):
                    return cx, cy, _tail_for(
                        cpos, cx, cy, bubble_width, bubble_height,
                        speaker_x, speaker_y, speaker_w, speaker_h
                    )

            # 回退：仍以人物头顶为锚，优先向上错层，再做很小的左右调整。
            # 这样即使人物贴近画面边缘，气泡也不会被过早送到画面中央。
            best = None
            for cpos, base_x, base_y in candidates:
                search_dx = [0, -12, 12, -24, 24, -36, 36]
                search_dy = [
                    0, -10, 10, -20, 20, -30, 30,
                    -40, -50, -60, -70, -80, -90, -100,
                    -110, -120, -130, -140, -150, -160,
                ]
                for dx in search_dx:
                    for dy in search_dy:
                        cx, cy = _clamp(base_x + dx, base_y + dy)
                        if not check_overlap(cx, cy, bubble_width, bubble_height, char_margin=15):
                            return cx, cy, _tail_for(
                                cpos, cx, cy, bubble_width, bubble_height,
                                speaker_x, speaker_y, speaker_w, speaker_h
                            )
                        area = _overlap_area(cx, cy, bubble_width, bubble_height, char_margin=15)
                        if best is None or area < best[2]:
                            best = (cx, cy, area)

            # 头顶纵向空间确实被占满时，继续沿说话者所在的局部列搜索整幅画面。
            # 可落到人物侧下方的安全空白，但横向必须留在说话者附近，不能跑到别人身边。
            near_positions = []
            near_x_offsets = [0, -24, 24, -48, 48, -72, 72]
            desired_y = speaker_y - bubble_height - 26
            # 局部空间不足时可使用人物脸部以下的画面，底边只保留 12px。
            # 仍严格不压脸，但不再把整个人体都当成禁区。
            max_y = int(bg_height - bubble_height - 12)
            for dx in near_x_offsets:
                cx = max(
                    safe_margin_x,
                    min(bubble_center_x - bubble_width / 2 + dx,
                        bg_width - bubble_width - safe_margin_x)
                )
                for cy in range(int(safe_margin_top), max_y + 1, 10):
                    if not check_overlap(
                        cx, cy, bubble_width, bubble_height,
                        char_margin=4
                    ):
                        horizontal_distance = abs(
                            (cx + bubble_width / 2) - bubble_center_x
                        )
                        vertical_distance = abs(cy - desired_y)
                        # 横向关联优先级高于上下位置，保证一眼能看出属于谁。
                        score = horizontal_distance * 3 + vertical_distance
                        near_positions.append((score, cx, cy))
            if near_positions:
                _, cx, cy = min(near_positions)
                return cx, cy, _tail_for(
                    'near_speaker', cx, cy, bubble_width, bubble_height,
                    speaker_x, speaker_y, speaker_w, speaker_h
                )

            # 兜底：先按说话者所在的左/中/右区域横向分布，再逐行换位。
            # 旧实现只增加 Y，多个气泡会在画面中央排成一条竖列。
            if bubble_center_x < bg_width * 0.33:
                anchor_x = bg_width * 0.20
            elif bubble_center_x > bg_width * 0.67:
                anchor_x = bg_width * 0.80
            else:
                anchor_x = bg_width * 0.50
            tc_x = int(anchor_x - bubble_width / 2)
            tc_y = safe_margin_top
            tc_x, tc_y = _clamp(tc_x, tc_y)
            bubble_margin = 20
            horizontal_slots = [
                int(bg_width * 0.20 - bubble_width / 2),
                int(bg_width * 0.50 - bubble_width / 2),
                int(bg_width * 0.80 - bubble_width / 2),
            ]
            horizontal_slots.sort(key=lambda x: abs((x + bubble_width / 2) - bubble_center_x))
            found = None
            for row_y in (
                safe_margin_top,
                safe_margin_top + bubble_height + bubble_margin,
                bg_height - bubble_height - safe_margin_bottom,
            ):
                for slot_x in horizontal_slots:
                    cx, cy = _clamp(slot_x, row_y)
                    if not check_overlap(cx, cy, bubble_width, bubble_height,
                                         margin=bubble_margin, char_margin=8):
                        found = (cx, cy)
                        break
                if found:
                    break
            if found:
                tc_x, tc_y = found
            logger.warning(
                f"[气泡避让] idx={idx} speaker={current_speaker} 找不到无重叠位置，"
                f"使用错位兜底 pos=({tc_x},{tc_y}) area={best[2] if best else 'NA'}"
            )
            return tc_x, tc_y, _tail_for(
                'top_center', tc_x, tc_y, bubble_width, bubble_height,
                speaker_x, speaker_y, speaker_w, speaker_h
            )

        for idx, (dialogue_text, current_speaker) in enumerate(valid_pairs):
            text_length = len(dialogue_text)

            # 自动识别格子尺寸：按面板短边缩放气泡基础字号，
            # 保证大格子字够大、小格子字不爆框
            panel_scale = max(0.85, min(1.8, min(bg_width, bg_height) / 768.0))

            if font_size is None:
                if text_length <= 6:
                    base_font_size = 42
                elif text_length <= 10:
                    base_font_size = 38
                elif text_length <= 15:
                    base_font_size = 34
                elif text_length <= 20:
                    base_font_size = 28
                else:
                    base_font_size = 24
                current_font_size = int(base_font_size * panel_scale)
            else:
                current_font_size = font_size

            font = ImageFont.truetype(font_path, current_font_size) if font_path else ImageFont.load_default()
            line_height = current_font_size + 4

            # 内容驱动的安全留白：图片气泡和异形气泡实际可用区域小于外接矩形，
            # 所以使用比文字描边更大的内边距，不能让字贴着弧边/锯齿边。
            # 气泡不仅要“刚好装下”文字，还应当留出可读的呼吸空间。单人气泡
            # 尤其常用异形/椭圆素材，安全边距再大一些，避免文字贴到弧边。
            # 预留的是文字到“可见气泡边缘”的安全区，而不只是外接矩形的空白。
            # 方形、圆角方形等素材的描边/尾巴会侵占可用空间，因此再扩大一档。
            padding = max(24, int(current_font_size * (1.30 if dialogue_count > 1 else 1.50)))

            # 与 render_text 共用断行逻辑，保证气泡尺寸和实际文字行数完全一致
            lines = _wrap_dialogue_lines(draw, dialogue_text, font, max_width, 11)

            text_width = max(draw.textlength(line, font=font) for line in lines) if lines else 0
            bubble_width = max(text_width + padding * 2, 60)
            bubble_height = len(lines) * line_height + padding * 2

            available_height = bg_height - safe_margin_top - safe_margin_bottom

            while bubble_height > available_height and current_font_size > 14:
                current_font_size -= 2
                font = ImageFont.truetype(font_path, current_font_size) if font_path else ImageFont.load_default()
                line_height = current_font_size + 4
                padding = max(24, int(current_font_size * (1.30 if dialogue_count > 1 else 1.50)))
                lines = _wrap_dialogue_lines(draw, dialogue_text, font, max_width, 11)

                text_width = max(draw.textlength(line, font=font) for line in lines) if lines else 0
                bubble_width = max(text_width + padding * 2, 60)
                bubble_height = len(lines) * line_height + padding * 2
            
            # 计算气泡位置（多个气泡错开显示）
            tail_pos = None
            speaker_target = None
            use_speaker_position = current_speaker and not is_guide and not skip_ip_overlay

            if use_speaker_position:
                if current_speaker in speaker_names:
                    spkr_idx = speaker_names.index(current_speaker)
                else:
                    spkr_char = current_speaker[-1] if current_speaker else ''
                    if spkr_char >= 'A' and spkr_char <= 'Z':
                        spkr_idx = ord(spkr_char) - ord('A')
                    else:
                        spkr_idx = -1

                if spkr_idx >= 0 and spkr_idx < len(ip_images):
                    if spkr_idx in display_indices:
                        display_idx = display_indices.index(spkr_idx)
                        if display_idx < len(ip_positions):
                            speaker_x, speaker_y = ip_positions[display_idx]
                            speaker_w, speaker_h = ip_transform_params[display_idx]['size']
                            # 使用最终合成后的精确坐标和缩放尺寸。
                            if display_idx < len(meta['characters']):
                                placed_char = meta['characters'][display_idx]
                                speaker_x = placed_char['x']
                                speaker_y = placed_char['y']
                                speaker_w = placed_char['width']
                                speaker_h = placed_char['height']
                            speaker_target = (speaker_x, speaker_y, speaker_w, speaker_h)
                            # Bug2 修复：依次尝试头顶/右侧/脚下，避开人物与已绘气泡（含回退策略）
                            bubble_x, bubble_y, tail_pos = _place_speaker_bubble(
                                speaker_x, speaker_y, speaker_w, speaker_h,
                                bubble_width, bubble_height, speaker_x + speaker_w // 2
                            )

            # 完整场景模式：气泡放在画面顶部/底部边缘，避免遮挡中央人物
            if tail_pos is None and skip_ip_overlay:
                cols = 2
                col_idx = idx % cols
                row_idx = idx // cols
                bubble_x = safe_margin_x + col_idx * (bg_width // cols - bubble_width - 10)
                # 优先放顶部，放不下则放底部
                if idx < 2:
                    bubble_y = safe_margin_top + row_idx * (bubble_height + 10)
                    tail_pos = (bubble_x + bubble_width // 2, bubble_y + bubble_height, 'down')
                else:
                    bubble_y = bg_height - bubble_height - safe_margin_bottom - row_idx * 20
                    tail_pos = (bubble_x + bubble_width // 2, bubble_y, 'up')
                bubble_x = max(safe_margin_x, min(bubble_x, bg_width - bubble_width - safe_margin_x))

            # 如果没有找到说话者位置，使用默认位置（错开显示）
            if tail_pos is None:
                cols = 2
                col_idx = idx % cols
                row_idx = idx // cols
                bubble_x = safe_margin_x + col_idx * (bg_width // cols - bubble_width - 10)
                bubble_y = safe_margin_top + row_idx * (bubble_height + 10)

                if bubble_y + bubble_height > bg_height - safe_margin_bottom:
                    bubble_y = bg_height - bubble_height - safe_margin_bottom - row_idx * 20
                tail_pos = (bubble_x + bubble_width // 2, bubble_y + bubble_height, 'down')
            
            # 全画布候选搜索：绝不保留与已有气泡重叠的位置。
            # 最终复核只检查真实脸部边界。局部避让已经确认不遮脸时，
            # 不能仅因不足 15px 的装饰性留白又触发全画布重排。
            if check_overlap(
                bubble_x, bubble_y, bubble_width, bubble_height,
                char_margin=0
            ):
                candidates = []
                step_x = int(max(30, bubble_width // 3))
                step_y = int(max(24, bubble_height // 2))
                for cy in range(int(safe_margin_top), int(max(safe_margin_top + 1, bg_height - bubble_height - safe_margin_bottom + 1)), step_y):
                    for cx in range(int(safe_margin_x), int(max(safe_margin_x + 1, bg_width - bubble_width - safe_margin_x + 1)), step_x):
                        overlap = _overlap_area(cx, cy, bubble_width, bubble_height)
                        bubble_overlap = sum(
                            max(0, min(cx + bubble_width, bx + bw) - max(cx, bx)) *
                            max(0, min(cy + bubble_height, by + bh) - max(cy, by))
                            for bx, by, bw, bh in drawn_bubbles
                        )
                        distance = abs((cx + bubble_width / 2) - (speaker_target[0] + speaker_target[2] / 2)) if speaker_target else 0
                        candidates.append((bubble_overlap > 0, bubble_overlap, distance, overlap, cx, cy))
                if candidates:
                    _, _, _, _, bubble_x, bubble_y = min(candidates)
                    if speaker_target:
                        sx, sy, sw, sh = speaker_target
                        tail_pos = _tail_for(
                            'top_center', bubble_x, bubble_y,
                            bubble_width, bubble_height,
                            sx, sy, sw, sh
                        )
            
            fill_color = bubble_colors[idx % len(bubble_colors)]
            outline_color = outline_colors[idx % len(outline_colors)]
            
            # ===== 气泡图来源：三优先级（显式人物指定 > 该格气泡 > 默认）=====
            # ① 人物显式指定 character_bubble_map[current_speaker] 命中且非「自动」→ 用人物指定气泡
            # ② 该格（panel）逐句手动指定 bubble_images[idx] 非空 → 用该格气泡
            # ③ 否则 → analyze_emotion 自动按语气匹配
            emotion_bubble_path = None

            # 人物气泡：精确匹配当前说话人；并兼容前后端命名/空白差异做一次 trim 兜底
            _char_bubble_value = None
            if isinstance(character_bubble_map, dict) and current_speaker:
                if current_speaker in character_bubble_map:
                    _char_bubble_value = character_bubble_map.get(current_speaker)
                else:
                    _stripped = current_speaker.strip()
                    for _k, _v in character_bubble_map.items():
                        if _k is not None and _k.strip() == _stripped:
                            _char_bubble_value = _v
                            break
            _per_character = (
                _char_bubble_value is not None
                and _char_bubble_value not in ('auto', '', None)
            )

            _per_dialogue = (
                bubble_images is not None
                and isinstance(bubble_images, list)
                and idx < len(bubble_images)
                and bool(bubble_images[idx])
            )

            if _per_character:
                _char_bubble = _char_bubble_value
                # 值为气泡图片 url（/ 或 http 开头）→ 解析为本地路径；否则当作情绪 key
                if isinstance(_char_bubble, str) and (
                    _char_bubble.startswith('/') or _char_bubble.startswith('http')
                ):
                    if _char_bubble.startswith('/'):
                        emotion_bubble_path = os.path.join(
                            app.root_path, _char_bubble.lstrip('/')
                        )
                    else:
                        # http(s) 远程 url：下载到本地临时文件
                        emotion_bubble_path = _download_bubble_to_local(_char_bubble)
                    logger.info(f"[人物气泡] {current_speaker} -> {_char_bubble} (url)")
                else:
                    # 值为情绪 key（happy/sad/angry/surprised/neutral/custom）→ 随机取该情绪气泡
                    emotion_bubble_path = get_random_bubble_image(_char_bubble)
                    logger.info(f"[人物气泡] {current_speaker} -> {_char_bubble} (emotion)")
            elif _per_dialogue:
                emotion_bubble_path = bubble_images[idx]
                logger.info(f"[逐句气泡] idx={idx} speaker={current_speaker} -> {bubble_images[idx]}")
            else:
                emotion = analyze_emotion(dialogue_text)
                emotion_bubble_path = get_random_bubble_image(emotion)

            native_tail_kind = bubble_native_tail_kind(emotion_bubble_path)
            content_scale_x, content_scale_y = bubble_content_scale(emotion_bubble_path)
            # 自带尾巴的 PNG 以“整张图”缩放，尾巴会占据下方一部分像素。
            # 因此扩展图片总高度，并把新增部分留给尾巴；文字仍只在上方主体区域排版。
            # 所有图片气泡都留一点下方安全区；带尾巴的素材再扩大到足够避开尾巴。
            # 这会让气泡整体随内容变大，而不是牺牲字号或把文字压在尾巴上。
            tail_content_inset = int(current_font_size * 0.35) if emotion_bubble_path else 0
            if native_tail_kind:
                tail_content_inset = max(tail_content_inset, max(18, int(current_font_size * 0.9)))

            # 思考气泡的正文主体只占 PNG 上半部。以原中心为锚放大外框，
            # 并在放大后重新执行位置避让，保证增加容量不会换来遮脸或压住其他气泡。
            if content_scale_x != 1.0 or content_scale_y != 1.0:
                old_center_x = bubble_x + bubble_width / 2
                old_center_y = bubble_y + bubble_height / 2
                max_bubble_width = bg_width - safe_margin_x * 2
                max_bubble_height = bg_height - safe_margin_top - safe_margin_bottom
                bubble_width = min(max_bubble_width, bubble_width * content_scale_x)
                bubble_height = min(max_bubble_height, bubble_height * content_scale_y)
                bubble_x = old_center_x - bubble_width / 2
                bubble_y = old_center_y - bubble_height / 2

            if tail_content_inset:
                bubble_height += tail_content_inset
            bubble_width = min(bubble_width, bg_width - safe_margin_x * 2)
            bubble_height = min(
                bubble_height, bg_height - safe_margin_top - safe_margin_bottom
            )

            if speaker_target and (content_scale_x != 1.0 or content_scale_y != 1.0):
                sx, sy, sw, sh = speaker_target
                bubble_x, bubble_y, tail_pos = _place_speaker_bubble(
                    sx, sy, sw, sh, bubble_width, bubble_height, sx + sw / 2
                )
            else:
                bubble_x = max(safe_margin_x, min(
                    bubble_x, bg_width - bubble_width - safe_margin_x
                ))
                bubble_y = max(safe_margin_top, min(
                    bubble_y, bg_height - bubble_height - safe_margin_bottom
                ))
            if native_tail_kind and speaker_target:
                sx, sy, sw, sh = speaker_target
                target_x = sx + sw / 2
                # 02 的原生三角在底边中央；07 的思考泡在左下角。
                # 反推气泡主体位置，让素材自己的尾巴对准人物，而不是另画一个尾巴。
                anchor_ratio = 0.50 if native_tail_kind == 'bottom_center' else 0.18
                native_x = target_x - bubble_width * anchor_ratio
                native_y = sy - bubble_height - 18
                native_x = max(safe_margin_x, min(
                    native_x, bg_width - bubble_width - safe_margin_x
                ))
                native_aligned = False
                for vertical_offset in (0, -16, -32, -48):
                    candidate_y = max(
                        safe_margin_top,
                        min(native_y + vertical_offset,
                            bg_height - bubble_height - safe_margin_bottom)
                    )
                    if not check_overlap(
                        native_x, candidate_y, bubble_width, bubble_height,
                        char_margin=15
                    ):
                        bubble_x, bubble_y = native_x, candidate_y
                        native_aligned = True
                        break
                if native_aligned:
                    # 原生尾巴已经属于气泡图片，渲染器与编辑元数据均不再生成第二条尾巴。
                    tail_pos = None
                else:
                    # 原生尾巴只能朝下。若人物上方没有安全位置，不能把该素材放到
                    # 人物下方导致尾巴反向；改用同情绪的无尾巴圆角素材并保留正确尾巴。
                    source_dir = os.path.dirname(str(emotion_bubble_path))
                    source_name = os.path.basename(str(emotion_bubble_path))
                    prefix = source_name.split('_', 1)[0]
                    fallback_bubble = os.path.join(
                        source_dir, f'{prefix}_01_rounded.png'
                    )
                    if os.path.exists(fallback_bubble):
                        emotion_bubble_path = fallback_bubble
                    native_tail_kind = None
            
            # 气泡与文字始终先画到独立透明图层。它既与最终成品合成，
            # 也保存给编辑器直接复用，避免浏览器再次排版中文后发生偏差。
            bubble_layer = Image.new('RGBA', background.size, (0, 0, 0, 0))
            bubble_draw = ImageDraw.Draw(bubble_layer)
            render_bubble(bubble_draw, bubble_x, bubble_y, bubble_width, bubble_height, 
                          style=bubble_style, fill_color=fill_color, 
                          outline_color=outline_color, opacity=opacity, tail_pos=tail_pos,
                          bubble_image=emotion_bubble_path)
            
            # 根据气泡底色自动选择文字和描边颜色，深色气泡用白字黑边，浅色气泡用黑字白边
            if emotion_bubble_path:
                text_color, stroke_fill = bubble_text_colors(bubble_image_path=emotion_bubble_path)
            else:
                text_color, stroke_fill = bubble_text_colors(fill_color=fill_color)
            
            text_render_info = render_text(
                bubble_draw, dialogue_text, bubble_x, bubble_y, bubble_width, bubble_height,
                font, line_height, padding=padding, align=text_align, font_path=font_path,
                text_color=text_color, stroke_fill=stroke_fill,
                content_bottom_inset=tail_content_inset
            ) or {}

            background = Image.alpha_composite(background, bubble_layer)
            draw = ImageDraw.Draw(background)
            bubble_layer_url = None
            bubble_layer_box = bubble_layer.getbbox()
            if bubble_layer_box:
                try:
                    output_stem = os.path.splitext(os.path.basename(output_path))[0]
                    bubble_layer_name = f"{output_stem}_bubble_{len(meta['bubbles'])}.png"
                    bubble_layer_path = os.path.join(os.path.dirname(output_path), bubble_layer_name)
                    bubble_layer.crop(bubble_layer_box).save(bubble_layer_path, format='PNG')
                    bubble_layer_url = '/' + os.path.relpath(bubble_layer_path, app.root_path).replace('\\', '/')
                except Exception as e:
                    logger.warning(f"保存气泡编辑图层失败: {e}")
            
            meta['bubbles'].append({
                'text': dialogue_text,
                'x': bubble_x,
                'y': bubble_y,
                'width': bubble_width,
                'height': bubble_height,
                'style': bubble_style,
                'bubble_url': _bubble_path_to_url(emotion_bubble_path),
                'fill_color': fill_color[:3] if fill_color else None,
                'outline_color': outline_color,
                'tail_pos': tail_pos,
                'native_tail': native_tail_kind,
                'font_size': text_render_info.get('font_size', current_font_size),
                'font_family': font_family,
                'font_weight': 'bold' if font_family == 'msyhbd' else 'normal',
                'line_height': text_render_info.get('line_height', line_height),
                'padding': padding,
                'content_bottom_inset': tail_content_inset,
                'text_align': text_align,
                'text_color': list(text_color[:3]),
                'stroke_fill': list(stroke_fill[:3]) if stroke_fill else None,
                'stroke_width': text_render_info.get('stroke_width', 0),
                'rendered_text': text_render_info.get('rendered_text', dialogue_text),
                'rendered_layer_url': bubble_layer_url,
                'rendered_layer_x': bubble_layer_box[0] if bubble_layer_box else bubble_x,
                'rendered_layer_y': bubble_layer_box[1] if bubble_layer_box else bubble_y,
                'rendered_layer_width': (bubble_layer_box[2] - bubble_layer_box[0]) if bubble_layer_box else bubble_width,
                'rendered_layer_height': (bubble_layer_box[3] - bubble_layer_box[1]) if bubble_layer_box else bubble_height,
                'speaker': current_speaker
            })
            
            drawn_bubbles.append((bubble_x, bubble_y, bubble_width, bubble_height))
        
        background = _apply_logo_to_image(background, logo_url, logo_position, logo_size)
        background.save(output_path, format='PNG')
        logger.info(f"图片合成成功: {output_path}")
        meta['success'] = True
        return meta
        
    except Exception as e:
        logger.error(f"图片合成失败: {str(e)}")
        return {'success': False, 'error': str(e)}


@app.route('/')
def index():
    """首页路由"""
    return render_template('index.html', guide_config=ROLE_GUIDE_CONFIG)


@app.route('/test')
def test():
    """测试页面"""
    return render_template('test.html')


@app.route('/analyze-content', methods=['POST'])
def analyze_content():
    """分析内容，提取知识点（Step 1）"""
    try:
        content = request.form.get('content', '').strip()
        
        if not content:
            return jsonify({'error': '请输入内容'}), 400
        
        logger.info(f"正在分析内容...")
        
        analysis = analyze_pdf_content(content)
        
        if analysis:
            return jsonify({
                'success': True,
                'analysis': analysis
            })
        else:
            return jsonify({
                'success': True,
                'analysis': {
                    "topic": "课程重点",
                    "key_points": ["知识点1", "知识点2", "知识点3"],
                    "teaching_goal": "帮助孩子理解课程内容",
                    "difficulty": "小学",
                    "story_premise": "小精灵们在有趣的场景中学习知识"
                }
            })
            
    except Exception as e:
        logger.error(f"分析内容失败: {str(e)}")
        return jsonify({'error': f'分析内容失败: {str(e)}'}), 500


@app.route('/generate-story-plan', methods=['POST'])
def generate_story_plan_api():
    """生成故事大纲（Step 2）"""
    try:
        course_topic = request.form.get('topic', '').strip()
        
        if not course_topic:
            return jsonify({'error': '请输入课程重点'}), 400
        
        panel_count = int(request.form.get('panel_count', 6))
        character_count = int(request.form.get('character_count', 3))
        scenario_style = request.form.get('scenario_style', 'humorous')
        
        panel_count = max(4, min(panel_count, 8))
        character_count = max(1, min(character_count, 4))
        
        logger.info(f"生成故事大纲: {course_topic[:50]}..., 格子数: {panel_count}, 角色数: {character_count}, 风格: {scenario_style}")
        
        story_plan = generate_story_plan(
            {'topic': course_topic, 'key_points': [], 'teaching_goal': '', 'story_premise': ''},
            character_count=character_count,
            panel_count=panel_count,
            scenario_style=scenario_style
        )
        
        if story_plan:
            return jsonify({
                'success': True,
                'story_plan': story_plan
            })
        else:
            characters = [f'小精灵{chr(ord("A") + i)}' for i in range(character_count)]
            return jsonify({
                'success': True,
                'story_plan': {
                    "world_setting": {
                        "main_location": "神奇的知识森林",
                        "time_period": "现代",
                        "visual_style": "Q版卡通，明亮色彩，奇幻风格",
                        "atmosphere": "有趣、探索"
                    },
                    "characters": [{"name": char, "role": "小精灵", "personality": "可爱"} for char in characters],
                    "story_arc": [{"panel": i+1, "location": "神奇的知识森林", "action": f"第{i+1}格场景", "mood": "开心"} for i in range(panel_count)]
                }
            })
            
    except Exception as e:
        logger.error(f"生成故事大纲失败: {str(e)}")
        return jsonify({'error': f'生成故事大纲失败: {str(e)}'}), 500


@app.route('/generate-panel-script', methods=['POST'])
def generate_panel_script_api():
    """生成单格内容（Step 3）"""
    try:
        story_plan_str = request.form.get('story_plan', '')
        panel_index = int(request.form.get('panel_index', 0))
        scenario_style = request.form.get('scenario_style', 'humorous')
        
        if not story_plan_str:
            return jsonify({'error': '请提供故事大纲'}), 400
        
        import json
        story_plan = json.loads(story_plan_str)
        
        panel_content = generate_panel_content(story_plan, panel_index, scenario_style)
        
        if panel_content:
            return jsonify({
                'success': True,
                'panel_content': panel_content,
                'panel_index': panel_index
            })
        else:
            return jsonify({
                'success': True,
                'panel_content': {
                    "scene": "cartoon background, Q version style, bright colors",
                    "scene_zh": "卡通背景",
                    "dialogue": "小精灵在说话",
                    "speaker": "小精灵A",
                    "hint": "思考一下这个问题"
                },
                'panel_index': panel_index
            })
            
    except Exception as e:
        logger.error(f"生成单格内容失败: {str(e)}")
        return jsonify({'error': f'生成单格内容失败: {str(e)}'}), 500


@app.route('/generate', methods=['POST'])
def generate_comic():
    """生成漫画的API接口（支持多格数、多角色、情景风格）"""
    try:
        course_topic = request.form.get('topic', '').strip()
        
        comic_type = request.form.get('comic_type', 'preview')
        
        file_content = request.form.get('file_content', '').strip()
        if file_content and not course_topic:
            course_topic = file_content[:500]
        elif file_content and course_topic:
            course_topic = f"{course_topic}\n\n参考内容：{file_content[:500]}"
        
        if not course_topic:
            return jsonify({'error': '请输入课程重点或上传文件'}), 400
        
        panel_count = int(request.form.get('panel_count', 4))
        character_count = int(request.form.get('character_count', 1))
        scenario_style = request.form.get('scenario_style', 'humorous')
        
        panel_count = max(4, min(panel_count, 8))
        character_count = max(1, min(character_count, 4))
        
        logger.info(f"收到课程重点: {course_topic[:100]}..., 漫画类型: {comic_type}, 格子数: {panel_count}, 角色数: {character_count}, 风格: {scenario_style}")
        
        style_seed = str(uuid.uuid4())[:8]
        
        world_setting = None
        
        story_plan_str = request.form.get('story_plan', '')
        if story_plan_str:
            try:
                import json
                story_plan = json.loads(story_plan_str)
                world_setting = story_plan.get('world_setting', None)
                logger.info("使用预先生成的故事大纲")
            except json.JSONDecodeError as e:
                logger.warning(f"解析故事大纲失败: {str(e)}")
                story_plan = None
        
        script_str = request.form.get('script', '')
        script = None
        
        if script_str:
            try:
                import json
                script = json.loads(script_str)
                logger.info("使用预先生成的脚本")
            except json.JSONDecodeError as e:
                logger.warning(f"解析预先生成的脚本失败: {str(e)}")
        
        if not script or 'scenes' not in script or 'dialogues' not in script:
            logger.info("生成新脚本...")
            script = generate_comic_script(
                course_topic, 
                comic_type, 
                panel_count=panel_count, 
                character_count=character_count, 
                scenario_style=scenario_style
            )
            
            if not script or 'scenes' not in script or 'dialogues' not in script:
                return jsonify({'error': '生成剧本失败'}), 500
        
        scenes = script['scenes']
        dialogues = script['dialogues']
        speakers = script.get('speakers', [])
        sound_effects = script.get('sound_effects', [])
        angles = script.get('angles', [])
        
        while len(scenes) < panel_count:
            scenes.append("cute cartoon background, Q version style, bright colors")
        while len(dialogues) < panel_count:
            dialogues.append("学习真有趣！")
        while len(speakers) < panel_count:
            speakers.append(f"小精灵{chr(ord('A') + len(speakers) % character_count)}")
        while len(sound_effects) < panel_count:
            sound_effects.append("")
        while len(angles) < panel_count:
            angles.append("eye level")
        
        use_uploaded_image = request.form.get('use_image', 'false') == 'true'
        image_file = request.form.get('image_file', '')
        
        logger.info("步骤2: 获取IP图片...")
        ip_paths_str = request.form.get('ip_paths', '')
        ip_paths = []
        
        if ip_paths_str:
            selected_ip_filenames = ip_paths_str.split(',')
            for filename in selected_ip_filenames:
                filename = filename.strip()
                if filename:
                    ip_path = os.path.join(IP_DIR, filename)
                    if os.path.exists(ip_path):
                        ip_paths.append(ip_path)
                    else:
                        uploaded_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                        if os.path.exists(uploaded_path):
                            ip_paths.append(uploaded_path)
        
        if not ip_paths:
            for i in range(character_count):
                ip_path = get_random_ip_image()
                if ip_path:
                    ip_paths.append(ip_path)
        
        if not ip_paths:
            return jsonify({'error': '未找到IP图片，请确保static/ip/目录下有PNG图片或上传IP图片'}), 500
        
        logger.info(f"使用的IP图片: {len(ip_paths)}个")
        
        poses_str = request.form.get('poses', '')
        poses = []
        if poses_str:
            try:
                import json
                poses = json.loads(poses_str)
            except json.JSONDecodeError:
                poses = []
        
        logger.info("步骤3: 生成并合成图片...")
        result_images = []
        
        for i in range(panel_count):
            scene_prompt = scenes[i] if i < len(scenes) else "cute cartoon background"
            angle = angles[i] if i < len(angles) else None
            
            environment_prompt = extract_environment_only(scene_prompt)
            if environment_prompt:
                logger.info(f"场景描述清洗: '{scene_prompt[:50]}...' -> '{environment_prompt[:50]}...'")
                scene_prompt = environment_prompt
            
            if use_uploaded_image and image_file:
                uploaded_image_path = os.path.join(app.config['UPLOAD_FOLDER'], image_file)
                if os.path.exists(uploaded_image_path):
                    try:
                        with open(uploaded_image_path, 'rb') as f:
                            background_data = f.read()
                    except Exception as e:
                        logger.error(f"读取上传图片失败: {str(e)}")
                        background_data = generate_image(scene_prompt, i, world_setting, style_seed, angle=angle, force_pure_background=True, block_names=list(dict.fromkeys(speakers)) if speakers else None)
                else:
                    background_data = generate_image(scene_prompt, i, world_setting, style_seed, angle=angle, force_pure_background=True, block_names=list(dict.fromkeys(speakers)) if speakers else None)
            else:
                background_data = generate_image(scene_prompt, i, world_setting, style_seed, angle=angle, force_pure_background=True, block_names=list(dict.fromkeys(speakers)) if speakers else None)
            
            output_filename = f"{uuid.uuid4().hex}_{i}.png"
            output_path = os.path.join(OUTPUT_DIR, output_filename)
            
            dialogue = dialogues[i] if i < len(dialogues) else ""
            speaker = speakers[i] if i < len(speakers) else ""
            sound_effect = sound_effects[i] if i < len(sound_effects) else ""
            
            current_poses = poses[i] if i < len(poses) else None
            pose_list = [current_poses] * len(ip_paths) if current_poses else None
            
            success = composite_image(
                background_data, 
                ip_paths, 
                dialogue, 
                output_path, 
                speaker=speaker,
                poses=pose_list
            )
            
            if success:
                result_images.append(f"/static/output/{output_filename}")
            else:
                logger.warning(f"第{i+1}张图片合成失败")
        
        if not result_images:
            return jsonify({'error': '图片合成失败'}), 500
        
        logger.info(f"漫画生成完成，共{len(result_images)}张图片")
        
        response_data = {
            'success': True,
            'images': result_images,
            'topic': course_topic[:100],
            'panel_count': panel_count,
            'character_count': character_count,
            'scenario_style': scenario_style
        }
        
        if 'summary' in script:
            response_data['summary'] = script['summary']
        
        return jsonify(response_data)
        
    except Exception as e:
        logger.error(f"生成漫画时发生错误: {str(e)}")
        return jsonify({'error': f'生成漫画失败: {str(e)}'}), 500


@app.route('/static/output/<filename>')
def serve_output(filename):
    """提供生成的漫画图片"""
    return send_from_directory(OUTPUT_DIR, filename)


def detect_file_type(file_stream, filename):
    """通过文件头检测文件类型"""
    file_stream.seek(0)
    header = file_stream.read(128)
    file_stream.seek(0)
    
    magic_numbers = {
        b'%PDF-': 'pdf',
        b'\x50\x4B\x03\x04': 'docx',
        b'\xD0\xCF\x11\xE0': 'doc',
        b'GIF87a': 'gif',
        b'GIF89a': 'gif',
        b'\xFF\xD8\xFF': 'jpg',
        b'\x89\x50\x4E\x47': 'png',
        b'BM': 'bmp',
        b'RIFF': 'bmp',
        b'{\rtf': 'rtf',
        b'\\rtf': 'rtf',
    }
    
    for magic, ext in magic_numbers.items():
        if header.startswith(magic):
            return ext
    
    if header[:4] == b'\x50\x4B\x03\x04':
        content_type = file_stream.content_type if hasattr(file_stream, 'content_type') else ''
        if 'word' in content_type.lower():
            return 'docx'
        return 'docx'
    
    return None


def allowed_file(filename):
    """检查文件扩展名是否允许"""
    if not filename or '.' not in filename or filename.startswith('.'):
        return False
    parts = filename.rsplit('.', 1)
    if len(parts) < 2:
        return False
    return parts[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_pdf(pdf_path):
    """从PDF中提取文本内容（支持文字版和图片版PDF）"""
    try:
        from pdf_extractor import extract_text_with_mineru_style, RAPIDOCRPDF_AVAILABLE, PDF_OXIDE_AVAILABLE
        
        if RAPIDOCRPDF_AVAILABLE:
            logger.info(f"🚀 使用RapidOCRPDF引擎处理PDF: {pdf_path}")
        elif PDF_OXIDE_AVAILABLE:
            logger.info(f"🚀 使用pdf_oxide引擎处理PDF: {pdf_path}")
        else:
            logger.info(f"使用PyMuPDF引擎处理PDF: {pdf_path}")
            
        text = extract_text_with_mineru_style(pdf_path)
        
        if text:
            logger.info(f"✅ PDF提取成功，共{len(text)}个字符")
            return text
        
        logger.warning("所有提取方法均失败，PDF可能是扫描版且OCR不可用")
        return None
        
    except Exception as e:
        logger.error(f"PDF读取失败: {str(e)}")
        return None


def extract_text_from_scanned_pdf(pdf_path):
    """使用OCR从扫描版PDF中提取文本（支持多种OCR工具）"""
    if not PDF2IMAGE_AVAILABLE:
        logger.warning("PyMuPDF未安装，无法处理图片版PDF")
        return None
    
    try:
        logger.info("正在使用PyMuPDF将PDF转换为图片...")
        
        doc = fitz.open(pdf_path)
        total_pages = doc.page_count
        logger.info(f"PDF共{total_pages}页")
        
        max_pages = min(total_pages, 10)
        if total_pages > max_pages:
            logger.info(f"为加快速度，仅识别前{max_pages}页")
        
        # 获取可用的OCR工具列表（豆包API优先，识别率最高）
        available_ocr_tools = []
        if OCR_DOUBAO_AVAILABLE:
            available_ocr_tools.append('doubao')
        if OCR_EASYOCR_AVAILABLE:
            available_ocr_tools.append('easyocr')
        if OCR_TESSERACT_AVAILABLE:
            available_ocr_tools.append('tesseract')
        if OCR_BAIDU_AVAILABLE:
            available_ocr_tools.append('baidu')
        
        if not available_ocr_tools:
            logger.warning("没有可用的OCR工具")
            return None
        
        logger.info(f"可用OCR工具: {available_ocr_tools}")
        
        # 优先使用的OCR工具（豆包API识别率最高）
        primary_ocr = available_ocr_tools[0]
        logger.info(f"使用{primary_ocr}进行OCR识别")
        
        def ocr_page(page_index):
            try:
                logger.info(f"正在识别第{page_index+1}页...")
                page = doc[page_index]
                
                pix = page.get_pixmap(dpi=150)
                
                if pix.width * pix.height < 10000:
                    return None
                
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                gray = img.convert('L')
                pixels = gray.getdata()
                avg_brightness = sum(pixels) / len(pixels)
                if avg_brightness > 240:
                    logger.info(f"第{page_index+1}页为空白页，跳过")
                    return None
                
                page_text = ""
                
                if primary_ocr == 'doubao':
                    try:
                        page_text = doubao_ocr(img)
                    except Exception as e:
                        logger.error(f"豆包OCR识别失败: {str(e)}")
                
                elif primary_ocr == 'easyocr':
                    try:
                        global easyocr_reader
                        if easyocr_reader is None:
                            logger.info("正在加载EasyOCR模型...")
                            easyocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False, 
                                model_storage_directory=os.path.join(os.path.dirname(__file__), 'models'))
                            logger.info("EasyOCR模型加载成功")
                        results = easyocr_reader.readtext(img)
                        page_text = "\n".join([result[1] for result in results])
                    except Exception as e:
                        logger.error(f"EasyOCR识别失败: {str(e)}")
                
                elif primary_ocr == 'baidu':
                    try:
                        page_text = baidu_ocr(img)
                    except Exception as e:
                        logger.error(f"百度OCR识别失败: {str(e)}")
                
                elif primary_ocr == 'tesseract':
                    try:
                        page_text = pytesseract.image_to_string(img, lang='chi_sim+eng', 
                            config='--psm 6 --oem 3')
                    except Exception as e:
                        logger.error(f"Tesseract识别失败: {str(e)}")
                
                if page_text and len(page_text.strip()) > 10:
                    return f"--- 第{page_index+1}页 ---\n{page_text}\n\n"
                return None
            except Exception as e:
                logger.error(f"第{page_index+1}页识别失败: {str(e)}")
                return None
        
        text = ""
        
        # 豆包API、EasyOCR和百度OCR不是线程安全的，需要顺序执行
        if primary_ocr in ['doubao', 'easyocr', 'baidu']:
            for i in range(max_pages):
                result = ocr_page(i)
                if result:
                    text += result
        else:
            # Tesseract是线程安全的，可以并行执行
            with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
                # ContextVar 默认不会自动传入线程池；逐页复制上下文，确保 OCR 使用
                # 发起请求的用户密钥，而不是其他账号或服务器全局密钥。
                futures = {
                    executor.submit(copy_context().run, ocr_page, i): i
                    for i in range(max_pages)
                }
                
                for future in concurrent.futures.as_completed(futures):
                    result = future.result()
                    if result:
                        text += result
        
        text = text.strip()
        
        if text:
            logger.info(f"OCR成功提取文本，共{len(text)}个字符")
            return text
        else:
            logger.warning("OCR未能识别出文本")
            return None
            
    except Exception as e:
        logger.error(f"OCR识别失败: {str(e)}")
        return None


def baidu_ocr(image):
    """使用百度OCR API识别图片中的文字"""
    try:
        import base64
        
        # 获取access_token
        token_url = 'https://aip.baidubce.com/oauth/2.0/token'
        params = {
            'grant_type': 'client_credentials',
            'client_id': BAIDU_OCR_API_KEY,
            'client_secret': BAIDU_OCR_SECRET_KEY
        }
        response = requests.get(token_url, params=params)
        response.raise_for_status()
        access_token = response.json().get('access_token')
        
        # 将图片转为base64
        from io import BytesIO
        buffer = BytesIO()
        image.save(buffer, format='PNG')
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        
        # 调用OCR API
        ocr_url = 'https://aip.baidubce.com/rest/2.0/ocr/v1/general_basic'
        headers = {'Content-Type': 'application/x-www-form-urlencoded'}
        data = {
            'image': image_base64,
            'access_token': access_token
        }
        response = requests.post(ocr_url, headers=headers, data=data)
        response.raise_for_status()
        
        result = response.json()
        words_result = result.get('words_result', [])
        
        text = "\n".join([item.get('words', '') for item in words_result])
        return text
        
    except Exception as e:
        logger.error(f"百度OCR API调用失败: {str(e)}")
        return None


def doubao_ocr(image):
    """使用豆包API视觉理解识别图片中的文字"""
    if not OCR_DOUBAO_AVAILABLE:
        logger.warning("豆包API未配置")
        return None
    
    try:
        import base64
        
        from io import BytesIO
        buffer = BytesIO()
        image.save(buffer, format='PNG')
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        
        url = "https://ark.cn-beijing.volces.com/api/v3/responses"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {ARK_API_KEY}"
        }
        
        data = {
            "model": DOUBAO_MODEL_ID,
            "input": [
                {"type": "text", "text": "请识别这张图片中的所有文字，并以纯文本形式返回，不要添加任何解释说明。"},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_base64}"}}
            ]
        }
        
        response = requests.post(url, headers=headers, json=data, timeout=60)
        response.raise_for_status()
        
        result = response.json()
        content = ""
        for output_item in result.get('output', []):
            if output_item.get('type') == 'message':
                for content_item in output_item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        content = content_item.get('text', '')
                        break
                if content:
                    break
        
        if content:
            logger.info(f"豆包OCR识别成功，提取{len(content)}个字符")
            return content.strip()
        else:
            logger.warning("豆包OCR未能识别出文本")
            return None
            
    except Exception as e:
        logger.error(f"豆包OCR API调用失败: {str(e)}")
        return None


def generate_comic_script_with_doubao(pdf_text, ip_images, comic_type='preview'):
    """使用豆包API生成四格漫画脚本（支持文件内容和IP图片）"""
    if not OCR_DOUBAO_AVAILABLE:
        logger.warning("豆包API未配置")
        return None
    
    try:
        import base64
        
        url = "https://ark.cn-beijing.volces.com/api/v3/responses"
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {ARK_API_KEY}"
        }
        
        # 准备内容
        input_content = []
        
        # 添加IP图片
        for ip_path in ip_images[:4]:
            if os.path.exists(ip_path):
                with open(ip_path, 'rb') as f:
                    ip_base64 = base64.b64encode(f.read()).decode('utf-8')
                    input_content.append({
                        "type": "image_url", 
                        "image_url": {"url": f"data:image/png;base64,{ip_base64}"}
                    })
        
        # 添加提示词
        type_desc = "课前预习" if comic_type == 'preview' else "课后复习"
        prompt = f"""我要制作{type_desc}漫画，请识别这个PDF内容，然后为我创作一个四格漫画脚本。

【PDF内容】
{pdf_text[:5000]}

【要求】
1. 图片内是漫画主人公，要严格用上漫画主人公作为主角，不要添加其他角色
2. 制作四格漫画
3. 请以JSON格式输出，包含以下字段：
   - scenes: 4个场景的中文描述（用于生成背景图）
   - scenes_zh: 4个场景的中文描述
   - dialogues: 4句中文对白（包含知识点的俏皮话）
4. 确保输出是合法的JSON格式，不要包含任何其他文字

例如：
{{
  "scenes": [
    "明亮的教室，学生们在读书",
    "老师在讲解课程",
    "学生们分组讨论",
    "开心的学生们举手"
  ],
  "scenes_zh": [
    "明亮的教室，学生们在读书",
    "老师在讲解课程",
    "学生们分组讨论",
    "开心的学生们举手"
  ],
  "dialogues": [
    "同学们，今天我们学习借景抒情！",
    "借景抒情就是通过景物表达情感哦~",
    "大家一起来找找课文里的景物吧！",
    "原来借景抒情这么有趣呀！"
  ]
}}
"""
        
        input_content.append({"type": "text", "text": prompt})
        
        data = {
            "model": DOUBAO_MODEL_ID,
            "input": input_content if len(input_content) > 1 else input_content[0]["text"]
        }
        
        logger.info("正在调用豆包API生成漫画脚本...")
        response = requests.post(url, headers=headers, json=data, timeout=120)
        response.raise_for_status()
        
        result = response.json()
        
        content = ""
        for output_item in result.get('output', []):
            if output_item.get('type') == 'message':
                for content_item in output_item.get('content', []):
                    if content_item.get('type') == 'output_text':
                        content = content_item.get('text', '')
                        break
                if content:
                    break
        
        if content:
            logger.info(f"豆包API生成脚本成功")
            # 提取JSON部分
            try:
                import json
                start = content.find('{')
                end = content.rfind('}') + 1
                if start != -1 and end != -1:
                    json_str = content[start:end]
                    script = json.loads(json_str)
                    return script
                else:
                    logger.warning("未能提取JSON格式的脚本")
                    return None
            except json.JSONDecodeError as e:
                logger.error(f"JSON解析失败: {str(e)}")
                return None
        else:
            logger.warning("豆包API未能生成脚本")
            return None
            
    except Exception as e:
        logger.error(f"豆包API生成脚本失败: {str(e)}")
        return None


def generate_comic_script_with_doubao_v2(course_topic, comic_type='preview', panel_count=4, character_count=1, scenario_style='humorous', ip_images=None, world_setting=None, character_system='student'):
    """使用豆包API生成高质量教育漫画剧本（支持多格数、多角色、多种情景风格、图片输入）
    
    参数:
        course_topic: 课程重点内容
        comic_type: 'preview'（课前预习）或 'summary'（课后复习）
        panel_count: 漫画格子数（4/6/8）
        character_count: 主角数量（1-4）
        scenario_style: 情景风格
        ip_images: IP图片路径列表（可选，用于让模型识别角色形象）
        world_setting: 世界设定字典（包含地点、时间、视觉风格、氛围）
        character_system: 角色体系（student/explorer/time_traveler/classic/elves）
    """
    if not OCR_DOUBAO_AVAILABLE:
        logger.warning("豆包API未配置")
        return None
    
    import random
    import time
    
    creative_directions = [
        '太空冒险', '海底世界', '魔法学院', '童话森林', 
        '未来城市', '恐龙乐园', '糖果王国', '机器人世界',
        '神秘岛屿', '古代城堡', '梦幻仙境', '星际旅行',
        '热带雨林', '冰雪世界', '昆虫王国', '音乐星球'
    ]
    random_direction = random.choice(creative_directions)
    random_seed = int(time.time() * 1000)
    
    url = "https://ark.cn-beijing.volces.com/api/v3/responses"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {ARK_API_KEY}"
    }
    
    input_content = []
    
    if ip_images and isinstance(ip_images, list) and len(ip_images) > 0:
        import base64
        for ip_path in ip_images[:4]:
            if os.path.exists(ip_path):
                with open(ip_path, 'rb') as f:
                    ip_base64 = base64.b64encode(f.read()).decode('utf-8')
                    input_content.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{ip_base64}"}
                    })
    
    style_desc = SCENARIO_STYLES.get(scenario_style, '轻松有趣的风格')
    
    system = CHARACTER_SYSTEMS.get(character_system, CHARACTER_SYSTEMS['student'])
    system_roles = system['roles'][:character_count]
    characters = [f'{system_roles[i]}{chr(ord("A") + i)}' for i in range(character_count)]
    
    panel_structure_desc = generate_panel_structure_desc(panel_count)
    
    style_instructions = {
        'humorous': f'搞笑风格：{system["name"]}有滑稽的行为，对话幽默搞笑',
        'adventure': f'冒险风格：{system["name"]}在探险场景中发现知识',
        'friendship': f'友情风格：通过友情故事学习知识，{system["name"]}互相帮助',
        'magic': f'魔法风格：用魔法比喻解释知识，{system["name"]}使用魔法道具',
        'school': f'校园风格：课堂上的师生互动，角色扮演学生和老师',
        'learning': f'学习风格：通过探索和发现学习知识',
        'animal': f'动物风格：角色拟人化为可爱的动物',
        'family': f'家庭风格：在家庭场景中学习知识'
    }
    
    style_instruction = style_instructions.get(scenario_style, style_instructions['humorous'])
    
    if comic_type == 'preview':
        prompt = f"""你是专业教育漫画编剧，擅长创作小学生爱看的搞笑漫画！

【课程重点】{course_topic}

【目标】让小学生看完后超级好奇，迫不及待想听课！

【要求】
- 类型：埋下伏笔型，绝对不要直接给答案！
- 风格：{style_desc}，{style_instruction}，搞笑有趣，像《阿衰》一样好玩
- 角色：{', '.join(characters)} + 知识小精灵（仅最后一格角落出场）
- 格子：{panel_count}格，{system['visual_style']}

【四段式结构】
{panel_structure_desc}

【创作要点】
1. 用日常场景引入，出现奇怪现象（比如气球自己飞、杯子倒过来水不漏）
2. 角色对话要搞笑、夸张，像好朋友聊天一样
3. 多用疑问句制造悬念："为什么？""怎么回事？""这不可能吧！"
4. 最后一格知识小精灵提出有趣的思考题，引发好奇心

【对话风格】
- 每格对话≤15字，简短有力！
- 语言像小学生说话，用口语化表达
- 可以用感叹号增强情绪！
- 多用网络流行词或童趣表达

【场景描述要求】
每个场景描述必须包含角色的具体动作和位置，例如："公园草地上，{characters[0]}蹲在地上好奇地看着自己飞起来的气球，{characters[1]}站在旁边惊讶地指着气球"

【示例（以"牛顿第三定律"为例）】
第1格：公园草地上，{characters[0]}蹲在地上看着自己飞起来的气球，{characters[1]}站在旁边惊讶地指着
第2格：{characters[0]}说："气球自己会飞？！"
第3格：{characters[1]}说："谁在推它呀？"
第4格：知识小精灵说："到底是什么力量？"

【规则】
- 前{panel_count-1}格说话者为{', '.join(characters)}之一
- 第{panel_count}格说话者为"知识小精灵"
- 输出JSON：scenes, scenes_zh, dialogues, speakers, hints"""
    else:
        prompt = f"""你是专业教育漫画编剧，擅长创作小学生易懂的复习漫画！

【课程重点】{course_topic}

【目标】让小学生轻松复习知识点，记得牢、用得上！

【要求】
- 类型：点题型，明确讲解知识点，给出清晰答案！
- 风格：{style_desc}，{style_instruction}，轻松有趣，像《父与子》一样温馨
- 角色：{', '.join(characters)} + 知识小精灵（仅最后一格角落出场）
- 格子：{panel_count}格，{system['visual_style']}

【四段式结构】
{panel_structure_desc}

【创作要点】
1. 回顾课堂内容，用简单例子解释核心知识点
2. 对话要明确、易懂，像老师在讲解但更有趣
3. 用生活中的例子让知识变得亲切
4. 最后一格知识小精灵提出复习思考题，巩固记忆

【对话风格】
- 每格对话≤15字，清晰明了！
- 语言简单，不用专业术语
- 多用肯定句："对！""没错！""记住了！"
- 可以加入鼓励的话增强信心

【场景描述要求】
每个场景描述必须包含角色的具体动作和位置，例如："明亮的教室里，{characters[0]}站在黑板前指着上面的公式，{characters[1]}坐在课桌前认真听讲"

【示例（以"牛顿第三定律"为例）】
第1格：明亮的教室里，{characters[0]}站在黑板前指着公式，{characters[1]}坐在课桌前认真听讲
第2格：{characters[0]}说："作用力和反作用力！"
第3格：{characters[1]}说："划船就是这个原理！"
第4格：知识小精灵说："推墙时墙也推你！"

【规则】
- 前{panel_count-1}格说话者为{', '.join(characters)}之一
- 第{panel_count}格说话者为"知识小精灵"
- 输出JSON：scenes, scenes_zh, dialogues, speakers, summary"""
    
    input_content.append({"type": "text", "text": prompt})
    
    data = {
        "model": DOUBAO_MODEL_ID,
        "input": input_content if len(input_content) > 1 else input_content[0]["text"],
        "temperature": 0.7,
        "seed": random_seed,
        "max_tokens": 2000
    }
    
    max_retries = 2
    retry_delay = 2
    
    for attempt in range(max_retries + 1):
        try:
            logger.info(f"正在调用豆包API生成漫画剧本（第{attempt+1}次尝试）...")
            response = requests.post(url, headers=headers, json=data, timeout=90)
            response.raise_for_status()
            
            result = response.json()
            
            content = ""
            for output_item in result.get('output', []):
                if output_item.get('type') == 'message':
                    for content_item in output_item.get('content', []):
                        if content_item.get('type') == 'output_text':
                            content = content_item.get('text', '')
                            break
                    if content:
                        break
            
            if content:
                logger.info(f"豆包API生成剧本成功")
                try:
                    import json
                    start = content.find('{')
                    end = content.rfind('}') + 1
                    if start != -1 and end != -1:
                        json_str = content[start:end]
                        script = json.loads(json_str)
                        
                        if comic_type == 'preview' and 'hints' not in script:
                            script['hints'] = ["带着问题去学习吧！"] * panel_count
                        if comic_type == 'summary' and 'summary' not in script:
                            script['summary'] = f"【知识点总结】{course_topic}"
                        
                        if 'dialogues' in script and isinstance(script['dialogues'], list):
                            processed_dialogues = []
                            processed_speakers = []
                            for d in script['dialogues']:
                                if isinstance(d, list):
                                    processed_dialogues.append(' '.join(d))
                                    processed_speakers.append(script.get('speakers', ['小精灵A'] * panel_count)[0])
                                else:
                                    processed_dialogues.append(str(d))
                                    processed_speakers.append(script.get('speakers', ['小精灵A'] * panel_count)[0])
                            script['dialogues'] = processed_dialogues
                            if len(processed_speakers) == len(processed_dialogues):
                                script['speakers'] = processed_speakers
                        
                        return script
                    else:
                        logger.warning("未能提取JSON格式的脚本")
                except json.JSONDecodeError as e:
                    logger.error(f"JSON解析失败: {str(e)}")
                if attempt < max_retries:
                    logger.info(f"等待{retry_delay}秒后重试...")
                    time.sleep(retry_delay)
                    continue
            else:
                logger.warning("豆包API未能生成脚本")
                if attempt < max_retries:
                    logger.info(f"等待{retry_delay}秒后重试...")
                    time.sleep(retry_delay)
                    continue
                
        except requests.exceptions.Timeout:
            logger.error(f"豆包API调用超时（第{attempt+1}次）")
            if attempt < max_retries:
                logger.info(f"等待{retry_delay}秒后重试...")
                time.sleep(retry_delay)
                continue
        except requests.exceptions.RequestException as e:
            logger.error(f"豆包API请求失败（第{attempt+1}次）: {str(e)}")
            if attempt < max_retries:
                logger.info(f"等待{retry_delay}秒后重试...")
                time.sleep(retry_delay)
                continue
        except Exception as e:
            logger.error(f"豆包API生成剧本异常（第{attempt+1}次）: {str(e)}")
            if attempt < max_retries:
                logger.info(f"等待{retry_delay}秒后重试...")
                time.sleep(retry_delay)
                continue
    
    logger.error("豆包API生成剧本失败，已达到最大重试次数")
    return None


def preprocess_image(image):
    """预处理图片以提高OCR识别率"""
    try:
        # 转换为灰度图
        img = image.convert('L')
        
        # 二值化处理
        threshold = 128
        img = img.point(lambda x: 0 if x < threshold else 255, '1')
        
        # 去除噪点（可选）
        # 使用简单的中值滤波
        img = img.filter(Image.MedianFilter(size=3))
        
        return img
    except Exception as e:
        logger.error(f"图片预处理失败: {str(e)}")
        return image


def extract_text_from_docx(docx_path):
    """从Word文档中提取文本"""
    try:
        from docx import Document
        doc = Document(docx_path)
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text.append(cell.text)
        return '\n'.join(text)
    except Exception as e:
        logger.error(f"读取Word文档失败: {str(e)}")
        return None


def extract_text_from_text_file(filepath):
    """从文本文件中提取文本（txt, md, rtf）"""
    try:
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        logger.error(f"无法解码文件: {filepath}")
        return None
    except Exception as e:
        logger.error(f"读取文本文件失败: {str(e)}")
        return None


def extract_text_from_image(image_path):
    """使用OCR从图片中提取文本（优先使用豆包API）"""
    available_tools = []
    if OCR_DOUBAO_AVAILABLE:
        available_tools.append('doubao')
    if OCR_EASYOCR_AVAILABLE:
        available_tools.append('easyocr')
    if OCR_TESSERACT_AVAILABLE:
        available_tools.append('tesseract')
    
    if not available_tools:
        logger.warning("OCR功能不可用")
        return None
    
    try:
        logger.info(f"正在识别图片: {image_path}")
        
        img = Image.open(image_path)
        
        text = ""
        primary_tool = available_tools[0]
        
        if primary_tool == 'doubao':
            text = doubao_ocr(img)
        elif primary_tool == 'easyocr':
            global easyocr_reader
            if easyocr_reader is None:
                logger.info("正在加载EasyOCR模型...")
                easyocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False, 
                    model_storage_directory=os.path.join(os.path.dirname(__file__), 'models'))
            results = easyocr_reader.readtext(img)
            text = "\n".join([result[1] for result in results])
        else:
            processed_img = preprocess_image(img)
            text = pytesseract.image_to_string(processed_img, lang='chi_sim+eng')
        
        text = text.strip()
        
        if text:
            logger.info(f"图片OCR成功提取文本，共{len(text)}个字符")
            return text
        else:
            logger.warning("图片OCR未能识别出文本")
            return None
            
    except Exception as e:
        logger.error(f"图片OCR识别失败: {str(e)}")
        return None


@app.route('/upload', methods=['POST'])
def upload_file():
    """上传PDF或图片文件（快速保存，不执行OCR）"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': '没有文件被上传'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': '没有选择文件'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': '不支持的文件类型，仅支持PDF和图片格式'}), 400
        
        # 生成唯一文件名
        ext = ''
        if '.' in file.filename and not file.filename.startswith('.'):
            parts = file.filename.rsplit('.', 1)
            if len(parts) > 1:
                ext = parts[1].lower()
        filename = f"{uuid.uuid4().hex}.{ext}"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # 保存文件
        file.save(filepath)
        logger.info(f"文件上传成功: {filepath}")
        
        return jsonify({
            'success': True,
            'type': 'pdf' if ext == 'pdf' else 'image',
            'filename': filename,
            'preview_url': f'/static/uploads/{filename}'
        })
        
    except Exception as e:
        logger.error(f"文件上传失败: {str(e)}")
        return jsonify({'error': f'文件处理失败: {str(e)}'}), 500


@app.route('/api/upload_logo', methods=['POST'])
def api_upload_logo():
    """上传 Logo 水印图片，保存到 static/logos/ 目录"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': '没有选择文件'}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '没有选择文件'}), 400

        # 仅允许 PNG / JPEG
        allowed_logo_exts = {'png', 'jpg', 'jpeg'}
        ext = ''
        if '.' in file.filename and not file.filename.startswith('.'):
            parts = file.filename.rsplit('.', 1)
            if len(parts) > 1:
                ext = parts[1].lower()
        if ext not in allowed_logo_exts:
            return jsonify({'error': '仅支持 PNG / JPEG 格式的 Logo'}), 400

        filename = f"{uuid.uuid4().hex}.{ext}"
        filepath = os.path.join(LOGO_DIR, filename)
        file.save(filepath)
        logger.info(f"Logo 上传成功: {filepath}")

        return jsonify({
            'success': True,
            'url': f'/static/logos/{filename}'
        })

    except Exception as e:
        logger.error(f"Logo 上传失败: {str(e)}")
        return jsonify({'error': f'Logo 上传失败: {str(e)}'}), 500


@app.route('/process-file/<filename>', methods=['POST'])
def process_file(filename):
    """处理文件（提取文本，带进度）"""
    try:
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        if not os.path.exists(filepath):
            return jsonify({'error': '文件不存在'}), 404
        
        ext = ''
        if '.' in filename and not filename.startswith('.'):
            parts = filename.rsplit('.', 1)
            if len(parts) > 1:
                ext = parts[1].lower()
        
        if ext == 'pdf':
            # 提取PDF文本（支持文字版和图片版）
            text = extract_text_from_pdf(filepath)
            if text:
                return jsonify({
                    'success': True,
                    'text': text[:5000],
                    'filename': filename
                })
            else:
                return jsonify({
                    'error': 'PDF文本提取失败，可能PDF是扫描件且OCR不可用，或PDF已加密',
                    'ocr_available': OCR_TESSERACT_AVAILABLE or OCR_EASYOCR_AVAILABLE,
                    'filename': filename
                }), 500
        elif ext in ['docx', 'doc']:
            # 提取Word文档文本
            text = extract_text_from_docx(filepath)
            if text:
                return jsonify({
                    'success': True,
                    'text': text[:5000],
                    'filename': filename
                })
            else:
                return jsonify({
                    'error': 'Word文档文本提取失败',
                    'filename': filename
                }), 500
        elif ext in ['txt', 'md', 'rtf']:
            # 提取文本文件内容
            text = extract_text_from_text_file(filepath)
            if text:
                return jsonify({
                    'success': True,
                    'text': text[:5000],
                    'filename': filename
                })
            else:
                return jsonify({
                    'error': '文本文件读取失败',
                    'filename': filename
                }), 500
        else:
            # 尝试从图片中提取文本
            text = extract_text_from_image(filepath)
            
            result = {
                'success': True,
                'filename': filename
            }
            
            if text:
                result['text'] = text[:5000]
            
            return jsonify(result)
            
    except Exception as e:
        logger.error(f"文件处理失败: {str(e)}")
        return jsonify({'error': f'文件处理失败: {str(e)}'}), 500


@app.route('/static/uploads/<filename>')
def serve_upload(filename):
    """提供上传的文件"""
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)


@app.route('/ip-images')
def list_ip_images():
    """获取所有可用的IP图片列表（内置+上传）"""
    logger.info(f"收到 /ip-images 请求")
    try:
        ip_images = []
        
        # 获取内置IP图片
        if os.path.exists(IP_DIR):
            for filename in os.listdir(IP_DIR):
                if filename.lower().endswith('.png'):
                    ip_images.append({
                        'filename': filename,
                        'type': 'built-in',
                        'url': f'/static/ip/{filename}'
                    })
        
        # 获取用户上传的IP图片
        if os.path.exists(app.config['UPLOAD_FOLDER']):
            for filename in os.listdir(app.config['UPLOAD_FOLDER']):
                if filename.lower().endswith('.png'):
                    ip_images.append({
                        'filename': filename,
                        'type': 'uploaded',
                        'url': f'/static/uploads/{filename}'
                    })
        
        return jsonify({'success': True, 'images': ip_images})
    
    except Exception as e:
        logger.error(f"获取IP图片列表失败: {str(e)}")
        return jsonify({'error': f'获取IP图片列表失败: {str(e)}'}), 500


@app.route('/upload-ip', methods=['POST'])
def upload_ip_image():
    """上传自定义IP图片"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': '没有选择文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '没有选择文件'}), 400
        
        # 仅允许PNG格式
        if not file.filename.lower().endswith('.png'):
            return jsonify({'error': '仅支持PNG格式的透明图片'}), 400
        
        # 生成唯一文件名
        ext = 'png'
        filename = f"{uuid.uuid4().hex}.{ext}"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # 保存文件
        file.save(filepath)
        logger.info(f"IP图片上传成功: {filepath}")
        
        return jsonify({
            'success': True,
            'filename': filename,
            'url': f'/static/uploads/{filename}'
        })
        
    except Exception as e:
        logger.error(f"IP图片上传失败: {str(e)}")
        return jsonify({'error': f'IP图片上传失败: {str(e)}'}), 500


@app.route('/api/upload_background', methods=['POST'])
def upload_background_image():
    """上传自定义背景图片（支持任意格式和清晰度）"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': '没有选择文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '没有选择文件'}), 400
        
        index = int(request.form.get('index', 0))
        
        ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else 'png'
        filename = f"bg_{index}_{uuid.uuid4().hex}.{ext}"
        filepath = os.path.join(OUTPUT_DIR, filename)
        
        file.save(filepath)
        logger.info(f"背景图片上传成功: {filepath}")
        
        return jsonify({
            'success': True,
            'image_url': f'/static/output/{filename}'
        })
        
    except Exception as e:
        logger.error(f"背景图片上传失败: {str(e)}")
        return jsonify({'error': f'背景图片上传失败: {str(e)}'}), 500


@app.route('/uploads')
def list_uploads():
    """列出已上传的文件（不执行耗时的OCR）"""
    logger.info(f"收到 /uploads 请求")
    try:
        files = []
        for filename in os.listdir(UPLOAD_FOLDER):
            filepath = os.path.join(UPLOAD_FOLDER, filename)
            if os.path.isfile(filepath):
                ext = ''
                if '.' in filename and not filename.startswith('.'):
                    parts = filename.rsplit('.', 1)
                    if len(parts) > 1:
                        ext = parts[1].lower()
                file_info = {
                    'filename': filename,
                    'type': 'pdf' if ext == 'pdf' else 'image',
                    'preview_url': f'/static/uploads/{filename}'
                }
                files.append(file_info)
        
        return jsonify({'success': True, 'files': files})
    except Exception as e:
        logger.error(f"获取文件列表失败: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/delete/<filename>', methods=['DELETE'])
def delete_file(filename):
    """删除上传的文件"""
    try:
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        if os.path.exists(filepath):
            os.remove(filepath)
            if filename in uploaded_files:
                del uploaded_files[filename]
            logger.info(f"文件删除成功: {filepath}")
            return jsonify({'success': True})
        else:
            return jsonify({'error': '文件不存在'}), 404
    except Exception as e:
        logger.error(f"文件删除失败: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/generate-script', methods=['POST'])
def generate_script_api():
    """使用豆包API生成漫画脚本（支持文件内容和IP图片）"""
    try:
        # 获取文件内容
        file_content = request.form.get('file_content', '').strip()
        filename = request.form.get('filename', '')
        
        if not file_content and not filename:
            return jsonify({'error': '请提供文件内容或文件名'}), 400
        
        # 如果没有直接提供文件内容，尝试从文件中提取
        if not file_content and filename:
            filepath = os.path.join(UPLOAD_FOLDER, filename)
            if os.path.exists(filepath):
                ext = ''
                if '.' in filename and not filename.startswith('.'):
                    parts = filename.rsplit('.', 1)
                    if len(parts) > 1:
                        ext = parts[1].lower()
                if ext == 'pdf':
                    file_content = extract_text_from_pdf(filepath)
                else:
                    file_content = extract_text_from_image(filepath)
            
            if not file_content:
                return jsonify({'error': '无法从文件中提取内容'}), 500
        
        # 获取漫画类型
        comic_type = request.form.get('comic_type', 'preview')
        
        # 获取用户选择的IP图片
        ip_paths_str = request.form.get('ip_paths', '')
        ip_paths = []
        
        if ip_paths_str:
            selected_ip_filenames = ip_paths_str.split(',')
            for ip_filename in selected_ip_filenames:
                ip_filename = ip_filename.strip()
                if ip_filename:
                    # 先检查IP目录
                    ip_path = os.path.join(IP_DIR, ip_filename)
                    if os.path.exists(ip_path):
                        ip_paths.append(ip_path)
                    else:
                        # 检查上传目录
                        uploaded_path = os.path.join(UPLOAD_FOLDER, ip_filename)
                        if os.path.exists(uploaded_path):
                            ip_paths.append(uploaded_path)
        
        if not ip_paths:
            ip_path = get_random_ip_image()
            if ip_path:
                ip_paths = [ip_path]
        
        if not ip_paths:
            return jsonify({'error': '未找到IP图片'}), 500
        
        logger.info(f"使用豆包API生成脚本，IP图片: {len(ip_paths)}个")
        
        # 调用豆包API生成脚本
        script = generate_comic_script_with_doubao(file_content, ip_paths, comic_type)
        
        if script and 'scenes' in script and 'dialogues' in script:
            speakers = script.get('speakers', [])
            unique_speakers = list(dict.fromkeys(speakers))
            characters = []
            for speaker in unique_speakers:
                if speaker == '知识小精灵':
                    characters.append({
                        'name': speaker,
                        'role': ROLE_GUIDE_CONFIG.get('role', '知识引导者'),
                        'appearance': ROLE_GUIDE_CONFIG.get('appearance', '')
                    })
                else:
                    characters.append({
                        'name': speaker,
                        'role': '',
                        'appearance': ''
                    })
            script['characters'] = characters
            return jsonify({
                'success': True,
                'script': script
            })
        else:
            # 如果豆包API失败，使用默认脚本
            logger.warning("豆包API生成脚本失败，使用默认脚本")
            default_script = get_default_script(file_content[:100], comic_type)
            default_script['characters'] = [{
                'name': '知识小精灵',
                'role': ROLE_GUIDE_CONFIG.get('role', '知识引导者'),
                'appearance': ROLE_GUIDE_CONFIG.get('appearance', '')
            }]
            return jsonify({
                'success': True,
                'script': default_script
            })
            
    except Exception as e:
        logger.error(f"生成脚本失败: {str(e)}")
        return jsonify({'error': f'生成脚本失败: {str(e)}'}), 500


@app.route('/api/create_task', methods=['POST'])
def api_create_task():
    """创建新任务"""
    try:
        task_id = create_task()
        logger.info(f"创建新任务: {task_id}")
        return jsonify({'task_id': task_id})
    except Exception as e:
        logger.error(f"创建任务失败: {str(e)}")
        return jsonify({'error': '创建任务失败'}), 500


@app.route('/api/task/<task_id>', methods=['GET'])
def api_get_task(task_id):
    """获取任务状态"""
    if task_id not in tasks:
        return jsonify({'error': '任务不存在'}), 404
    return jsonify(tasks[task_id].to_dict())


@app.route('/api/extract_pdf', methods=['POST'])
def api_extract_pdf():
    """上传并提取PDF内容，生成初始剧本"""
    try:
        task_id = request.form.get('task_id')
        if not task_id or task_id not in tasks:
            task_id = create_task()
        
        task = tasks[task_id]
        
        file_content = request.form.get('file_content', '').strip()
        
        panel_count = int(request.form.get('panel_count', 4))
        character_count = int(request.form.get('character_count', 1))
        character_system = request.form.get('character_system', 'student')
        scenario_style = request.form.get('scenario_style', 'humorous')
        
        if 'file' in request.files:
            file = request.files['file']
            if file.filename:
                from werkzeug.utils import secure_filename
                filename = secure_filename(file.filename)
                
                ext = ''
                if '.' in filename and not filename.startswith('.'):
                    parts = filename.rsplit('.', 1)
                    if len(parts) > 1:
                        ext = parts[1].lower()
                
                if not ext:
                    ext = detect_file_type(file.stream, filename)
                    if ext:
                        filename = f"{filename}.{ext}"
                    else:
                        return jsonify({'error': '无法识别文件类型，请确保文件有正确的扩展名（如.pdf、.docx等）'}), 400
                
                filepath = os.path.join(UPLOAD_FOLDER, filename)
                file.save(filepath)
                
                logger.info(f"上传文件: {filename}, 类型: {ext}")
                
                if ext == 'pdf':
                    file_content = extract_text_from_pdf(filepath)
                elif ext in ['docx', 'doc']:
                    file_content = extract_text_from_docx(filepath)
                elif ext in ['txt', 'md', 'rtf']:
                    file_content = extract_text_from_text_file(filepath)
                else:
                    file_content = extract_text_from_image(filepath)
        
        if not file_content:
            return jsonify({'error': '无法从文件中提取内容'}), 400
        
        task.pdf_text = file_content
        task.comic_type = request.form.get('comic_type', 'preview')
        
        ip_paths_str = request.form.get('ip_paths', '')
        ip_paths = []
        if ip_paths_str:
            selected_ip_filenames = ip_paths_str.split(',')
            for filename in selected_ip_filenames:
                filename = filename.strip()
                if filename:
                    ip_path = os.path.join(IP_DIR, filename)
                    if os.path.exists(ip_path):
                        ip_paths.append(ip_path)
                    else:
                        uploaded_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                        if os.path.exists(uploaded_path):
                            ip_paths.append(uploaded_path)
        task.ip_paths = ip_paths
        
        script = generate_comic_script(file_content, task.comic_type, panel_count, character_count, scenario_style, ip_paths, character_system)
        task.script = script
        
        return jsonify({
            'task_id': task_id,
            'pdf_text': file_content,
            'script': script
        })
    except Exception as e:
        logger.error(f"提取PDF失败: {str(e)}")
        error_msg = f'提取PDF失败: {str(e)}'
        if 'list index out of range' in str(e):
            error_msg = '文件名格式不正确，请确保文件有正确的扩展名（如.pdf、.docx等）'
        elif '识别失败' in str(e) or '未获取到文本' in str(e):
            error_msg = '无法提取PDF内容，可能原因：\n1. PDF是扫描版图片，需要OCR识别\n2. 当前环境缺少中文OCR语言包\n3. 建议使用可编辑的PDF或Word文档\n\n您可以尝试：\n- 手动输入课程重点文本\n- 使用已配置好OCR的环境\n- 提供文字版PDF而非扫描版'
        return jsonify({'error': error_msg}), 500


@app.route('/api/regenerate_script', methods=['POST'])
def api_regenerate_script():
    """重新生成剧本（支持多格数、多角色、多种情景风格）"""
    try:
        task_id = request.form.get('task_id')
        if not task_id or task_id not in tasks:
            return jsonify({'error': '任务不存在'}), 404
        
        task = tasks[task_id]
        course_topic = request.form.get('topic', task.pdf_text)
        comic_type = request.form.get('comic_type', task.comic_type)
        
        panel_count = int(request.form.get('panel_count', 4))
        character_count = int(request.form.get('character_count', 1))
        scenario_style = request.form.get('scenario_style', 'humorous')
        
        ip_paths_str = request.form.get('ip_paths', '')
        ip_paths = []
        if ip_paths_str:
            selected_ip_filenames = ip_paths_str.split(',')
            for filename in selected_ip_filenames:
                filename = filename.strip()
                if filename:
                    ip_path = os.path.join(IP_DIR, filename)
                    if os.path.exists(ip_path):
                        ip_paths.append(ip_path)
                    else:
                        uploaded_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                        if os.path.exists(uploaded_path):
                            ip_paths.append(uploaded_path)
        if ip_paths:
            task.ip_paths = ip_paths
        
        script = generate_comic_script(course_topic, comic_type, panel_count, character_count, scenario_style, task.ip_paths)
        task.script = script
        
        return jsonify({'script': script})
    except Exception as e:
        logger.error(f"重新生成剧本失败: {str(e)}")
        return jsonify({'error': f'重新生成剧本失败: {str(e)}'}), 500


@app.route('/api/update_script', methods=['POST'])
def api_update_script():
    """更新剧本（用户编辑）"""
    try:
        task_id = request.form.get('task_id') or request.args.get('task_id')
        if not task_id or task_id not in tasks:
            return jsonify({'error': '任务不存在'}), 404
        
        task = tasks[task_id]
        new_script = None
        
        content_type = request.content_type or ''
        
        if 'application/json' in content_type:
            try:
                json_data = request.get_json(silent=True)
                if json_data and isinstance(json_data, dict):
                    new_script = json_data.get('script')
            except:
                new_script = None
        else:
            script_str = request.form.get('script')
            if script_str:
                try:
                    new_script = json.loads(script_str)
                except:
                    new_script = None
        
        if new_script and isinstance(new_script, dict):
            task.script = new_script
        
        return jsonify({'script': task.script})
    except Exception as e:
        logger.error(f"更新剧本失败: {str(e)}\n{traceback.format_exc()}")
        return jsonify({'error': f'更新剧本失败: {str(e)}'}), 500


@app.route('/api/text_input', methods=['POST'])
def api_text_input():
    """接收用户输入的纯文字课程重点，直接生成剧本"""
    try:
        text = request.form.get('text', '').strip()
        
        if not text:
            return jsonify({'error': '请输入课程重点内容'}), 400
        
        logger.info(f"收到纯文字输入，长度: {len(text)}")
        
        comic_type = request.form.get('comic_type', 'preview')
        panel_count = int(request.form.get('panel_count', 4))
        character_count = int(request.form.get('character_count', 1))
        character_system = request.form.get('character_system', 'student')
        scenario_style = request.form.get('scenario_style', 'humorous')
        
        task_id = create_task()
        task = tasks[task_id]
        task.pdf_text = text
        task.comic_type = comic_type
        
        script = generate_comic_script(text, comic_type, panel_count, character_count, scenario_style, None, character_system)
        
        if script and 'scenes' in script and 'dialogues' in script:
            task.script = script
            logger.info(f"纯文字输入生成剧本成功，任务ID: {task_id}")
            
            return jsonify({
                'success': True,
                'task_id': task_id,
                'script': script
            })
        else:
            logger.warning("纯文字输入生成剧本失败")
            return jsonify({'error': '生成剧本失败'}), 500
            
    except Exception as e:
        logger.error(f"纯文字输入处理失败: {str(e)}")
        return jsonify({'error': f'处理失败: {str(e)}'}), 500


def calc_panel_size(aspect_ratio_str, base_size=768):
    """根据宽高比字符串计算面板尺寸，保持长边=base_size"""
    try:
        if aspect_ratio_str and ':' in aspect_ratio_str:
            w_ratio, h_ratio = aspect_ratio_str.split(':')
            w_ratio = float(w_ratio)
            h_ratio = float(h_ratio)
            if w_ratio >= h_ratio:
                width = base_size
                height = int(base_size * h_ratio / w_ratio)
            else:
                height = base_size
                width = int(base_size * w_ratio / h_ratio)
            return width, height
    except:
        pass
    return base_size, base_size


@app.route('/api/generate_backgrounds', methods=['POST'])
def api_generate_backgrounds():
    """生成背景图片"""
    try:
        task_id = request.form.get('task_id')
        layout = request.form.get('layout', 'four-grid')
        if not task_id or task_id not in tasks:
            return jsonify({'error': '任务不存在'}), 404
        
        task = tasks[task_id]
        
        if not task.script:
            return jsonify({'error': '请先生成剧本'}), 400
        
        scenes = task.script.get('scenes', [])
        backgrounds = []
        
        # 汇总剧本里的角色名，作为纯背景屏蔽名单（即使场景描述里写了角色也不出现在背景）
        _bg_block_names = []
        for _c in (task.script.get('characters', []) or []):
            if isinstance(_c, dict):
                _nm = _c.get('name')
                if _nm:
                    _bg_block_names.append(str(_nm))
            elif isinstance(_c, str) and _c.strip():
                _bg_block_names.append(_c.strip())
        for _p in (task.script.get('panels', []) or []):
            for _sp in (_p.get('speakers', []) or []):
                if _sp and _sp not in _bg_block_names:
                    _bg_block_names.append(str(_sp))
            _sp = _p.get('speaker')
            if _sp and _sp not in _bg_block_names:
                _bg_block_names.append(str(_sp))
        
        world_setting = task.script.get('story_plan', {}).get('world_setting', None)
        style_seed = f"{task_id}_{random.randint(1000, 9999)}"
        
        layout_config = get_layout_config(layout)
        grid_specs = layout_config.get('grid_specs', [])
        is_freeform = layout_config.get('freeform', False)
        
        panels_info = {}
        if is_freeform:
            for i, spec in enumerate(grid_specs):
                panels_info[i] = {
                    'width': spec.get('width', 768),
                    'height': spec.get('height', 768)
                }
        else:
            panels_info, _ = calc_layout_positions(grid_specs, 720, 10)
        
        for i, scene in enumerate(scenes):
            try:
                panel_spec = grid_specs[i] if i < len(grid_specs) else {}
                info = panels_info.get(i, {'width': 768, 'height': 768})
                width = info['width']
                height = info['height']
                
                background_data = generate_image(
                    scene, i, world_setting, style_seed,
                    width=width, height=height, panel_spec=panel_spec,
                    is_freeform=is_freeform, force_pure_background=True,
                    block_names=_bg_block_names,
                    allow_default_fallback=False
                )
                if background_data is not None and not background_visual_style_matches_prompt(background_data, scene):
                    raise RuntimeError('生图结果与场景要求的暖色卡通风格明显不符，请重新生成')
                if background_data is None:
                    raise RuntimeError('生图服务当前限流或暂不可用，未生成与场景描述匹配的背景')
                output_filename = f"{task_id}_bg_{i}.png"
                output_path = os.path.join(OUTPUT_DIR, output_filename)
                
                from io import BytesIO
                img = Image.open(BytesIO(background_data))
                img = normalize_background_for_panel(img, width, height)
                img.save(output_path)
                
                backgrounds.append(f'/static/output/{output_filename}')
            except Exception as e:
                logger.error(f"生成背景图{i}失败: {str(e)}")
                backgrounds.append(None)
        
        task.backgrounds = backgrounds
        
        return jsonify({'backgrounds': backgrounds})
    except Exception as e:
        logger.error(f"生成背景图失败: {str(e)}")
        return jsonify({'error': f'生成背景图失败: {str(e)}'}), 500


@app.route('/api/regenerate_background', methods=['POST'])
def api_regenerate_background():
    """重新生成单个背景图片"""
    try:
        task_id = request.form.get('task_id')
        index = int(request.form.get('index', 0))
        layout = request.form.get('layout', 'four-grid')
        
        if not task_id or task_id not in tasks:
            return jsonify({'error': '任务不存在'}), 404
        
        task = tasks[task_id]
        
        if not task.script:
            return jsonify({'error': '请先生成剧本'}), 400
        
        scenes = task.script.get('scenes', [])
        if index < 0 or index >= len(scenes):
            return jsonify({'error': '无效的索引'}), 400
        
        scene = scenes[index]
        world_setting = task.script.get('story_plan', {}).get('world_setting', None)
        # 从剧本提取角色名作为屏蔽名单
        _rb_names = []
        for _c in (task.script.get('characters', []) or []):
            if isinstance(_c, dict):
                _nm = _c.get('name')
                if _nm:
                    _rb_names.append(str(_nm))
            elif isinstance(_c, str) and _c.strip():
                _rb_names.append(_c.strip())
        for _p in (task.script.get('panels', []) or []):
            for _sp in (_p.get('speakers', []) or []):
                if _sp and _sp not in _rb_names:
                    _rb_names.append(str(_sp))
            _sp = _p.get('speaker')
            if _sp and _sp not in _rb_names:
                _rb_names.append(str(_sp))

        layout_config = get_layout_config(layout)
        grid_specs = layout_config.get('grid_specs', [])
        panel_spec = grid_specs[index] if index < len(grid_specs) else {}
        if layout_config.get('freeform', False):
            target_width = int(panel_spec.get('width', 768))
            target_height = int(panel_spec.get('height', 768))
        else:
            panels_info, _ = calc_layout_positions(grid_specs, 720, 10)
            info = panels_info.get(index, {'width': 768, 'height': 768})
            target_width = int(info['width'])
            target_height = int(info['height'])
            cell_ar = calculate_cell_aspect_ratio(layout_config, panel_spec)
            panel_spec = {'aspect_ratio': cell_ar}

        background_data = generate_image(
            scene, index, world_setting,
            width=target_width, height=target_height,
            panel_spec=panel_spec,
            force_pure_background=True, block_names=_rb_names,
            allow_default_fallback=False
        )
        if background_data is None:
            return jsonify({
                'error_code': 'IMAGE_PROVIDER_UNAVAILABLE',
                'error': '生图服务当前限流或暂不可用，未生成与场景描述匹配的背景，请稍后重试。'
            }), 503
        output_filename = f"{task_id}_bg_{index}.png"
        output_path = os.path.join(OUTPUT_DIR, output_filename)
        
        from io import BytesIO
        img = Image.open(BytesIO(background_data))
        img = normalize_background_for_panel(
            img, target_width, target_height
        )
        img.save(output_path)
        
        task.backgrounds[index] = f'/static/output/{output_filename}'
        
        return jsonify({
            'index': index,
            'background': task.backgrounds[index]
        })
    except Exception as e:
        logger.error(f"重新生成背景图失败: {str(e)}")
        return jsonify({'error': f'重新生成背景图失败: {str(e)}'}), 500


@app.route('/api/compose_comic', methods=['POST'])
def api_compose_comic():
    """合成漫画"""
    try:
        task_id = request.form.get('task_id')
        if not task_id or task_id not in tasks:
            return jsonify({'error': '任务不存在'}), 404
        
        panel_count = int(request.form.get('panel_count', 4))
        font_family = request.form.get('font_family', 'msyh')
        font_size = int(request.form.get('font_size', 0))
        font_size = font_size if font_size > 0 else None
        
        layout = request.form.get('layout', 'four-grid')
        layout_config = get_layout_config(layout)
        
        task = tasks[task_id]
        
        if not task.script:
            return jsonify({'error': '请先生成剧本'}), 400
        
        if not task.backgrounds:
            return jsonify({'error': '请先生成背景图'}), 400
        
        ip_paths_str = request.form.get('ip_paths', '')
        ip_paths = []
        
        if ip_paths_str:
            selected_ip_filenames = ip_paths_str.split(',')
            for ip_filename in selected_ip_filenames:
                ip_filename = ip_filename.strip()
                if ip_filename:
                    ip_path = os.path.join(IP_DIR, ip_filename)
                    if os.path.exists(ip_path):
                        ip_paths.append(ip_path)
                    else:
                        uploaded_path = os.path.join(UPLOAD_FOLDER, ip_filename)
                        if os.path.exists(uploaded_path):
                            ip_paths.append(uploaded_path)
        
        if not ip_paths:
            ip_path = get_random_ip_image()
            if ip_path:
                ip_paths = [ip_path]
        
        task.ip_paths = ip_paths
        
        dialogues = task.script.get('dialogues', [])
        comics = []
        
        custom_backgrounds = getattr(task, 'custom_backgrounds', {})
        
        speakers = task.script.get('speakers', [])
        
        poses_str = request.form.get('poses', '')
        poses = []
        if poses_str:
            try:
                import json
                poses = json.loads(poses_str)
            except:
                poses = poses_str.split(',')
        
        panel_dimensions = calculate_panel_dimensions(layout_config, 720, 900)
        
        for i, (bg_path, dialogue) in enumerate(zip(task.backgrounds, dialogues)):
            try:
                if i in custom_backgrounds and os.path.exists(custom_backgrounds[i]):
                    bg_full_path = custom_backgrounds[i]
                    logger.info(f"使用自定义背景: {bg_full_path}")
                elif not bg_path:
                    comics.append(None)
                    continue
                else:
                    bg_full_path = os.path.join(app.root_path, bg_path.lstrip('/'))
                
                with open(bg_full_path, 'rb') as f:
                    bg_data = f.read()
                
                output_filename = f"{task_id}_comic_{i}.png"
                output_path = os.path.join(OUTPUT_DIR, output_filename)
                
                speaker = speakers[i] if i < len(speakers) else ''
                
                panel_poses = poses[i] if isinstance(poses, list) and i < len(poses) else None
                
                target_width = panel_dimensions[i]['width'] if i < len(panel_dimensions) else None
                target_height = panel_dimensions[i]['height'] if i < len(panel_dimensions) else None
                
                composite_image(bg_data, ip_paths, dialogue, output_path, speaker=speaker, font_family=font_family, font_size=font_size, poses=panel_poses, target_width=target_width, target_height=target_height)
                comics.append(f'/static/output/{output_filename}')
            except Exception as e:
                logger.error(f"合成漫画{i}失败: {str(e)}")
                comics.append(None)
        
        task.comics = comics
        
        layout = request.form.get('layout', 'four-grid')
        comic_page_path = create_comic_page(task_id, comics, panel_count, layout)
        
        return jsonify({'comics': comics, 'comic_page': comic_page_path})
    except Exception as e:
        logger.error(f"合成漫画失败: {str(e)}")
        return jsonify({'error': f'合成漫画失败: {str(e)}'}), 500


def parse_css_grid(css_grid):
    """解析 CSS grid 字符串，提取行和列的比例"""
    col_ratios = []
    row_ratios = []
    
    if not css_grid:
        return col_ratios, row_ratios
    
    parts = css_grid.split(';')
    for part in parts:
        part = part.strip()
        if part.startswith('grid-template-columns'):
            col_str = part.replace('grid-template-columns:', '').strip()
            for col in col_str.split():
                if col.endswith('fr'):
                    try:
                        fr_match = re.search(r'(\d*\.?\d+)fr', col)
                        col_ratios.append(float(fr_match.group(1)) if fr_match else 1.0)
                    except:
                        col_ratios.append(1.0)
        elif part.startswith('grid-template-rows'):
            row_str = part.replace('grid-template-rows:', '').strip()
            for row in row_str.split():
                if row.endswith('fr'):
                    try:
                        fr_match = re.search(r'(\d*\.?\d+)fr', row)
                        row_ratios.append(float(fr_match.group(1)) if fr_match else 1.0)
                    except:
                        row_ratios.append(1.0)
    
    return col_ratios, row_ratios


def parse_aspect_ratio(ar_str):
    """解析宽高比字符串，返回 width/height 比值"""
    try:
        if ar_str and ':' in ar_str:
            w, h = ar_str.split(':')
            return float(w) / float(h)
    except:
        pass
    return 1.0


def calc_layout_positions(grid_specs, total_content_width, gap):
    """根据每个格子的 aspect_ratio 计算布局位置和尺寸
    按行排列，每行高度由该行所有格子的宽高比和总宽度共同决定
    """
    rows = {}
    for i, spec in enumerate(grid_specs):
        row = spec.get('row', 0)
        if row not in rows:
            rows[row] = []
        rows[row].append((i, spec))
    
    row_list = sorted(rows.keys())
    panels_info = {}
    current_y = 0
    
    for row_idx in row_list:
        row_panels = rows[row_idx]
        col_count = len(row_panels)
        total_gap = (col_count - 1) * gap
        
        ratios = []
        for i, spec in row_panels:
            ar = parse_aspect_ratio(spec.get('aspect_ratio', '1:1'))
            ratios.append(ar)
        
        row_height = (total_content_width - total_gap) / sum(ratios)
        
        current_x = 0
        for idx, (i, spec) in enumerate(row_panels):
            ar = ratios[idx]
            w = int(ar * row_height)
            h = int(row_height)
            panels_info[i] = {
                'x': current_x,
                'y': current_y,
                'width': w,
                'height': h
            }
            current_x += w + gap
        
        current_y += int(row_height) + gap
    
    total_height = current_y - gap if current_y > 0 else 0
    return panels_info, total_height


def create_comic_page(task_id, comic_paths, panel_count, layout='four-grid'):
    """创建漫画分镜框布局页面，支持网格布局和自由布局（倾斜旋转、重叠）"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        border_width = 6
        gap = 10
        padding = 16
        panel_border = 3
        total_content_width = 720
        
        layout_config = get_layout_config(layout)
        grid_specs = layout_config.get('grid_specs', [])
        is_freeform = layout_config.get('freeform', False)
        
        if is_freeform:
            max_x = 0
            max_y = 0
            for spec in grid_specs:
                x = spec.get('x', 0)
                y = spec.get('y', 0)
                w = spec.get('width', 0)
                h = spec.get('height', 0)
                angle = abs(spec.get('tilt_angle', 0))
                radians = angle * 3.14159 / 180
                diag = (w*w + h*h)**0.5
                expand = diag * (1 - math.cos(radians)) / 2
                max_x = max(max_x, x + w + expand)
                max_y = max(max_y, y + h + expand)
            
            total_width = int(max_x) + 2 * padding + 2 * border_width + 50
            total_height = int(max_y) + 2 * padding + 2 * border_width + 50
            
            comic_page = Image.new('RGBA', (total_width, total_height), color=(255, 255, 255, 255))
            
            panels_with_z = []
            for i, spec in enumerate(grid_specs):
                panels_with_z.append((spec.get('z_order', i), i, spec))
            
            panels_with_z.sort(key=lambda x: x[0])
            
            for z_order, i, spec in panels_with_z:
                if i >= len(comic_paths) or not comic_paths[i]:
                    continue
                
                panel_width = spec.get('width', 768)
                panel_height = spec.get('height', 768)
                x = padding + border_width + spec.get('x', 0)
                y = padding + border_width + spec.get('y', 0)
                tilt_angle = spec.get('tilt_angle', 0)
                
                comic_full_path = os.path.join(app.root_path, comic_paths[i].lstrip('/'))
                if os.path.exists(comic_full_path):
                    panel_img = Image.open(comic_full_path).convert('RGB')
                    if panel_img.size != (panel_width, panel_height):
                        panel_img = panel_img.resize((panel_width, panel_height), Image.Resampling.BICUBIC)
                    
                    panel_img_with_alpha = panel_img.convert('RGBA')
                    
                    if tilt_angle != 0:
                        rotated = panel_img_with_alpha.rotate(tilt_angle, expand=True, resample=Image.Resampling.BICUBIC)
                        rotated_w, rotated_h = rotated.size
                        offset_x = (rotated_w - panel_width) // 2
                        offset_y = (rotated_h - panel_height) // 2
                        paste_x = x - offset_x
                        paste_y = y - offset_y
                        comic_page.paste(rotated, (paste_x, paste_y), rotated)
                    else:
                        comic_page.paste(panel_img_with_alpha, (x, y), panel_img_with_alpha)
            
            comic_page_rgb = comic_page.convert('RGB')
            
        else:
            panels_info, content_height = calc_layout_positions(grid_specs, total_content_width, gap)
            
            total_width = total_content_width + 2 * padding + 2 * border_width
            total_height = content_height + 2 * padding + 2 * border_width
            
            comic_page_rgb = Image.new('RGB', (total_width, total_height), color=(255, 255, 255))
            draw = ImageDraw.Draw(comic_page_rgb)
            
            draw.rounded_rectangle([0, 0, total_width - 1, total_height - 1], radius=20, outline=(0, 0, 0), width=border_width)
            
            for i, comic_path in enumerate(comic_paths):
                if not comic_path or i not in panels_info:
                    continue
                
                info = panels_info[i]
                panel_width = info['width']
                panel_height = info['height']
                x = padding + border_width + info['x']
                y = padding + border_width + info['y']
                
                comic_full_path = os.path.join(app.root_path, comic_path.lstrip('/'))
                if os.path.exists(comic_full_path):
                    panel_bg = Image.new('RGB', (panel_width, panel_height), color=(255, 255, 255))
                    panel_draw = ImageDraw.Draw(panel_bg)
                    
                    panel_img = Image.open(comic_full_path).convert('RGB')
                    if panel_img.size != (panel_width, panel_height):
                        panel_img = panel_img.resize((panel_width, panel_height), Image.Resampling.BICUBIC)
                    panel_bg.paste(panel_img, (0, 0))
                    
                    panel_draw.rounded_rectangle([0, 0, panel_width - 1, panel_height - 1], radius=12, outline=(0, 0, 0), width=panel_border)
                    
                    comic_page_rgb.paste(panel_bg, (x, y))
        
        output_filename = f"{task_id}_comic_page.png"
        output_path = os.path.join(OUTPUT_DIR, output_filename)
        comic_page_rgb.save(output_path)
        
        return f'/static/output/{output_filename}'
    
    except Exception as e:
        logger.error(f"创建漫画页面失败: {str(e)}")
        return None


@app.route('/api/regenerate_panel', methods=['POST'])
def api_regenerate_panel():
    """重新生成单个漫画格子"""
    try:
        task_id = request.form.get('task_id')
        index = int(request.form.get('index', 0))
        
        if not task_id or task_id not in tasks:
            return jsonify({'error': '任务不存在'}), 404
        
        task = tasks[task_id]
        
        if not task.script or not task.backgrounds:
            return jsonify({'error': '请先生成剧本和背景图'}), 400
        
        bg_path = task.backgrounds[index]
        dialogue = task.script.get('dialogues', [])[index]
        
        if not bg_path:
            return jsonify({'error': '背景图不存在'}), 400
        
        bg_full_path = os.path.join(app.root_path, bg_path.lstrip('/'))
        with open(bg_full_path, 'rb') as f:
            bg_data = f.read()
        
        output_filename = f"{task_id}_comic_{index}.png"
        output_path = os.path.join(OUTPUT_DIR, output_filename)
        
        composite_image(bg_data, task.ip_paths, dialogue, output_path)
        task.comics[index] = f'/static/output/{output_filename}'
        
        return jsonify({
            'index': index,
            'comic': task.comics[index]
        })
    except Exception as e:
        logger.error(f"重新生成漫画格子失败: {str(e)}")
        return jsonify({'error': f'重新生成漫画格子失败: {str(e)}'}), 500


def compose_aggregate_image(comic_images, layout='four-grid'):
    """将所有漫画格子按布局合并为一张汇总大图，严格按 CSS Grid 比例填满，无白边/黑边。"""
    try:
        from PIL import Image
        from io import BytesIO

        if not comic_images:
            logger.error("没有可合成的漫画图片")
            return None

        def _trim_black_borders(img, threshold=15):
            """裁掉图片四周的纯黑/近黑边缘，避免汇总图出现黑边。"""
            try:
                import numpy as np
                arr = np.array(img)
                if arr.ndim == 3:
                    mask = arr.sum(axis=2) > threshold * 3
                else:
                    mask = arr > threshold
                if not mask.any():
                    return img
                rows = mask.any(axis=1)
                cols = mask.any(axis=0)
                top = rows.argmax()
                bottom = len(rows) - rows[::-1].argmax()
                left = cols.argmax()
                right = len(cols) - cols[::-1].argmax()
                # 至少保留 90% 的画面，避免过度裁剪
                min_w = int(img.width * 0.9)
                min_h = int(img.height * 0.9)
                if right - left < min_w:
                    left = max(0, (img.width - min_w) // 2)
                    right = left + min_w
                if bottom - top < min_h:
                    top = max(0, (img.height - min_h) // 2)
                    bottom = top + min_h
                if left > 0 or top > 0 or right < img.width or bottom < img.height:
                    return img.crop((left, top, right, bottom))
                return img
            except Exception:
                return img

        images = []
        for i, img_url in enumerate(comic_images):
            if not img_url:
                continue
            # 修复：以前端传入的 URL 为准（它已是当前最新——编辑过的是 panel_i_edited.png?t=时间戳，
            # 未编辑的是本批原始成品图）。不再无脑优先读磁盘上固定名的 panel_i_edited.png，
            # 否则上一批残留的旧编辑图会污染合集图，导致"合集没用每格最新图片"。
            # 必须 strip 掉 URL 上的 ?t= 查询串，否则拼出的路径带 ? 会找不到文件。
            clean_url = img_url.split('?')[0]
            img_path = os.path.join(app.root_path, clean_url.lstrip('/'))
            if os.path.exists(img_path):
                raw_img = Image.open(img_path).convert('RGB')
                images.append(_trim_black_borders(raw_img))
                logger.info(f"合集图第{i+1}格使用: {clean_url}")
            else:
                # 兜底：URL 指向的文件不存在时，才回退到固定名的编辑图
                edited_path = os.path.join(app.root_path, 'static', 'output', f'panel_{i}_edited.png')
                if os.path.exists(edited_path):
                    raw_img = Image.open(edited_path).convert('RGB')
                    images.append(_trim_black_borders(raw_img))
                    logger.warning(f"合集图第{i+1}格 URL 缺失，回退编辑图: panel_{i}_edited.png")

        if not images:
            logger.error("无法找到任何漫画图片文件")
            return None

        panel_count = len(images)
        # 修复：直接吃完整 layout_config 字典（与链接版同一数据源），不再只靠字符串 id 查表 fallback。
        # 链接版布局正确，正是因为它直接用了这个含 css_grid/grid_specs 的完整对象。
        if isinstance(layout, dict) and (layout.get('css_grid') or layout.get('grid_specs')):
            layout_config = layout
            logger.info(f"合集图直接使用前端透传的完整布局配置: {layout.get('detected_layout') or layout.get('name') or '(无名)'}")
        else:
            layout_config = get_layout_config(layout) or LAYOUT_CONFIG.get('four-grid', LAYOUT_CONFIG['four-grid'])
        is_freeform = layout_config.get('freeform', False)

        target_total_width = 2400  # 汇总图目标宽度，足够清晰
        gap = 20

        if is_freeform:
            # 自由布局：直接按 grid_specs 里的 x/y/width/height 拼贴，并缩放到目标宽度
            raw_width = max([s.get('x', 0) + s.get('width', 100) for s in layout_config.get('grid_specs', [])] + [100])
            scale = target_total_width / max(raw_width, 1)
            total_width = int(raw_width * scale)
            total_height = int(max([s.get('y', 0) + s.get('height', 100) for s in layout_config.get('grid_specs', [])] + [100]) * scale)
        else:
            # 网格布局：按 CSS Grid fr 值计算总高度
            rows, cols = parse_fr_values(layout_config.get('css_grid', ''))
            if not rows:
                rows = [1.0] * layout_config.get('rows', 2)
            if not cols:
                cols = [1.0] * layout_config.get('cols', 2)
            total_col_fr = sum(cols)
            total_row_fr = sum(rows)
            # 总高度 = 总宽度 * (row_fr / col_fr) + 间隙
            total_height = int((target_total_width - (len(cols) - 1) * gap) * total_row_fr / total_col_fr + (len(rows) - 1) * gap)
            total_width = target_total_width

        grid_specs = layout_config.get('grid_specs', [])
        # 兜底：若图片数量多于布局格数，为多余图片追加 1x1 占位格顺序排在现有格之后，
        # 避免汇总图只拼前几格（用户反馈"汇总并没有汇总所有的格子"）。
        if len(grid_specs) < len(images) and not is_freeform:
            ncols = max(len(cols), 1)
            start = len(grid_specs)
            base_row = (max([s.get('row', 0) for s in grid_specs], default=-1) + 1) if grid_specs else 0
            for k in range(start, len(images)):
                m = k - start
                grid_specs.append({
                    'row': base_row + m // ncols,
                    'col': m % ncols,
                    'row_span': 1, 'col_span': 1
                })
            logger.info(f"汇总图：图片数({len(images)})>格数，已扩展布局到 {len(grid_specs)} 格")

        # 确保行列 fr 数组能覆盖所有 grid_specs（否则扩展出来的格会被算成 0 尺寸而崩溃）
        max_row = max([s.get('row', 0) + s.get('row_span', 1) - 1 for s in grid_specs], default=0)
        max_col = max([s.get('col', 0) + s.get('col_span', 1) - 1 for s in grid_specs], default=0)
        rows_fr = list(rows)
        cols_fr = list(cols)
        while len(rows_fr) <= max_row:
            rows_fr.append(1.0)
        while len(cols_fr) <= max_col:
            cols_fr.append(1.0)
        total_row_fr = sum(rows_fr)
        total_col_fr = sum(cols_fr)

        if is_freeform:
            # 自由布局：直接按 grid_specs 里的 x/y/width/height 拼贴，并缩放到目标宽度
            raw_width = max([s.get('x', 0) + s.get('width', 100) for s in grid_specs] + [100])
            scale = target_total_width / max(raw_width, 1)
            total_width = int(raw_width * scale)
            total_height = int(max([s.get('y', 0) + s.get('height', 100) for s in grid_specs] + [100]) * scale)
        else:
            total_height = int((target_total_width - (len(cols_fr) - 1) * gap) * total_row_fr / total_col_fr + (len(rows_fr) - 1) * gap)
            total_width = target_total_width

        aggregate_img = Image.new('RGB', (total_width, total_height), (255, 255, 255))
        draw = ImageDraw.Draw(aggregate_img)

        # 预计算每行的高度 / 每列的宽度（像素），并对极端情况做 >=1 守卫
        row_heights = []
        rem_h = total_height - (len(rows_fr) - 1) * gap
        for r in rows_fr:
            row_heights.append(max(1, int(rem_h * r / total_row_fr)))
        col_widths = []
        rem_w = total_width - (len(cols_fr) - 1) * gap
        for c in cols_fr:
            col_widths.append(max(1, int(rem_w * c / total_col_fr)))

        for i, img in enumerate(images):
            if i >= len(grid_specs):
                break

            spec = grid_specs[i]

            if is_freeform:
                x = int(spec.get('x', 0) * scale)
                y = int(spec.get('y', 0) * scale)
                cw = max(1, int(spec.get('width', 100) * scale))
                ch = max(1, int(spec.get('height', 100) * scale))
            else:
                row = spec.get('row', 0)
                col = spec.get('col', 0)
                rs = spec.get('row_span', 1)
                cs = spec.get('col_span', 1)
                x = sum(col_widths[:col]) + col * gap
                y = sum(row_heights[:row]) + row * gap
                cw = sum(col_widths[col:col + cs]) + (cs - 1) * gap
                ch = sum(row_heights[row:row + rs]) + (rs - 1) * gap

            # 守卫：尺寸必须 > 0，否则跳过该格，避免 resize 抛 "height and width must be > 0"
            if cw <= 0 or ch <= 0:
                logger.warning(f"汇总图第{i+1}格尺寸非法 (w={cw}, h={ch})，跳过")
                continue

            # contain 填充：等比缩放并居中，白边填充，绝不裁剪、绝不出 0 尺寸
            img_w, img_h = img.size
            if img_w <= 0 or img_h <= 0:
                continue
            img_ratio = img_w / img_h
            target_ratio = cw / ch
            if img_ratio > target_ratio:
                nw = cw
                nh = max(1, int(cw / img_ratio))
            else:
                nh = ch
                nw = max(1, int(ch * img_ratio))
            try:
                resized = img.resize((nw, nh), Image.BICUBIC)
            except Exception as re:
                logger.warning(f"汇总图第{i+1}格缩放失败: {re}")
                continue
            px = x + (cw - nw) // 2
            py = y + (ch - nh) // 2
            aggregate_img.paste(resized, (px, py))
            draw.rectangle([x, y, x + cw, y + ch], outline=(255, 179, 198), width=3)

        buffer = BytesIO()
        aggregate_img.save(buffer, format='PNG')
        return buffer.getvalue()

    except Exception as e:
        logger.error(f"合成汇总图片失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return None


def create_comic_html_page(comic_images, script=None, base_url='', animation='fade', layout_config=None, inline=False, share_view='mobile'):
    """创建包含所有漫画图片的HTML页面，支持动画效果和布局配置"""
    speakers = []
    dialogues = []
    scenes_zh = []
    
    if script:
        if isinstance(script, dict):
            speakers = script.get('speakers', [])
            dialogues = script.get('dialogues', [])
            scenes_zh = script.get('scenes_zh', [])
        elif isinstance(script, list) and len(script) > 0 and isinstance(script[0], dict):
            for panel in script:
                speakers.append(panel.get('speakers', [''])[0] if isinstance(panel.get('speakers'), list) else '')
                dialogues.append(panel.get('dialogues', ''))
                scenes_zh.append(panel.get('scene', ''))
    
    topic = ''
    if isinstance(script, dict):
        topic = script.get('topic', '')[:50]

    import base64
    def _resolve_src(comic_url):
        """inline 模式：把图片转 base64 data URI 内联，生成自包含 HTML（无需服务器即可打开/分享）"""
        if inline and comic_url:
            p = os.path.join(app.root_path, comic_url.lstrip('/'))
            if os.path.exists(p):
                with open(p, 'rb') as f:
                    b64 = base64.b64encode(f.read()).decode('ascii')
                return 'data:image/png;base64,' + b64
        return base_url + comic_url.lstrip('/') if base_url else comic_url
    
    animation_css = '''
        @keyframes fadeIn {
            from { opacity: 0; }
            to { opacity: 1; }
        }
        @keyframes zoomIn {
            from { opacity: 0; transform: scale(0.3); }
            to { opacity: 1; transform: scale(1); }
        }
        @keyframes slideLeft {
            from { opacity: 0; transform: translateX(-100px); }
            to { opacity: 1; transform: translateX(0); }
        }
        @keyframes slideRight {
            from { opacity: 0; transform: translateX(100px); }
            to { opacity: 1; transform: translateX(0); }
        }
        @keyframes slideTop {
            from { opacity: 0; transform: translateY(-100px); }
            to { opacity: 1; transform: translateY(0); }
        }
        @keyframes slideBottom {
            from { opacity: 0; transform: translateY(100px); }
            to { opacity: 1; transform: translateY(0); }
        }
        @keyframes bounceIn {
            0% { opacity: 0; transform: scale(0.3); }
            50% { opacity: 1; transform: scale(1.05); }
            70% { transform: scale(0.9); }
            100% { opacity: 1; transform: scale(1); }
        }
        @keyframes flipIn {
            from { opacity: 0; transform: perspective(400px) rotateY(90deg); }
            to { opacity: 1; transform: perspective(400px) rotateY(0); }
        }
        @keyframes rotateIn {
            from { opacity: 0; transform: rotate(-20deg); }
            to { opacity: 1; transform: rotate(0); }
        }
        .comic-panel {
            opacity: 0;
            transform-origin: center center;
        }
        .animate-fade { animation: fadeIn 0.6s ease-out forwards; }
        .animate-zoom { animation: zoomIn 0.6s ease-out forwards; }
        .animate-slide-left { animation: slideLeft 0.6s ease-out forwards; }
        .animate-slide-right { animation: slideRight 0.6s ease-out forwards; }
        .animate-slide-top { animation: slideTop 0.6s ease-out forwards; }
        .animate-slide-bottom { animation: slideBottom 0.6s ease-out forwards; }
        .animate-bounce { animation: bounceIn 0.8s ease-out forwards; }
        .animate-flip { animation: flipIn 0.6s ease-out forwards; }
        .animate-rotate { animation: rotateIn 0.6s ease-out forwards; }
    ''' if animation != 'none' else ''
    
    html_content = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{topic or '漫画标题'}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            background: linear-gradient(135deg, #FFF5E6, #FFE4E1);
            min-height: 100vh;
            padding: 20px;
        }}
        .container {{
            max-width: 1400px;
            margin: 0 auto;
        }}
        .header {{
            text-align: center;
            margin-bottom: 30px;
        }}
        .header h1 {{
            font-size: 2.5rem;
            color: #FFB3C6;
            text-shadow: 2px 2px 0px #FF8FA3;
            margin-bottom: 10px;
        }}
        .header p {{
            color: #888;
            font-size: 1.1rem;
        }}
        .comic-grid {{
            display: grid;
            gap: 16px;
            margin-bottom: 30px;
            width: 100%;
            max-width: 900px;
            margin-left: auto;
            margin-right: auto;
            align-items: stretch;
            min-height: 0;
        }}
        /* “手机阅读版”不依赖设备宽度判断：微信/内嵌浏览器也始终单列大图。 */
        body.mobile-reader .comic-grid {{
            grid-template-columns: 1fr !important;
            grid-template-rows: none !important;
            gap: 16px;
        }}
        body.mobile-reader .container {{ max-width: 680px; padding: 18px 12px; }}
        body.mobile-reader .comic-panel {{ min-height: 0 !important; }}
        .comic-panel {{
            background: white;
            border-radius: 12px;
            box-shadow: 0 4px 12px rgba(0,0,0,0.15);
            overflow: hidden;
            position: relative;
            border: 3px solid #222222;
            height: auto;
            min-height: 250px;
        }}
        .comic-panel:hover {{
            box-shadow: 0 8px 24px rgba(255, 179, 198, 0.4);
        }}
        .comic-image {{
            width: 100%;
            height: 100%;
            object-fit: cover;
            display: block;
        }}
        .comic-image-scaled {{
            width: 100%;
            height: auto;
            max-height: 500px;
            object-fit: contain;
            display: block;
        }}
        .panel-number {{
            display: inline-block;
            background: linear-gradient(135deg, #FFB3C6, #FF8FA3);
            color: white;
            font-weight: bold;
            font-size: 0.85rem;
            padding: 4px 10px;
            border-radius: 15px;
            margin-bottom: 8px;
        }}
        .speaker {{
            color: #FF6B9D;
            font-weight: bold;
            font-size: 1rem;
            margin-bottom: 5px;
        }}
        .dialogue {{
            color: #333;
            font-size: 1rem;
            margin-bottom: 5px;
            line-height: 1.6;
        }}
        .scene {{
            color: #888;
            font-size: 0.85rem;
            margin-top: 10px;
            padding-top: 10px;
            border-top: 1px dashed #FFB3C6;
        }}
        .footer {{
            text-align: center;
            color: #888;
            font-size: 0.85rem;
            padding: 25px;
            background: rgba(255, 255, 255, 0.5);
            border-radius: 10px;
        }}
        @media (max-width: 768px) {{
            .comic-grid {{
                /* 动态布局写在 inline style 中，手机端必须覆盖为单列。 */
                grid-template-columns: 1fr !important;
                grid-template-rows: none !important;
                gap: 14px;
            }}
            .container {{ padding: 16px 12px; }}
            .header {{ margin-bottom: 16px; }}
            .comic-panel {{ min-height: 0 !important; }}
            .header h1 {{
                font-size: 1.35rem;
            }}
        }}
        {animation_css}
    </style>
</head>
<body class="{'mobile-reader' if share_view == 'mobile' else 'desktop-reader'}">
    <div class="container">
        <div class="header">
            <h1>{topic or '漫画标题'}</h1>
        </div>
        <div class="comic-grid" id="comicGrid">
'''
    
    animation_class_map = {
        'fade': 'animate-fade',
        'zoom': 'animate-zoom',
        'slide-left': 'animate-slide-left',
        'slide-right': 'animate-slide-right',
        'slide-top': 'animate-slide-top',
        'slide-bottom': 'animate-slide-bottom',
        'bounce': 'animate-bounce',
        'flip': 'animate-flip',
        'rotate': 'animate-rotate',
        'none': ''
    }
    animate_class = animation_class_map.get(animation, 'animate-fade')
    
    if layout_config and layout_config.get('grid_specs') and len(layout_config['grid_specs']) > 0:
        specs = layout_config['grid_specs']
        css_grid = layout_config.get('css_grid', '')
        
        grid_style = f'''
        <style>
            #comicGrid {{
                display: grid;
                {css_grid}
                gap: 16px;
                margin-bottom: 30px;
                width: 100%;
                max-width: 900px;
                margin-left: auto;
                margin-right: auto;
                align-items: stretch;
                min-height: 0;
            }}
        </style>
        '''
        html_content += grid_style
        
        for i, comic_url in enumerate(comic_images):
            if comic_url and i < len(specs):
                spec = specs[i]
                speaker = speakers[i] if i < len(speakers) else ''
                dialogue = dialogues[i] if i < len(dialogues) else ''
                scene_zh = scenes_zh[i] if i < len(scenes_zh) else ''
                image_url = _resolve_src(comic_url)
                
                panel_class = f'comic-panel {animate_class}' if animate_class else 'comic-panel'
                delay_style = f'animation-delay: {round(i * 0.3, 1)}s;' if animate_class else ''
                
                aspect_ratio_style = ''
                if spec.get('target_aspect_ratio'):
                    aspect_ratio_style = f'aspect-ratio: {spec["target_aspect_ratio"].replace(":", " / ")};'
                elif spec.get('target_width') and spec.get('target_height'):
                    aspect_ratio_style = f'aspect-ratio: {spec["target_width"]} / {spec["target_height"]};'
                elif spec.get('aspect_ratio'):
                    aspect_ratio_style = f'aspect-ratio: {spec["aspect_ratio"]};'
                elif spec.get('col_span') and spec.get('row_span'):
                    aspect_ratio_style = f'aspect-ratio: {spec["col_span"]} / {spec["row_span"]};'
                
                panel_style = f'''
                    grid-row: {spec["row"] + 1} / span {spec.get("row_span", 1)};
                    grid-column: {spec["col"] + 1} / span {spec.get("col_span", 1)};
                    {aspect_ratio_style}
                    {delay_style}
                '''
                
                html_content += f'''<div class="{panel_class}" style="{panel_style}">
                    <img src="{image_url}" alt="漫画第{i+1}格" class="comic-image">
                    <div class="panel-number" style="position:absolute;bottom:8px;left:8px;background:rgba(0,0,0,0.6);color:#fff;padding:2px 8px;border-radius:12px;font-size:12px;">第{i+1}格</div>
                    {f'<div class="speaker" style="position:absolute;top:8px;left:8px;background:rgba(255,255,255,0.9);padding:4px 10px;border-radius:8px;font-size:12px;color:#FF6B9D;font-weight:bold;">🎭 {speaker}</div>' if speaker else ''}
                    {f'<div class="dialogue" style="position:absolute;top:35px;left:8px;background:rgba(255,255,255,0.9);padding:4px 10px;border-radius:8px;font-size:12px;color:#333;max-width:80%;">{dialogue}</div>' if dialogue else ''}
                </div>
'''
        
        if len(comic_images) > len(specs):
            for i in range(len(specs), len(comic_images)):
                comic_url = comic_images[i]
                if comic_url:
                    speaker = speakers[i] if i < len(speakers) else ''
                    dialogue = dialogues[i] if i < len(dialogues) else ''
                    scene_zh = scenes_zh[i] if i < len(scenes_zh) else ''
                    image_url = _resolve_src(comic_url)
                    
                    panel_class = f'comic-panel {animate_class}' if animate_class else 'comic-panel'
                    delay_style = f' style="animation-delay: {round(i * 0.3, 1)}s;"' if animate_class else ''
                    
                    html_content += f'''<div class="{panel_class}" style="height:auto;min-height:250px;aspect-ratio:1/1;align-self:stretch;{delay_style}">
                        <img src="{image_url}" alt="漫画第{i+1}格" class="comic-image">
                        <div class="panel-number" style="position:absolute;bottom:8px;left:8px;background:rgba(0,0,0,0.6);color:#fff;padding:2px 8px;border-radius:12px;font-size:12px;">第{i+1}格</div>
                        {f'<div class="speaker" style="position:absolute;top:8px;left:8px;background:rgba(255,255,255,0.9);padding:4px 10px;border-radius:8px;font-size:12px;color:#FF6B9D;font-weight:bold;">🎭 {speaker}</div>' if speaker else ''}
                        {f'<div class="dialogue" style="position:absolute;top:35px;left:8px;background:rgba(255,255,255,0.9);padding:4px 10px;border-radius:8px;font-size:12px;color:#333;max-width:80%;">{dialogue}</div>' if dialogue else ''}
                    </div>
'''
    else:
        for i, comic_url in enumerate(comic_images):
            if comic_url:
                speaker = speakers[i] if i < len(speakers) else ''
                dialogue = dialogues[i] if i < len(dialogues) else ''
                scene_zh = scenes_zh[i] if i < len(scenes_zh) else ''
                image_url = _resolve_src(comic_url)
                
                panel_class = f'comic-panel {animate_class}' if animate_class else 'comic-panel'
                delay_style = f' style="animation-delay: {round(i * 0.3, 1)}s;"' if animate_class else ''
                
                html_content += f'''<div class="{panel_class}"{delay_style}>
                    <div class="panel-number">第{i+1}格</div>
                    <img src="{image_url}" alt="漫画第{i+1}格" class="comic-image-scaled">
                    {f'<div class="speaker">🎭 {speaker}</div>' if speaker else ''}
                    {f'<div class="dialogue">💬 {dialogue}</div>' if dialogue else ''}
                    {f'<div class="scene">📍 {scene_zh}</div>' if scene_zh else ''}
                </div>
'''
    
    html_content += '''</div>
        <div class="footer"><p>真读书·有深度·用得上</p></div>
    </div>
</body>
</html>'''
    
    return html_content


@ app.route('/api/export', methods=['GET', 'POST'])
def api_export():
    """导出漫画（图片版、ZIP版、链接版、PPT版）"""
    try:
        if request.method == 'POST':
            data = request.get_json(silent=True) or {}

            def _get_field(name, default=None):
                """优先从 JSON 体读取，缺失则从表单字段读取（兼容 multipart/form-data 上传）。"""
                if name in data and data[name] is not None:
                    return data[name]
                return request.form.get(name, default)

            # 复杂 JSON 字段可能以 form 字符串形式到达，尝试反序列化
            raw_images = _get_field('images', [])
            if isinstance(raw_images, str):
                try:
                    raw_images = json.loads(raw_images)
                except Exception:
                    raw_images = []
            raw_panels = _get_field('panels', [])
            if isinstance(raw_panels, str):
                try:
                    raw_panels = json.loads(raw_panels)
                except Exception:
                    raw_panels = []
            raw_layout = _get_field('layout_config', None)
            if isinstance(raw_layout, str):
                try:
                    raw_layout = json.loads(raw_layout)
                except Exception:
                    raw_layout = None

            export_format = _get_field('format', 'images')
            comic_images = raw_images or []
            script = _get_field('script', {})
            if isinstance(script, str):
                try:
                    script = json.loads(script)
                except Exception:
                    script = {}
            if not isinstance(script, dict):
                script = {}
            export_title = str(_get_field('title', '') or '').strip()
            animation = _get_field('animation', 'fade')
            layout_config = raw_layout
            share_mode = _get_field('share_mode', 'server')
            panels_data = raw_panels or []  # 分层数据，GIF/MP4/PPT 动态效果使用
        else:
            task_id = request.args.get('task_id')
            export_format = request.args.get('format', 'images')
            animation = request.args.get('animation', 'fade')
            
            if not task_id or task_id not in tasks:
                return jsonify({'error': '任务不存在'}), 404
            
            task = tasks[task_id]
            comic_images = task.comics
            script = task.script
            export_title = str(getattr(task, 'title', '') or '').strip()
            layout_config = None  # GET 路径无布局配置，汇总图按默认 four-grid（与原行为一致）

        if not isinstance(script, dict):
            script = {}
        
        # 部署时设置 PUBLIC_BASE_URL（例如 https://comic.example.com）。本机开发仍可
        # 正常使用 localhost；分享页中的所有图片也使用同一公开域名，避免外部访问者加载不到图。
        public_base_url = (
            os.getenv('PUBLIC_BASE_URL', '') or
            os.getenv('COMIC_SHARE_BASE_URL', '')
        ).strip().rstrip('/')
        if public_base_url:
            base_url = public_base_url
        else:
            # 从本机 127.0.0.1 打开时，分享链接至少应换成局域网 IP，供同一
            # Wi-Fi/办公网中的同事访问；真正外网访问则由 PUBLIC_BASE_URL 覆盖。
            request_host = request.host.split(':', 1)[0]
            if request_host in ('127.0.0.1', 'localhost', '::1'):
                try:
                    import socket
                    host_ips = socket.gethostbyname_ex(socket.gethostname())[2]
                    lan_candidates = [
                        ip for ip in host_ips
                        if not (ip.startswith('127.') or ip.startswith('169.254.') or
                                ip.startswith('198.18.') or ip.startswith('0.'))
                    ]
                    lan_host = lan_candidates[0] if lan_candidates else None
                    if not lan_host:
                        probe = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
                        probe.connect(('8.8.8.8', 80))
                        lan_host = probe.getsockname()[0]
                        probe.close()
                    request_port = request.environ.get('SERVER_PORT', '5000')
                    base_url = f'{request.scheme}://{lan_host}:{request_port}'
                except Exception:
                    base_url = request.host_url.rstrip('/')
            else:
                base_url = request.host_url.rstrip('/')
        if not base_url.endswith('/'):
            base_url += '/'
        
        if export_format == 'images':
            return jsonify({
                'success': True,
                'format': 'images',
                'comics': comic_images
            })
        
        elif export_format == 'zip':
            try:
                import zipfile
                from io import BytesIO
                from flask import send_file

                # 汇总大图：直接把完整 layout_config 字典交给 compose_aggregate_image（与链接版同源）
                aggregate_bytes = None
                layout_arg = layout_config if isinstance(layout_config, (dict, str)) and layout_config else 'four-grid'
                try:
                    aggregate_bytes = compose_aggregate_image(comic_images, layout=layout_arg)
                except Exception as agg_e:
                    logger.warning(f"汇总图生成失败（不影响单格导出）: {agg_e}")

                zip_buffer = BytesIO()
                zip_title = export_title or str(script.get('title') or '').strip()
                zip_basename = sanitize_export_filename(zip_title or PPT_DEFAULT_TITLE)

                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    for i, comic_url in enumerate(comic_images):
                        if comic_url:
                            # 修复：以传入 URL 为准（含编辑后 panel_i_edited.png?t=时间戳），strip 查询串；
                            # 不再无脑优先旧残留 panel_i_edited.png
                            img_path = os.path.join(app.root_path, comic_url.split('?')[0].lstrip('/'))
                            if os.path.exists(img_path):
                                zipf.write(img_path, f'{zip_basename}_第{i+1}格.png')
                            else:
                                edited_path = os.path.join(app.root_path, 'static', 'output', f'panel_{i}_edited.png')
                                if os.path.exists(edited_path):
                                    zipf.write(edited_path, f'{zip_basename}_第{i+1}格.png')
                    if aggregate_bytes:
                        zipf.writestr(f'{zip_basename}_合并图.png', aggregate_bytes)

                zip_buffer.seek(0)

                # 压缩包文件名 = 漫画名，与 PPT / MP4 导出保持一致的命名规则
                return send_file(
                    zip_buffer,
                    mimetype='application/zip',
                    as_attachment=True,
                    download_name=f'{zip_basename}_图片.zip'
                )

            except Exception as e:
                logger.error(f"ZIP打包失败: {str(e)}")
                return jsonify({'success': False, 'error': f'打包失败: {str(e)}'}), 500
        
        elif export_format == 'mp4':
            try:
                import subprocess
                import tempfile
                import shutil
                from io import BytesIO
                from flask import send_file
                # Image/ImageDraw 已在模块顶部导入，此处不再重复导入
                # （重复导入会让 Image 成为函数局部变量，导致 PPT 分支 UnboundLocalError）

                # 优先使用项目依赖自带的 ffmpeg；若环境尚未安装该 Python 包，
                # 仍可使用系统 PATH 中的 ffmpeg，避免因单个可选模块直接中断导出。
                try:
                    import imageio_ffmpeg
                    FFMPEG = imageio_ffmpeg.get_ffmpeg_exe()
                except ImportError:
                    FFMPEG = shutil.which('ffmpeg')
                    if not FFMPEG:
                        raise RuntimeError(
                            '视频组件未安装，请在项目虚拟环境中安装 imageio-ffmpeg'
                        )

                # 读取 MP4 专用参数（兼容 JSON 体与 multipart 表单）
                def _mp4_field(name, default=None):
                    if name in data and data[name] is not None:
                        return data[name]
                    return request.form.get(name, default)

                resolution = _mp4_field('resolution', 'sd')
                bgm = _mp4_field('bgm', 'none')
                volume_raw = _mp4_field('volume', 0.7)
                fade_in_raw = _mp4_field('fade_in', 1.0)
                fade_out_raw = _mp4_field('fade_out', 1.0)
                try:
                    volume = float(volume_raw)
                except (TypeError, ValueError):
                    volume = 0.7
                try:
                    fade_in = float(fade_in_raw)
                except (TypeError, ValueError):
                    fade_in = 1.0
                try:
                    fade_out = float(fade_out_raw)
                except (TypeError, ValueError):
                    fade_out = 1.0
                volume = max(0.0, min(1.0, volume))
                fade_in = max(0.0, fade_in)
                fade_out = max(0.0, fade_out)
                if resolution not in ('sd', 'hd'):
                    resolution = 'sd'

                # 优先用前端透传的 layout_config（与链接版同一份布局）
                layout_config = data.get('layout_config', layout_config)

                # 收集可导出的图片（兼容仅传 panels 的情况）
                export_urls = list(comic_images or [])
                if not export_urls and panels_data:
                    export_urls = [p.get('url') or p.get('background_url')
                                   for p in panels_data if isinstance(p, dict)]
                export_urls = [u for u in export_urls if u]

                loaded_images = []
                for u in export_urls:
                    # 编辑后 URL 会携带缓存参数；文件系统路径必须剥离查询串，
                    # 否则该格在 MP4 中无法读取，导致整格空白或被跳过。
                    clean_u = str(u).split('?', 1)[0]
                    p = os.path.join(app.root_path, clean_u.lstrip('/'))
                    if os.path.exists(p):
                        loaded_images.append(Image.open(p).convert('RGB'))

                if not loaded_images:
                    return jsonify({'success': False, 'error': '没有可导出的图片'}), 400

                def _cover_resize(img, tw, th):
                    iw, ih = img.size
                    tr = tw / th if th else 1
                    ir = iw / ih if ih else 1
                    if ir > tr:
                        nh = ih; nw = int(nh * tr)
                    else:
                        nw = iw; nh = int(nw / tr)
                    return img.resize((max(1, nw), max(1, nh)), Image.BICUBIC)

                def _ken_burns_frame(img, tw, th, progress):
                    iw, ih = img.size
                    scale = 1.0 + 0.10 * progress
                    crop_w = tw / scale; crop_h = th / scale
                    cx = iw * (0.5 + (progress - 0.5) * 0.10)
                    cy = ih * (0.5 + (progress - 0.5) * 0.06)
                    left = max(0, int(cx - crop_w / 2)); top = max(0, int(cy - crop_h / 2))
                    right = min(iw, int(left + crop_w)); bottom = min(ih, int(top + crop_h))
                    if right - left < crop_w: left = max(0, int(right - crop_w))
                    if bottom - top < crop_h: top = max(0, int(bottom - crop_h))
                    return img.crop((left, top, right, bottom)).resize((tw, th), Image.BICUBIC)

                def _render_kenburns():
                    """退化方案：无布局时整图 Ken Burns + 交叉淡化转场"""
                    W = H = 768
                    frames = []
                    for i, img in enumerate(loaded_images):
                        for k in range(8):
                            frames.append(_ken_burns_frame(img, W, H, k / 7))
                        if i < len(loaded_images) - 1:
                            a = _cover_resize(loaded_images[i], W, H).convert('RGBA')
                            b = _cover_resize(loaded_images[i + 1], W, H).convert('RGBA')
                            for t in range(2):
                                frames.append(Image.blend(a, b, (t + 1) / 3).convert('RGB'))
                    return frames

                def _render_grid_reveal():
                    """漫画网格逐格揭示：格子按布局出现，出现后留在原地（与链接版一致）。"""
                    grid_specs = layout_config.get('grid_specs', []) or []
                    css_grid = layout_config.get('css_grid', '')
                    rows_fr, cols_fr = parse_fr_values(css_grid) if css_grid else ([], [])
                    if not cols_fr:
                        maxc = max([s.get('col', 0) for s in grid_specs], default=0)
                        cols_fr = [1.0] * (maxc + 1)
                    if not rows_fr:
                        maxr = max([s.get('row', 0) for s in grid_specs], default=0)
                        rows_fr = [1.0] * (maxr + 1)

                    gap = 10
                    col_total = sum(cols_fr) or 1
                    row_total = sum(rows_fr) or 1
                    # 分辨率驱动画布尺寸：sd=820×1230，hd=1640×2460
                    if resolution == 'hd':
                        RES_W, RES_H = 1640, 2460
                    else:
                        RES_W, RES_H = 820, 1230
                    avail_w = RES_W - gap * (len(cols_fr) - 1)
                    avail_h = RES_H - gap * (len(rows_fr) - 1)
                    unit = min(avail_w / col_total, avail_h / row_total)
                    cw = int(col_total * unit + gap * (len(cols_fr) - 1))
                    ch = int(row_total * unit + gap * (len(rows_fr) - 1))

                    def cell_rect(spec):
                        r = spec.get('row', 0); c = spec.get('col', 0)
                        rs = spec.get('row_span', 1); cs = spec.get('col_span', 1)
                        x = int(sum(cols_fr[:c]) * unit + c * gap)
                        y = int(sum(rows_fr[:r]) * unit + r * gap)
                        w = int(sum(cols_fr[c:c + cs]) * unit + (cs - 1) * gap)
                        h = int(sum(rows_fr[r:r + rs]) * unit + (rs - 1) * gap)
                        return x, y, w, h

                    def _anim_params(anim, t):
                        """返回 (scale, dx_frac, dy_frac, angle_deg, alpha) 用于单格入场特效。
                        t ∈ [0,1] 为揭示进度；dx_frac/dy_frac 为相对单元宽/高的平移比例（slide 用）。
                        未知类型降级为 fade。"""
                        if anim == 'fade':
                            return (1.0, 0.0, 0.0, 0.0, t)
                        if anim == 'zoom':
                            scale = 0.4 + 0.6 * t
                            return (scale, 0.0, 0.0, 0.0, 1.0)
                        if anim == 'bounce':
                            # 0.3 -> 1.05 -> 0.9 -> 1 的 overshoot 关键帧
                            if t < 0.5:
                                s = 0.3 + (1.05 - 0.3) * (t / 0.5)
                            elif t < 0.7:
                                s = 1.05 + (0.9 - 1.05) * ((t - 0.5) / 0.2)
                            else:
                                s = 0.9 + (1.0 - 0.9) * ((t - 0.7) / 0.3)
                            return (s, 0.0, 0.0, 0.0, 1.0)
                        if anim == 'slide-left':
                            return (1.0, -0.4 * (1 - t), 0.0, 0.0, 1.0)
                        if anim == 'slide-right':
                            return (1.0, 0.4 * (1 - t), 0.0, 0.0, 1.0)
                        if anim == 'slide-top':
                            return (1.0, 0.0, -0.4 * (1 - t), 0.0, 1.0)
                        if anim == 'slide-bottom':
                            return (1.0, 0.0, 0.4 * (1 - t), 0.0, 0.0)
                        if anim == 'flip':
                            # 近似 scaleX：横向缩放模拟翻面（GIF/PIL 难以做真正 3D 翻转）
                            scale = 0.2 + 0.8 * t
                            return (scale, 0.0, 0.0, 0.0, 1.0)
                        if anim == 'rotate':
                            angle = -20.0 * (1 - t)
                            return (1.0, 0.0, 0.0, angle, t)
                        if anim == 'none':
                            return (1.0, 0.0, 0.0, 0.0, 1.0)
                        # 未知类型降级为 fade
                        return (1.0, 0.0, 0.0, 0.0, t)

                    def draw_panel(canvas, img, rect, t):
                        """按 animation 类型把 img 以"入场特效"合成到 canvas 的 rect 单元内。
                        t ∈ [0,1] 为当前格的揭示进度；已完全揭示的格子用 t=1.0 调用。"""
                        x, y, w, h = rect
                        scale, dxf, dyf, angle, alpha = _anim_params(animation, t)
                        scale = max(0.01, float(scale))
                        iw, ih = img.size
                        tw = max(1, int(round(w * scale)))
                        th = max(1, int(round(h * scale)))
                        resized = img.resize((tw, th), Image.Resampling.BICUBIC)
                        if angle:
                            resized = resized.rotate(angle, expand=True, resample=Image.BICUBIC)
                        rw, rh = resized.size
                        # 居中到单元中心，再按方向偏移
                        cx = x + w // 2 + int(round(dxf * w))
                        cy = y + h // 2 + int(round(dyf * h))
                        px = cx - rw // 2
                        py = cy - rh // 2
                        if alpha >= 0.995:
                            canvas.paste(resized, (px, py))
                        else:
                            layer = Image.new('RGBA', canvas.size, (0, 0, 0, 0))
                            rgba = resized.convert('RGBA')
                            rgba.putalpha(int(255 * alpha))
                            layer.paste(rgba, (px, py))
                            canvas = Image.alpha_composite(canvas.convert('RGBA'), layer).convert('RGB')
                        return canvas

                    def draw_border(canvas, rect):
                        d = ImageDraw.Draw(canvas)
                        d.rectangle([rect[0], rect[1], rect[0] + rect[2] - 1, rect[1] + rect[3] - 1],
                                    outline=(30, 30, 30), width=3)

                    # 入场动画帧数：none 直接定格（1 帧），其余每格 4 帧逐步揭示（兼顾流畅与文件体积，4 帧恰好命中 bounce 过冲峰值）
                    reveal = 4 if (animation and animation != 'none') else 1
                    frames = []
                    n = len(loaded_images)
                    # 兜底：若图片数量多于布局格数（如布局配置异常/未带布局），
                    # 为多余图片追加 1x1 占位格顺序排在现有格之后，避免 fallback 成整张画布
                    # 叠放造成互相遮挡（尤其是后面几个格子）。
                    if len(grid_specs) < n:
                        ncols = max(len(cols_fr), 1)
                        start = len(grid_specs)
                        base_row = (max([s.get('row', 0) for s in grid_specs], default=-1) + 1) if grid_specs else 0
                        for k in range(start, n):
                            m = k - start
                            grid_specs.append({
                                'row': base_row + m // ncols,
                                'col': m % ncols,
                                'row_span': 1, 'col_span': 1
                            })
                        max_row = max([s.get('row', 0) for s in grid_specs])
                        while len(rows_fr) <= max_row:
                            rows_fr.append(1.0)
                    for i in range(n):
                        spec = grid_specs[i] if i < len(grid_specs) else None
                        rect = cell_rect(spec) if spec else (0, 0, cw, ch)
                        # 揭示本格：按 animation 施加入场特效，每格 reveal 帧，确保相邻帧不重复
                        # （GIF 编码器会丢弃连续相同帧，故每帧 t 递增，保证有差异）
                        for f in range(reveal):
                            t = (f + 1) / reveal
                            canvas = Image.new('RGB', (cw, ch), (255, 255, 255))
                            for j in range(i):
                                js = grid_specs[j] if j < len(grid_specs) else None
                                jr = cell_rect(js) if js else (0, 0, cw, ch)
                                canvas = draw_panel(canvas, loaded_images[j], jr, 1.0)
                                draw_border(canvas, jr)
                            canvas = draw_panel(canvas, loaded_images[i], rect, t)
                            draw_border(canvas, rect)
                            frames.append(canvas)
                    return frames

                if layout_config and isinstance(layout_config, dict) and layout_config.get('grid_specs'):
                    frames = _render_grid_reveal()
                else:
                    frames = _render_kenburns()

                if not frames:
                    return jsonify({'success': False, 'error': '没有可渲染的帧'}), 400

                # 编码参数
                fps = 8 if resolution == 'hd' else 10
                crf = 30 if resolution == 'hd' else 28

                # 封面定格在最后一帧：在帧序列最前插入若干张末帧副本，
                # 使文件管理器/微信等场景取到的首帧(=封面)为完整漫画而非动画早期帧。
                if frames:
                    _cover_hold = max(1, int(fps * 0.5))
                    _last_frame = frames[-1]
                    frames = [_last_frame] * _cover_hold + frames

                tmp_dir = tempfile.mkdtemp(prefix='comic_mp4_')
                try:
                    # 1) 帧写入临时 PNG
                    for idx, fr in enumerate(frames):
                        fr.save(os.path.join(tmp_dir, f'frame_{idx:04d}.png'))

                    video_path = os.path.join(tmp_dir, 'video.mp4')

                    def _run_ffmpeg(args, label):
                        """运行 ffmpeg：带 300s 超时，失败时把 stderr 透出以便排查（修复静默失败/卡死）。"""
                        try:
                            subprocess.run(args, check=True,
                                           stdout=subprocess.DEVNULL,
                                           stderr=subprocess.PIPE,
                                           timeout=300)
                        except subprocess.CalledProcessError as e:
                            stderr_text = (e.stderr.decode('utf-8', 'replace')
                                           if isinstance(e.stderr, bytes) else (e.stderr or ''))
                            raise RuntimeError(f"{label}失败: {stderr_text}") from e
                        except subprocess.TimeoutExpired:
                            raise RuntimeError(f"{label}超时（>300s），ffmpeg 可能被资源耗尽卡死") from None

                    def encode_video(crf_val, fps_val, out_path):
                        cmd = [
                            FFMPEG, '-y',
                            '-framerate', str(fps_val),
                            '-i', os.path.join(tmp_dir, 'frame_%04d.png'),
                            # H.264/yuv420p 要求宽高均为偶数。
                            '-vf', 'scale=trunc(iw/2)*2:trunc(ih/2)*2,format=yuv420p',
                            '-c:v', 'libx264',
                            '-crf', str(crf_val),
                            '-pix_fmt', 'yuv420p',
                            '-movflags', '+faststart',
                            out_path,
                        ]
                        _run_ffmpeg(cmd, '视频编码')

                    encode_video(crf, fps, video_path)

                    final_path = video_path
                    audio_path = None

                    # 2) 背景音乐混流
                    if bgm != 'none':
                        if bgm == 'upload':
                            up = request.files.get('bgm_file')
                            if not up:
                                raise ValueError('bgm=upload 但未收到 bgm_file')
                            audio_path = os.path.join(tmp_dir, 'upload_bgm.mp3')
                            up.save(audio_path)
                        else:
                            bgm_map = {
                                'calm': 'static/bgm/bgm-soothing-calm.mp3',
                                'peaceful': 'static/bgm/bgm-soothing-peaceful.mp3',
                                'piano': 'static/bgm/bgm-soothing-piano.mp3',
                            }
                            rel = bgm_map.get(bgm)
                            if not rel:
                                raise ValueError(f'未知背景音乐类型: {bgm}')
                            audio_path = os.path.join(app.root_path, rel)
                            if not os.path.exists(audio_path):
                                raise FileNotFoundError(f'背景音乐文件不存在: {audio_path}')

                        dur = len(frames) / fps
                        audio_out = os.path.join(tmp_dir, 'with_audio.mp4')
                        af = (f"[1:a]volume={volume},"
                              f"afade=t=in:st=0:d={fade_in},"
                              f"afade=t=out:st={max(0, dur - fade_out):.2f}:d={fade_out},"
                              f"apad[a]")
                        _run_ffmpeg([
                            FFMPEG, '-y',
                            '-i', video_path,
                            '-stream_loop', '-1', '-i', audio_path,
                            '-filter_complex', af,
                            '-map', '0:v', '-map', '[a]',
                            '-t', f'{dur:.2f}',          # 绑定输出到视频时长，修复 -stream_loop -1 + -shortest 死锁
                            '-c:v', 'copy',
                            '-c:a', 'aac', '-b:a', '96k',
                            '-shortest',
                            '-movflags', '+faststart',
                            audio_out,
                        ], '背景音乐混流')
                        final_path = audio_out

                    # 3) ≤ 2MB 体积保护
                    max_bytes = 2 * 1024 * 1024
                    best_path = final_path
                    best_size = os.path.getsize(best_path)
                    for attempt in range(3):
                        if best_size <= max_bytes:
                            break
                        crf = min(38, crf + 4)
                        fps = max(5, fps - 2)
                        re_video = os.path.join(tmp_dir, f'video_r{attempt + 1}.mp4')
                        encode_video(crf, fps, re_video)
                        if bgm != 'none':
                            dur = len(frames) / fps
                            re_audio = os.path.join(tmp_dir, f'with_audio_r{attempt + 1}.mp4')
                            af2 = (f"[1:a]volume={volume},"
                                   f"afade=t=in:st=0:d={fade_in},"
                                   f"afade=t=out:st={max(0, dur - fade_out):.2f}:d={fade_out},"
                                   f"apad[a]")
                            _run_ffmpeg([
                                FFMPEG, '-y',
                                '-i', re_video,
                                '-stream_loop', '-1', '-i', audio_path,
                                '-filter_complex', af2,
                                '-map', '0:v', '-map', '[a]',
                                '-t', f'{dur:.2f}',          # 绑定输出到视频时长，修复 -stream_loop -1 + -shortest 死锁
                                '-c:v', 'copy',
                                '-c:a', 'aac', '-b:a', '96k',
                                '-shortest',
                                '-movflags', '+faststart',
                                re_audio,
                            ], '背景音乐混流(重试)')
                            candidate = re_audio
                        else:
                            candidate = re_video
                        cand_size = os.path.getsize(candidate)
                        if cand_size < best_size:
                            best_size = cand_size
                            best_path = candidate

                    if best_size > max_bytes:
                        logger.warning(f"MP4 仍超过 2MB 限制: {best_size} bytes")

                    # 4) 读取到内存后再清理临时目录，确保安全返回
                    with open(best_path, 'rb') as _fh:
                        _buf = BytesIO(_fh.read())
                    mp4_title = export_title or str(script.get('title') or '').strip()
                    mp4_basename = sanitize_export_filename(mp4_title or PPT_DEFAULT_TITLE)
                    return send_file(_buf, mimetype='video/mp4', as_attachment=True,
                                     download_name=f'{mp4_basename}.mp4')
                finally:
                    shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception as e:
                logger.error(f"MP4生成失败: {str(e)}")
                import traceback
                traceback.print_exc()
                return jsonify({'success': False, 'error': f'MP4生成失败: {str(e)}'}), 500

        elif export_format == 'links':
            # inline 模式：生成自包含 HTML（图片 base64 内联），直接作为文件下载，无需服务器即可分享
            if share_mode == 'inline':
                html_content = create_comic_html_page(comic_images, script, base_url, animation, layout_config, inline=True,
                                                      share_view=data.get('share_view', 'mobile'))
                from io import BytesIO
                from flask import send_file
                buf = BytesIO(html_content.encode('utf-8'))
                buf.seek(0)
                return send_file(buf, mimetype='text/html', as_attachment=True,
                                download_name=f'comic_share_{str(uuid.uuid4())[:8]}.html')

            html_content = create_comic_html_page(comic_images, script, base_url, animation, layout_config,
                                                  share_view=data.get('share_view', 'mobile'))

            aggregate_filename = f"{str(uuid.uuid4())[:8]}_comic.html"
            aggregate_path = os.path.join(OUTPUT_DIR, aggregate_filename)
            with open(aggregate_path, 'w', encoding='utf-8') as f:
                f.write(html_content)

            aggregate_link = base_url + f'static/output/{aggregate_filename}'

            return jsonify({
                'success': True,
                'format': 'links',
                'links': [base_url + comic.lstrip('/') for comic in comic_images if comic],
                'aggregate_link': aggregate_link,
                'share_url': aggregate_link,
                'is_public': bool(public_base_url),
                'message': 'HTML展示页面已生成，可直接在浏览器中打开'
            })
        
        elif export_format == 'ppt':
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt, Emu
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN
                import tempfile
                
                prs = Presentation()
                
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = "大阅读精灵漫画家"
                try:
                    title.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
                    title.text_frame.paragraphs[0].font.size = Pt(36)
                    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x6B, 0x9D)
                except Exception:
                    pass
                
                topic = ''
                if isinstance(script, dict):
                    topic = script.get('topic', '')[:50]
                subtitle.text = f"主题: {topic}..." if topic else "漫画作品"
                try:
                    subtitle.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
                    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
                except Exception:
                    pass
                
                speakers = []
                dialogues = []
                scenes_zh = []
                backgrounds = []
                ip_paths = []
                panels_data = []
                
                if isinstance(script, dict):
                    speakers = script.get('speakers', [])
                    dialogues = script.get('dialogues', [])
                    scenes_zh = script.get('scenes_zh', [])
                    
                    panels = script.get('panels', [])
                    if panels and isinstance(panels, list):
                        for panel in panels:
                            if isinstance(panel, dict):
                                panel_speakers = panel.get('speakers', [])
                                panel_dialogues = panel.get('dialogues', [])
                                panel_scene = panel.get('scene', '')
                                panel_speaker = panel.get('speaker', '')
                                panel_dialogue = panel.get('dialogue', '')
                                
                                if not isinstance(panel_speakers, list):
                                    panel_speakers = [panel_speakers] if panel_speakers else []
                                if not isinstance(panel_dialogues, list):
                                    panel_dialogues = [panel_dialogues] if panel_dialogues else []
                                
                                panels_data.append({
                                    'speakers': panel_speakers,
                                    'dialogues': panel_dialogues,
                                    'scene': panel_scene,
                                    'speaker': panel_speaker,
                                    'dialogue': panel_dialogue
                                })
                
                if request.method == 'POST':
                    data = request.get_json()
                    backgrounds = data.get('backgrounds', [])
                    characters = data.get('characters', [])
                    # 优先使用前端传来的完整面板元数据（分层PPT需要）
                    if data.get('panels'):
                        panels_data = data.get('panels')
                    
                    for char in characters:
                        if isinstance(char, dict):
                            img_url = char.get('image_url') or char.get('url')
                            if img_url:
                                ip_paths.append(img_url)
                        elif isinstance(char, str) and char:
                            ip_paths.append(char)

                    # Logo 水印参数（PPT 分层模式第一格叠加）
                    ppt_logo_url = data.get('logoUrl', None)
                    ppt_logo_position = data.get('logoPosition', 'top-right')
                    ppt_logo_size = data.get('logoSize', 0.08)
                    try:
                        ppt_logo_size = float(ppt_logo_size)
                    except (TypeError, ValueError):
                        ppt_logo_size = 0.08
                
                task_id = None
                if request.method == 'POST':
                    data = request.get_json()
                    task_id = data.get('task_id')
                
                if task_id and task_id in tasks:
                    task = tasks[task_id]
                    if not backgrounds:
                        backgrounds = task.backgrounds
                    if not ip_paths:
                        ip_paths = task.ip_paths
                    if not speakers:
                        speakers = task.script.get('speakers', [])
                    if not dialogues:
                        dialogues = task.script.get('dialogues', [])
                    if not scenes_zh:
                        scenes_zh = task.script.get('scenes_zh', [])
                
                slide_width = Inches(10)
                slide_height = Inches(7.5)
                slide_width_emu = slide_width.emu
                slide_height_emu = slide_height.emu
                
                for i, comic_url in enumerate(comic_images):
                    slide_layout = prs.slide_layouts[5]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 辅助：把 /static/... 或相对路径转成绝对路径
                    def _path_from_url(url):
                        if not url:
                            return None
                        if url.startswith('/'):
                            return os.path.join(app.root_path, url.lstrip('/'))
                        elif url.startswith('http'):
                            return None
                        return os.path.join(app.root_path, url)
                    
                    # 辅助：把图片按 cover 模式裁剪到目标尺寸，填满幻灯片/格子
                    def _cover_image_path(img_path, target_w, target_h):
                        if not img_path or not os.path.exists(img_path):
                            return img_path
                        try:
                            img = Image.open(img_path)
                            img_w, img_h = img.size
                            if img_w <= 0 or img_h <= 0:
                                return img_path
                            target_ratio = target_w / target_h if target_h > 0 else 1
                            img_ratio = img_w / img_h
                            if img_ratio > target_ratio:
                                new_h = img_h
                                new_w = int(new_h * target_ratio)
                                left = (img_w - new_w) // 2
                                crop = (left, 0, left + new_w, new_h)
                            else:
                                new_w = img_w
                                new_h = int(new_w / target_ratio)
                                top = (img_h - new_h) // 2
                                crop = (0, top, new_w, top + new_h)
                            cropped = img.crop(crop)
                            tmp_path = os.path.join(tempfile.gettempdir(), f"ppt_cover_{uuid.uuid4().hex[:8]}.png")
                            cropped.save(tmp_path, format='PNG')
                            return tmp_path
                        except Exception as e:
                            logger.warning(f"cover 裁剪失败: {e}")
                            return img_path
                    
                    # 判断当前面板是否有完整的分层元数据
                    panel = panels_data[i] if isinstance(panels_data, list) and i < len(panels_data) else None
                    # PPT 的首要承诺是与漫画预览完全一致。原“分层重建”会由
                    # 默认导出为分层 PPT：背景、人物、气泡和文字均保留为独立对象，
                    # 以便教师继续移动人物、改气泡及字体。只有缺少面板元数据时才降级为成品图。
                    panel_meta = panels_data[i] if i < len(panels_data) and isinstance(panels_data[i], dict) else {}
                    has_layered = bool(panel_meta.get('background_url') and
                                       (panel_meta.get('characters') or panel_meta.get('bubbles')))
                    
                    slide_width = Inches(10)
                    slide_height = Inches(7.5)
                    slide_width_emu = slide_width.emu
                    slide_height_emu = slide_height.emu
                    
                    if has_layered:
                        # ===== 分层模式：背景 + 人物 + 气泡 + 文字，全部可编辑 =====
                        try:
                            bg_url = panel['background_url']
                            bg_path = _path_from_url(bg_url)
                            bg_img = Image.open(bg_path)
                            bg_w, bg_h = bg_img.size
                            # contain 填充：完整显示背景，保持角色/气泡相对位置（避免 cover 裁剪导致人物被截到幻灯片外）
                            scale = min(slide_width_emu / bg_w, slide_height_emu / bg_h)
                            new_w = int(bg_w * scale)
                            new_h = int(bg_h * scale)
                            left = int((slide_width_emu - new_w) / 2)
                            top = int((slide_height_emu - new_h) / 2)
                            
                            # 背景（最底层）
                            bg_pic = slide.shapes.add_picture(bg_path, left, top, width=new_w, height=new_h)
                            slide.shapes._spTree.insert(0, bg_pic._element)
                            logger.info(f"第{i+1}页：分层模式，背景 {bg_w}x{bg_h}，cover 缩放 scale={scale:.3f}")
                            
                            # 人物
                            for char in panel.get('characters', []):
                                char_url = char.get('url') or char.get('path')
                                char_path = _path_from_url(char_url)
                                if not char_path or not os.path.exists(char_path):
                                    continue
                                cx = int(char.get('x', 0) * scale)
                                cy = int(char.get('y', 0) * scale)
                                cw = max(1, int(char.get('width', 100) * scale))
                                ch = max(1, int(char.get('height', 100) * scale))
                                try:
                                    char_pic = slide.shapes.add_picture(char_path, left + cx, top + cy, width=cw, height=ch)
                                    logger.info(f"第{i+1}页：添加人物图层 {char_url}")
                                except Exception as e:
                                    logger.warning(f"添加人物图层失败: {e}")
                            
                            # 气泡 + 文字
                            for bubble in panel.get('bubbles', []):
                                bubble_url = bubble.get('bubble_url')
                                if bubble_url:
                                    bubble_path = _path_from_url(bubble_url)
                                    if bubble_path and os.path.exists(bubble_path):
                                        bx = int(bubble.get('x', 0) * scale)
                                        by = int(bubble.get('y', 0) * scale)
                                        bw = max(1, int(bubble.get('width', 100) * scale))
                                        bh = max(1, int(bubble.get('height', 100) * scale))
                                        try:
                                            bubble_pic = slide.shapes.add_picture(bubble_path, left + bx, top + by, width=bw, height=bh)
                                            logger.info(f"第{i+1}页：添加气泡图层")
                                        except Exception as e:
                                            logger.warning(f"添加气泡图层失败: {e}")
                                
                                text = bubble.get('text', '')
                                if text:
                                    bx = int(bubble.get('x', 0) * scale)
                                    by = int(bubble.get('y', 0) * scale)
                                    bw = max(1, int(bubble.get('width', 100) * scale))
                                    bh = max(1, int(bubble.get('height', 100) * scale))
                                    try:
                                        text_box = slide.shapes.add_textbox(left + bx, top + by, bw, bh)
                                        tf = text_box.text_frame
                                        tf.word_wrap = True
                                        tf.margin_left = 0
                                        tf.margin_right = 0
                                        tf.margin_top = 0
                                        tf.margin_bottom = 0
                                        tf.text = text
                                        p = tf.paragraphs[0]
                                        # Pillow 的 font_size 是像素，scale 已把像素映射到 EMU；
                                        # 1pt = 12700 EMU。此前按气泡高度猜字号并固定最小 14pt，
                                        # 导致短气泡内文字严重放大、换行超出页面。
                                        source_font_px = float(bubble.get('font_size') or bubble.get('fontSize') or 18)
                                        font_size_pt = max(7, min(28, int(source_font_px * scale / 12700)))
                                        p.font.size = Pt(font_size_pt)
                                        p.font.name = 'Microsoft YaHei'
                                        p.font.bold = True
                                        p.font.color.rgb = RGBColor(50, 50, 50)
                                        p.alignment = PP_ALIGN.CENTER
                                    except Exception as e:
                                        logger.warning(f"添加气泡文字失败: {e}")

                            # Logo 水印（仅第一格，PPT 分层模式额外叠加）
                            if i == 0 and ppt_logo_url:
                                try:
                                    if ppt_logo_url.startswith('/'):
                                        logo_local_path = os.path.join(app.root_path, ppt_logo_url.lstrip('/'))
                                    elif os.path.isabs(ppt_logo_url):
                                        logo_local_path = ppt_logo_url
                                    else:
                                        logo_local_path = os.path.join(app.root_path, ppt_logo_url)

                                    if os.path.exists(logo_local_path):
                                        logo_img = Image.open(logo_local_path)
                                        logo_w_px, logo_h_px = logo_img.size
                                        logo_w_emu = max(1, int(new_w * ppt_logo_size))
                                        logo_h_emu = max(1, int(logo_w_emu * logo_h_px / logo_w_px)) if logo_w_px > 0 else logo_w_emu
                                        logo_margin = int(new_w * 0.02)

                                        if ppt_logo_position == 'top-left':
                                            logo_left = left + logo_margin
                                            logo_top = top + logo_margin
                                        elif ppt_logo_position == 'bottom-right':
                                            logo_left = left + new_w - logo_w_emu - logo_margin
                                            logo_top = top + new_h - logo_h_emu - logo_margin
                                        elif ppt_logo_position == 'bottom-left':
                                            logo_left = left + logo_margin
                                            logo_top = top + new_h - logo_h_emu - logo_margin
                                        else:  # top-right
                                            logo_left = left + new_w - logo_w_emu - logo_margin
                                            logo_top = top + logo_margin

                                        slide.shapes.add_picture(logo_local_path, logo_left, logo_top,
                                                                 width=logo_w_emu, height=logo_h_emu)
                                        logger.info(f"PPT第{i+1}页：Logo 水印叠加成功 position={ppt_logo_position}")
                                    else:
                                        logger.warning(f"PPT Logo 文件不存在: {logo_local_path}")
                                except Exception as e:
                                    logger.warning(f"PPT Logo 叠加失败: {e}")
                        except Exception as e:
                            logger.error(f"分层PPT第{i+1}页失败，降级为合成图: {e}")
                            has_layered = False
                    
                    if not has_layered:
                        # ===== 降级模式：直接插入合成图 =====
                        if comic_url:
                            comic_url_clean = comic_url.split('?', 1)[0]
                            if comic_url_clean.startswith('/'):
                                img_path = os.path.join(app.root_path, comic_url_clean.lstrip('/'))
                            elif comic_url_clean.startswith('http'):
                                try:
                                    import tempfile
                                    img_response = requests.get(comic_url_clean, timeout=30)
                                    img_response.raise_for_status()
                                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                                        tmp.write(img_response.content)
                                        img_path = tmp.name
                                except Exception as e:
                                    logger.error(f"下载漫画图片失败: {str(e)}")
                                    continue
                            else:
                                img_path = os.path.join(app.root_path, comic_url_clean)

                            # 仅当当前成品路径失效时，才回退到旧编辑缓存；不能覆盖用户当前预览。
                            if not os.path.exists(img_path):
                                edited_path = os.path.join(app.root_path, 'static', 'output', f'panel_{i}_edited.png')
                                if os.path.exists(edited_path):
                                    img_path = edited_path
                                    logger.info(f"PPT导出第{i+1}页回退使用编辑缓存: panel_{i}_edited.png")
                            
                            if os.path.exists(img_path):
                                try:
                                    img = Image.open(img_path)
                                    img_w, img_h = img.size
                                    # contain 填充：完整显示图片，居中留白，避免人物被裁到幻灯片外
                                    scale = min(slide_width_emu / img_w, slide_height_emu / img_h)
                                    new_w = int(img_w * scale)
                                    new_h = int(img_h * scale)
                                    left = int((slide_width_emu - new_w) / 2)
                                    top = int((slide_height_emu - new_h) / 2)
                                    
                                    pic = slide.shapes.add_picture(img_path, left, top, width=new_w, height=new_h)
                                    try:
                                        pic.line.color.rgb = RGBColor(0x22, 0x22, 0x22)
                                        pic.line.width = Pt(3)
                                    except Exception:
                                        pass
                                    slide.shapes._spTree.insert(0, pic._element)
                                    logger.info(f"第{i+1}页：使用合成图作为背景（cover 填充）")
                                except Exception as e:
                                    logger.error(f"添加漫画图片失败: {str(e)}")
                    
                    # 场景描述（可选）
                    panel_speakers = []
                    panel_dialogues = []
                    panel_scene = ''
                    if panel and isinstance(panel, dict):
                        panel_speakers = panel.get('speakers', [])
                        panel_dialogues = panel.get('dialogues', [])
                        panel_scene = panel.get('scene', '')
                    
                    if not panel_speakers:
                        panel_speakers = [speakers[i]] if i < len(speakers) else []
                    if not panel_dialogues:
                        panel_dialogues = [dialogues[i]] if i < len(dialogues) else []
                    if not panel_scene:
                        panel_scene = scenes_zh[i] if i < len(scenes_zh) else ''
                    
                    if panel_scene:
                        try:
                            scene_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.4))
                            scene_tf = scene_box.text_frame
                            scene_tf.text = panel_scene
                            for paragraph in scene_tf.paragraphs:
                                paragraph.font.size = Pt(12)
                                paragraph.font.color.rgb = RGBColor(100, 100, 100)
                                paragraph.font.name = 'Microsoft YaHei'
                        except Exception as e:
                            logger.warning(f"添加场景描述失败: {e}")
                
                if isinstance(script, dict):
                    if script.get('summary'):
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        content = slide.placeholders[1]
                        title.text = "学习总结"
                        content.text = script['summary']
                    elif script.get('hints'):
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        content = slide.placeholders[1]
                        title.text = "学习总结"
                        content.text = "\n".join(script['hints'])
                
                output_filename = f"{str(uuid.uuid4())[:8]}_comic.pptx"
                output_path = os.path.join(OUTPUT_DIR, output_filename)
                prs.save(output_path)
                
                from flask import send_file
                return send_file(output_path, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True, download_name=output_filename)
                
            except ImportError:
                return jsonify({'error': '需要安装python-pptx库'}), 500
        
        else:
            return jsonify({'error': '不支持的导出格式'}), 400
        
    except Exception as e:
        logger.error(f"导出失败: {str(e)}")
        return jsonify({'error': f'导出失败: {str(e)}'}), 500


# ---------------------------------------------------------------------------
# 视频直链导出（Railway 永久链接 + 本地临时链接兼容）
# ---------------------------------------------------------------------------
# 背景：本地 Flask 进程会被系统回收，导致 /api/export 生成的分享链接失效。
# 方案：把 MP4 + 自包含播放页写入一个独立部署目录，由主理人上传到 CloudStudio，
#       得到不依赖本地服务的公网直链。本模块只负责"准备部署目录"，不做部署。

VIDEO_SHARE_DEPLOY_ROOT = os.path.normpath(
    os.environ.get(
        'VIDEO_SHARE_STORAGE_ROOT',
        '/data/video_shares' if os.environ.get('RAILWAY_ENVIRONMENT')
        else r'C:/Users/matiancheng/WorkBuddy/2026-07-19-21-49-54/cloudstudio_deploy'
    )
)

BUNDLED_VIDEO_SHARE_ROOT = os.path.join(app.root_path, 'seed_video_shares')
VIDEO_SHARE_MANIFEST_PATH = os.path.join(BUNDLED_VIDEO_SHARE_ROOT, 'manifest.json')

VIDEO_SHARE_PAGE_TITLE = '漫画视频'

# 临时视频公网分享：每次分享只启动一个指向“该视频独立目录”的静态服务器，
# 再由 cloudflared Quick Tunnel 暴露该服务器；不会把 Flask 主站暴露到公网。
CLOUDFLARED_PATH = os.path.join(app.root_path, 'tools', 'cloudflared.exe')
VIDEO_SHARE_TUNNELS = {}


def _load_bundled_video_manifest():
    """读取随部署发布的历史视频清单。"""
    try:
        with open(VIDEO_SHARE_MANIFEST_PATH, 'r', encoding='utf-8') as fh:
            value = json.load(fh)
        return value if isinstance(value, dict) else {}
    except (OSError, ValueError) as exc:
        logger.warning(f'[video-share] 历史视频清单读取失败: {exc}')
        return {}


BUNDLED_VIDEO_SHARE_MANIFEST = _load_bundled_video_manifest()


def _safe_video_share_id(share_id):
    """只允许短横线、下划线与 ASCII 字母数字，避免目录穿越。"""
    value = str(share_id or '').strip().lower()
    if not re.fullmatch(r'[a-z0-9_-]{3,64}', value):
        abort(404)
    return value


def _resolve_video_share(share_id):
    """先找 Volume 中的新视频，再找随程序发布的 40 份历史视频。"""
    share_id = _safe_video_share_id(share_id)
    storage_dir = os.path.join(VIDEO_SHARE_DEPLOY_ROOT, f'video_share_{share_id}')
    storage_video = os.path.join(storage_dir, 'video.mp4')
    storage_meta = os.path.join(storage_dir, 'metadata.json')
    if os.path.isfile(storage_video):
        title = VIDEO_SHARE_PAGE_TITLE
        try:
            with open(storage_meta, 'r', encoding='utf-8') as fh:
                title = str((json.load(fh) or {}).get('title') or title)
        except (OSError, ValueError, TypeError):
            pass
        return storage_video, title

    bundled_meta = BUNDLED_VIDEO_SHARE_MANIFEST.get(share_id)
    bundled_video = os.path.join(BUNDLED_VIDEO_SHARE_ROOT, f'{share_id}.mp4')
    if bundled_meta and os.path.isfile(bundled_video):
        title = bundled_meta.get('title') if isinstance(bundled_meta, dict) else bundled_meta
        return bundled_video, str(title or VIDEO_SHARE_PAGE_TITLE)
    abort(404)


@app.route('/videos/<share_id>/')
def permanent_video_share_page(share_id):
    """Railway 上稳定不变的视频播放页。"""
    _video_path, title = _resolve_video_share(share_id)
    return build_video_share_page(
        video_filename=url_for('permanent_video_share_media', share_id=share_id),
        page_title=title
    )


@app.route('/videos/<share_id>/media')
def permanent_video_share_media(share_id):
    """支持浏览器 Range 请求的 MP4 文件响应。"""
    video_path, title = _resolve_video_share(share_id)
    return send_file(
        video_path,
        mimetype='video/mp4',
        conditional=True,
        download_name=f'{sanitize_export_filename(title)}.mp4'
    )


def _reserve_local_port():
    """获取一个可供临时静态视频服务器使用的本机端口。"""
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
        sock.bind(('127.0.0.1', 0))
        return sock.getsockname()[1]


def start_temporary_video_share(deploy_dir, share_id):
    """启动仅包含当前视频的本机静态服务与 Cloudflare 临时公网隧道。

    返回播放页公网基地址。Quick Tunnel 不需用户域名或账号，
    但链接只在本机和本进程保持运行时有效。
    """
    if not os.path.isfile(CLOUDFLARED_PATH):
        raise RuntimeError('未找到 cloudflared 客户端，请先安装后再生成临时公网链接。')

    port = _reserve_local_port()
    creation_flags = getattr(subprocess, 'CREATE_NO_WINDOW', 0)
    null_output = subprocess.DEVNULL
    http_process = subprocess.Popen(
        [sys.executable, '-m', 'http.server', str(port), '--bind', '127.0.0.1',
         '--directory', deploy_dir],
        stdin=null_output, stdout=null_output, stderr=null_output,
        creationflags=creation_flags
    )

    # Quick Tunnel 偶尔会受网络抖动或 Cloudflare 首次分配延迟影响。此前只
    # 尝试一次、等待 25 秒，因此看起来像“偶发点击报错”。静态服务只启动一次，
    # 隧道则最多重试一次，每次等待 40 秒。
    log_handle = None
    tunnel_process = None
    public_base = None
    last_error = ''
    try:
        time.sleep(0.4)
        if http_process.poll() is not None:
            raise RuntimeError('本机视频临时服务启动失败。')

        for attempt in range(2):
            log_path = os.path.join(deploy_dir, f'cloudflared_{attempt + 1}.log')
            log_handle = open(log_path, 'w+', encoding='utf-8')
            tunnel_process = subprocess.Popen(
                [CLOUDFLARED_PATH, 'tunnel', '--no-autoupdate', '--url', f'http://127.0.0.1:{port}'],
                stdin=null_output, stdout=log_handle, stderr=subprocess.STDOUT,
                creationflags=creation_flags
            )
            # cloudflared 会把随机 trycloudflare 地址写进启动日志；每次最多等待 40 秒。
            for _ in range(80):
                if tunnel_process.poll() is not None:
                    break
                log_handle.flush()
                log_handle.seek(0)
                match = re.search(r'https://[a-zA-Z0-9-]+\.trycloudflare\.com', log_handle.read())
                if match:
                    public_base = match.group(0).rstrip('/')
                    break
                time.sleep(0.5)
            if public_base:
                break

            log_handle.flush()
            log_handle.seek(0)
            last_error = log_handle.read()[-500:].strip()
            if tunnel_process.poll() is None:
                tunnel_process.terminate()
                try:
                    tunnel_process.wait(timeout=3)
                except subprocess.TimeoutExpired:
                    tunnel_process.kill()
            log_handle.close()
            log_handle = None
            tunnel_process = None
            if attempt == 0:
                time.sleep(1)
        if not public_base:
            detail = f'（{last_error[-160:]}）' if last_error else ''
            raise RuntimeError(f'Cloudflare 临时隧道未能启动，请检查网络后重试。{detail}')

        VIDEO_SHARE_TUNNELS[share_id] = {
            'http_process': http_process,
            'tunnel_process': tunnel_process,
            'log_handle': log_handle,
            'public_base': public_base,
            'started_at': datetime.now().isoformat()
        }
        return public_base
    except Exception:
        for process in (tunnel_process, http_process):
            if process is not None and process.poll() is None:
                process.terminate()
        if log_handle is not None:
            log_handle.close()
        raise


def build_video_share_page(video_filename='comic.mp4', page_title=VIDEO_SHARE_PAGE_TITLE):
    """构建自包含的响应式视频播放页。

    要求：不引用任何外部 CDN，全部内联样式，离线/公网均可直接打开；
    页面占满视口无滚动条，视频等比居中铺满。

    Args:
        video_filename: 与 index.html 同目录的视频文件名，默认 ``comic.mp4``。
        page_title: 页面标题文案，默认「漫画视频」。

    Returns:
        str: 完整的 HTML 文本。
    """
    safe_title = str(page_title or VIDEO_SHARE_PAGE_TITLE)
    safe_video = str(video_filename or 'comic.mp4')
    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<meta name="format-detection" content="telephone=no">
<title>{safe_title}</title>
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  html, body {{
    width: 100%;
    height: 100%;
    overflow: hidden;
    background: #000;
    font-family: -apple-system, BlinkMacSystemFont, "PingFang SC",
                 "Microsoft YaHei", "Helvetica Neue", Arial, sans-serif;
  }}
  .stage {{
    position: fixed;
    top: 0;
    left: 0;
    width: 100%;
    height: 100%;
    display: flex;
    align-items: center;
    justify-content: center;
    background: #000;
  }}
  video {{
    width: 100%;
    height: 100%;
    max-width: 100vw;
    max-height: 100vh;
    object-fit: contain;
    display: block;
    background: #000;
  }}
  .page-title {{
    position: fixed;
    top: 12px;
    left: 0;
    width: 100%;
    text-align: center;
    color: #fff;
    font-size: 15px;
    font-weight: 600;
    letter-spacing: 1px;
    text-shadow: 0 1px 4px rgba(0, 0, 0, 0.8);
    pointer-events: none;
    z-index: 2;
  }}
  .tip {{
    position: fixed;
    bottom: 10px;
    left: 0;
    width: 100%;
    text-align: center;
    color: rgba(255, 255, 255, 0.6);
    font-size: 12px;
    text-shadow: 0 1px 3px rgba(0, 0, 0, 0.8);
    pointer-events: none;
    z-index: 2;
  }}
</style>
</head>
<body>
  <div class="page-title">{safe_title}</div>
  <div class="stage">
    <video
      src="{safe_video}"
      controls
      muted
      autoplay
      playsinline
      webkit-playsinline
      preload="auto">
      您的浏览器不支持 HTML5 视频播放。
    </video>
  </div>
  <div class="tip">如无法播放，请检查网络</div>
  <script>
    // 不循环播放；ended 时保持浏览器已解码的最后一帧。
    document.querySelector('video').addEventListener('ended', function () {{ this.pause(); }});
  </script>
</body>
</html>
"""


def generate_share_mp4_bytes(payload):
    """复用 ``/api/export`` 的 ``format=mp4`` 分支生成 MP4 二进制数据。

    通过 ``app.test_request_context`` 直接调用 :func:`api_export`，
    既完整复用既有渲染/编码/压缩逻辑，又不改动原接口任何一行代码。

    Args:
        payload: 与 ``/api/export`` 一致的请求体字典（images/panels/script 等）。

    Returns:
        bytes: 生成好的 MP4 文件字节流。

    Raises:
        RuntimeError: MP4 生成失败时抛出，附带原接口的错误文案。
    """
    mp4_payload = dict(payload or {})
    mp4_payload['format'] = 'mp4'

    # 内部调用无法携带 multipart 文件，自定义上传 BGM 降级为无背景音乐，避免必然失败。
    if mp4_payload.get('bgm') == 'upload':
        app.logger.warning('[video-share] bgm=upload 无法内部转发文件，已降级为 none')
        mp4_payload['bgm'] = 'none'

    with app.test_request_context('/api/export', method='POST', json=mp4_payload):
        result = api_export()

    status_code = 200
    response_obj = result
    if isinstance(result, tuple):
        response_obj = result[0]
        if len(result) > 1 and isinstance(result[1], int):
            status_code = result[1]

    if status_code != 200:
        error_text = 'MP4 生成失败'
        try:
            payload_json = response_obj.get_json(silent=True) or {}
            error_text = payload_json.get('error') or payload_json.get('message') or error_text
        except Exception:
            pass
        raise RuntimeError(error_text)

    # send_file 返回的 Response 默认 direct_passthrough=True，直接 get_data 会报错。
    try:
        response_obj.direct_passthrough = False
    except Exception:
        pass

    mp4_bytes = response_obj.get_data()
    if not mp4_bytes:
        raise RuntimeError('MP4 生成结果为空')

    content_type = (getattr(response_obj, 'mimetype', '') or '').lower()
    if content_type and 'video' not in content_type:
        # 例如误返回 application/json 错误体
        try:
            payload_json = response_obj.get_json(silent=True) or {}
            raise RuntimeError(payload_json.get('error') or f'MP4 生成返回了非视频内容: {content_type}')
        except RuntimeError:
            raise
        except Exception:
            raise RuntimeError(f'MP4 生成返回了非视频内容: {content_type}')

    return mp4_bytes


@app.route('/api/export-video-share', methods=['POST'])
def api_export_video_share():
    """生成"视频直链包"：MP4 + 响应式播放页，输出到 CloudStudio 部署目录。

    请求体与 ``/api/export``（format=mp4）完全一致，例如::

        {"images": [...], "panels": [...], "script": {...},
         "layout_config": {...}, "resolution": "sd", "bgm": "none"}

    Returns:
        flask.Response: 成功返回 deploy_dir / mp4_size / mp4_filename；
        失败返回 ``{"success": false, "error": "..."}`` 与 HTTP 500。
    """
    import shutil  # noqa: F401  # 保留：与既有 mp4 分支一致的局部导入风格

    share_id = uuid.uuid4().hex[:8]
    app.logger.info(f'[video-share] 开始生成视频直链包 share_id={share_id}')

    try:
        data = request.get_json(silent=True) or {}

        # 1) 复用既有 MP4 生成逻辑
        mp4_bytes = generate_share_mp4_bytes(data)
        app.logger.info(f'[video-share] MP4 生成完成 share_id={share_id} size={len(mp4_bytes)} bytes')

        # 2) 创建部署目录（Windows 反斜杠安全）
        deploy_dir = os.path.normpath(
            os.path.join(VIDEO_SHARE_DEPLOY_ROOT, f'video_share_{share_id}')
        )
        os.makedirs(deploy_dir, exist_ok=True)

        # 3) 写入 MP4。固定文件名让永久路由不受中文标题和改名影响。
        share_script = data.get('script') or {}
        if isinstance(share_script, str):
            try:
                share_script = json.loads(share_script)
            except Exception:
                share_script = {}
        if not isinstance(share_script, dict):
            share_script = {}
        share_title = str(
            data.get('title') or share_script.get('title') or PPT_DEFAULT_TITLE
        ).strip() or PPT_DEFAULT_TITLE
        mp4_filename = 'video.mp4'
        mp4_path = os.path.normpath(os.path.join(deploy_dir, mp4_filename))
        with open(mp4_path, 'wb') as fh:
            fh.write(mp4_bytes)
        mp4_size = os.path.getsize(mp4_path)

        with open(os.path.join(deploy_dir, 'metadata.json'), 'w', encoding='utf-8') as fh:
            json.dump({'title': share_title, 'created_at': datetime.now().isoformat()}, fh,
                      ensure_ascii=False, indent=2)

        # 4) 写入自包含响应式播放页（中文文件名需 URL 编码后再写进 <source src>）
        from urllib.parse import quote as _url_quote

        index_path = os.path.normpath(os.path.join(deploy_dir, 'index.html'))
        with open(index_path, 'w', encoding='utf-8') as fh:
            fh.write(build_video_share_page(
                video_filename=_url_quote(mp4_filename),
                page_title=share_title
            ))

        # 5) Railway 直接通过固定路由提供服务；本地仍保留 Quick Tunnel 兼容。
        is_railway = bool(os.environ.get('RAILWAY_ENVIRONMENT'))
        if is_railway:
            public_url = url_for('permanent_video_share_page', share_id=share_id, _external=True)
            video_url = url_for('permanent_video_share_media', share_id=share_id, _external=True)
            has_volume = bool(os.environ.get('RAILWAY_VOLUME_MOUNT_PATH'))
            share_mode = 'railway_persistent' if has_volume else 'railway_ephemeral'
            share_notice = (
                '永久链接已保存到 Railway 持久化存储。' if has_volume else
                '固定链接已生成，但 Railway 尚未挂载持久化 Volume，重新部署后文件可能丢失。'
            )
        else:
            public_base = start_temporary_video_share(deploy_dir, share_id)
            public_url = public_base + '/'
            video_url = public_base + '/' + _url_quote(mp4_filename)
            share_mode = 'temporary_tunnel'
            share_notice = '本地临时链接：本机和本程序保持运行期间可访问。'

        deploy_dir_out = deploy_dir.replace('\\', '/').rstrip('/') + '/'
        app.logger.info(
            f'[video-share] 视频直链包已生成 share_id={share_id} '
            f'deploy_dir={deploy_dir_out} mp4_size={mp4_size}'
        )

        return jsonify({
            'success': True,
            'message': '视频直链已生成',
            'deploy_dir': deploy_dir_out,
            'mp4_size': mp4_size,
            'mp4_filename': mp4_filename,
            'public_url': public_url,
            'video_url': video_url,
            'share_mode': share_mode,
            'share_notice': share_notice
        })

    except Exception as e:
        app.logger.error(f'[video-share] 生成失败 share_id={share_id}: {e}')
        logger.error(f'视频直链包生成失败: {e}')
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': f'视频直链包生成失败: {e}'}), 500


# ---------------------------------------------------------------------------
# PPT 导出（合并漫画页 · 每格独立可编辑）
# ---------------------------------------------------------------------------
# 与 /api/export?format=ppt 的区别：
#   - 旧接口：一格一页，直接 send_file 返回二进制流；
#   - 新接口：一页承载全部格子（2x3 等接近正方形的网格），每格是**独立的**
#     图片形状 + 独立文本框，用户在 PowerPoint 里可单独拖动/删除/替换/改字；
#     并且返回 JSON（pptx_url），便于前端弹窗提示或新标签页下载。
# ---------------------------------------------------------------------------
PPT_SLIDE_WIDTH_INCHES = 13.333   # 16:9 宽屏，给 2x3 网格更充裕的横向空间
PPT_SLIDE_HEIGHT_INCHES = 7.5
PPT_DEFAULT_TITLE = '大阅读精灵漫画'


def resolve_local_asset_path(url):
    """把前端传入的图片 URL 转换成服务器本地绝对路径。

    兼容以下几种形态：
      * ``/static/output/xxx.png``（站内绝对路径）
      * ``static/output/xxx.png``（站内相对路径）
      * ``http://host/static/output/xxx.png``（带域名，截取 /static/ 之后的部分）
      * 携带缓存参数的 ``xxx.png?t=123``（编辑后前端会追加时间戳）

    Args:
        url: 前端传入的图片地址，允许为 None/空串。

    Returns:
        str | None: 存在的本地绝对路径；无法解析或文件不存在时返回 None。
    """
    if not url:
        return None
    raw = str(url).strip()
    if not raw:
        return None

    # 去掉查询串与锚点，否则拼出的文件系统路径必然不存在
    raw = raw.split('?', 1)[0].split('#', 1)[0]
    if not raw:
        return None

    if raw.startswith('http://') or raw.startswith('https://'):
        marker = '/static/'
        idx = raw.find(marker)
        if idx == -1:
            return None
        raw = raw[idx:]

    candidate = os.path.normpath(os.path.join(app.root_path, raw.lstrip('/\\')))
    return candidate if os.path.exists(candidate) else None


def compute_ppt_grid(panel_count):
    """按"尽量接近正方形"的原则计算合并页的网格行列数。

    对照关系：4 格 → 2x2，6 格 → 2x3，9 格 → 3x3；
    3 格及以下直接排成一行，避免出现大片空格。

    Args:
        panel_count: 漫画格数量。

    Returns:
        tuple[int, int]: ``(rows, cols)``，两者均 >= 1。
    """
    count = max(int(panel_count or 0), 0)
    if count <= 0:
        return 1, 1
    if count <= 3:
        return 1, count
    cols = int(math.ceil(math.sqrt(count)))
    rows = int(math.ceil(count / float(cols)))
    return max(rows, 1), max(cols, 1)


def extract_panel_caption(panel, index=0):
    """提取某一格的默认可编辑文案。

    取值优先级：第一个气泡文本 > 台词 dialogues > caption > scene/场景描述；
    全部缺失时回退为 ``第N格``，保证文本框始终有内容可编辑。

    Args:
        panel: 单格数据字典。
        index: 该格在整部漫画中的序号（从 0 开始），用于兜底文案。

    Returns:
        str: 供 PPT 文本框使用的文案，已裁剪到 120 字以内。
    """
    fallback = f'第{int(index) + 1}格'
    if not isinstance(panel, dict):
        return fallback

    def _pick(value):
        """从字符串或字典中取出非空文本。"""
        if isinstance(value, str):
            return value.strip()
        if isinstance(value, dict):
            for key in ('text', 'content', 'dialogue', 'caption'):
                text = value.get(key)
                if isinstance(text, str) and text.strip():
                    return text.strip()
        return ''

    for field in ('bubbles', 'dialogues'):
        items = panel.get(field) or []
        if isinstance(items, (list, tuple)):
            for item in items:
                text = _pick(item)
                if text:
                    return text[:120]

    for field in ('caption', 'scene_zh', 'scene', 'description'):
        text = _pick(panel.get(field))
        if text:
            return text[:120]

    return fallback


def add_editable_comic_grid_slide(prs, panels, slide_title=''):
    """向演示文稿追加一页"合并漫画"，每格都是可独立编辑的对象。

    每格产出两个形状：
      1. 一个独立的图片形状（可单独拖动 / 删除 / 右键替换图片）；
      2. 图片正下方一个独立文本框（默认填充台词，可直接改字）。

    Args:
        prs: ``pptx.Presentation`` 实例。
        panels: 已归一化的格子列表，元素为含 ``url`` 键的字典。
        slide_title: 该页标题；传空串则不生成标题、网格占满整页。

    Returns:
        tuple: ``(slide, inserted_count)``，``inserted_count`` 为成功插入的图片数。
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # slide_layouts[6] 是空白版式：不带任何占位符，方便自由摆放
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.35)
    top_start = Inches(0.18)

    if slide_title:
        title_box = slide.shapes.add_textbox(
            margin, top_start, slide_w - margin * 2, Inches(0.55)
        )
        title_box.name = 'ComicSlideTitle'
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = str(slide_title)
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.name = 'Microsoft YaHei'
        title_run.font.size = Pt(22)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(0xFF, 0x6B, 0x9D)
        usable_top = top_start + Inches(0.70)
    else:
        usable_top = top_start

    usable_w = slide_w - margin * 2
    usable_h = slide_h - usable_top - margin

    rows, cols = compute_ppt_grid(len(panels))
    gap = Inches(0.12)
    cell_w = int((usable_w - gap * (cols - 1)) / cols)
    cell_h = int((usable_h - gap * (rows - 1)) / rows)
    caption_h = int(min(Inches(0.42), cell_h * 0.24))
    image_h = max(cell_h - caption_h, 1)

    inserted_count = 0
    for idx, panel in enumerate(panels):
        row_i = idx // cols
        col_i = idx % cols
        cell_left = int(margin + col_i * (cell_w + gap))
        cell_top = int(usable_top + row_i * (cell_h + gap))

        image_url = panel.get('url') or panel.get('image') or panel.get('background_url')
        image_path = resolve_local_asset_path(image_url)

        if image_path:
            # 按原始宽高比等比缩放并居中，避免变形；图片保持独立形状便于单独操作
            try:
                with Image.open(image_path) as img_obj:
                    img_w, img_h = img_obj.size
            except Exception as img_err:
                logger.warning(f'[export-ppt] 读取图片尺寸失败({image_path}): {img_err}')
                img_w, img_h = (4, 3)
            if img_w <= 0 or img_h <= 0:
                img_w, img_h = (4, 3)

            scale = min(cell_w / float(img_w), image_h / float(img_h))
            draw_w = max(int(img_w * scale), 1)
            draw_h = max(int(img_h * scale), 1)
            pic_left = int(cell_left + (cell_w - draw_w) / 2)
            pic_top = int(cell_top + (image_h - draw_h) / 2)

            picture = slide.shapes.add_picture(
                image_path, pic_left, pic_top, width=draw_w, height=draw_h
            )
            picture.name = f'ComicPanel_{idx + 1}'
            inserted_count += 1
        else:
            # 图片缺失时留一个占位文本框，保证网格结构完整且可手动替换
            holder = slide.shapes.add_textbox(cell_left, cell_top, cell_w, image_h)
            holder.name = f'ComicPanelMissing_{idx + 1}'
            holder_para = holder.text_frame.paragraphs[0]
            holder_para.text = f'第{idx + 1}格图片缺失'
            holder_para.alignment = PP_ALIGN.CENTER
            holder_run = holder_para.runs[0]
            holder_run.font.name = 'Microsoft YaHei'
            holder_run.font.size = Pt(12)
            holder_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        caption_box = slide.shapes.add_textbox(
            cell_left, int(cell_top + image_h), cell_w, caption_h
        )
        caption_box.name = f'ComicCaption_{idx + 1}'
        caption_frame = caption_box.text_frame
        caption_frame.word_wrap = True
        caption_para = caption_frame.paragraphs[0]
        caption_para.text = extract_panel_caption(panel, idx)
        caption_para.alignment = PP_ALIGN.CENTER
        caption_run = caption_para.runs[0]
        caption_run.font.name = 'Microsoft YaHei'
        caption_run.font.size = Pt(11)
        caption_run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    return slide, inserted_count


def normalize_export_panels(data):
    """把请求体中的 panels/images 归一化成统一的格子列表。

    兼容三种输入：``panels`` 为字典数组、``panels`` 为 URL 字符串数组、
    仅提供 ``images`` URL 数组；``panels`` 缺 url 时用同下标的 images 补齐。

    Args:
        data: ``/api/export-ppt`` 的请求体字典。

    Returns:
        list[dict]: 每个元素至少含有 ``url`` 或 ``background_url``。
    """
    raw_panels = data.get('panels') or []
    if isinstance(raw_panels, str):
        try:
            raw_panels = json.loads(raw_panels)
        except Exception:
            raw_panels = []
    if not isinstance(raw_panels, list):
        raw_panels = []

    raw_images = data.get('images') or []
    if isinstance(raw_images, str):
        try:
            raw_images = json.loads(raw_images)
        except Exception:
            raw_images = []
    if not isinstance(raw_images, list):
        raw_images = []

    normalized = []
    for item in raw_panels:
        if isinstance(item, dict):
            normalized.append(dict(item))
        elif isinstance(item, str) and item.strip():
            normalized.append({'url': item.strip()})

    if not normalized:
        normalized = [{'url': u} for u in raw_images if isinstance(u, str) and u.strip()]

    # panels 只有分层元数据、没有成图 URL 时，用 images 同下标补齐
    for i, panel in enumerate(normalized):
        if not panel.get('url') and i < len(raw_images) and isinstance(raw_images[i], str):
            panel['url'] = raw_images[i]

    return [p for p in normalized if p.get('url') or p.get('background_url')]


def group_panels_by_sections(panels, raw_sections):
    """按前端「节」配置把格子分组。

    ``panel_ids`` 支持两种写法：前端默认的 ``"p0"/"p1"`` 下标式 ID，
    以及直接的整数下标；越界与重复引用都会被忽略。

    Args:
        panels: 归一化后的格子列表。
        raw_sections: 请求体中的 ``sections``，元素形如
            ``{"id": "sec_x", "title": "第一节", "panel_ids": ["p0", "p1"]}``。

    Returns:
        list[dict]: ``[{'title': str, 'panels': list[dict]}]``；无有效分节时返回空列表。
    """
    if isinstance(raw_sections, str):
        try:
            raw_sections = json.loads(raw_sections)
        except Exception:
            raw_sections = []
    if not isinstance(raw_sections, list) or not raw_sections:
        return []

    total = len(panels)
    grouped = []
    for order, section in enumerate(raw_sections):
        if not isinstance(section, dict):
            continue
        title = str(section.get('title') or f'第{order + 1}节').strip() or f'第{order + 1}节'
        panel_ids = section.get('panel_ids') or []
        if not isinstance(panel_ids, (list, tuple)):
            continue

        picked = []
        used = set()
        for pid in panel_ids:
            index = -1
            if isinstance(pid, int):
                index = pid
            elif isinstance(pid, str):
                digits = re.sub(r'\D', '', pid)
                if digits:
                    index = int(digits)
            if 0 <= index < total and index not in used:
                used.add(index)
                picked.append(panels[index])

        if picked:
            grouped.append({'title': title, 'panels': picked})

    return grouped


# ---------------------------------------------------------------------------
# PPT 矢量化导出：背景图 + 角色图 + 气泡图 + 原生可编辑文字框
# ---------------------------------------------------------------------------
# 与旧的"整格一张图"方案的根本区别：
#   - 旧方案：一格 = 一张合成好的 PNG，PPT 里只能整体拖动，文字改不了；
#   - 新方案：一格 = 若干独立形状（背景 / 每个角色 / 每个气泡 / 每段文字），
#     文字用 PowerPoint 原生 TextBox 承载，字体、字号、颜色、粗斜体、对齐
#     全部还原成网页预览的取值，用户可直接双击改字。
#
# 坐标换算（所有元素共用一套公式，保证分层严丝合缝）：
#   canvas_w/canvas_h : 该格在网页画布上的像素尺寸（bubbles/characters 的坐标基准）
#   scale             = min(cell_w / canvas_w, cell_h / canvas_h)   # EMU per canvas px
#   draw_w/draw_h     = canvas_w * scale, canvas_h * scale          # 等比后的实际绘制尺寸
#   origin_x/origin_y = cell_left + (cell_w - draw_w) / 2, 同理 y   # 居中留白
#   任意元素(x, y, w, h) → (origin_x + x*scale, origin_y + y*scale, w*scale, h*scale)
#   字号 pt           = font_size_px * scale / 12700                # EMU→pt，等价于 px*0.75*比例
# ---------------------------------------------------------------------------
PPT_EMU_PER_PX = 9525            # 96dpi：1 CSS px = 9525 EMU
PPT_EMU_PER_PT = 12700           # 1 pt = 12700 EMU
PPT_DEFAULT_FONT = 'Microsoft YaHei'
PPT_DEFAULT_TEXT_RGB = (0x33, 0x33, 0x33)
PPT_MIN_FONT_PT = 1.0
PPT_MAX_FONT_PT = 400.0
PPT_FALLBACK_CANVAS = (1024.0, 1024.0)

# 后端合成时用的是字体文件名（msyh/msyhbd），PPT 需要真正的字体族名
PPT_FONT_ALIAS = {
    'msyh': 'Microsoft YaHei',
    'msyhbd': 'Microsoft YaHei',
    'msyhl': 'Microsoft YaHei Light',
    'simhei': 'SimHei',
    'simsun': 'SimSun',
    'simkai': 'KaiTi',
    'kaiti': 'KaiTi',
    'fangsong': 'FangSong',
    'yahei': 'Microsoft YaHei',
    'microsoft yahei': 'Microsoft YaHei',
    '微软雅黑': 'Microsoft YaHei',
    '黑体': 'SimHei',
    '宋体': 'SimSun',
    '楷体': 'KaiTi',
}


def sanitize_export_filename(title, fallback=None, max_length=60):
    """把漫画名清洗成跨平台合法的文件名主干。

    去掉 Windows/POSIX 都不允许的字符 ``\\ / : * ? " < > |`` 以及换行、
    首尾空格与点号（Windows 不允许文件名以点结尾）。

    Args:
        title: 原始标题，允许为 None / 空串。
        fallback: 标题为空时使用的兜底名称。
        max_length: 截断长度，避免超出文件系统上限。

    Returns:
        str: 永不为空的合法文件名主干。
    """
    raw = str(title or '').strip()
    if not raw:
        raw = str(fallback if fallback is not None else PPT_DEFAULT_TITLE).strip()
    cleaned = re.sub(r'[\\/:*?"<>|\r\n\t]', '_', raw)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
    cleaned = cleaned.strip(' .')[:max_length].strip(' .')
    return cleaned or 'comic'


def _ppt_pick(mapping, keys, default=None):
    """在字典中按优先级取第一个"有值"的键（兼容 camelCase / snake_case 两套字段）。

    Args:
        mapping: 任意字典，非字典直接返回 default。
        keys: 候选键名序列，靠前优先。
        default: 全部缺失时的返回值。

    Returns:
        Any: 命中的值或 default。
    """
    if not isinstance(mapping, dict):
        return default
    for key in keys:
        value = mapping.get(key)
        if value is not None and value != '':
            return value
    return default


def _ppt_num(value, default=0.0):
    """安全地把任意输入转成有限浮点数，NaN/Inf/非法输入统一回退。

    Args:
        value: 待转换的值。
        default: 转换失败时的返回值。

    Returns:
        float: 有限浮点数。
    """
    try:
        number = float(value)
    except (TypeError, ValueError):
        return float(default)
    if math.isnan(number) or math.isinf(number):
        return float(default)
    return number


def _ppt_font_name(raw):
    """把字体标识归一化成 PPT 可识别的字体族名。

    Args:
        raw: 'msyh' / 'Microsoft YaHei' / '微软雅黑' / "'Microsoft YaHei', sans-serif" 等。

    Returns:
        str: 字体族名；无法识别时原样返回（PowerPoint 会自行 fallback）。
    """
    name = str(raw or '').strip()
    if not name:
        return PPT_DEFAULT_FONT
    # CSS font stack 只取第一段，并剥掉引号
    name = name.split(',')[0].strip().strip('"').strip("'").strip()
    if not name:
        return PPT_DEFAULT_FONT
    return PPT_FONT_ALIAS.get(name.lower(), name)


def _ppt_rgb(value, default=PPT_DEFAULT_TEXT_RGB):
    """把多种颜色写法解析成 ``(r, g, b)`` 三元组。

    支持 ``'#333'`` / ``'#333333'`` / ``'333333'`` / ``'rgb(51,51,51)'`` /
    ``'rgba(51,51,51,0.8)'`` / ``[51, 51, 51]`` / ``(51, 51, 51)``。

    Args:
        value: 颜色值。
        default: 解析失败时的兜底 RGB。

    Returns:
        tuple[int, int, int]: 每个分量都在 0..255。
    """
    def _clamp(channel):
        return max(0, min(255, int(round(_ppt_num(channel, 0)))))

    if isinstance(value, (list, tuple)) and len(value) >= 3:
        return (_clamp(value[0]), _clamp(value[1]), _clamp(value[2]))

    text = str(value or '').strip().lower()
    if not text:
        return tuple(default)

    match = re.match(r'rgba?\s*\(([^)]*)\)', text)
    if match:
        parts = [p.strip() for p in match.group(1).split(',') if p.strip()]
        if len(parts) >= 3:
            return (_clamp(parts[0]), _clamp(parts[1]), _clamp(parts[2]))
        return tuple(default)

    hex_text = text.lstrip('#')
    if re.fullmatch(r'[0-9a-f]{3}', hex_text):
        return tuple(int(ch * 2, 16) for ch in hex_text)
    if re.fullmatch(r'[0-9a-f]{6}', hex_text):
        return (int(hex_text[0:2], 16), int(hex_text[2:4], 16), int(hex_text[4:6], 16))
    if re.fullmatch(r'[0-9a-f]{8}', hex_text):  # #RRGGBBAA，丢弃 alpha
        return (int(hex_text[0:2], 16), int(hex_text[2:4], 16), int(hex_text[4:6], 16))

    return tuple(default)


def _ppt_is_transparent_color(value):
    """判断网页颜色是否明确表示透明，供 PPT 气泡底色使用。"""
    text = str(value or '').strip().lower().replace(' ', '')
    if text in ('', 'transparent', 'none'):
        return True
    match = re.fullmatch(r'rgba\(([^)]*)\)', text)
    if match:
        parts = [part.strip() for part in match.group(1).split(',')]
        if len(parts) >= 4:
            try:
                return float(parts[3]) <= 0.001
            except (TypeError, ValueError):
                return False
    if re.fullmatch(r'#?[0-9a-f]{8}', text):
        return int(text.lstrip('#')[-2:], 16) == 0
    return False


def _ppt_align(value):
    """把 textAlign / text_align 映射成 ``PP_ALIGN`` 枚举。

    Args:
        value: 'left' / 'center' / 'right' / 'justify'，大小写不敏感。

    Returns:
        PP_ALIGN: 默认 LEFT。
    """
    from pptx.enum.text import PP_ALIGN

    text = str(value or '').strip().lower()
    return {
        'left': PP_ALIGN.LEFT,
        'start': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'centre': PP_ALIGN.CENTER,
        'middle': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'end': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY,
    }.get(text, PP_ALIGN.LEFT)


def _ppt_apply_east_asian_font(run, font_name):
    """给 run 同时写入 ``a:latin`` 与 ``a:ea``，保证中文也用指定字体。

    python-pptx 的 ``run.font.name`` 只写 latin 字形；中文字符走的是
    East Asian 字形，不补 ``a:ea`` 时 PowerPoint 会退回主题字体（宋体），
    导致"字体与网页预览不一致"。

    Args:
        run: ``pptx.text.text._Run`` 实例。
        font_name: 字体族名。
    """
    try:
        from pptx.oxml.ns import qn

        # 先让 python-pptx 写好 a:latin（同时保证 rPr 元素存在）
        run.font.name = font_name

        # Font._element 本身就是 <a:rPr>；python-pptx 没有 a:ea 的封装，手写即可。
        # DrawingML schema 要求 a:ea 紧跟在 a:latin 之后，否则 PowerPoint 会判定文件损坏。
        rPr = run.font._element
        existing_ea = rPr.find(qn('a:ea'))
        if existing_ea is not None:
            existing_ea.set('typeface', font_name)
            return
        ea = rPr.makeelement(qn('a:ea'), {'typeface': font_name})
        latin = rPr.find(qn('a:latin'))
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    except Exception as font_err:  # 字体细节失败不应影响整份 PPT
        logger.warning(f'[export-ppt] 设置东亚字体失败({font_name}): {font_err}')


def resolve_panel_canvas_size(panel):
    """推断某一格在网页画布上的像素尺寸（bubbles/characters 坐标的基准）。

    优先取前端显式传来的 ``canvas_width`` / ``canvas_height``；缺失时读成图
    或背景图的真实像素尺寸兜底；再不行用 1024x1024。

    Args:
        panel: 单格数据字典。

    Returns:
        tuple[float, float]: ``(canvas_width, canvas_height)``，两者恒 > 0。
    """
    width = _ppt_num(_ppt_pick(panel, ('canvas_width', 'canvasWidth')), 0)
    height = _ppt_num(_ppt_pick(panel, ('canvas_height', 'canvasHeight')), 0)
    if width > 0 and height > 0:
        return width, height

    for key in ('url', 'background_url', 'image'):
        image_path = resolve_local_asset_path(_ppt_pick(panel, (key,)))
        if not image_path:
            continue
        try:
            with Image.open(image_path) as img_obj:
                img_w, img_h = img_obj.size
            if img_w > 0 and img_h > 0:
                return float(img_w), float(img_h)
        except Exception as size_err:
            logger.warning(f'[export-ppt] 读取画布尺寸失败({image_path}): {size_err}')

    return PPT_FALLBACK_CANVAS


def resolve_character_asset(character):
    """解析角色图的本地路径，兼容 ``url`` 与后端记录的绝对 ``path``。

    Args:
        character: 单个角色元数据字典。

    Returns:
        str | None: 存在的本地绝对路径，找不到返回 None。
    """
    if not isinstance(character, dict):
        return None
    asset_path = resolve_local_asset_path(
        _ppt_pick(character, ('url', 'image', 'image_url', 'src'))
    )
    if asset_path:
        return asset_path
    raw_path = _ppt_pick(character, ('path', 'local_path'))
    if raw_path and os.path.exists(str(raw_path)):
        return os.path.normpath(str(raw_path))
    return None


def normalize_ppt_characters(panel):
    """把一格里的角色列表归一化成 PPT 可直接使用的结构。

    Args:
        panel: 单格数据字典。

    Returns:
        list[dict]: 每项含 ``path/x/y/width/height/speaker``，图片缺失的项被丢弃。
    """
    raw_items = _ppt_pick(panel, ('characters', 'chars'), []) or []
    if isinstance(raw_items, str):
        try:
            raw_items = json.loads(raw_items)
        except Exception:
            raw_items = []
    if not isinstance(raw_items, (list, tuple)):
        return []

    normalized = []
    for item in raw_items:
        if not isinstance(item, dict):
            continue
        asset_path = resolve_character_asset(item)
        if not asset_path:
            continue
        width = _ppt_num(_ppt_pick(item, ('width', 'w')), 0)
        height = _ppt_num(_ppt_pick(item, ('height', 'h')), 0)
        if width <= 0 or height <= 0:
            # 没有尺寸就按图片原始像素铺，至少不丢角色
            try:
                with Image.open(asset_path) as img_obj:
                    width, height = (float(img_obj.size[0]), float(img_obj.size[1]))
            except Exception:
                continue
        if width <= 0 or height <= 0:
            continue
        normalized.append({
            'path': asset_path,
            'speaker': str(_ppt_pick(item, ('speaker', 'name'), '') or ''),
            'x': _ppt_num(_ppt_pick(item, ('x', 'left')), 0),
            'y': _ppt_num(_ppt_pick(item, ('y', 'top')), 0),
            'width': width,
            'height': height,
        })
    return normalized


def normalize_ppt_bubble(raw, canvas_size):
    """把气泡数据统一成 PPT 排版结构（同时兼容前端编辑器与后端合成两套字段）。

    字段来源差异::

        前端编辑器保存: bubbleImage / fontSize / fontFamily / fill / textAlign / text_align
        后端合成 meta : bubble_url / font_size / font_family / text_color / text_align
                        + rendered_layer_url（气泡与文字**已烘焙**的成品层）

    因此气泡底图优先取"无字气泡"（``bubbleImage`` / ``bubble_url``），这样才能
    在上面叠加原生可编辑文字框；只有在无字气泡缺失时才退回 ``rendered_layer_url``，
    此时把 ``baked`` 置 True，调用方需跳过文字框以免文字重影。

    Args:
        raw: 单个气泡数据字典。
        canvas_size: ``(canvas_width, canvas_height)``，用于兜底尺寸。

    Returns:
        dict | None: 归一化结果；既无图又无文字时返回 None。
    """
    if not isinstance(raw, dict):
        return None

    canvas_w, canvas_h = canvas_size
    # 合成预览时 Pillow 已经按中文标点、气泡可用宽度计算出最终换行。
    # PPT 必须优先复用它，不能把原始台词交给 PowerPoint 再断一次行。
    rendered_text = _ppt_pick(raw, ('rendered_text', 'renderedText'), '')
    text = rendered_text or _ppt_pick(raw, ('text', 'content', 'dialogue', 'caption'), '')
    text = str(text or '').strip()

    # 编辑器中气泡外形（红框）与文字（蓝框）可以分别移动；导出时必须
    # 优先读各自的坐标，不能再把文字硬塞回气泡位置。
    box_x = _ppt_num(_ppt_pick(raw, ('bubble_x', 'bubbleX', 'x', 'left')), 0)
    box_y = _ppt_num(_ppt_pick(raw, ('bubble_y', 'bubbleY', 'y', 'top')), 0)
    box_w = _ppt_num(_ppt_pick(raw, ('bubble_width', 'bubbleWidth', 'width', 'w')), 0)
    box_h = _ppt_num(_ppt_pick(raw, ('bubble_height', 'bubbleHeight', 'height', 'h')), 0)
    if box_w <= 0:
        box_w = canvas_w * 0.35
    if box_h <= 0:
        box_h = canvas_h * 0.14

    text_x = _ppt_num(_ppt_pick(raw, ('text_x', 'textX')), box_x)
    text_y = _ppt_num(_ppt_pick(raw, ('text_y', 'textY')), box_y)
    text_w = _ppt_num(_ppt_pick(raw, ('text_width', 'textWidth')), box_w)
    if text_w <= 0:
        text_w = box_w

    plain_path = resolve_local_asset_path(
        _ppt_pick(raw, ('bubbleImage', 'bubble_image', 'bubble_url', 'bubbleUrl'))
    )
    baked_path = resolve_local_asset_path(
        _ppt_pick(raw, ('rendered_layer_url', 'renderedLayerUrl'))
    )

    if plain_path:
        image_path, baked = plain_path, False
        img_x, img_y, img_w, img_h = box_x, box_y, box_w, box_h
    elif baked_path:
        image_path, baked = baked_path, True
        img_x = _ppt_num(_ppt_pick(raw, ('rendered_layer_x', 'renderedLayerX')), box_x)
        img_y = _ppt_num(_ppt_pick(raw, ('rendered_layer_y', 'renderedLayerY')), box_y)
        img_w = _ppt_num(_ppt_pick(raw, ('rendered_layer_width', 'renderedLayerWidth')), box_w)
        img_h = _ppt_num(_ppt_pick(raw, ('rendered_layer_height', 'renderedLayerHeight')), box_h)
        if img_w <= 0:
            img_w = box_w
        if img_h <= 0:
            img_h = box_h
    else:
        image_path, baked = None, False
        img_x, img_y, img_w, img_h = box_x, box_y, box_w, box_h

    if not image_path and not text:
        return None

    font_size = _ppt_num(_ppt_pick(raw, ('fontSize', 'font_size', 'logicalFontSize')), 0)
    if font_size <= 0:
        font_size = max(canvas_h * 0.028, 12.0)

    font_weight = str(_ppt_pick(raw, ('fontWeight', 'font_weight'), 'normal') or 'normal').lower()
    font_style = str(_ppt_pick(raw, ('fontStyle', 'font_style'), 'normal') or 'normal').lower()
    color_value = _ppt_pick(raw, ('fill', 'text_color', 'textColor', 'color'))

    return {
        'text': text,
        'has_exact_line_breaks': bool(rendered_text),
        'image_path': image_path,
        'baked': baked,
        'image_x': img_x,
        'image_y': img_y,
        'image_width': img_w,
        'image_height': img_h,
        'x': text_x,
        'y': text_y,
        'width': text_w,
        'height': box_h,
        'padding': max(_ppt_num(_ppt_pick(raw, ('padding', 'logicalPadding')), 8.0), 0.0),
        'font_size': font_size,
        'line_height': max(_ppt_num(_ppt_pick(raw, ('line_height', 'lineHeight', 'lineHeightPx')), font_size * 1.16), 1.0),
        'font_family': _ppt_font_name(_ppt_pick(raw, ('fontFamily', 'font_family'))),
        'rgb': _ppt_rgb(color_value),
        'bold': font_weight in ('bold', 'bolder', '600', '700', '800', '900') or font_weight == 'true',
        'italic': font_style in ('italic', 'oblique'),
        'align': _ppt_pick(raw, ('textAlign', 'text_align', 'align'), 'center'),
        'background_rgb': _ppt_rgb(_ppt_pick(raw, ('backgroundColor', 'background_color')), (255, 255, 255)),
        'background_transparent': _ppt_is_transparent_color(_ppt_pick(raw, ('backgroundColor', 'background_color'))),
        'border_rgb': _ppt_rgb(_ppt_pick(raw, ('borderColor', 'border_color')), (34, 34, 34)),
        'border_transparent': _ppt_is_transparent_color(_ppt_pick(raw, ('borderColor', 'border_color'))),
        'border_width': max(_ppt_num(_ppt_pick(raw, ('borderWidth', 'border_width')), 2.0), 0.0),
    }


def normalize_ppt_bubbles(panel, canvas_size):
    """归一化一格中的全部气泡。

    Args:
        panel: 单格数据字典。
        canvas_size: ``(canvas_width, canvas_height)``。

    Returns:
        list[dict]: ``normalize_ppt_bubble`` 结果列表（已过滤 None）。
    """
    raw_items = _ppt_pick(panel, ('bubbles',), []) or []
    if isinstance(raw_items, str):
        try:
            raw_items = json.loads(raw_items)
        except Exception:
            raw_items = []
    if not isinstance(raw_items, (list, tuple)):
        return []

    normalized = []
    for item in raw_items:
        bubble = normalize_ppt_bubble(item, canvas_size)
        if bubble:
            normalized.append(bubble)
    return normalized


def add_ppt_text_shape(slide, left, top, width, height, bubble, scale, shape_name):
    """在幻灯片上放一个还原网页样式的原生可编辑文本框。

    文本框本身不设填充与边框，直接"浮"在气泡图之上，视觉与网页一致，
    同时保留 PowerPoint 的全部文字编辑能力。

    Args:
        slide: 目标幻灯片。
        left/top/width/height: 文本框矩形（EMU，int）。
        bubble: ``normalize_ppt_bubble`` 的结果。
        scale: EMU per canvas px 的缩放比，用于换算字号与内边距。
        shape_name: 形状名，便于用户在选择窗格里辨认。

    Returns:
        pptx.shapes.autoshape.Shape: 新建的文本框。
    """
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR

    textbox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    try:
        textbox.name = shape_name
    except Exception:
        pass

    text_frame = textbox.text_frame
    # 已保存了预览最终换行时，关闭 PowerPoint 二次换行，避免中文粗体和
    # 书名号/标点被 PPT 拆成与漫画预览不同的行；旧数据仍保留自动换行兼容。
    text_frame.word_wrap = not bubble.get('has_exact_line_breaks', False)

    # 内边距按网页 padding 等比换算，并留出安全上限，避免小气泡把文字挤没
    inset = int(max(min(bubble['padding'] * scale, width * 0.25, height * 0.25), 0))
    text_frame.margin_left = Emu(inset)
    text_frame.margin_right = Emu(inset)
    text_frame.margin_top = Emu(0)
    text_frame.margin_bottom = Emu(0)
    try:
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = _ppt_align(bubble['align'])
    # Pillow 预览是固定像素行高；PPT 默认行距会随字体替换变化，导致文字块
    # 在同一气泡内变高或变矮，因此同样换算为固定 pt 行距。
    try:
        line_pt = bubble['line_height'] * scale / float(PPT_EMU_PER_PT)
        paragraph.line_spacing = Pt(max(0.5, round(line_pt, 2)))
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
    except Exception:
        pass

    run = paragraph.add_run()
    run.text = bubble['text']

    font_pt = bubble['font_size'] * scale / float(PPT_EMU_PER_PT)
    font_pt = max(PPT_MIN_FONT_PT, min(PPT_MAX_FONT_PT, font_pt))
    run.font.size = Pt(round(font_pt, 1))
    run.font.bold = bool(bubble['bold'])
    run.font.italic = bool(bubble['italic'])
    red, green, blue = bubble['rgb']
    run.font.color.rgb = RGBColor(red, green, blue)
    _ppt_apply_east_asian_font(run, bubble['font_family'])

    return textbox


def add_layered_panel_to_slide(slide, panel, panel_index, cell_left, cell_top,
                               cell_w, cell_h):
    """把一格漫画以"分层可编辑"的方式铺进幻灯片的指定矩形区域。

    产出的形状（每个都能在 PowerPoint 里单独选中 / 拖动 / 删除 / 换图 / 改字）::

        Bg_<i>            背景图（等比居中）
        Char_<i>_<j>      第 j 个角色图
        BubbleImg_<i>_<j> 第 j 个气泡底图（无字）
        Text_<i>_<j>      第 j 段台词（原生文本框，字体/字号/颜色还原网页）

    若该格没有 ``background_url``（说明前端只给了成图、没有分层数据），
    自动退化为"整格一张图"模式，避免角色/气泡与成图内容重复叠加。

    Args:
        slide: 目标幻灯片。
        panel: 单格数据字典。
        panel_index: 该格序号（从 0 开始），用于形状命名。
        cell_left/cell_top/cell_w/cell_h: 该格在幻灯片上的矩形（EMU）。

    Returns:
        dict: ``{'layered': bool, 'background': int, 'characters': int,
        'bubble_images': int, 'textboxes': int, 'missing': int}`` 统计信息。
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    stats = {
        'layered': False, 'background': 0, 'characters': 0,
        'bubble_images': 0, 'textboxes': 0, 'missing': 0,
    }

    cell_left = int(cell_left)
    cell_top = int(cell_top)
    cell_w = max(int(cell_w), 1)
    cell_h = max(int(cell_h), 1)
    tag = panel_index + 1

    canvas_w, canvas_h = resolve_panel_canvas_size(panel)
    scale = min(cell_w / canvas_w, cell_h / canvas_h)
    draw_w = max(int(canvas_w * scale), 1)
    draw_h = max(int(canvas_h * scale), 1)
    origin_x = int(cell_left + (cell_w - draw_w) / 2)
    origin_y = int(cell_top + (cell_h - draw_h) / 2)

    def _to_slide_rect(x, y, width, height):
        """canvas 像素矩形 → 幻灯片 EMU 矩形。"""
        return (
            int(origin_x + x * scale),
            int(origin_y + y * scale),
            max(int(width * scale), 1),
            max(int(height * scale), 1),
        )

    # 是否走分层模式由"前端有没有给 background_url"决定，而不是由图片能否解析决定：
    # 背景图临时失效时仍要把角色/气泡/文字铺上去，否则这一格的台词会整段丢失。
    raw_background = _ppt_pick(panel, ('background_url', 'backgroundUrl'))
    layered = bool(raw_background)
    if layered:
        background_path = resolve_local_asset_path(raw_background)
    else:
        # 没有分层数据：退回"整格一张成图"，此时不再叠加角色/气泡，避免内容重影
        background_path = resolve_local_asset_path(_ppt_pick(panel, ('url', 'image')))
    stats['layered'] = layered

    if background_path:
        try:
            picture = slide.shapes.add_picture(
                background_path, origin_x, origin_y, width=draw_w, height=draw_h
            )
            picture.name = f'Bg_{tag}'
            stats['background'] = 1
        except Exception as bg_err:
            logger.warning(f'[export-ppt] 背景图插入失败({background_path}): {bg_err}')
            stats['missing'] += 1
    else:
        holder = slide.shapes.add_textbox(cell_left, cell_top, cell_w, cell_h)
        holder.name = f'BgMissing_{tag}'
        holder_para = holder.text_frame.paragraphs[0]
        holder_para.text = f'第{tag}格背景缺失'
        holder_para.alignment = PP_ALIGN.CENTER
        holder_run = holder_para.runs[0]
        holder_run.font.size = Pt(12)
        holder_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        _ppt_apply_east_asian_font(holder_run, PPT_DEFAULT_FONT)
        stats['missing'] += 1

    if not layered:
        # 成图已经包含角色与气泡，直接返回，避免重影
        return stats

    for char_i, character in enumerate(normalize_ppt_characters(panel)):
        left, top, width, height = _to_slide_rect(
            character['x'], character['y'], character['width'], character['height']
        )
        try:
            picture = slide.shapes.add_picture(
                character['path'], left, top, width=width, height=height
            )
            picture.name = f'Char_{tag}_{char_i + 1}'
            stats['characters'] += 1
        except Exception as char_err:
            logger.warning(f"[export-ppt] 角色图插入失败({character['path']}): {char_err}")
            stats['missing'] += 1

    bubbles = normalize_ppt_bubbles(panel, (canvas_w, canvas_h))

    # 先铺所有气泡底图，再铺所有文字框：保证文字层始终压在气泡层之上
    for bubble_i, bubble in enumerate(bubbles):
        left, top, width, height = _to_slide_rect(
            bubble['image_x'], bubble['image_y'],
            bubble['image_width'], bubble['image_height']
        )
        if bubble['image_path']:
            try:
                picture = slide.shapes.add_picture(
                    bubble['image_path'], left, top, width=width, height=height
                )
                picture.name = (
                    f'BubbleRendered_{tag}_{bubble_i + 1}' if bubble['baked']
                    else f'BubbleImg_{tag}_{bubble_i + 1}'
                )
                stats['bubble_images'] += 1
            except Exception as bubble_err:
                logger.warning(f"[export-ppt] 气泡图插入失败({bubble['image_path']}): {bubble_err}")
                stats['missing'] += 1
            continue

        # 普通气泡在网页中本来就是 Fabric 的矩形而非图片文件。旧代码
        # 因为找不到 bubbleImage 而直接跳过，造成下载 PPT 只有文字没有
        # 气泡。这里改为原生可编辑圆角矩形，位置、颜色、边框均取保存数据。
        try:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            shape.name = f'BubbleShape_{tag}_{bubble_i + 1}'
            fill_r, fill_g, fill_b = bubble['background_rgb']
            line_r, line_g, line_b = bubble['border_rgb']
            if bubble['background_transparent']:
                # rgba(..., 0) 在网页是透明；若解析成 RGB(0,0,0) 后再 solid，
                # PowerPoint 就会显示成黑色气泡。
                shape.fill.background()
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(fill_r, fill_g, fill_b)
            if bubble['border_transparent'] or bubble['border_width'] <= 0:
                shape.line.fill.background()
            else:
                shape.line.color.rgb = RGBColor(line_r, line_g, line_b)
                shape.line.width = max(int(bubble['border_width'] * scale), 1)
            stats['bubble_images'] += 1
        except Exception as bubble_shape_err:
            logger.warning(f'[export-ppt] 原生气泡形状插入失败: {bubble_shape_err}')
            stats['missing'] += 1

    for bubble_i, bubble in enumerate(bubbles):
        # baked=True 说明底图里已经烘焙了文字，再加文本框会出现双重文字
        if not bubble['text'] or bubble['baked']:
            continue
        left, top, width, height = _to_slide_rect(
            bubble['x'], bubble['y'], bubble['width'], bubble['height']
        )
        add_ppt_text_shape(
            slide, left, top, width, height, bubble, scale,
            f'Text_{tag}_{bubble_i + 1}'
        )
        stats['textboxes'] += 1

    return stats


def add_ppt_cover_slide(prs, main_title, subtitle_text=''):
    """生成封面页（标题 = 漫画名）。

    Args:
        prs: ``pptx.Presentation`` 实例。
        main_title: 主标题文本。
        subtitle_text: 副标题文本，空串则留空。

    Returns:
        pptx.slide.Slide: 封面幻灯片。
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = str(main_title)
    try:
        cover_title_para = cover.shapes.title.text_frame.paragraphs[0]
        cover_title_para.font.size = Pt(40)
        cover_title_para.font.bold = True
        cover_title_para.font.color.rgb = RGBColor(0xFF, 0x6B, 0x9D)
        for run in cover_title_para.runs:
            _ppt_apply_east_asian_font(run, PPT_DEFAULT_FONT)
    except Exception as style_err:
        logger.warning(f'[export-ppt] 封面标题样式设置失败: {style_err}')

    try:
        cover.placeholders[1].text = str(subtitle_text or '')
        cover_sub_para = cover.placeholders[1].text_frame.paragraphs[0]
        cover_sub_para.font.size = Pt(18)
        for run in cover_sub_para.runs:
            _ppt_apply_east_asian_font(run, PPT_DEFAULT_FONT)
    except Exception as style_err:
        logger.warning(f'[export-ppt] 封面副标题设置失败: {style_err}')

    return cover


def create_comic_presentation():
    """新建一个 16:9 空白演示文稿，尺寸取自全局 PPT 常量。

    Returns:
        pptx.Presentation: 已设置好画布尺寸的演示文稿。
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(PPT_SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(PPT_SLIDE_HEIGHT_INCHES)
    return prs


def build_merged_vector_pptx(panels, main_title, subtitle_text=''):
    """构建「合并版」：封面 1 页 + 全部格子挤在同一页的网格页。

    Args:
        panels: 归一化后的格子列表。
        main_title: 漫画名（封面标题）。
        subtitle_text: 副标题。

    Returns:
        tuple: ``(prs, summary_dict)``；summary 含 grid / 各类形状计数。
    """
    from pptx.util import Inches

    prs = create_comic_presentation()
    add_ppt_cover_slide(prs, main_title, subtitle_text)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    margin = Inches(0.3)
    gap = Inches(0.1)
    usable_w = prs.slide_width - margin * 2
    usable_h = prs.slide_height - margin * 2

    rows, cols = compute_ppt_grid(len(panels))
    cell_w = int((usable_w - gap * (cols - 1)) / cols)
    cell_h = int((usable_h - gap * (rows - 1)) / rows)

    totals = {'background': 0, 'characters': 0, 'bubble_images': 0,
              'textboxes': 0, 'missing': 0, 'layered_panels': 0}
    for idx, panel in enumerate(panels):
        row_i = idx // cols
        col_i = idx % cols
        cell_left = int(margin + col_i * (cell_w + gap))
        cell_top = int(margin + row_i * (cell_h + gap))
        stats = add_layered_panel_to_slide(
            slide, panel, idx, cell_left, cell_top, cell_w, cell_h
        )
        for key in ('background', 'characters', 'bubble_images', 'textboxes', 'missing'):
            totals[key] += stats[key]
        totals['layered_panels'] += 1 if stats['layered'] else 0

    totals['grid'] = f'{rows}x{cols}'
    totals['slide_count'] = len(prs.slides._sldIdLst)
    return prs, totals


def build_split_vector_pptx(panels, main_title, subtitle_text=''):
    """构建「分开版」：封面 1 页 + 每格独占一页（整页即画布）。

    Args:
        panels: 归一化后的格子列表。
        main_title: 漫画名（封面标题）。
        subtitle_text: 副标题。

    Returns:
        tuple: ``(prs, summary_dict)``。
    """
    prs = create_comic_presentation()
    add_ppt_cover_slide(prs, main_title, subtitle_text)

    totals = {'background': 0, 'characters': 0, 'bubble_images': 0,
              'textboxes': 0, 'missing': 0, 'layered_panels': 0}
    for idx, panel in enumerate(panels):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        stats = add_layered_panel_to_slide(
            slide, panel, idx, 0, 0, prs.slide_width, prs.slide_height
        )
        for key in ('background', 'characters', 'bubble_images', 'textboxes', 'missing'):
            totals[key] += stats[key]
        totals['layered_panels'] += 1 if stats['layered'] else 0

    totals['grid'] = f'{len(panels)}x1'
    totals['slide_count'] = len(prs.slides._sldIdLst)
    return prs, totals


def save_presentation_with_name(prs, base_name, suffix, export_id):
    """把演示文稿保存到 OUTPUT_DIR，文件名以漫画名命名。

    同名文件被占用（例如用户正开着上一次导出的 PPT）时，自动追加 export_id
    重试一次，避免整个导出因 PermissionError 失败。

    Args:
        prs: ``pptx.Presentation`` 实例。
        base_name: 已清洗的漫画名。
        suffix: 版本后缀，如 ``'合并版'``。
        export_id: 本次导出的短 ID，用于重名兜底。

    Returns:
        tuple[str, str, int]: ``(绝对路径, 文件名, 字节大小)``。
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    filename = f'{base_name}_{suffix}.pptx'
    output_path = os.path.normpath(os.path.join(OUTPUT_DIR, filename))
    try:
        prs.save(output_path)
    except Exception as save_err:
        logger.warning(f'[export-ppt] 以漫画名保存失败({filename}): {save_err}，改用带 ID 的文件名')
        filename = f'{base_name}_{suffix}_{export_id}.pptx'
        output_path = os.path.normpath(os.path.join(OUTPUT_DIR, filename))
        prs.save(output_path)
    return output_path, filename, os.path.getsize(output_path)


@app.route('/api/export-ppt', methods=['POST'])
def api_export_ppt():
    """导出**真可编辑**的矢量 PPT：背景图 + 角色图 + 气泡图 + 原生文字框。

    一次生成两份文件，文件名均以漫画名命名：

      * ``<漫画名>_合并版.pptx``：封面 1 页 + 全部格子铺在同一页的网格页；
      * ``<漫画名>_分开版.pptx``：封面 1 页 + 每格独占一页（整页即画布）。

    两份都不是"整格截图"：每个角色、每个气泡、每段台词都是独立形状，
    在 PowerPoint 里可单独选中、拖动、替换、改字，字体/字号/颜色与网页预览一致。

    请求体（与旧版兼容）::

        {
          "panels": [{
            "url": "/static/output/a.png",
            "background_url": "/static/output/bg.png",
            "characters": [{"url": "...", "x": 100, "y": 200, "width": 300, "height": 500}],
            "bubbles": [{"text": "你好", "bubbleImage": "...", "x": 40, "y": 30,
                         "width": 300, "height": 120, "fontSize": 28,
                         "fontFamily": "Microsoft YaHei", "fill": "#333333"}],
            "canvas_width": 1024, "canvas_height": 1024
          }],
          "images": ["/static/output/a.png"],
          "script": {"title": "小猫历险记", "topic": "勇气"},
          "title": "小猫历险记"
        }

    Returns:
        flask.Response: 成功返回
        ``{"success": true, "merged_pptx_url": "...", "split_pptx_url": "...",
        "merged_filename": "...", "split_filename": "...", "panel_count": 6}``；
        失败返回 ``{"success": false, "error": "..."}``。
    """
    export_id = uuid.uuid4().hex[:8]
    app.logger.info(f'[export-ppt] 开始生成矢量可编辑 PPT export_id={export_id}')

    try:
        import pptx  # noqa: F401  # 仅做依赖存在性校验，实际导入在各构建函数内
    except ImportError:
        app.logger.error('[export-ppt] 缺少 python-pptx 依赖')
        return jsonify({
            'success': False,
            'error': '服务器缺少 python-pptx 库，请执行 pip install python-pptx 后重试'
        }), 500

    try:
        data = request.get_json(silent=True) or {}
        panels = normalize_export_panels(data)
        if not panels:
            app.logger.warning(f'[export-ppt] 无可导出的漫画数据 export_id={export_id}')
            return jsonify({'success': False, 'error': '没有可导出的漫画数据'}), 400

        script = data.get('script') or {}
        if isinstance(script, str):
            try:
                script = json.loads(script)
            except Exception:
                script = {}
        if not isinstance(script, dict):
            script = {}

        main_title = str(
            data.get('title') or script.get('title') or PPT_DEFAULT_TITLE
        ).strip() or PPT_DEFAULT_TITLE
        subtitle_text = str(script.get('topic') or script.get('summary') or '').strip()[:60]
        base_name = sanitize_export_filename(main_title)

        # ---- 合并版：封面 + 一页网格（每格分层可编辑）----
        merged_prs, merged_stats = build_merged_vector_pptx(panels, main_title, subtitle_text)
        merged_path, merged_filename, merged_size = save_presentation_with_name(
            merged_prs, base_name, '合并版', export_id
        )
        app.logger.info(
            f'[export-ppt] 合并版完成 export_id={export_id} file={merged_filename} '
            f'grid={merged_stats["grid"]} chars={merged_stats["characters"]} '
            f'bubbles={merged_stats["bubble_images"]} texts={merged_stats["textboxes"]} '
            f'missing={merged_stats["missing"]}'
        )

        # ---- 分开版：封面 + 每格一页 ----
        split_prs, split_stats = build_split_vector_pptx(panels, main_title, subtitle_text)
        split_path, split_filename, split_size = save_presentation_with_name(
            split_prs, base_name, '分开版', export_id
        )
        app.logger.info(
            f'[export-ppt] 分开版完成 export_id={export_id} file={split_filename} '
            f'slides={split_stats["slide_count"]} chars={split_stats["characters"]} '
            f'bubbles={split_stats["bubble_images"]} texts={split_stats["textboxes"]} '
            f'missing={split_stats["missing"]}'
        )

        from urllib.parse import quote

        return jsonify({
            'success': True,
            'message': 'PPT 已生成（合并版 + 分开版），文字可直接编辑',
            'title': main_title,
            'merged_pptx_url': f'/static/output/{quote(merged_filename)}',
            'split_pptx_url': f'/static/output/{quote(split_filename)}',
            'merged_filename': merged_filename,
            'split_filename': split_filename,
            'merged_file_size': merged_size,
            'split_file_size': split_size,
            'merged_pptx_path': merged_path.replace('\\', '/'),
            'split_pptx_path': split_path.replace('\\', '/'),
            # 兼容旧前端：pptx_url / filename 指向合并版
            'pptx_url': f'/static/output/{quote(merged_filename)}',
            'filename': merged_filename,
            'pptx_path': merged_path.replace('\\', '/'),
            'file_size': merged_size,
            'panel_count': len(panels),
            'layered_panel_count': merged_stats['layered_panels'],
            'character_count': merged_stats['characters'],
            'bubble_count': merged_stats['bubble_images'],
            'textbox_count': merged_stats['textboxes'],
            'missing_asset_count': merged_stats['missing'],
            'grid': merged_stats['grid'],
            'merged_slide_count': merged_stats['slide_count'],
            'split_slide_count': split_stats['slide_count'],
            'slide_count': merged_stats['slide_count'],
        })

    except Exception as e:
        app.logger.error(f'[export-ppt] 生成失败 export_id={export_id}: {e}')
        logger.error(f'PPT 导出失败: {e}')
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': f'PPT 生成失败: {e}'}), 500


@app.route('/api/generate_character_image', methods=['POST'])
def api_generate_character_image():
    """生成角色图片"""
    import traceback
    try:
        import json
        data = request.get_json()
        
        name = data.get('name', '')
        role = data.get('role', '')
        appearance = data.get('appearance', '')
        dialogues = data.get('dialogues', [])
        
        logger.info(f"收到角色图片生成请求: name={name}, role={role}, appearance={appearance}")
        
        if not name:
            logger.warning("角色名称为空")
            return jsonify({'success': False, 'error': '角色名称不能为空', 'error_code': 'EMPTY_NAME'}), 400

        # === 预设 IP 角色（许多/莉莉/高远/哈哈/...）直接返回静态姿势库，不走 Pollinations 生图 ===
        # 用户在页面手动指定的形象必须优先于脚本中的角色原名。
        # 否则“许多/莉莉/知识小精灵”等已有名字会先命中默认库，覆盖用户刚选的其他小精灵。
        preset_poses = get_preset_poses(data.get('preset')) or get_preset_poses(name)
        if preset_poses:
            logger.info(f"角色[{name}] 命中预设 IP 姿势库({len(preset_poses)}个)，直接返回静态姿势图")
            return jsonify({
                'success': True,
                'image_url': preset_poses.get('stand') or next(iter(preset_poses.values())),
                'poses': preset_poses,
                'preset': True
            })

        # === 非预设角色：使用知识小精灵姿势库（catch-all）===
        # 所有非预设角色（如李白、小明）自动使用精灵姿势，不再走 Pollinations
        sprite_poses = get_preset_poses(SPRITE_FOLDER)
        if sprite_poses:
            logger.info(f"角色[{name}] 非预设角色，使用知识小精灵姿势库({len(sprite_poses)}个)")
            return jsonify({
                'success': True,
                'image_url': sprite_poses.get('stand') or next(iter(sprite_poses.values())),
                'poses': sprite_poses,
                'preset': True,
                'sprite': True  # 标记为精灵模式
            })

        char_index = sum(ord(c) for c in name) % 9
        char_seed = data.get('seed') or ((sum(ord(c) for c in name) % 900000) + 100000)
        logger.info(f"角色[{name}] 使用 seed: {char_seed}")

        base_parts = [
            f"角色：{name}",
            "isolated cartoon character on solid white background",
            "white background only",
            "no shadows, no ground, no gradient background",
            "full body character",
            "isolated character"
        ]
        if role:
            base_parts.append(f"角色定位：{role}")
        if appearance:
            base_parts.append(f"外观描述：{appearance}")

        needed_poses = {'stand'}  # 标准姿势始终生成，作为兜底

        if dialogues and isinstance(dialogues, list):
            for dialogue in dialogues:
                if isinstance(dialogue, str) and dialogue:
                    _, action = extract_action(dialogue)
                    if action:
                        pose_key = map_action_to_pose(action)
                        if pose_key in POSE_LIBRARY and pose_key != 'stand':
                            needed_poses.add(pose_key)
                    else:
                        # 台词无动作标记：按语义自动挑姿势并预生成，避免漫画里全员 stand
                        pk = auto_pose_from_text(dialogue)
                        if pk in POSE_LIBRARY and pk != 'stand':
                            needed_poses.add(pk)
        
        pose_items = [(pk, POSE_LIBRARY[pk]) for pk in needed_poses if pk in POSE_LIBRARY]
        logger.info(f"角色[{name}] 分析脚本需要的姿势: {sorted(needed_poses)}")

        provider = data.get('provider', 'pollinations')
        use_face = (provider == 'pollinations_face')
        engine_name = 'Pollinations(保脸换姿势)' if use_face else 'Pollinations'
        logger.info(f"角色[{name}]生图引擎: {engine_name}")

        poses = {}
        stand_url = None
        from io import BytesIO
        import time

        CHAR_SIZE = 768
        MAX_TIME_PER_POSE = 45
        total_start_time = time.time()

        standard_char_bytes = None
        if use_face:
            standard_prompt = "，".join(base_parts)
            standard_char_bytes = generate_image(standard_prompt, char_index, is_character=True,
                                                 seed=char_seed, width=CHAR_SIZE, height=CHAR_SIZE)
            if standard_char_bytes is None:
                logger.warning("保脸模式标准角色生成失败，整体降级为逐姿势文生图")
                use_face = False
            else:
                logger.info(f"保脸模式标准角色生成成功（{len(standard_char_bytes)} bytes）")

        def _gen_pose_with_timeout(pose_key, pose_cn):
            import signal
            
            class TimeoutException(Exception):
                pass
            
            def timeout_handler(signum, frame):
                raise TimeoutException("姿势生成超时")
            
            try:
                if use_face and standard_char_bytes:
                    en_pose = POSE_LIBRARY_EN.get(pose_key, 'standing naturally')
                    img2img_prompt = (f"same exact cartoon character, {en_pose}, "
                                      f"keep identical face, eyes, hair, outfit, hair color and clothes color, "
                                      f"only change body pose, do not change face, "
                                      f"isolated on solid white background, white background only, no shadows, no ground")
                    image_data = generate_image_pollinations_img2img(standard_char_bytes, img2img_prompt,
                                                                    seed=char_seed, width=CHAR_SIZE, height=CHAR_SIZE)
                    if image_data is None:
                        logger.warning(f"姿势[{pose_key}] 保脸图生图失败，降级为普通文生图")
                        full_prompt = "，".join(base_parts) + f"，当前姿势：{pose_cn}"
                        image_data = generate_image(full_prompt, char_index, is_character=True,
                                                    seed=char_seed, width=CHAR_SIZE, height=CHAR_SIZE)
                else:
                    full_prompt = "，".join(base_parts) + f"，当前姿势：{pose_cn}"
                    image_data = generate_image(full_prompt, char_index, is_character=True,
                                                seed=char_seed, width=CHAR_SIZE, height=CHAR_SIZE)
                return pose_key, image_data
            except TimeoutException:
                logger.error(f"姿势[{pose_key}] 生成超时")
                return pose_key, None
            except Exception as e:
                logger.error(f"姿势[{pose_key}] 生成异常: {e}")
                return pose_key, None

        raw_results = {}
        POSE_SPACING = 1.2
        logger.info(f"角色[{name}] 串行生成 {len(pose_items)} 个姿势（间隔 {POSE_SPACING}s）")
        
        for pose_key, pose_cn in pose_items:
            elapsed_total = time.time() - total_start_time
            if elapsed_total > 360:
                logger.error(f"角色[{name}] 生成总时间已超过360秒，停止生成")
                break
                
            try:
                time.sleep(POSE_SPACING)
                pk, img_data = _gen_pose_with_timeout(pose_key, pose_cn)
                raw_results[pk] = img_data
            except Exception as e:
                logger.error(f"姿势[{pose_key}] 生成异常: {e}")
                raw_results[pose_key] = None

        for pose_key, pose_cn in pose_items:
            image_data = raw_results.get(pose_key)
            if image_data is None:
                logger.warning(f"姿势[{pose_key}]生成失败，跳过")
                continue
            removed = remove_background(image_data)
            if removed:
                image_data = removed
            output_filename = f"char_{pose_key}_{str(uuid.uuid4())[:8]}.png"
            output_path = os.path.join(OUTPUT_DIR, output_filename)
            img = Image.open(BytesIO(image_data))
            if img.mode != 'RGBA':
                img = img.convert('RGBA')
            img.save(output_path, format='PNG')
            url = f'/static/output/{output_filename}'
            poses[pose_key] = url
            if pose_key == 'stand' or stand_url is None:
                stand_url = url
            logger.info(f"角色姿势图生成成功: {name}/{pose_key} -> {url}")

        # stand 是角色“标准姿势/头像”，免费端点偶发限流会跳过它，缺失时单独重试一次
        if 'stand' not in poses:
            logger.warning(f"角色[{name}] stand 标准姿势缺失，单独重试一次")
            try:
                full_prompt = "，".join(base_parts) + "，当前姿势：自然站立，isolated on solid white background"
                img_data = generate_image(full_prompt, char_index, is_character=True,
                                          seed=char_seed, width=CHAR_SIZE, height=CHAR_SIZE)
                if img_data:
                    removed = remove_background(img_data)
                    if removed:
                        img_data = removed
                    output_filename = f"char_stand_{str(uuid.uuid4())[:8]}.png"
                    output_path = os.path.join(OUTPUT_DIR, output_filename)
                    img = Image.open(BytesIO(img_data))
                    if img.mode != 'RGBA':
                        img = img.convert('RGBA')
                    img.save(output_path, format='PNG')
                    poses['stand'] = f'/static/output/{output_filename}'
                    stand_url = poses['stand']
                    logger.info(f"角色[{name}] stand 重试成功 -> {stand_url}")
            except Exception as e:
                logger.error(f"角色[{name}] stand 重试失败: {e}")

        # 最终回退：若 stand 仍不可用，用最中性的 smile 或任一成功姿势作为主图
        if stand_url is None and poses:
            stand_url = poses.get('smile') or next(iter(poses.values()))
            logger.warning(f"角色[{name}] stand 不可用，主图回退为: {stand_url}")

        if not poses:
            logger.error("所有姿势生成均失败")
            return jsonify({'success': False, 'error': '图片生成失败', 'error_code': 'IMAGE_GENERATION_FAILED'}), 500

        total_time = time.time() - total_start_time
        logger.info(f"角色图片生成成功(姿势库): {name} -> {stand_url}, 共{len(poses)}个姿势, 耗时: {total_time:.2f}秒")
        return jsonify({
            'success': True,
            'image_url': stand_url,
            'poses': poses
        })
        
    except requests.exceptions.ConnectionError as e:
        error_msg = f"连接错误: {str(e)}"
        logger.error(f"生成角色图片连接错误: {error_msg}")
        return jsonify({'success': False, 'error': error_msg, 'error_code': 'CONNECTION_ERROR'}), 503
    except requests.exceptions.Timeout as e:
        error_msg = f"请求超时: {str(e)}"
        logger.error(f"生成角色图片超时: {error_msg}")
        return jsonify({'success': False, 'error': error_msg, 'error_code': 'TIMEOUT'}), 504
    except requests.exceptions.RequestException as e:
        error_msg = f"HTTP请求异常: {str(e)}"
        logger.error(f"生成角色图片HTTP异常: {error_msg}")
        return jsonify({'success': False, 'error': error_msg, 'error_code': 'HTTP_ERROR'}), 502
    except Exception as e:
        error_msg = f"生成角色图片失败: {str(e)}"
        logger.error(f"{error_msg}\n{traceback.format_exc()}")
        return jsonify({'success': False, 'error': error_msg, 'error_code': 'INTERNAL_ERROR'}), 500


@app.route('/api/generate_background', methods=['POST'])
def api_generate_background():
    """生成单个背景图片（支持将角色融入背景）"""
    try:
        import json
        data = request.get_json()
        
        scene = data.get('scene', '')
        index = data.get('index', 0)
        characters = data.get('characters', [])
        world_setting_data = data.get('world_setting', {})
        
        if not scene:
            return jsonify({'error': '场景描述不能为空'}), 400
        
        logger.info(f"生成背景图片: 场景='{scene[:50]}...', 索引={index}, 角色数={len(characters)}")
        
        world_setting = {
            'visual_style': world_setting_data.get('visual_style', 'Q version cartoon style, bright colors'),
            'main_location': world_setting_data.get('main_location', 'comic scene'),
            'atmosphere': world_setting_data.get('atmosphere', 'fun, educational'),
            'time_period': world_setting_data.get('time_period', '')
        }
        style_seed = f"bg_{str(uuid.uuid4())[:8]}"
        
        if not characters or len(characters) == 0:
            _char_names = []
            if characters and isinstance(characters, list):
                _char_names = [str(c.get('name') if isinstance(c, dict) else c) for c in characters]
            environment_prompt = extract_environment_only(scene, character_names=_char_names)
            if environment_prompt:
                logger.info(f"场景描述清洗: '{scene[:50]}...' -> '{environment_prompt[:50]}...'")
                scene = environment_prompt

        # 读取布局配置（容错：可能不存在），用于按每格比例"框定"生成纯场景背景
        layout_config = data.get('layout_config')
        panel_spec = None
        target_width = 768
        target_height = 768
        if isinstance(layout_config, dict):
            grid_specs = layout_config.get('grid_specs') or []
            if isinstance(grid_specs, list) and len(grid_specs) > index:
                spec = grid_specs[index]
                if isinstance(spec, dict):
                    # 关键修复：用「单元格真实比例」（由 CSS Grid 的 fr 跨度推导，
                    # 见 calculate_cell_aspect_ratio），而非 grid_spec 里声明的 aspect_ratio。
                    # 两者常不一致（如 six-cinematic 顶部大格声明 16:9，但实际单元格是 5:2），
                    # 之前按声明比例生成 → 汇总图 cover 裁剪会切掉人物/文字左右。
                    cell_ar = calculate_cell_aspect_ratio(layout_config, spec)
                    if cell_ar and ':' in str(cell_ar):
                        panel_spec = {'aspect_ratio': cell_ar}
                        target_width, target_height = calculate_image_size_for_panel(panel_spec)
                        logger.info(f"[背景] 第{index + 1}格按单元格真实比例框定生成: {cell_ar}")
                    else:
                        aspect = spec.get('aspect_ratio', '1:1')
                        if aspect and ':' in str(aspect):
                            panel_spec = {'aspect_ratio': aspect}
                            target_width, target_height = calculate_image_size_for_panel(panel_spec)
                            logger.info(f"[背景] 第{index + 1}格按布局比例框定生成: {aspect}")

        # 背景必须严格纯场景、无人物、无人群（force_pure_background 绕过 is_crowd 背影分支）
        # 免费图像 API 间歇性失败：加最多 3 次自动重试（含首次），降低偶发假图（占位渐变图）概率
        max_bg_retries = 3
        bg_image_data = None
        # 把本背景关联的角色名作为屏蔽名单，确保即使场景描述里写了角色也不会出现在纯背景中
        _bg_block_names = []
        for _c in (characters or []):
            if isinstance(_c, dict):
                _nm = _c.get('name')
                if _nm:
                    _bg_block_names.append(str(_nm))
            elif isinstance(_c, str) and _c.strip():
                _bg_block_names.append(_c.strip())
        for bg_attempt in range(max_bg_retries):
            bg_image_data = generate_image(
                scene, index, world_setting, style_seed,
                is_character=False, characters=[],
                width=target_width, height=target_height,
                panel_spec=panel_spec,
                force_pure_background=True,
                block_names=_bg_block_names,
                # 服务限流/故障时不能把内置示意图当成 AI 生成结果返回；那会造成
                # “提示词正确、图片完全无关”的假成功。
                allow_default_fallback=False
            )
            if bg_image_data is not None and not background_visual_style_matches_prompt(bg_image_data, scene):
                logger.warning(f"[背景质检] 第{index + 1}格第{bg_attempt + 1}次结果偏离暖色卡通提示词，自动重试")
                bg_image_data = None
            if bg_image_data is not None:
                if bg_attempt > 0:
                    logger.info(f"[背景重试] 第{index+1}格背景在第{bg_attempt+1}次尝试成功")
                break
            else:
                logger.warning(f"[背景重试] 第{index+1}格背景第{bg_attempt+1}次尝试失败，{'重试中...' if bg_attempt < max_bg_retries - 1 else '已达最大重试次数'}")
                if bg_attempt < max_bg_retries - 1:
                    time.sleep(2)  # 重试间隔 2 秒

        if bg_image_data is None:
            return jsonify({
                'success': False,
                'error_code': 'IMAGE_PROVIDER_UNAVAILABLE',
                'error': '生图服务当前限流或暂不可用，未生成与场景描述匹配的背景。请稍后点击“重新生成”。'
            }), 503
        
        output_filename = f"bg_{str(uuid.uuid4())[:8]}_{index}.png"
        output_path = os.path.join(OUTPUT_DIR, output_filename)
        
        from io import BytesIO
        img = Image.open(BytesIO(bg_image_data))
        img = normalize_background_for_panel(
            img, target_width, target_height
        )
        img.save(output_path)
        
        image_url = f'/static/output/{output_filename}'
        
        logger.info(f"背景图片生成成功: {image_url}")
        
        return jsonify({
            'success': True,
            'image_url': image_url
        })
        
    except Exception as e:
        logger.error(f"生成背景图片失败: {str(e)}")
        return jsonify({'error': f'生成背景图片失败: {str(e)}'}), 500


def generate_single_background(args):
    """生成单张背景图片（用于并行执行）"""
    is_translated = False
    if len(args) == 6:
        i, scene, characters, world_setting, style_seed, is_translated = args
    else:
        i, scene, characters, world_setting, style_seed = args
    
    try:
        logger.info(f"开始生成背景{i+1}: {scene[:50]}...")
        
        # 汇总角色名（即便 characters 为空也保留），作为纯背景时屏蔽名单
        _char_names = []
        if characters and isinstance(characters, list):
            _char_names = [str(c.get('name') if isinstance(c, dict) else c) for c in characters]
        if not characters or len(characters) == 0:
            environment_prompt = extract_environment_only(scene, character_names=_char_names)
            if environment_prompt:
                logger.info(f"场景描述清洗: '{scene[:50]}...' -> '{environment_prompt[:50]}...'")
                scene = environment_prompt
        
        bg_image_data = generate_image(
            scene, i, world_setting, style_seed,
            is_character=False, characters=[], is_translated=is_translated,
            force_pure_background=True, block_names=_char_names,
            allow_default_fallback=False
        )
        if bg_image_data is None:
            raise RuntimeError('生图服务当前限流或暂不可用，未生成与场景描述匹配的背景')
        if not background_visual_style_matches_prompt(bg_image_data, scene):
            raise RuntimeError('生图结果与场景要求的暖色卡通风格明显不符，请重新生成')
        
        output_filename = f"bg_{style_seed}_{i}.png"
        output_path = os.path.join(OUTPUT_DIR, output_filename)
        
        from io import BytesIO
        img = Image.open(BytesIO(bg_image_data))
        img.save(output_path)
        
        url = f'/static/output/{output_filename}'
        
        logger.info(f"背景{i+1}生成成功: {url}")
        
        return {
            'index': i,
            'url': url,
            'success': True,
            'is_default': False
        }
        
    except Exception as e:
        logger.error(f"生成背景{i+1}失败: {str(e)}")
        return {
            'index': i,
            'url': None,
            'success': False,
            'is_default': False,
            'error': str(e)
        }


@app.route('/api/generate_all_backgrounds', methods=['POST'])
def api_generate_all_backgrounds():
    """批量生成所有背景图片（优化版：先批量翻译，再并行生成，带速率限制）"""
    try:
        import json
        import time
        import threading
        
        logger.info("="*50)
        logger.info("接收到批量生成背景图片请求")
        start_total = time.time()
        
        data = request.get_json()
        logger.info(f"请求数据长度: {len(json.dumps(data, ensure_ascii=False))}")
        
        scenes = data.get('scenes', [])
        characters = data.get('characters', [])
        world_setting_data = data.get('world_setting', {})
        
        if not scenes or len(scenes) == 0:
            return jsonify({'error': '场景列表不能为空'}), 400
        
        logger.info(f"批量生成背景图片: 共{len(scenes)}个场景, 角色数={len(characters)}")
        logger.info(f"场景列表: {scenes}")
        
        world_setting = {
            'visual_style': world_setting_data.get('visual_style', 'Q version cartoon style, bright colors'),
            'main_location': world_setting_data.get('main_location', 'comic scene'),
            'atmosphere': world_setting_data.get('atmosphere', 'fun, educational'),
            'time_period': world_setting_data.get('time_period', '')
        }
        style_seed = f"bg_batch_{str(uuid.uuid4())[:8]}"
        
        images = [None] * len(scenes)
        results = []
        
        start_time = time.time()
        
        translated_scenes = []
        for i, scene in enumerate(scenes):
            try:
                translated = translate_to_english(scene)
                translated_scenes.append(translated)
                logger.info(f"场景{i+1}翻译完成: '{scene[:30]}...' -> '{translated[:30]}...'")
            except Exception as e:
                logger.warning(f"场景{i+1}翻译失败，使用原文: {str(e)}")
                translated_scenes.append(scene)
        
        translate_time = time.time()
        logger.info(f"批量翻译完成，耗时{translate_time - start_time:.2f}秒")
        
        logger.info(f"使用串行方式生成背景图片（避免Pollinations.AI限流）")
        
        min_request_interval = 5.0
        
        for i in range(len(scenes)):
            if i > 0:
                elapsed = time.time() - last_request_time[0]
                wait_time = min_request_interval - elapsed
                if wait_time > 0:
                    logger.info(f"等待{wait_time:.2f}秒后开始生成背景{i+1}")
                    time.sleep(wait_time)
            
            last_request_time = [time.time()]
            
            try:
                logger.info(f"开始生成背景{i+1}/{len(scenes)}")
                result = generate_single_background((i, translated_scenes[i], [], world_setting, style_seed, True))
                images[i] = result['url']
                results.append(result)
                logger.info(f"背景{i+1}/{len(scenes)}完成")
            except Exception as e:
                logger.error(f"背景{i+1}处理异常: {str(e)}")
                images[i] = None
                results.append({
                    'index': i,
                    'url': None,
                    'success': False,
                    'is_default': False,
                    'error': str(e)
                })
        
        results.sort(key=lambda x: x['index'])
        
        success_count = sum(1 for r in results if r['success'])
        total_time = time.time() - start_total
        generate_time = time.time() - translate_time
        
        logger.info(f"批量背景生成完成: 成功{success_count}/{len(scenes)}个")
        logger.info(f"总耗时: {total_time:.2f}秒, 翻译耗时: {translate_time - start_total:.2f}秒, 生成耗时: {generate_time:.2f}秒")
        logger.info(f"平均每个背景耗时: {total_time/len(scenes):.2f}秒")
        logger.info("="*50)
        
        return jsonify({
            'success': True,
            'images': images,
            'results': results,
            'success_count': success_count,
            'total_count': len(scenes),
            'total_time': round(total_time, 2),
            'translate_time': round(translate_time - start_total, 2),
            'generate_time': round(generate_time, 2)
        })
        
    except Exception as e:
        total_time = time.time() - start_total
        logger.error(f"批量生成背景图片失败: {str(e)}")
        logger.error(f"失败时已耗时: {total_time:.2f}秒")
        logger.info("="*50)
        return jsonify({'error': f'批量生成背景图片失败: {str(e)}'}), 500


@app.route('/api/generate_comic_from_script', methods=['POST'])
def api_generate_comic_from_script():
    """从脚本生成漫画"""
    try:
        import json
        data = request.get_json()
        
        logger.info(f"接收到漫画生成请求，数据长度: {len(json.dumps(data, ensure_ascii=False))}")
        
        script = data.get('script', {})
        characters = data.get('characters', {})
        backgrounds = data.get('backgrounds', {})
        layout_prompts = data.get('layoutPrompts', [])
        template_id = data.get('template_id', '')
        bubble_images_data = data.get('bubbleImages', [])
        character_descriptions = data.get('characterDescriptions', script.get('characters', []))
        integrated_scene = data.get('integrated_scene', False)
        character_poses = data.get('characterPoses', {})  # 角色名 -> {poseKey: url}
        selected_guide_sprite = data.get('selectedGuideSprite')
        # 按人物批量指定气泡：人物名 -> 气泡标识（emotion key 或气泡图片 url）
        character_bubble_map = data.get('characterBubbleMap', {})
        if not isinstance(character_bubble_map, dict):
            character_bubble_map = {}

        # Logo 水印参数（仅第一格叠加）
        logo_url = data.get('logoUrl', None)
        logo_position = data.get('logoPosition', 'top-right')
        logo_size = data.get('logoSize', 0.08)
        try:
            logo_size = float(logo_size)
        except (TypeError, ValueError):
            logo_size = 0.08

        # 汇总本漫画所有角色名，作为背景纯净化时要屏蔽的名字（避免角色名出现在背景里）
        comic_character_names = []
        for _cd in (character_descriptions or []):
            if isinstance(_cd, dict):
                _nm = _cd.get('name')
                if _nm:
                    comic_character_names.append(str(_nm))
            elif isinstance(_cd, str) and _cd.strip():
                comic_character_names.append(_cd.strip())
        # 也把脚本里出现的说话者/角色名补进来
        for _p in (script.get('panels', []) or []):
            for _sp in (_p.get('speakers', []) or []):
                if _sp and _sp not in comic_character_names:
                    comic_character_names.append(str(_sp))
            _sp = _p.get('speaker')
            if _sp and _sp not in comic_character_names:
                comic_character_names.append(str(_sp))

        logger.info(f"布局提示词数量: {len(layout_prompts)}, 角色数量: {len(characters)}, "
                    f"panel数量: {len(script.get('panels', []))}, 模板ID: {template_id}, "
                    f"气泡数据: {len(bubble_images_data)}, "
                    f"人物气泡映射: {character_bubble_map}, "
                    f"完整场景模式: {integrated_scene}")
        
        if not script or not script.get('panels'):
            return jsonify({'error': '脚本数据无效'}), 400
        
        panels = script.get('panels', [])
        
        if len(panels) == 0:
            return jsonify({'error': '脚本中没有分镜数据'}), 400
        
        # ===== 解析布局：优先使用前端指定的模板ID，其次从提示词推断 =====
        layout_config = None
        if template_id and template_id in LAYOUT_CONFIG:
            layout_config = get_layout_config(template_id)
            logger.info(f"使用前端指定布局模板: {template_id}")
        elif layout_prompts:
            layout_config = parse_layout_prompts(layout_prompts)
        
        if layout_config:
            logger.info(f"检测到布局模板: {layout_config.get('detected_layout', layout_config.get('name', 'unknown'))}, "
                        f"grid规格数: {len(layout_config.get('grid_specs', []))}")
        
        comic_images = []
        panels_meta = []
        
        world_setting = {
            'visual_style': 'Q version cartoon style, bright colors',
            'main_location': 'comic scene',
            'atmosphere': 'fun, educational'
        }
        style_seed = f"{str(uuid.uuid4())[:8]}"
        
        for i, panel in enumerate(panels):
            scene = panel.get('scene', '')
            dialogue = panel.get('dialogue', '')
            dialogues = panel.get('dialogues', [])
            speaker = panel.get('speaker', '')
            speaker_list = panel.get('speakers', [])
            
            layout_prompt = layout_prompts[i] if i < len(layout_prompts) else ''
            
            # 获取当前格的布局规格
            panel_spec = None
            if layout_config and layout_config.get('grid_specs'):
                grid_specs = layout_config['grid_specs']
                if i < len(grid_specs):
                    panel_spec = grid_specs[i]
                    # 用 CSS Grid fr 值计算单元格真实宽高比，覆盖默认的 1:1
                    if panel_spec and not layout_config.get('freeform'):
                        actual_ratio = calculate_cell_aspect_ratio(layout_config, panel_spec)
                        panel_spec['aspect_ratio'] = actual_ratio
                        logger.info(f"分镜{i+1}布局规格: aspect={actual_ratio} (原={panel_spec.get('aspect_ratio', '1:1')}), "
                                    f"shot={panel_spec.get('shot_type')}, angle={panel_spec.get('angle')}")
                    else:
                        logger.info(f"分镜{i+1}布局规格: aspect={panel_spec.get('aspect_ratio', '1:1')}, "
                                    f"shot={panel_spec.get('shot_type')}, angle={panel_spec.get('angle')}")

            # 根据当前布局规格计算目标尺寸（必须在背景生成分支之前完成，避免 UnboundLocalError）
            # 确保合成图与前端卡片比例一致，避免 cover 裁剪
            target_width, target_height = None, None
            if panel_spec:
                try:
                    # 合成画布用更高基准(1152)，文字/气泡画在高分辨率画布上，前端放大不糊
                    target_width, target_height = calculate_image_size_for_panel(panel_spec, base_size=1152)
                    logger.info(f"分镜{i+1}目标尺寸: {target_width}x{target_height} (aspect={panel_spec.get('aspect_ratio', '1:1')})")
                except Exception as e:
                    logger.warning(f"计算分镜{i+1}目标尺寸失败: {e}")

            logger.info(f"处理分镜{i+1}: 场景='{scene}', 说话者={speaker_list}, 对话={dialogues}, 布局提示词='{layout_prompt}'")
            
            bg_image_data = None
            bg_url_used = None
            
            if backgrounds and str(i) in backgrounds:
                bg_url = backgrounds[str(i)]
                bg_path = os.path.join(app.root_path, bg_url.lstrip('/'))
                if os.path.exists(bg_path):
                    logger.info(f"使用预先生成的背景: {bg_url}")
                    # 如果 target 已定义，确保背景图尺寸与合成图一致，避免编辑/PPT 中人物坐标空间错位
                    if target_width and target_height:
                        try:
                            with Image.open(bg_path) as img:
                                # 即便尺寸已经相同也必须规范化：模型可能在图片内部画了
                                # 上下深灰/黑色电影边，普通 resize 无法消除。
                                resized = normalize_background_for_panel(
                                    img, target_width, target_height
                                )
                                resized.save(bg_path, 'PNG')
                                logger.info(
                                    f"预生成背景已规范化到 {target_width}x{target_height}"
                                )
                        except Exception as e:
                            logger.warning(f"预生成背景 resize 失败: {e}")
                    with open(bg_path, 'rb') as f:
                        bg_image_data = f.read()
                    bg_url_used = bg_url
            
            # 根据当前格的说话者筛选角色描述，用于生成完整场景
            panel_character_descs = []
            if integrated_scene and character_descriptions:
                panel_speakers = set(speaker_list) if speaker_list else ({speaker} if speaker else set())
                for char_desc in character_descriptions:
                    if isinstance(char_desc, dict):
                        char_name = char_desc.get('name', '')
                        if char_name in panel_speakers:
                            desc_parts = [char_name]
                            if char_desc.get('role'):
                                desc_parts.append(char_desc['role'])
                            if char_desc.get('appearance'):
                                desc_parts.append(char_desc['appearance'])
                            panel_character_descs.append(', '.join(desc_parts))
                    elif isinstance(char_desc, str):
                        panel_character_descs.append(char_desc)

            if bg_image_data is None:
                logger.info(f"生成新背景: {scene}")
                # 把panel_spec的angle追加到angle参数（如果有的话）
                angle = None
                if panel_spec and panel_spec.get('angle'):
                    angle = panel_spec['angle']
                # 把panel_spec的shot_type追加到composition
                composition = None
                if panel_spec and panel_spec.get('shot_type'):
                    composition = panel_spec['shot_type']

                bg_image_data = generate_image(
                    scene, i, world_setting, style_seed,
                    layout_prompt=layout_prompt,
                    panel_spec=panel_spec,
                    angle=angle,
                    composition=composition,
                    characters=None,
                    force_pure_background=True,
                    block_names=comic_character_names
                )
                # 保存新生成的背景并记录 URL，供编辑面板使用
                bg_filename = f"bg_{str(uuid.uuid4())[:8]}_{i}.png"
                bg_path = os.path.join(OUTPUT_DIR, bg_filename)
                with open(bg_path, 'wb') as f:
                    f.write(bg_image_data)
                # 如果 target 已定义，确保背景图尺寸与合成图一致
                if target_width and target_height:
                    try:
                        with Image.open(bg_path) as img:
                            if img.size != (target_width, target_height):
                                resized = normalize_background_for_panel(
                                    img, target_width, target_height
                                )
                                resized.save(bg_path, 'PNG')
                                logger.info(f"新背景已 cover 规范化到 {target_width}x{target_height}")
                    except Exception as e:
                        logger.warning(f"新背景 resize 失败: {e}")
                # 重新读取 resize 后的数据供 composite_image 使用
                with open(bg_path, 'rb') as f:
                    bg_image_data = f.read()
                bg_url_used = f'/static/output/{bg_filename}'
                logger.info(f"新背景已保存: {bg_url_used}")
            
            output_filename = f"comic_{str(uuid.uuid4())[:8]}_{i}.png"
            output_path = os.path.join(OUTPUT_DIR, output_filename)
            
            ip_paths = []
            speaker_names = []
            
            if speaker_list:
                for sp in speaker_list:
                    if sp in characters:
                        ip_paths.append(os.path.join(app.root_path, characters[sp].lstrip('/')))
                        speaker_names.append(sp)
                    elif sp == '知识小精灵':
                        default_ip = get_selected_guide_image(characters, selected_guide_sprite)
                        if default_ip:
                            ip_paths.append(default_ip)
                            speaker_names.append(sp)
            elif speaker and speaker in characters:
                ip_paths = [os.path.join(app.root_path, characters[speaker].lstrip('/'))]
                speaker_names = [speaker]
            elif speaker == '知识小精灵':
                default_ip = get_selected_guide_image(characters, selected_guide_sprite)
                if default_ip:
                    ip_paths = [default_ip]
                    speaker_names = [speaker]
            
            # 如果panel_spec指定了guide_sprite，但当前没有知识小精灵，自动加上
            if panel_spec and panel_spec.get('guide_sprite') and '知识小精灵' not in speaker_names:
                default_ip = get_selected_guide_image(characters, selected_guide_sprite)
                if default_ip:
                    ip_paths.append(default_ip)
                    speaker_names.append('知识小精灵')
                    if not speaker:
                        speaker = '知识小精灵'
                    logger.info(f"分镜{i+1}自动添加知识小精灵（布局指定）")
            
            # 音效文字
            sound_effect_text = None
            
            panel_bubble_images = []
            if i < len(bubble_images_data) and bubble_images_data[i]:
                for bubble_url in bubble_images_data[i]:
                    if bubble_url:
                        bubble_path = os.path.join(app.root_path, bubble_url.lstrip('/'))
                        if os.path.exists(bubble_path):
                            panel_bubble_images.append(bubble_path)
                            logger.info(f"分镜{i+1}使用用户选择的气泡: {bubble_path}")
            
            composite_result = composite_image(
                bg_image_data,
                ip_paths,
                dialogue,
                output_path,
                speaker=speaker,
                speakers=speaker_names,
                dialogues=dialogues,
                bubble_images=panel_bubble_images if panel_bubble_images else None,
                character_bubble_map=character_bubble_map,
                skip_ip_overlay=integrated_scene,
                character_poses=character_poses,
                font_family='msyhbd',
                target_width=target_width,
                target_height=target_height,
                logo_url=logo_url if i == 0 else None,
                logo_position=logo_position,
                logo_size=logo_size
            )
            
            panel_meta = composite_result if isinstance(composite_result, dict) else {
                'success': bool(composite_result),
                'characters': [],
                'bubbles': [],
                'sound_effect': None
            }

            # 合成失败时直接返回错误，避免返回不存在的图片 URL
            if not panel_meta.get('success', False):
                error_msg = panel_meta.get('error', '未知合成错误')
                logger.error(f"分镜{i+1}合成失败: {error_msg}")
                return jsonify({'error': f'分镜{i+1}合成失败: {error_msg}'}), 500

            panel_meta['url'] = f'/static/output/{output_filename}'
            panel_meta['background_url'] = bg_url_used
            panel_meta['index'] = i
            panel_meta['dialogues'] = dialogues
            panel_meta['speakers'] = speaker_names
            panel_meta['bubble_images'] = panel_bubble_images
            panels_meta.append(panel_meta)
            
            comic_images.append(f'/static/output/{output_filename}')
        
        logger.info(f"漫画生成成功，共{len(comic_images)}格")
        
        result = {
            'success': True,
            'images': comic_images,
            'panels': panels_meta,
            'layout_prompts': layout_prompts
        }
        
        # 返回布局配置给前端，用于CSS Grid渲染
        if layout_config:
            grid_specs_out = []
            for spec in layout_config.get('grid_specs', []):
                spec_copy = dict(spec)
                # 补充真实目标宽高（根据 aspect_ratio 计算，便于前端锁定卡片比例）
                try:
                    tw, th = calculate_image_size_for_panel(spec_copy, base_size=768)
                    spec_copy['target_width'] = tw
                    spec_copy['target_height'] = th
                    spec_copy['target_aspect_ratio'] = spec_copy.get('aspect_ratio', '1:1')
                except Exception as e:
                    logger.warning(f"计算目标宽高失败: {e}")
                    spec_copy['target_width'] = 768
                    spec_copy['target_height'] = 768
                grid_specs_out.append(spec_copy)
            result['layout_config'] = {
                'css_grid': layout_config.get('css_grid', ''),
                'grid_specs': grid_specs_out,
                'rows': layout_config.get('rows', 2),
                'cols': layout_config.get('cols', 2),
                'detected_layout': layout_config.get('detected_layout', ''),
                'name': layout_config.get('name', '')
            }
        
        return jsonify(result)
        
    except Exception as e:
        logger.error(f"生成漫画失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'生成漫画失败: {str(e)}'}), 500


@app.route('/api/layout_templates', methods=['GET'])
def api_layout_templates():
    """获取所有可用的布局模板"""
    try:
        panel_count = request.args.get('panel_count', type=int)
        templates = []
        
        for key, config in LAYOUT_CONFIG.items():
            if panel_count and config.get('panel_count') != panel_count:
                continue
            templates.append({
                'id': key,
                'name': config.get('name', key),
                'description': config.get('description', ''),
                'panel_count': config.get('panel_count', 0),
                'rows': config.get('rows', 1),
                'cols': config.get('cols', 1),
                'css_grid': config.get('css_grid', ''),
                'grid_specs': config.get('grid_specs', [])
            })
        
        return jsonify({
            'success': True,
            'templates': templates
        })
    except Exception as e:
        logger.error(f"获取布局模板失败: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/bubbles', methods=['GET'])
def api_get_bubbles():
    """获取可用的气泡列表"""
    try:
        emotion = request.args.get('emotion', 'all')
        bubbles = []
        
        base_dir = os.path.join(app.root_path, 'static', 'bubbles')
        
        emotions_list = ['happy', 'sad', 'angry', 'surprised', 'neutral']
        if emotion != 'all' and emotion in emotions_list:
            emotions_list = [emotion]
        
        for emo in emotions_list:
            emo_dir = os.path.join(base_dir, emo)
            if os.path.exists(emo_dir):
                for filename in os.listdir(emo_dir):
                    if filename.lower().endswith('.png'):
                        bubbles.append({
                            'id': filename,
                            'emotion': emo,
                            'name': os.path.splitext(filename)[0],
                            'url': f'/static/bubbles/{emo}/{filename}'
                        })
        
        return jsonify({
            'success': True,
            'bubbles': bubbles
        })
    except Exception as e:
        logger.error(f"获取气泡列表失败: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/bubbles/upload', methods=['POST'])
def api_upload_bubble():
    """上传自定义气泡图片"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': '请选择上传的文件'}), 400
        
        file = request.files['file']
        emotion = request.form.get('emotion', 'custom')
        
        if file.filename == '':
            return jsonify({'error': '文件名不能为空'}), 400
        
        if not file.filename.lower().endswith('.png'):
            return jsonify({'error': '只支持PNG格式'}), 400
        
        custom_dir = os.path.join(app.root_path, 'static', 'bubbles', 'custom')
        os.makedirs(custom_dir, exist_ok=True)
        
        import uuid
        unique_name = f"custom_{uuid.uuid4().hex[:8]}_{file.filename}"
        filepath = os.path.join(custom_dir, unique_name)
        
        file.save(filepath)
        
        return jsonify({
            'success': True,
            'bubble': {
                'id': unique_name,
                'emotion': 'custom',
                'name': os.path.splitext(unique_name)[0],
                'url': f'/static/bubbles/custom/{unique_name}'
            }
        })
    except Exception as e:
        logger.error(f"上传气泡失败: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/analyze_emotion', methods=['POST'])
def api_analyze_emotion():
    """分析文本情绪"""
    try:
        text = request.json.get('text', '')
        emotion = analyze_emotion(text)
        
        return jsonify({
            'success': True,
            'emotion': emotion
        })
    except Exception as e:
        logger.error(f"情绪分析失败: {str(e)}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/export_comic', methods=['POST'])
def api_export_comic():
    """导出漫画（批量打包为ZIP）"""
    try:
        import json
        from urllib.parse import quote
        data = request.get_json()
        
        images = data.get('images', [])
        export_title = str(data.get('title') or '').strip()
        export_basename = sanitize_export_filename(export_title or PPT_DEFAULT_TITLE)
        
        if not images:
            return jsonify({'error': '没有可导出的图片'}), 400
        
        import zipfile
        from io import BytesIO
        from flask import send_file
        
        zip_buffer = BytesIO()
        
        # 汇总大图：把各格按布局拼成一张（复用已有 compose_aggregate_image）
        # 修复：直接把前端透传的完整 layout_config 字典交给 compose_aggregate_image，
        # 与链接版同一数据源（含 css_grid/grid_specs），彻底避免"只抠 id 抠不到就 fallback four-grid"。
        aggregate_bytes = None
        layout_config = data.get('layout_config', None)
        layout_arg = layout_config if isinstance(layout_config, (dict, str)) and layout_config else 'four-grid'
        try:
            aggregate_bytes = compose_aggregate_image(images, layout=layout_arg)
        except Exception as agg_e:
            logger.warning(f"汇总图生成失败（不影响单格导出）: {agg_e}")
        
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for i, img_url in enumerate(images):
                if img_url:
                    # strip ?t= 查询串，保证编辑过的格子也能正确写入 comic_N.png
                    img_path = os.path.join(app.root_path, img_url.split('?')[0].lstrip('/'))
                    if os.path.exists(img_path):
                        zipf.write(img_path, f'{export_basename}_第{i+1}格.png')
            if aggregate_bytes:
                zipf.writestr(f'{export_basename}_合并图.png', aggregate_bytes)
        
        zip_buffer.seek(0)
        if data.get('prepare_download'):
            # 让浏览器以普通 GET 文件下载，而非 fetch/blob 接收大文件；后者在部分
            # Chromium 环境会无提示中断并显示 net::ERR_FAILED。
            export_name = f'comic_export_{datetime.now().strftime("%Y%m%d_%H%M%S")}_{uuid.uuid4().hex[:6]}.zip'
            export_path = os.path.join(OUTPUT_DIR, export_name)
            with open(export_path, 'wb') as export_file:
                export_file.write(zip_buffer.getvalue())
            return jsonify({
                'success': True,
                'filename': f'{export_basename}_图片.zip',
                'download_url': f'/api/download_export/{export_name}?download_name={quote(f"{export_basename}_图片.zip")}'
            })
        
        return send_file(
            zip_buffer,
            mimetype='application/zip',
            as_attachment=True,
            download_name=f'{export_basename}_图片.zip'
        )
        
    except Exception as e:
        logger.error(f"导出漫画失败: {str(e)}")
        return jsonify({'error': f'导出漫画失败: {str(e)}'}), 500


@app.route('/api/download_export/<filename>', methods=['GET'])
def download_export_file(filename):
    """下载已准备好的 ZIP；只允许导出目录内的 comic_export 文件。"""
    safe_name = os.path.basename(filename)
    if safe_name != filename or not safe_name.startswith('comic_export_') or not safe_name.endswith('.zip'):
        return jsonify({'error': '无效的下载文件'}), 400
    export_path = os.path.join(OUTPUT_DIR, safe_name)
    if not os.path.isfile(export_path):
        return jsonify({'error': '下载文件不存在或已过期'}), 404
    requested_name = str(request.args.get('download_name') or '').strip()
    if requested_name.lower().endswith('.zip'):
        requested_name = f'{sanitize_export_filename(requested_name[:-4])}.zip'
    else:
        requested_name = safe_name
    return send_file(export_path, mimetype='application/zip', as_attachment=True,
                     download_name=requested_name, conditional=True)


@app.errorhandler(404)
def page_not_found(e):
    """404错误处理"""
    return render_template('index.html', guide_config=ROLE_GUIDE_CONFIG, error='页面未找到'), 404


@app.errorhandler(413)
def request_entity_too_large(e):
    """413文件过大错误处理"""
    logger.error(f"文件上传过大: {str(e)}")
    return jsonify({'error': '文件太大了！请上传小于500MB的文件'}), 413


@app.route('/api/upload_character_image', methods=['POST'])
def api_upload_character_image():
    """上传角色图片"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未选择文件'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        if file:
            allowed_extensions = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
            ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
            
            if ext not in allowed_extensions:
                return jsonify({'error': '不支持的文件格式'}), 400
            
            output_filename = f"char_upload_{str(uuid.uuid4())[:8]}.{ext}"
            output_path = os.path.join(OUTPUT_DIR, output_filename)
            
            file.save(output_path)
            
            image_url = f'/static/output/{output_filename}'
            
            logger.info(f"角色图片上传成功: {image_url}")
            
            return jsonify({
                'success': True,
                'image_url': image_url
            })
            
    except Exception as e:
        logger.error(f"上传角色图片失败: {str(e)}")
        return jsonify({'error': f'上传角色图片失败: {str(e)}'}), 500


@app.errorhandler(500)
def internal_server_error(e):
    """500错误处理"""
    import traceback
    logger.error(f"服务器内部错误: {str(e)}\n{traceback.format_exc()}")
    
    if request.path.startswith('/api/'):
        return jsonify({'error': f'服务器内部错误: {str(e)}'}), 500
    
    return render_template('index.html', guide_config=ROLE_GUIDE_CONFIG, error='服务器出错了，请稍后重试'), 500



@app.route('/api/save_panel_edit', methods=['POST'])
def save_panel_edit():
    """保存前端编辑器修改后的一格漫画图片"""
    import traceback
    try:
        panel_index = request.form.get('panel_index', '0')
        # 编辑后的成图不能再使用 panel_0_edited.png 这类全局固定名。
        # 不同任务编辑同一格时会互相覆盖该文件，造成任务预览串图。
        task_id = request.form.get('task_id', '').strip()
        image_data = request.form.get('image_data', '')
        bubbles_data = request.form.get('bubbles', '')
        
        if isinstance(panel_index, str) and panel_index.isdigit():
            panel_index = int(panel_index)
        else:
            panel_index = 0
        
        logger.info(f"收到保存面板编辑请求: task_id={task_id or 'unsaved'}, panel_index={panel_index}, image_data长度={len(image_data) if image_data else 0}, bubbles长度={len(bubbles_data) if bubbles_data else 0}")
        
        if not image_data or len(image_data) < 100:
            logger.error(f"缺少图片数据或图片数据过短: {len(image_data) if image_data else 0}")
            return jsonify({'error': '缺少图片数据'}), 400

        import base64
        if ',' in image_data:
            image_data = image_data.split(',')[1]
            logger.info(f"移除data前缀后，image_data长度: {len(image_data)}")
        
        if len(image_data) % 4 != 0:
            padding = 4 - (len(image_data) % 4)
            image_data += '=' * padding
            logger.info(f"添加{padding}个=填充，长度变为: {len(image_data)}")
        
        try:
            img_bytes = base64.b64decode(image_data)
            logger.info(f"Base64解码成功，图片字节数: {len(img_bytes)}")
            logger.info(f"图片前10字节: {img_bytes[:10].hex()}")
            if img_bytes[:8] == b'\x89PNG\r\n\x1a\n':
                logger.info("图片是有效的PNG格式")
            else:
                logger.warning(f"图片可能不是有效的PNG，魔术字节: {img_bytes[:8].hex()}")
        except Exception as decode_err:
            logger.error(f"Base64解码失败: {str(decode_err)}")
            return jsonify({'error': '图片数据格式错误'}), 400
        
        if len(img_bytes) < 100:
            logger.error(f"图片数据过短，可能是无效图片: {len(img_bytes)}字节")
            return jsonify({'error': '无效的图片数据'}), 400

        import os
        output_dir = os.path.join(app.root_path, 'static', 'output')
        os.makedirs(output_dir, exist_ok=True)
        logger.info(f"输出目录: {output_dir}")
        logger.info(f"目录是否可写: {os.access(output_dir, os.W_OK)}")
        
        # task_id 仅用于可读性；UUID 使同一任务反复保存同一格也不会覆盖旧任务。
        safe_task_id = ''.join(ch for ch in task_id if ch.isalnum() or ch in ('-', '_'))[:48] or 'unsaved'
        filename = f'panel_{safe_task_id}_{panel_index}_{uuid.uuid4().hex[:12]}_edited.png'
        output_path = os.path.join(output_dir, filename)
        logger.info(f"输出路径: {output_path}")

        try:
            from io import BytesIO
            logger.info("开始打开图片...")
            img = Image.open(BytesIO(img_bytes))
            logger.info(f"图片打开成功，格式: {img.format}, 模式: {img.mode}, 尺寸: {img.size}")
            
            if img.mode != 'RGBA':
                logger.info(f"转换图片模式从{img.mode}到RGBA")
                img = img.convert('RGBA')
            
            logger.info(f"开始保存图片到: {output_path}")
            
            for attempt in range(3):
                try:
                    img.save(output_path, format='PNG')
                    logger.info(f"面板 {panel_index} 编辑已保存: {output_path}, 尺寸: {img.size}")
                    break
                except Exception as save_err:
                    if attempt < 2:
                        logger.warning(f"保存尝试 {attempt + 1} 失败，重试中: {str(save_err)}")
                        time.sleep(0.1)
                    else:
                        raise save_err
        except Exception as img_err:
            logger.error(f"图片处理失败: {str(img_err)}\n{traceback.format_exc()}")
            return jsonify({'error': f'图片处理失败: {str(img_err)}'}), 500

        url = f'/static/output/{filename}?t={int(time.time())}'
        return jsonify({'success': True, 'url': url})
    except Exception as e:
        logger.error(f"保存面板编辑失败: {str(e)}\n{traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500


@app.route('/api/client_log', methods=['POST'])
def api_client_log():
    """接收前端日志上报"""
    try:
        data = request.get_json(silent=True) or {}
        level = data.get('level', 'info')
        message = data.get('message', '')
        source = data.get('source', 'unknown')
        
        if level.lower() == 'error':
            logger.error(f"[CLIENT-ERROR] [{source}] {message}")
        elif level.lower() == 'warning':
            logger.warning(f"[CLIENT-WARN] [{source}] {message}")
        elif level.lower() == 'debug':
            logger.debug(f"[CLIENT-DEBUG] [{source}] {message}")
        else:
            logger.info(f"[CLIENT-INFO] [{source}] {message}")
        
        return jsonify({'success': True})
    except Exception as e:
        logger.error(f"处理客户端日志失败: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500


# ===================== 用户与任务系统（工号+姓名登录，SQLite） =====================
import sqlite3

DB_DIR = os.path.normpath(os.environ.get('APP_DATA_DIR') or os.path.join(app.root_path, 'data'))
DB_PATH = os.path.join(DB_DIR, 'app.db')
PASSWORD_RESET_ATTEMPTS = {}


def _credential_cipher():
    """获取用户凭证加密器；Railway 必须显式配置独立主密钥。"""
    configured = os.environ.get('USER_SECRET_ENCRYPTION_KEY', '').strip()
    if configured:
        try:
            return Fernet(configured.encode('ascii'))
        except (ValueError, TypeError):
            raise RuntimeError('USER_SECRET_ENCRYPTION_KEY 配置无效')
    if os.environ.get('RAILWAY_ENVIRONMENT'):
        raise RuntimeError('Railway 尚未配置用户凭证加密主密钥')
    # 仅供本地开发兼容；生产环境不会走到这里。
    material = hashlib.sha256(app.config['SECRET_KEY'].encode('utf-8')).digest()
    return Fernet(base64.urlsafe_b64encode(material))


def encrypt_user_secret(value):
    return _credential_cipher().encrypt(str(value).encode('utf-8')).decode('ascii')


def decrypt_user_secret(value):
    if not value:
        return ''
    try:
        return _credential_cipher().decrypt(str(value).encode('ascii')).decode('utf-8')
    except InvalidToken:
        logger.error('[user-api-key] 用户凭证无法解密，可能是主密钥发生变化')
        return ''

def get_db_conn():
    os.makedirs(DB_DIR, exist_ok=True)
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_user_db():
    conn = get_db_conn()
    try:
        conn.execute('''CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            work_id TEXT UNIQUE NOT NULL,
            name TEXT NOT NULL,
            password_hash TEXT,
            auth_version INTEGER NOT NULL DEFAULT 0,
            ark_api_key_encrypted TEXT,
            ark_api_key_last4 TEXT,
            image_provider TEXT,
            image_api_key_encrypted TEXT,
            image_api_key_last4 TEXT,
            image_base_url TEXT,
            image_model TEXT,
            created_at TEXT NOT NULL
        )''')
        user_columns = {row['name'] for row in conn.execute('PRAGMA table_info(users)').fetchall()}
        for column, definition in (
            ('password_hash', 'TEXT'),
            ('auth_version', 'INTEGER NOT NULL DEFAULT 0'),
            ('ark_api_key_encrypted', 'TEXT'),
            ('ark_api_key_last4', 'TEXT'),
            ('image_provider', 'TEXT'),
            ('image_api_key_encrypted', 'TEXT'),
            ('image_api_key_last4', 'TEXT'),
            ('image_base_url', 'TEXT'),
            ('image_model', 'TEXT'),
        ):
            if column not in user_columns:
                conn.execute(f'ALTER TABLE users ADD COLUMN {column} {definition}')
        conn.execute('''CREATE TABLE IF NOT EXISTS tasks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            title TEXT NOT NULL,
            pinned INTEGER NOT NULL DEFAULT 0,
            sort_order INTEGER NOT NULL DEFAULT 0,
            data_json TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            FOREIGN KEY (user_id) REFERENCES users(id)
        )''')
        # 兼容已经存在的本地数据库：任务位置必须独立于 updated_at，保存编辑时不应跳到第一项。
        columns = {row['name'] for row in conn.execute('PRAGMA table_info(tasks)').fetchall()}
        if 'sort_order' not in columns:
            conn.execute('ALTER TABLE tasks ADD COLUMN sort_order INTEGER NOT NULL DEFAULT 0')
            for user_row in conn.execute('SELECT DISTINCT user_id FROM tasks').fetchall():
                task_rows = conn.execute(
                    'SELECT id FROM tasks WHERE user_id=? ORDER BY pinned DESC, updated_at DESC, id DESC',
                    (user_row['user_id'],),
                ).fetchall()
                for order, task_row in enumerate(task_rows, start=1):
                    conn.execute('UPDATE tasks SET sort_order=? WHERE id=?', (order, task_row['id']))
        # 「节」是任务级分组，必须随账号持久化，不能只依赖浏览器 localStorage。
        conn.execute('''CREATE TABLE IF NOT EXISTS task_sections (
            user_id INTEGER PRIMARY KEY,
            sections_json TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            FOREIGN KEY (user_id) REFERENCES users(id)
        )''')
        conn.commit()
    finally:
        conn.close()

init_user_db()


def get_user_ark_api_key(user_id):
    conn = get_db_conn()
    try:
        row = conn.execute('SELECT ark_api_key_encrypted FROM users WHERE id=?', (user_id,)).fetchone()
    finally:
        conn.close()
    return decrypt_user_secret(row['ark_api_key_encrypted']) if row else ''


AI_KEY_REQUIRED_PATHS = {
    '/analyze-content', '/generate-story-plan', '/generate-panel-script', '/generate',
    '/generate-script', '/process-file/', '/api/extract_pdf', '/api/regenerate_script',
    '/api/text_input', '/api/generate_backgrounds', '/api/regenerate_background',
    '/api/regenerate_panel', '/api/generate_character_image', '/api/generate_background',
    '/api/generate_all_backgrounds', '/api/generate_comic_from_script', '/api/analyze_emotion',
}


@app.before_request
def bind_user_api_key_to_request():
    """验证会话并绑定当前用户密钥；公网 AI 功能必须登录且完成绑定。"""
    uid = session.get('user_id')
    api_key = ''
    valid_user = False
    _request_image_provider.set(None)
    if uid:
        conn = get_db_conn()
        try:
            row = conn.execute('''SELECT auth_version, ark_api_key_encrypted,
                                  image_provider, image_api_key_encrypted,
                                  image_base_url, image_model
                                  FROM users WHERE id=?''', (uid,)).fetchone()
        finally:
            conn.close()
        valid_user = bool(
            row and int(session.get('auth_version', -1)) == int(row['auth_version'] or 0)
        )
        if valid_user:
            api_key = decrypt_user_secret(row['ark_api_key_encrypted'])
            image_key = decrypt_user_secret(row['image_api_key_encrypted'])
            provider = (row['image_provider'] or '').strip().lower()
            if provider and image_key:
                _request_image_provider.set({
                    'provider': provider,
                    'api_key': image_key,
                    'base_url': (row['image_base_url'] or '').strip(),
                    'model': (row['image_model'] or '').strip(),
                })
            elif api_key:
                # 老用户无感迁移：仍使用原豆包密钥和默认图片模型。
                _request_image_provider.set({
                    'provider': 'doubao', 'api_key': api_key,
                    'base_url': 'https://ark.cn-beijing.volces.com/api/v3/images/generations',
                    'model': ARK_IMAGE_MODEL,
                })
            else:
                _request_image_provider.set(None)
        else:
            session.pop('user_id', None)
            session.pop('auth_version', None)
    _request_ark_key.set(api_key)

    needs_ark_key = any(
        request.path == path or (path.endswith('/') and request.path.startswith(path))
        for path in AI_KEY_REQUIRED_PATHS
    )
    if (os.environ.get('RAILWAY_ENVIRONMENT') and request.method == 'POST' and needs_ark_key):
        if not valid_user:
            return jsonify({'success': False, 'error': '请先登录账号'}), 401
        if not api_key and not get_request_image_provider():
            return jsonify({'success': False, 'error': '请先在左侧账号区域配置 AI 服务'}), 428

def get_current_user():
    uid = session.get('user_id')
    if not uid:
        return None
    conn = get_db_conn()
    try:
        row = conn.execute('''SELECT id, work_id, name, created_at, auth_version,
                              ark_api_key_last4, image_provider, image_api_key_last4,
                              image_base_url, image_model FROM users WHERE id=?''', (uid,)).fetchone()
    finally:
        conn.close()
    if not row:
        return None
    if int(session.get('auth_version', -1)) != int(row['auth_version'] or 0):
        session.pop('user_id', None)
        session.pop('auth_version', None)
        return None
    image_provider = (row['image_provider'] or ('doubao' if row['ark_api_key_last4'] else '')).strip()
    image_last4 = row['image_api_key_last4'] or row['ark_api_key_last4'] or ''
    return {'id': row['id'], 'work_id': row['work_id'], 'name': row['name'],
            'created_at': row['created_at'], 'ark_key_configured': bool(row['ark_api_key_last4']),
            'ark_key_last4': row['ark_api_key_last4'] or '',
            'image_provider': image_provider, 'image_key_configured': bool(image_last4),
            'image_key_last4': image_last4, 'image_model': row['image_model'] or '',
            'image_base_url': row['image_base_url'] or ''}

def _now():
    return datetime.now().strftime('%Y-%m-%d %H:%M:%S')

@app.route('/api/register', methods=['POST'])
def api_register():
    try:
        data = request.get_json(force=True, silent=True) or {}
        work_id = (data.get('work_id') or '').strip()
        name = (data.get('name') or '').strip()
        password = str(data.get('password') or '')
        if not work_id or not name or not password:
            return jsonify({'success': False, 'error': '工号、姓名和密码都不能为空'}), 400
        if len(password) < 6:
            return jsonify({'success': False, 'error': '密码至少需要 6 位'}), 400
        conn = get_db_conn()
        try:
            existing = conn.execute('SELECT id FROM users WHERE work_id=?', (work_id,)).fetchone()
            if existing:
                return jsonify({'success': False, 'error': '该工号已注册，请直接登录'}), 409
            now = _now()
            cur = conn.execute('''INSERT INTO users
                (work_id, name, password_hash, created_at)
                VALUES (?,?,?,?)''',
                (work_id, name, generate_password_hash(password), now))
            uid = cur.lastrowid
            conn.commit()
        finally:
            conn.close()
        session['user_id'] = uid
        session['auth_version'] = 0
        return jsonify({'success': True, 'user': {'id': uid, 'work_id': work_id, 'name': name}})
    except Exception as e:
        logger.error(f"注册失败: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/login', methods=['POST'])
def api_login():
    try:
        data = request.get_json(force=True, silent=True) or {}
        work_id = (data.get('work_id') or '').strip()
        name = (data.get('name') or '').strip()
        password = str(data.get('password') or '')
        if not work_id or not name or not password:
            return jsonify({'success': False, 'error': '工号、姓名和密码都不能为空'}), 400
        conn = get_db_conn()
        try:
            row = conn.execute('''SELECT id, work_id, name, password_hash, auth_version
                                  FROM users WHERE work_id=?''', (work_id,)).fetchone()
            if not row:
                return jsonify({'success': False, 'error': '工号未注册，请先注册'}), 404
            if row['name'] != name:
                return jsonify({'success': False, 'error': '姓名与注册信息不符'}), 403
            if not row['password_hash']:
                return jsonify({'success': False, 'error': '旧账号尚未设置密码，请联系管理员迁移'}), 403
            if not check_password_hash(row['password_hash'], password):
                return jsonify({'success': False, 'error': '密码错误'}), 403
            uid = row['id']
        finally:
            conn.close()
        session['user_id'] = uid
        session['auth_version'] = int(row['auth_version'] or 0)
        return jsonify({'success': True, 'user': {'id': uid, 'work_id': row['work_id'], 'name': row['name']}})
    except Exception as e:
        logger.error(f"登录失败: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/logout', methods=['POST'])
def api_logout():
    session.pop('user_id', None)
    session.pop('auth_version', None)
    return jsonify({'success': True})


@app.route('/api/profile', methods=['PUT'])
def api_update_profile():
    """允许已登录用户用当前密码修正自己的姓名。"""
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '请先登录'}), 401
    data = request.get_json(silent=True) or {}
    name = str(data.get('name') or '').strip()
    current_password = str(data.get('current_password') or '')
    if not name or not current_password:
        return jsonify({'success': False, 'error': '姓名和当前密码不能为空'}), 400
    conn = get_db_conn()
    try:
        row = conn.execute('SELECT password_hash FROM users WHERE id=?', (user['id'],)).fetchone()
        if not row or not row['password_hash'] or not check_password_hash(row['password_hash'], current_password):
            return jsonify({'success': False, 'error': '当前密码错误'}), 403
        conn.execute('UPDATE users SET name=? WHERE id=?', (name, user['id']))
        conn.commit()
    finally:
        conn.close()
    return jsonify({'success': True, 'message': '姓名已更新'})


@app.route('/api/password/reset', methods=['POST'])
def api_password_reset():
    data = request.get_json(silent=True) or {}
    work_id = str(data.get('work_id') or '').strip()
    name = str(data.get('name') or '').strip()
    new_password = str(data.get('new_password') or '')
    if not work_id or not name or len(new_password) < 6:
        return jsonify({'success': False, 'error': '请填写工号、姓名和至少 6 位的新密码'}), 400
    attempt_key = (request.remote_addr or 'unknown', work_id)
    now_ts = time.time()
    recent = [ts for ts in PASSWORD_RESET_ATTEMPTS.get(attempt_key, []) if now_ts - ts < 600]
    if len(recent) >= 5:
        return jsonify({'success': False, 'error': '尝试次数过多，请 10 分钟后再试'}), 429
    recent.append(now_ts)
    PASSWORD_RESET_ATTEMPTS[attempt_key] = recent
    conn = get_db_conn()
    try:
        row = conn.execute('SELECT id, name FROM users WHERE work_id=?', (work_id,)).fetchone()
        if not row or row['name'] != name:
            return jsonify({'success': False, 'error': '工号或姓名不正确'}), 403
        conn.execute('''UPDATE users SET password_hash=?, auth_version=auth_version+1
                        WHERE id=?''', (generate_password_hash(new_password), row['id']))
        conn.commit()
    finally:
        conn.close()
    PASSWORD_RESET_ATTEMPTS.pop(attempt_key, None)
    return jsonify({'success': True, 'message': '密码已重置，请重新登录'})


@app.route('/api/ark-key', methods=['GET', 'PUT', 'DELETE'])
def api_user_ark_key():
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '请先登录'}), 401
    if request.method == 'GET':
        return jsonify({'success': True, 'configured': user['ark_key_configured'],
                        'last4': user['ark_key_last4']})
    conn = get_db_conn()
    try:
        if request.method == 'DELETE':
            conn.execute('''UPDATE users SET ark_api_key_encrypted=NULL,
                            ark_api_key_last4=NULL WHERE id=?''', (user['id'],))
            conn.commit()
            _request_ark_key.set('')
            return jsonify({'success': True, 'message': '豆包 API Key 已解除绑定'})

        api_key = str((request.get_json(silent=True) or {}).get('api_key') or '').strip()
        if len(api_key) < 16 or re.search(r'\s', api_key):
            return jsonify({'success': False, 'error': 'API Key 格式不正确'}), 400
        encrypted = encrypt_user_secret(api_key)
        conn.execute('''UPDATE users SET ark_api_key_encrypted=?, ark_api_key_last4=?
                        WHERE id=?''', (encrypted, api_key[-4:], user['id']))
        conn.commit()
        _request_ark_key.set(api_key)
        return jsonify({'success': True, 'configured': True, 'last4': api_key[-4:],
                        'message': '豆包 API Key 已加密保存'})
    except RuntimeError as exc:
        return jsonify({'success': False, 'error': str(exc)}), 503
    finally:
        conn.close()


IMAGE_PROVIDER_DEFAULTS = {
    'doubao': {
        'label': '豆包（火山方舟）',
        'base_url': 'https://ark.cn-beijing.volces.com/api/v3/images/generations',
        'model': ARK_IMAGE_MODEL,
    },
    'aliyun': {
        'label': '阿里云百炼（通义万相）',
        'base_url': 'https://dashscope.aliyuncs.com/api/v1/services/aigc/text2image/image-synthesis',
        'model': 'wanx2.1-t2i-turbo',
    },
    'openai': {
        'label': 'OpenAI',
        'base_url': 'https://api.openai.com/v1/images/generations',
        'model': 'gpt-image-1',
    },
    'custom': {
        'label': '自定义 OpenAI 兼容接口',
        'base_url': '',
        'model': '',
    },
}


def _validate_public_https_url(raw_url):
    """Only allow public HTTPS URLs; block SSRF targets and embedded credentials."""
    from urllib.parse import urlparse
    import ipaddress
    import socket
    value = str(raw_url or '').strip().rstrip('/')
    if not value:
        raise ValueError('接口地址不能为空')
    parsed = urlparse(value)
    if (parsed.scheme != 'https' or not parsed.hostname or parsed.username or parsed.password
            or parsed.query or parsed.fragment):
        raise ValueError('自定义接口必须是公网 HTTPS 地址，且不能包含账号密码')
    host = parsed.hostname.lower()
    if host == 'localhost' or host.endswith('.local'):
        raise ValueError('不能使用本机或内网接口地址')
    try:
        addresses = {item[4][0] for item in socket.getaddrinfo(host, parsed.port or 443)}
        if any(ipaddress.ip_address(address).is_private or ipaddress.ip_address(address).is_loopback
               or ipaddress.ip_address(address).is_link_local or ipaddress.ip_address(address).is_reserved
               for address in addresses):
            raise ValueError('不能使用本机或内网接口地址')
    except socket.gaierror:
        raise ValueError('自定义接口域名无法解析')
    return value


def _validate_custom_image_url(raw_url):
    value = _validate_public_https_url(raw_url)
    if not value.endswith('/images/generations'):
        value += '/images/generations'
    return value


@app.route('/api/image-provider', methods=['GET', 'PUT', 'DELETE'])
def api_image_provider():
    """Per-user encrypted image API settings, independent between accounts."""
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '请先登录'}), 401
    if request.method == 'GET':
        providers = {key: {'label': value['label'], 'default_model': value['model'],
                           'default_base_url': value['base_url']}
                     for key, value in IMAGE_PROVIDER_DEFAULTS.items()}
        return jsonify({'success': True, 'configured': user['image_key_configured'],
                        'provider': user['image_provider'], 'last4': user['image_key_last4'],
                        'model': user['image_model'], 'base_url': user['image_base_url'],
                        'providers': providers})
    conn = get_db_conn()
    try:
        if request.method == 'DELETE':
            existing = conn.execute('SELECT image_provider FROM users WHERE id=?',
                                    (user['id'],)).fetchone()
            if existing and existing['image_provider']:
                conn.execute('''UPDATE users SET image_provider=NULL, image_api_key_encrypted=NULL,
                                image_api_key_last4=NULL, image_base_url=NULL, image_model=NULL
                                WHERE id=?''', (user['id'],))
            else:
                # 旧版账号把豆包同时作为文字和生图密钥；解除时保持旧操作语义。
                conn.execute('''UPDATE users SET ark_api_key_encrypted=NULL, ark_api_key_last4=NULL
                                WHERE id=?''', (user['id'],))
                _request_ark_key.set('')
            conn.commit()
            _request_image_provider.set(None)
            return jsonify({'success': True, 'message': '生图服务配置已解除'})

        body = request.get_json(silent=True) or {}
        provider = str(body.get('provider') or '').strip().lower()
        if provider not in IMAGE_PROVIDER_DEFAULTS:
            return jsonify({'success': False, 'error': '不支持的生图服务商'}), 400
        api_key = str(body.get('api_key') or '').strip()
        if len(api_key) < 16 or re.search(r'\s', api_key):
            return jsonify({'success': False, 'error': 'API Key 格式不正确'}), 400
        defaults = IMAGE_PROVIDER_DEFAULTS[provider]
        model = str(body.get('model') or defaults['model']).strip()
        if not model or len(model) > 120:
            return jsonify({'success': False, 'error': '模型名称不能为空或过长'}), 400
        base_url = defaults['base_url']
        if provider == 'custom':
            try:
                base_url = _validate_custom_image_url(body.get('base_url'))
            except ValueError as exc:
                return jsonify({'success': False, 'error': str(exc)}), 400
        encrypted = encrypt_user_secret(api_key)
        conn.execute('''UPDATE users SET image_provider=?, image_api_key_encrypted=?,
                        image_api_key_last4=?, image_base_url=?, image_model=? WHERE id=?''',
                     (provider, encrypted, api_key[-4:], base_url, model, user['id']))
        conn.commit()
        _request_image_provider.set({'provider': provider, 'api_key': api_key,
                                     'base_url': base_url, 'model': model})
        return jsonify({'success': True, 'configured': True, 'provider': provider,
                        'last4': api_key[-4:], 'message': f"{defaults['label']}生图服务已加密保存"})
    except RuntimeError as exc:
        return jsonify({'success': False, 'error': str(exc)}), 503
    finally:
        conn.close()

@app.route('/api/me', methods=['GET'])
def api_me():
    user = get_current_user()
    return jsonify({'success': True, 'user': user})

@app.route('/api/tasks', methods=['GET'])
def api_list_tasks():
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '未登录'}), 401
    conn = get_db_conn()
    try:
        rows = conn.execute(
            'SELECT id, title, pinned, sort_order, data_json, created_at, updated_at FROM tasks WHERE user_id=? ORDER BY pinned DESC, sort_order ASC, id ASC',
            (user['id'],)).fetchall()
    finally:
        conn.close()
    tasks = []
    for r in rows:
        try:
            data = json.loads(r['data_json'])
        except Exception:
            data = {}
        tasks.append({'id': r['id'], 'title': r['title'], 'pinned': bool(r['pinned']),
                      'sort_order': r['sort_order'], 'data': data, 'created_at': r['created_at'], 'updated_at': r['updated_at']})
    return jsonify({'success': True, 'tasks': tasks})


def _normalize_task_sections_payload(raw_sections):
    """清洗任务级「节」配置，避免错误输入破坏某个用户的全部分组。"""
    if not isinstance(raw_sections, list):
        return []
    normalized, seen_ids = [], set()
    for index, section in enumerate(raw_sections):
        if not isinstance(section, dict):
            continue
        section_id = str(section.get('id') or '').strip()[:80]
        if not section_id or section_id in seen_ids:
            continue
        seen_ids.add(section_id)
        title = str(section.get('title') or f'第{index + 1}节').strip()[:100] or f'第{index + 1}节'
        task_ids = section.get('task_ids') or []
        if not isinstance(task_ids, list):
            task_ids = []
        unique_task_ids, seen_tasks = [], set()
        for task_id in task_ids:
            try:
                clean_task_id = int(task_id)
            except (TypeError, ValueError):
                continue
            if clean_task_id > 0 and clean_task_id not in seen_tasks:
                seen_tasks.add(clean_task_id)
                unique_task_ids.append(clean_task_id)
        normalized.append({'id': section_id, 'title': title, 'task_ids': unique_task_ids})
    return normalized


@app.route('/api/task-sections', methods=['GET'])
def api_get_task_sections():
    """读取当前登录用户持久化的任务分节。"""
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '未登录'}), 401
    conn = get_db_conn()
    try:
        row = conn.execute(
            'SELECT sections_json FROM task_sections WHERE user_id=?', (user['id'],)
        ).fetchone()
    finally:
        conn.close()
    if not row:
        return jsonify({'success': True, 'exists': False, 'sections': []})
    try:
        sections = _normalize_task_sections_payload(json.loads(row['sections_json']))
    except Exception:
        sections = []
    return jsonify({'success': True, 'exists': True, 'sections': sections})


@app.route('/api/task-sections', methods=['PUT'])
def api_save_task_sections():
    """保存当前登录用户的任务分节；覆盖的是该用户自己的分节配置。"""
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '未登录'}), 401
    data = request.get_json(force=True, silent=True) or {}
    sections = _normalize_task_sections_payload(data.get('sections'))
    conn = get_db_conn()
    try:
        conn.execute(
            '''INSERT INTO task_sections (user_id, sections_json, updated_at)
               VALUES (?, ?, ?)
               ON CONFLICT(user_id) DO UPDATE SET
                 sections_json=excluded.sections_json,
                 updated_at=excluded.updated_at''',
            (user['id'], json.dumps(sections, ensure_ascii=False), _now())
        )
        conn.commit()
    finally:
        conn.close()
    return jsonify({'success': True, 'sections': sections})


@app.route('/api/tasks', methods=['POST'])
def api_create_comic_task():
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '未登录，无法保存任务'}), 401
    data = request.get_json(force=True, silent=True) or {}
    title = (data.get('title') or '未命名任务').strip() or '未命名任务'
    task_data = data.get('data', {})
    now = _now()
    conn = get_db_conn()
    try:
        next_order = conn.execute(
            'SELECT COALESCE(MAX(sort_order), 0) + 1 AS next_order FROM tasks WHERE user_id=?',
            (user['id'],),
        ).fetchone()['next_order']
        cur = conn.execute(
            # 7 个字段：user_id、title、固定的 pinned=0，以及 4 个动态值。
            # 原来这里多写了一个占位符，导致每次新建任务均触发 SQLite 500，
            # 表现为“点击保存任务后侧边栏没有新增任务”。
            'INSERT INTO tasks (user_id, title, pinned, sort_order, data_json, created_at, updated_at) VALUES (?,?,0,?,?,?,?)',
            (user['id'], title, next_order, json.dumps(task_data, ensure_ascii=False), now, now))
        tid = cur.lastrowid
        conn.commit()
    finally:
        conn.close()
    return jsonify({'success': True, 'task': {'id': tid, 'title': title, 'pinned': False, 'sort_order': next_order, 'data': task_data}})


@app.route('/api/tasks/order', methods=['PUT'])
def api_reorder_comic_tasks():
    """Persist a user-chosen task sequence without touching updated_at."""
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '未登录'}), 401
    payload = request.get_json(force=True, silent=True) or {}
    task_ids = payload.get('task_ids')
    if not isinstance(task_ids, list) or not task_ids:
        return jsonify({'success': False, 'error': '任务排序数据无效'}), 400
    try:
        ordered_ids = [int(item) for item in task_ids]
    except (TypeError, ValueError):
        return jsonify({'success': False, 'error': '任务编号无效'}), 400
    if len(set(ordered_ids)) != len(ordered_ids):
        return jsonify({'success': False, 'error': '任务排序包含重复项'}), 400
    conn = get_db_conn()
    try:
        owned = {row['id'] for row in conn.execute('SELECT id FROM tasks WHERE user_id=?', (user['id'],)).fetchall()}
        if set(ordered_ids) != owned:
            return jsonify({'success': False, 'error': '任务排序与当前账号不一致，请刷新后重试'}), 400
        for order, task_id in enumerate(ordered_ids, start=1):
            conn.execute('UPDATE tasks SET sort_order=? WHERE id=? AND user_id=?', (order, task_id, user['id']))
        conn.commit()
    finally:
        conn.close()
    return jsonify({'success': True})

@app.route('/api/tasks/<int:task_id>', methods=['PUT'])
def api_update_task(task_id):
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '未登录'}), 401
    data = request.get_json(force=True, silent=True) or {}
    conn = get_db_conn()
    try:
        row = conn.execute('SELECT id, user_id FROM tasks WHERE id=?', (task_id,)).fetchone()
        if not row or row['user_id'] != user['id']:
            return jsonify({'success': False, 'error': '任务不存在或无权限'}), 404
        fields = []
        params = []
        if 'title' in data:
            fields.append('title=?'); params.append((data['title'] or '未命名任务').strip())
        if 'pinned' in data:
            fields.append('pinned=?'); params.append(1 if data['pinned'] else 0)
        if 'data' in data:
            fields.append('data_json=?'); params.append(json.dumps(data['data'], ensure_ascii=False))
        fields.append('updated_at=?'); params.append(_now())
        params.append(task_id)
        conn.execute('UPDATE tasks SET ' + ', '.join(fields) + ' WHERE id=?', params)
        conn.commit()
    finally:
        conn.close()
    return jsonify({'success': True})

@app.route('/api/tasks/<int:task_id>', methods=['DELETE'])
def api_delete_task(task_id):
    user = get_current_user()
    if not user:
        return jsonify({'success': False, 'error': '未登录'}), 401
    conn = get_db_conn()
    try:
        row = conn.execute('SELECT id, user_id FROM tasks WHERE id=?', (task_id,)).fetchone()
        if not row or row['user_id'] != user['id']:
            return jsonify({'success': False, 'error': '任务不存在或无权限'}), 404
        conn.execute('DELETE FROM tasks WHERE id=?', (task_id,))
        conn.commit()
    finally:
        conn.close()
    return jsonify({'success': True})


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    logger.info(f"启动Flask应用，端口: {port}")
    # 开发时仍保留代码/模板热重载，但忽略运行时生成的图片、ZIP、PPT 和视频。
    # 否则用户一导出文件，Werkzeug 会中断正在发送的自动保存请求，造成任务丢失。
    app.run(
        host='0.0.0.0',
        port=port,
        debug=True,
        threaded=True,
        exclude_patterns=['static/output/*', 'static\\output\\*'],
    )
